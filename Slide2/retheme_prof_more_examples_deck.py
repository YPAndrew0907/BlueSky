#!/usr/bin/env python3
from __future__ import annotations

import argparse
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Literal

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


PIL_RGBA = "RGBA"


@dataclass(frozen=True)
class ThemeColors:
    card_fill_old: RGBColor
    card_fill_new: RGBColor
    card_line_old: RGBColor
    card_line_new: RGBColor
    muted_text_old: RGBColor
    muted_text_new: RGBColor


DV_THEME = ThemeColors(
    card_fill_old=RGBColor(0x16, 0x20, 0x33),
    card_fill_new=RGBColor(0x11, 0x1C, 0x2D),
    card_line_old=RGBColor(0x3A, 0x4D, 0x66),
    card_line_new=RGBColor(0x2B, 0x3A, 0x55),
    muted_text_old=RGBColor(0xC9, 0xD2, 0xE3),
    muted_text_new=RGBColor(0xB7, 0xC2, 0xD6),
)


def _extract_image_from_pptx(*, pptx_path: Path, media_name: str, out_path: Path) -> Path:
    with zipfile.ZipFile(pptx_path, "r") as z:
        candidates = [n for n in z.namelist() if n == f"ppt/media/{media_name}"]
        if not candidates:
            raise FileNotFoundError(f"missing media in template: ppt/media/{media_name} (pptx={pptx_path})")
        out_path.parent.mkdir(parents=True, exist_ok=True)
        out_path.write_bytes(z.read(candidates[0]))
        return out_path


def _generate_dv_bg_overlay(*, image6_jpg: Path, out_png: Path) -> Path:
    out_png.parent.mkdir(parents=True, exist_ok=True)

    canvas_w, canvas_h = 1920, 1080
    base = Image.new(PIL_RGBA, (canvas_w, canvas_h), "#0B1320")

    overlay = Image.open(image6_jpg).convert(PIL_RGBA)
    resampling = getattr(Image, "Resampling", Image)
    overlay = ImageOps.fit(overlay, (canvas_w, canvas_h), method=resampling.LANCZOS)

    overlay = ImageEnhance.Brightness(overlay).enhance(0.35)
    overlay = ImageEnhance.Color(overlay).enhance(0.0)
    overlay = overlay.filter(ImageFilter.GaussianBlur(radius=1.5))

    alpha = overlay.getchannel("A").point(lambda p: int(round(p * 0.22)))
    overlay.putalpha(alpha)

    out = Image.alpha_composite(base, overlay)
    out.save(out_png, format="PNG")
    return out_png


def _picture_shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh
    raise KeyError(f"picture shape_id={shape_id} not found")


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001
    blip.set(qn("r:embed"), rid)


def _slide_text(slide) -> str:
    tokens: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            t = (p.text or "").strip()
            if t:
                tokens.append(t)
    return " ".join(tokens)


def _is_locked_slide(slide) -> bool:
    text = _slide_text(slide)
    return ("Bluesky UI:" in text) or ("Data receipt:" in text)


def _maybe_swap_shape_fill(shape, *, old: RGBColor, new: RGBColor) -> bool:
    try:
        fill = shape.fill
    except Exception:
        return False
    if fill.type != MSO_FILL.SOLID:
        return False
    fore = fill.fore_color
    if fore.type != MSO_COLOR_TYPE.RGB:
        return False
    if fore.rgb != old:
        return False
    fore.rgb = new
    return True


def _maybe_swap_shape_line(shape, *, old: RGBColor, new: RGBColor) -> bool:
    try:
        line = shape.line
    except Exception:
        return False
    if line is None:
        return False
    color = line.color
    if color is None or color.type != MSO_COLOR_TYPE.RGB:
        return False
    if color.rgb != old:
        return False
    color.rgb = new
    return True


def _maybe_swap_text_colors(shape, *, old: RGBColor, new: RGBColor) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0
    swapped = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            color = run.font.color
            if color is None or color.type != MSO_COLOR_TYPE.RGB:
                continue
            if color.rgb != old:
                continue
            color.rgb = new
            swapped += 1
    return swapped


def _normalize_nonlocked_slide_colors(slide, *, theme: ThemeColors) -> None:
    for shape in slide.shapes:
        _maybe_swap_shape_fill(shape, old=theme.card_fill_old, new=theme.card_fill_new)
        _maybe_swap_shape_line(shape, old=theme.card_line_old, new=theme.card_line_new)
        _maybe_swap_text_colors(shape, old=theme.muted_text_old, new=theme.muted_text_new)


def retheme_deck(
    *,
    pptx_in: Path,
    pptx_out: Path,
    assert_dir: Path,
    cache_dir: Path,
    mode: Literal["safe", "full"],
) -> None:
    pptx_in = pptx_in.resolve()
    pptx_out = pptx_out.resolve()
    assert_dir = assert_dir.resolve()
    cache_dir = cache_dir.resolve()

    if not pptx_in.exists():
        raise FileNotFoundError(f"missing --in: {pptx_in}")

    template = assert_dir / "Data Visual by Slidesgo.pptx"
    if not template.exists():
        raise FileNotFoundError(f"missing template: {template}")

    image6 = _extract_image_from_pptx(pptx_path=template, media_name="image6.jpg", out_path=cache_dir / "image6.jpg")
    overlay_png = _generate_dv_bg_overlay(image6_jpg=image6, out_png=cache_dir / "dv_bg_overlay.png")

    prs = Presentation(str(pptx_in))
    for slide_num, slide in enumerate(prs.slides, start=1):
        overlay_pic = _picture_shape(slide, 3)
        _replace_picture_image(overlay_pic, slide, overlay_png)

        if mode == "full":
            if _is_locked_slide(slide):
                continue
            _normalize_nonlocked_slide_colors(slide, theme=DV_THEME)

    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_out))


def main() -> None:
    here = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(description="Retheme the professor deck by swapping the global overlay image.")
    parser.add_argument(
        "--in",
        dest="pptx_in",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v12.pptx",
        help="Input PPTX (source of truth).",
    )
    parser.add_argument(
        "--out",
        dest="pptx_out",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v13.base.pptx",
        help="Output PPTX.",
    )
    parser.add_argument(
        "--assert-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slides/assert"),
        help="Slides/assert directory containing the Slidesgo template PPTX.",
    )
    parser.add_argument(
        "--cache-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slide2/prof_build/assets_cache/theme_dv"),
        help="Cache directory for extracted template media and derived overlay.",
    )
    parser.add_argument(
        "--mode",
        type=str,
        choices=("safe", "full"),
        default="safe",
        help="safe: replace only picture id=3. full: also normalize fills/lines/text on non-locked slides.",
    )
    args = parser.parse_args()

    retheme_deck(
        pptx_in=args.pptx_in,
        pptx_out=args.pptx_out,
        assert_dir=args.assert_dir,
        cache_dir=args.cache_dir,
        mode=args.mode,
    )
    print(f"OK: wrote {args.pptx_out.resolve()}")


if __name__ == "__main__":
    main()
