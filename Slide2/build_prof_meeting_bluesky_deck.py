#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Literal

from PIL import Image, ImageEnhance, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


AssetRole = Literal["structural", "semantic"]

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
FONT_NAME = "Calibri"

COLORS = {
    "bg": "0B1320",
    "ink": "F7F9FC",
    "muted": "C9D2E3",
    "card": "162033",
    "card_line": "3A4D66",
    "accent_cyan": "58B7E6",
    "accent_coral": "FF7D77",
    "accent_amber": "F4C55D",
    "accent_green": "54C687",
}

REQUIRED_ARTIFACTS = [
    "02_csv_exports/feed_generators_index.csv",
    "02_csv_exports/starterpacks.csv",
    "02_csv_exports/starterpack_feeds.csv",
    "02_csv_exports/discovery_feed_inclusions.csv",
    "02_csv_exports/popular_feeds.csv",
    "02_csv_exports/feed_panel.csv",
    "02_csv_exports/feed_items.csv.gz",
    "02_csv_exports/posts.csv.gz",
    "02_csv_exports/authors.csv.gz",
    "02_csv_exports/post_labels.csv.gz",
    "state.db",
    "05_manifest/run_metadata.csv",
    "05_manifest/run_summary.csv",
    "05_manifest/manifest.csv",
    "05_manifest/data_dictionary.csv",
    "05_manifest/validation_report.csv",
]

_BskyUiMode = Literal["clear", "sanitized"]

# User request: "not even blur" - default to clear/raw screenshots (no redaction).
# If you later need a shareable/public version, switch to "sanitized".
BSKY_UI_MODE: _BskyUiMode = "clear"

# Bluesky UI screenshot assets.
BSKY_UI_ASSETS_BY_MODE: dict[_BskyUiMode, dict[str, str]] = {
    # Clear/raw UI screenshots (no redaction).
    "clear": {
        "B1": "bsky_raw_feeds.png",
        "B2": "bsky_raw_feed_for_science.png",
        "B3": "bsky_raw_starterpack.png",
    },
    # Privacy-sanitized UI screenshots (pixelation + blur in risky regions).
    "sanitized": {
        "B1": "feeds_market_v2.png",
        "B2": "feed_science_timeline_v2.png",
        "B3": "starterpack_science_minds_v2.png",
    },
}


@dataclass(frozen=True)
class AssetSource:
    asset_id: str
    source_pptx: Path
    media_path: str
    role: AssetRole


@dataclass(frozen=True)
class PreparedAsset:
    asset_id: str
    source_pptx: Path
    media_path: str
    role: AssetRole
    derived_path: Path


@dataclass(frozen=True)
class AssetRef:
    source_pptx: Path
    media_path: str
    derived_path: Path
    role: AssetRole
    usage: str


@dataclass(frozen=True)
class SlideSpec:
    num: int
    title: str
    hidden: bool
    takeaway: str
    purpose: str
    speaker_bullets: list[str]
    asset_ids: list[str]


@dataclass(frozen=True)
class ManifestEntry:
    slide: int
    title: str
    assets: list[AssetRef]


@dataclass(frozen=True)
class BuildConfig:
    assert_dir: Path
    run_dir: Path
    working_out: Path
    final_out: Path
    slide_index_out: Path
    asset_manifest_out: Path
    assets_cache_dir: Path


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _extract_media(*, source_pptx: Path, media_path: str, out_path: Path) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    out_path.parent.mkdir(parents=True, exist_ok=True)
    with zipfile.ZipFile(source_pptx, "r") as zf:
        try:
            payload = zf.read(media_path)
        except KeyError as err:
            raise RuntimeError(f"Missing media path {media_path} in {source_pptx}") from err
    out_path.write_bytes(payload)
    return out_path


def _prep_full_background(*, src_path: Path, out_path: Path, size_px: tuple[int, int]) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    base = Image.new("RGBA", size_px, tuple(int(COLORS["bg"][i : i + 2], 16) for i in (0, 2, 4)) + (255,))
    overlay = Image.open(src_path).convert("RGBA")
    overlay = ImageOps.fit(overlay, size_px, method=Image.LANCZOS)
    overlay = ImageEnhance.Brightness(overlay).enhance(0.42)
    overlay = ImageEnhance.Color(overlay).enhance(0.15)
    alpha = overlay.split()[-1].point(lambda a: int(a * 0.55))
    overlay.putalpha(alpha)
    out = Image.alpha_composite(base, overlay)
    out.save(out_path)
    return out_path


def _prep_card_texture(*, src_path: Path, out_path: Path, size_px: tuple[int, int]) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    texture = Image.open(src_path).convert("RGBA")
    texture = ImageOps.fit(texture, size_px, method=Image.LANCZOS)
    texture = ImageEnhance.Brightness(texture).enhance(0.25)
    texture = ImageEnhance.Color(texture).enhance(0.0)
    texture = texture.filter(ImageFilter.GaussianBlur(radius=1.1))
    texture.putalpha(texture.split()[-1].point(lambda a: int(a * 0.35)))
    texture.save(out_path)
    return out_path


def _prep_strip(*, src_path: Path, out_path: Path, size_px: tuple[int, int], alpha_mult: float) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    strip = Image.open(src_path).convert("RGBA")
    strip = ImageOps.fit(strip, size_px, method=Image.LANCZOS)
    strip.putalpha(strip.split()[-1].point(lambda a: int(a * alpha_mult)))
    strip.save(out_path)
    return out_path


def _prep_icon(*, src_path: Path, out_path: Path, size_px: int) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    icon = Image.open(src_path).convert("RGBA")
    icon = ImageOps.contain(icon, (size_px, size_px), method=Image.LANCZOS)
    canvas = Image.new("RGBA", (size_px, size_px), (0, 0, 0, 0))
    left = (size_px - icon.width) // 2
    top = (size_px - icon.height) // 2
    canvas.paste(icon, (left, top), icon)
    canvas.save(out_path)
    return out_path


def _prep_photo_overlay(*, src_path: Path, out_path: Path, size_px: tuple[int, int]) -> Path:
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    photo = Image.open(src_path).convert("RGBA")
    photo = ImageOps.fit(photo, size_px, method=Image.LANCZOS)
    photo = ImageEnhance.Color(photo).enhance(0.0)
    photo = ImageEnhance.Brightness(photo).enhance(0.42)
    photo = photo.filter(ImageFilter.GaussianBlur(radius=0.6))
    photo.putalpha(photo.split()[-1].point(lambda a: int(a * 0.35)))
    photo.save(out_path)
    return out_path


def _build_asset_sources(assert_dir: Path) -> dict[str, AssetSource]:
    return {
        "A1": AssetSource("A1", assert_dir / "Architect Infographics by Slidesgo.pptx", "ppt/media/image12.png", "structural"),
        "A2": AssetSource("A2", assert_dir / "Architect Infographics by Slidesgo.pptx", "ppt/media/image3.png", "structural"),
        "A3": AssetSource("A3", assert_dir / "How To Make An Infographic by Slidesgo.pptx", "ppt/media/image2.png", "structural"),
        "A4": AssetSource("A4", assert_dir / "How To Make An Infographic by Slidesgo.pptx", "ppt/media/image1.png", "structural"),
        "A5": AssetSource("A5", assert_dir / "EN What Is a Satellite_ by Slidesgo.pptx", "ppt/media/image2.png", "structural"),
        "A6": AssetSource("A6", assert_dir / "EN What Is a Satellite_ by Slidesgo.pptx", "ppt/media/image3.png", "structural"),
        "A7": AssetSource(
            "A7",
            assert_dir / "Space Exploration Mission Pitch Deck by Slidesgo.pptx",
            "ppt/media/image9.jpg",
            "structural",
        ),
        "A8": AssetSource("A8", assert_dir / "How To Make An Infographic by Slidesgo.pptx", "ppt/media/image3.png", "structural"),
        "A9": AssetSource("A9", assert_dir / "Copy of Stock Market Infographics by Slidesgo.pptx", "ppt/media/image1.png", "structural"),
        "S1": AssetSource("S1", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image1.png", "semantic"),
        "S2": AssetSource("S2", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image2.png", "semantic"),
        "S3": AssetSource("S3", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image3.png", "semantic"),
        "S4": AssetSource("S4", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image4.png", "semantic"),
        "S5": AssetSource("S5", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image5.png", "semantic"),
        "S6": AssetSource("S6", assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx", "ppt/media/image6.png", "semantic"),
    }


def _ensure_bsky_ui_screenshot_assets(*, bsky_dir: Path, mode: _BskyUiMode) -> dict[str, Path]:
    """Ensure Bluesky UI screenshots exist on disk (clear or sanitized).

    These images are used to explain core Bluesky concepts (feed marketplace, feed timeline, starter pack)
    to a non-Bluesky audience. In "sanitized" mode we redact handles/post text via pixelation+blur.
    """

    raw_dir = bsky_dir / "raw"
    sanitized_dir = bsky_dir / "sanitized_test3"
    filenames = BSKY_UI_ASSETS_BY_MODE[mode]
    base_dir = raw_dir if mode == "clear" else sanitized_dir
    expected = {asset_id: base_dir / filename for asset_id, filename in filenames.items()}
    if all(path.exists() and path.stat().st_size > 0 for path in expected.values()):
        return expected

    # Lazy import: capture requires Chrome + network, but most builds should be able to reuse
    # already-sanitized outputs.
    from bsky_screenshots import (  # noqa: PLC0415
        CaptureSpec,
        DerivedSpec,
        RectRel,
        capture_bsky_screenshots,
        derive_sanitized_images,
    )

    capture_specs = [
        CaptureSpec(shot_id="feeds", url="https://bsky.app/feeds"),
        CaptureSpec(
            shot_id="feed_for_science",
            url="https://bsky.app/profile/did:plc:jfhpnnst6flqway4eaeqzj2a/feed/for-science",
        ),
        CaptureSpec(
            shot_id="starterpack",
            url="https://bsky.app/starter-pack/did:plc:32kt7lsk5pgcjkmtptu6sfyc/3lk562bzxnd2u",
        ),
    ]
    raw_by_id = capture_bsky_screenshots(out_dir=raw_dir, specs=capture_specs)

    # If we're in clear mode, we only need the raw captures.
    if mode == "clear":
        missing = [p for p in expected.values() if not (p.exists() and p.stat().st_size > 0)]
        if missing:
            missing_text = "\n".join(f"- {p}" for p in missing)
            raise RuntimeError(f"Bluesky UI raw assets missing after capture:\n{missing_text}")
        return expected

    # Redaction rectangles are tuned for 1920x1080 desktop captures (keep UI chrome + key headings).
    derived_specs = [
        DerivedSpec(
            asset_id="feeds_market_v2",
            source_shot_id="feeds",
            crop_rel=None,
            redact_rels=(RectRel(0.28, 0.22, 0.72, 0.96),),
            out_size_px=None,
        ),
        DerivedSpec(
            asset_id="feed_science_timeline_v2",
            source_shot_id="feed_for_science",
            crop_rel=None,
            redact_rels=(
                RectRel(0.25, 0.14, 0.75, 0.96),  # posts area
                RectRel(0.42, 0.03, 0.75, 0.12),  # header byline/handle line
            ),
            out_size_px=None,
        ),
        DerivedSpec(
            asset_id="starterpack_science_minds_v2",
            source_shot_id="starterpack",
            crop_rel=None,
            redact_rels=(
                RectRel(0.40, 0.16, 0.70, 0.235),  # byline handle strip
                RectRel(0.25, 0.38, 0.75, 0.96),  # account list
            ),
            out_size_px=None,
        ),
    ]
    derive_sanitized_images(raw_by_id=raw_by_id, out_dir=sanitized_dir, derived=derived_specs)

    missing = [p for p in expected.values() if not (p.exists() and p.stat().st_size > 0)]
    if missing:
        missing_text = "\n".join(f"- {p}" for p in missing)
        raise RuntimeError(f"Bluesky UI sanitized assets missing after capture:\n{missing_text}")

    return expected


def _prepare_bsky_ui_assets(*, cache_dir: Path) -> dict[str, PreparedAsset]:
    bsky_dir = Path(__file__).resolve().parent / "prof_build" / "bsky"
    by_id = _ensure_bsky_ui_screenshot_assets(bsky_dir=bsky_dir, mode=BSKY_UI_MODE)

    prepared: dict[str, PreparedAsset] = {}
    for asset_id, src_path in by_id.items():
        derived_path = cache_dir / f"{asset_id}_ui.png"
        derived_path.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(src_path, derived_path)
        prepared[asset_id] = PreparedAsset(
            asset_id=asset_id,
            source_pptx=src_path,
            media_path=str(src_path.name),
            role="semantic",
            derived_path=derived_path,
        )
    return prepared


def prepare_assets(assert_dir: Path, cache_dir: Path) -> dict[str, PreparedAsset]:
    sources = _build_asset_sources(assert_dir)
    cache_dir.mkdir(parents=True, exist_ok=True)

    # Ensure deterministic regeneration when asset source mapping changes.
    generated_files = [
        "A1_bg_blueprint_overlay.png",
        "A2_card_texture.png",
        "A3_divider_strip.png",
        "A4_header_bar.png",
        "A5_illustration.png",
        "A6_illustration.png",
        "A7_photo_overlay.png",
        "A8_donut.png",
        "A9_chart.png",
        "S1_icon.png",
        "S2_icon.png",
        "S3_icon.png",
        "S4_icon.png",
        "S5_icon.png",
        "S6_icon.png",
    ]
    for filename in generated_files:
        target = cache_dir / filename
        if target.exists():
            target.unlink()

    for src in sources.values():
        if not src.source_pptx.exists():
            raise RuntimeError(f"Missing assert source pptx: {src.source_pptx}")

    raw_cache: dict[tuple[Path, str], Path] = {}
    for src in sources.values():
        key = (src.source_pptx, src.media_path)
        if key in raw_cache:
            continue
        raw_name = f"raw_{src.source_pptx.stem.replace(' ', '_')}__{Path(src.media_path).name}"
        raw_cache[key] = _extract_media(
            source_pptx=src.source_pptx,
            media_path=src.media_path,
            out_path=cache_dir / raw_name,
        )

    prepared: dict[str, PreparedAsset] = {}

    prepared["A1"] = PreparedAsset(
        asset_id="A1",
        source_pptx=sources["A1"].source_pptx,
        media_path=sources["A1"].media_path,
        role="structural",
        derived_path=_prep_full_background(
            src_path=raw_cache[(sources["A1"].source_pptx, sources["A1"].media_path)],
            out_path=cache_dir / "A1_bg_blueprint_overlay.png",
            size_px=(1920, 1080),
        ),
    )
    prepared["A2"] = PreparedAsset(
        asset_id="A2",
        source_pptx=sources["A2"].source_pptx,
        media_path=sources["A2"].media_path,
        role="structural",
        derived_path=_prep_card_texture(
            src_path=raw_cache[(sources["A2"].source_pptx, sources["A2"].media_path)],
            out_path=cache_dir / "A2_card_texture.png",
            size_px=(1400, 900),
        ),
    )
    prepared["A3"] = PreparedAsset(
        asset_id="A3",
        source_pptx=sources["A3"].source_pptx,
        media_path=sources["A3"].media_path,
        role="structural",
        derived_path=_prep_strip(
            src_path=raw_cache[(sources["A3"].source_pptx, sources["A3"].media_path)],
            out_path=cache_dir / "A3_divider_strip.png",
            size_px=(1600, 180),
            alpha_mult=0.96,
        ),
    )
    prepared["A4"] = PreparedAsset(
        asset_id="A4",
        source_pptx=sources["A4"].source_pptx,
        media_path=sources["A4"].media_path,
        role="structural",
        derived_path=_prep_strip(
            src_path=raw_cache[(sources["A4"].source_pptx, sources["A4"].media_path)],
            out_path=cache_dir / "A4_header_bar.png",
            size_px=(1300, 220),
            alpha_mult=0.95,
        ),
    )
    prepared["A5"] = PreparedAsset(
        asset_id="A5",
        source_pptx=sources["A5"].source_pptx,
        media_path=sources["A5"].media_path,
        role="structural",
        derived_path=_prep_icon(
            src_path=raw_cache[(sources["A5"].source_pptx, sources["A5"].media_path)],
            out_path=cache_dir / "A5_illustration.png",
            size_px=1200,
        ),
    )
    prepared["A6"] = PreparedAsset(
        asset_id="A6",
        source_pptx=sources["A6"].source_pptx,
        media_path=sources["A6"].media_path,
        role="structural",
        derived_path=_prep_icon(
            src_path=raw_cache[(sources["A6"].source_pptx, sources["A6"].media_path)],
            out_path=cache_dir / "A6_illustration.png",
            size_px=1200,
        ),
    )
    prepared["A7"] = PreparedAsset(
        asset_id="A7",
        source_pptx=sources["A7"].source_pptx,
        media_path=sources["A7"].media_path,
        role="structural",
        derived_path=_prep_photo_overlay(
            src_path=raw_cache[(sources["A7"].source_pptx, sources["A7"].media_path)],
            out_path=cache_dir / "A7_photo_overlay.png",
            size_px=(1500, 900),
        ),
    )
    prepared["A8"] = PreparedAsset(
        asset_id="A8",
        source_pptx=sources["A8"].source_pptx,
        media_path=sources["A8"].media_path,
        role="structural",
        derived_path=_prep_icon(
            src_path=raw_cache[(sources["A8"].source_pptx, sources["A8"].media_path)],
            out_path=cache_dir / "A8_donut.png",
            size_px=900,
        ),
    )
    prepared["A9"] = PreparedAsset(
        asset_id="A9",
        source_pptx=sources["A9"].source_pptx,
        media_path=sources["A9"].media_path,
        role="structural",
        derived_path=_prep_icon(
            src_path=raw_cache[(sources["A9"].source_pptx, sources["A9"].media_path)],
            out_path=cache_dir / "A9_chart.png",
            size_px=900,
        ),
    )

    for sid in ("S1", "S2", "S3", "S4", "S5", "S6"):
        src = sources[sid]
        prepared[sid] = PreparedAsset(
            asset_id=sid,
            source_pptx=src.source_pptx,
            media_path=src.media_path,
            role="semantic",
            derived_path=_prep_icon(
                src_path=raw_cache[(src.source_pptx, src.media_path)],
                out_path=cache_dir / f"{sid}_icon.png",
                size_px=560,
            ),
        )

    prepared.update(_prepare_bsky_ui_assets(cache_dir=cache_dir))
    return prepared


def _shape(slide, shape_id: int):
    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape
    raise KeyError(f"Missing shape_id={shape_id}")


def _set_shape_text(
    shape,
    lines: list[str],
    *,
    size_pt: float,
    color_hex: str = COLORS["ink"],
    bold_first: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(size_pt)
            run.font.bold = bool(idx == 0 and bold_first)
            run.font.color.rgb = _rgb(color_hex)


def _add_textbox(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: list[str],
    size_pt: float,
    bold_first: bool = False,
    color_hex: str = COLORS["ink"],
    align: PP_ALIGN = PP_ALIGN.LEFT,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_shape_text(
        shape,
        lines,
        size_pt=size_pt,
        color_hex=color_hex,
        bold_first=bold_first,
        align=align,
    )
    return shape


def _add_card(slide, *, x: float, y: float, w: float, h: float):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(COLORS["card"])
    card.fill.transparency = 0.08
    card.line.color.rgb = _rgb(COLORS["card_line"])
    card.line.width = Pt(1.4)
    return card


def _add_chip(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    accent_hex: str,
    size_pt: float = 20,
) -> None:
    chip = _add_card(slide, x=x, y=y, w=w, h=h)
    chip.fill.transparency = 0.14
    chip.line.color.rgb = _rgb(accent_hex)
    chip.line.width = Pt(2.2)
    _set_shape_text(chip, [text], size_pt=size_pt, bold_first=True)


def _add_marker_badge(
    slide,
    *,
    x: float,
    y: float,
    d: float,
    label: str,
    fill_hex: str,
    text_hex: str = "1A1A1A",
    size_pt: float = 20,
) -> None:
    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(x),
        Inches(y),
        Inches(d),
        Inches(d),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(fill_hex)
    badge.line.fill.background()
    _set_shape_text(badge, [label], size_pt=size_pt, color_hex=text_hex, bold_first=True, align=PP_ALIGN.CENTER)


def _add_highlight_rect(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    accent_hex: str,
    fill_alpha: float = 0.86,
) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    # NOTE: Filled translucent overlays don't render reliably in our PDF export pipeline
    # (LibreOffice tends to flatten transparency). Use outline-only highlights so the
    # screenshot stays readable in both PPTX and PDF.
    box.fill.background()
    box.line.color.rgb = _rgb(accent_hex)
    box.line.width = Pt(3.2)


def _add_full_picture(slide, prs: Presentation, image_path: Path):
    slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)


def _add_picture(slide, *, image_path: Path, x: float, y: float, w: float, h: float):
    return slide.shapes.add_picture(str(image_path), Inches(x), Inches(y), Inches(w), Inches(h))


def _record_asset(
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
    *,
    slide_num: int,
    asset: PreparedAsset,
    usage: str,
) -> None:
    bucket = manifest.setdefault(slide_num, [])
    bucket_ids = manifest_ids.setdefault(slide_num, [])
    key = (str(asset.derived_path), usage)
    for existing in bucket:
        if (str(existing.derived_path), existing.usage) == key:
            return
    bucket.append(
        AssetRef(
            source_pptx=asset.source_pptx,
            media_path=asset.media_path,
            derived_path=asset.derived_path,
            role=asset.role,
            usage=usage,
        )
    )
    bucket_ids.append(asset.asset_id)


def _apply_base_slide(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    bg = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        prs.slide_width,
        prs.slide_height,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(COLORS["bg"])
    bg.line.fill.background()
    _add_full_picture(slide, prs, assets["A1"].derived_path)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A1"], usage="full background overlay")


def _draw_system_link(slide, *, x1: float, y1: float, x2: float, y2: float):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    line.line.color.rgb = _rgb(COLORS["muted"])
    line.line.width = Pt(1.6)


def _set_slide_fade_transition(slide, *, speed: str = "slow") -> None:
    sld = slide._element
    for child in list(sld):
        if child.tag == qn("p:transition"):
            sld.remove(child)
    transition = OxmlElement("p:transition")
    transition.set("spd", speed)
    transition.append(OxmlElement("p:fade"))
    timing = sld.find(qn("p:timing"))
    if timing is not None:
        sld.insert(list(sld).index(timing), transition)
    else:
        sld.append(transition)


def _slide_specs() -> list[SlideSpec]:
    return [
        SlideSpec(
            num=1,
            title="Algorithmic choice relocates power",
            hidden=False,
            takeaway="Choice does not remove concentration risk.",
            purpose="Choice does not remove concentration risk.",
            speaker_bullets=[
                "State the thesis in one sentence.",
                "Set expectation: no results in this meeting.",
                "Focus the discussion on priorities and recollection.",
            ],
            asset_ids=["A1", "A2", "A3", "A5", "A8", "S6"],
        ),
        SlideSpec(
            num=2,
            title="System framing: pluggable services, new chokepoints",
            hidden=False,
            takeaway="Feed and labeler modularity shifts where power concentrates.",
            purpose="Feed and labeler modularity shifts where power concentrates.",
            speaker_bullets=[
                "Show key nodes in the ecosystem.",
                "Explain discovery and hosting chokepoints.",
                "Connect chokepoints to security and governance stakes.",
            ],
            asset_ids=["A1", "A2", "A3", "A4", "A6", "S2"],
        ),
        SlideSpec(
            num=3,
            title="Bluesky UI: feed marketplace (/feeds)",
            hidden=False,
            takeaway="A feed is a selectable algorithmic timeline; feeds are discoverable in a directory.",
            purpose="Ground the discussion in real UI: where feeds are discovered/selected.",
            speaker_bullets=[
                "Show the public feeds directory.",
                "Explain that each feed is an algorithmic timeline made by someone.",
                "Point out that discovery is centralized by directory + defaults.",
            ],
            asset_ids=["A1", "A2", "A4", "B1", "S2"],
        ),
        SlideSpec(
            num=4,
            title="Bluesky UI: what a feed looks like",
            hidden=False,
            takeaway="A feed turns an algorithm into a timeline (posts).",
            purpose="Show what a feed output looks like (timeline) and what users can do with it.",
            speaker_bullets=[
                "Show a feed page and its identity (name + likes).",
                "Explain the output is a timeline of posts.",
                "Tie back to research questions about exposure and manipulation.",
            ],
            asset_ids=["A1", "A2", "A4", "B2", "S1"],
        ),
        SlideSpec(
            num=5,
            title="Bluesky UI: starter packs (onboarding bundles)",
            hidden=False,
            takeaway="Starter packs bundle accounts (and sometimes feeds) for onboarding.",
            purpose="Explain starter packs as a discovery/onboarding surface worth measuring.",
            speaker_bullets=[
                "Show a starter pack page.",
                "Explain one-click onboarding: follow people + add feeds.",
                "Tie to discovery concentration + governance questions.",
            ],
            asset_ids=["A1", "A2", "A4", "B3", "S3"],
        ),
        SlideSpec(
            num=6,
            title="Research questions (R1-R5)",
            hidden=False,
            takeaway="The RQs define our security and privacy measurement scope.",
            purpose="The RQs define our security and privacy measurement scope.",
            speaker_bullets=[
                "Frame R1-R4 as measurable with current artifacts.",
                "Frame R5 as future work requiring extra guardrails.",
            ],
            asset_ids=["A1", "A2", "A8", "A9", "S1", "S2", "S3", "S4", "S5"],
        ),
        SlideSpec(
            num=7,
            title="Data artifacts collected (no results)",
            hidden=False,
            takeaway="Artifacts are joinable, reproducible, and auditable.",
            purpose="Artifacts are joinable, reproducible, and auditable.",
            speaker_bullets=[
                "Group artifacts by discovery, snapshots, and trust receipts.",
                "Show join spine for future analysis.",
                "Emphasize no empirical claims yet.",
            ],
            asset_ids=["A1", "A2", "A4", "A5", "A8", "S6"],
        ),
        SlideSpec(
            num=8,
            title="Collection credibility (one slide)",
            hidden=False,
            takeaway="Read-only collection and reproducible receipts support later inference.",
            purpose="Read-only collection and reproducible receipts support later inference.",
            speaker_bullets=[
                "Summarize read-only XRPC posture.",
                "Show seven-stage pipeline at high level.",
                "Anchor reproducibility in state and manifest artifacts.",
            ],
            asset_ids=["A1", "A3", "A7", "A9", "S4"],
        ),
        SlideSpec(
            num=9,
            title="What this enables next + advisor asks",
            hidden=False,
            takeaway="Next measurements are clear, and priority decisions are explicit.",
            purpose="Next measurements are clear, and priority decisions are explicit.",
            speaker_bullets=[
                "Map next analyses to RQ1-RQ4.",
                "Request priority guidance for venue fit.",
                "Ask what recollection to run next.",
            ],
            asset_ids=["A1", "A4", "A6", "A8", "S5"],
        ),
        SlideSpec(
            num=10,
            title="Backup: API surfaces (names only)",
            hidden=True,
            takeaway="API scope is clear without endpoint clutter.",
            purpose="API scope is clear without endpoint clutter.",
            speaker_bullets=[
                "Keep this table for Q and A only.",
                "Confirm read APIs vs auth-only APIs.",
            ],
            asset_ids=["A1", "A4", "A9", "S2"],
        ),
        SlideSpec(
            num=11,
            title="Backup: folder layout + trust stamps",
            hidden=True,
            takeaway="Run receipts are auditable and reproducible.",
            purpose="Run receipts are auditable and reproducible.",
            speaker_bullets=[
                "Show where each artifact class is stored.",
                "Point to trust stamps for integrity and schema.",
            ],
            asset_ids=["A1", "A2", "A5", "S3"],
        ),
    ]


def _write_storyboard(slides: list[SlideSpec], out_path: Path) -> None:
    lines = ["# Internal storyboard", ""]
    for spec in slides:
        hidden_tag = " [hidden]" if spec.hidden else ""
        lines.append(f"Slide {spec.num}: {spec.title}{hidden_tag}")
        lines.append(f"- takeaway: {spec.takeaway}")
        lines.append(f"- assets: {', '.join(spec.asset_ids)}")
        lines.append(f"- bullets: {'; '.join(spec.speaker_bullets)}")
        lines.append("")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("\n".join(lines), encoding="utf-8")


def _build_slide_1(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)

    _add_picture(slide, image_path=assets["A3"].derived_path, x=0.0, y=0.0, w=13.333, h=0.36)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A3"], usage="top divider strip")
    _add_picture(slide, image_path=assets["A4"].derived_path, x=8.85, y=0.06, w=3.6, h=0.28)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="accent strip")
    _add_picture(slide, image_path=assets["A8"].derived_path, x=11.02, y=0.18, w=1.45, h=1.45)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A8"], usage="hero donut accent")

    _add_textbox(
        slide,
        x=0.72,
        y=0.62,
        w=9.0,
        h=0.9,
        lines=["Algorithmic choice relocates power"],
        size_pt=40,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.55,
        w=8.9,
        h=0.66,
        lines=["Research questions plus data artifacts collected (no results yet)."],
        size_pt=22,
        color_hex=COLORS["muted"],
    )

    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.72, y=2.52, w=11.9, h=2.25)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="thesis card texture")
    _add_card(slide, x=0.72, y=2.52, w=11.9, h=2.25)
    _add_textbox(
        slide,
        x=1.02,
        y=2.95,
        w=11.3,
        h=1.45,
        lines=[
            "Algorithmic choice does not remove power - it relocates it.",
            "Meeting focus: prioritize RQs and define what to recollect next.",
        ],
        size_pt=24,
    )

    badge = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(9.9),
        Inches(1.78),
        Inches(2.7),
        Inches(0.62),
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(COLORS["accent_coral"])
    badge.line.fill.background()
    _set_shape_text(
        badge,
        ["NO RESULTS TODAY"],
        size_pt=20,
        color_hex="1A1A1A",
        bold_first=True,
        align=PP_ALIGN.CENTER,
    )

    _add_picture(slide, image_path=assets["A5"].derived_path, x=9.8, y=4.05, w=2.62, h=2.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A5"], usage="hero illustration")
    _add_picture(slide, image_path=assets["S6"].derived_path, x=11.25, y=0.5, w=1.3, h=1.3)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S6"], usage="title semantic icon")


def _build_slide_2(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.28, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A3"].derived_path, x=0.7, y=6.82, w=12.0, h=0.24)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A3"], usage="footer accent strip")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.94, w=12.08, h=3.34)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="diagram texture")

    _add_textbox(
        slide,
        x=0.72,
        y=0.72,
        w=10.6,
        h=0.82,
        lines=["System framing: pluggable services, new chokepoints"],
        size_pt=33,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.52,
        w=9.5,
        h=0.5,
        lines=["Takeaway: modular services relocate concentration."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    node_specs = [
        ("viewer", 0.72, 2.65, 1.9, 0.9),
        ("discovery surfaces", 2.92, 2.12, 2.3, 0.9),
        ("feed generators", 5.52, 2.12, 2.3, 0.9),
        ("provider hosting", 8.12, 2.12, 2.3, 0.9),
        ("exposure outcomes", 10.72, 2.65, 1.9, 0.9),
        ("labelers", 5.52, 4.05, 2.3, 0.9),
    ]
    for label, x, y, w, h in node_specs:
        card = _add_card(slide, x=x, y=y, w=w, h=h)
        _set_shape_text(card, [label], size_pt=20, align=PP_ALIGN.CENTER, bold_first=True)

    _draw_system_link(slide, x1=2.62, y1=3.10, x2=2.92, y2=2.57)
    _draw_system_link(slide, x1=5.22, y1=2.57, x2=5.52, y2=2.57)
    _draw_system_link(slide, x1=7.82, y1=2.57, x2=8.12, y2=2.57)
    _draw_system_link(slide, x1=10.42, y1=2.57, x2=10.72, y2=3.10)
    _draw_system_link(slide, x1=6.67, y1=3.02, x2=6.67, y2=4.05)
    _draw_system_link(slide, x1=7.37, y1=3.02, x2=7.37, y2=4.05)

    _add_textbox(
        slide,
        x=0.72,
        y=5.65,
        w=8.8,
        h=0.95,
        lines=["Discovery surfaces and provider hosting are key chokepoints."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    _add_picture(slide, image_path=assets["A6"].derived_path, x=9.4, y=4.32, w=2.86, h=2.22)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A6"], usage="supporting illustration")
    _add_picture(slide, image_path=assets["S2"].derived_path, x=11.28, y=0.52, w=1.2, h=1.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S2"], usage="diagram semantic icon")


def _build_slide_3_bsky_feeds_market(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.28, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.32, w=12.1, h=5.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="UI card texture")

    _add_textbox(
        slide,
        x=0.72,
        y=0.70,
        w=10.6,
        h=0.82,
        lines=["Bluesky UI: feed marketplace (/feeds)"],
        size_pt=33,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.52,
        w=10.8,
        h=0.42,
        lines=["Takeaway: feeds are discoverable + selectable algorithmic timelines."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    # Screenshot (sanitized).
    shot_x, shot_y, shot_w = 4.05, 1.86, 8.58
    shot_h = shot_w * 9.0 / 16.0
    _add_picture(slide, image_path=assets["B1"].derived_path, x=shot_x, y=shot_y, w=shot_w, h=shot_h)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["B1"], usage="Bluesky /feeds screenshot (sanitized)")
    _add_picture(slide, image_path=assets["S2"].derived_path, x=11.28, y=0.52, w=1.2, h=1.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S2"], usage="UI slide icon")

    chips = [
        ("A. /feeds directory", COLORS["accent_cyan"]),
        ("Search feeds", COLORS["accent_cyan"]),
        ("B. Feed cards", COLORS["accent_amber"]),
        ("Creator + likes", COLORS["accent_amber"]),
        ("C. Trending widget", COLORS["accent_coral"]),
        ("Defaults matter", COLORS["accent_coral"]),
        ("Discovery risk", COLORS["accent_coral"]),
    ]
    y = 2.06
    for text, accent in chips:
        _add_chip(slide, x=0.72, y=y, w=3.12, h=0.50, text=text, accent_hex=accent)
        y += 0.56

    # Highlights on the screenshot.
    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 0.32, d=0.34, label="A", fill_hex=COLORS["accent_cyan"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.52,
        y=shot_y + 0.30,
        w=shot_w * 0.62,
        h=0.52,
        accent_hex=COLORS["accent_cyan"],
        fill_alpha=0.90,
    )
    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 1.18, d=0.34, label="B", fill_hex=COLORS["accent_amber"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.52,
        y=shot_y + 1.14,
        w=shot_w * 0.62,
        h=shot_h * 0.62,
        accent_hex=COLORS["accent_amber"],
        fill_alpha=0.92,
    )
    _add_marker_badge(
        slide,
        x=shot_x + (shot_w * 0.70) + 0.18,
        y=shot_y + 0.32,
        d=0.34,
        label="C",
        fill_hex=COLORS["accent_coral"],
        text_hex=COLORS["bg"],
    )
    _add_highlight_rect(
        slide,
        x=shot_x + (shot_w * 0.70) + 0.46,
        y=shot_y + 0.30,
        w=shot_w * 0.24,
        h=1.22,
        accent_hex=COLORS["accent_coral"],
        fill_alpha=0.92,
    )

    _draw_system_link(slide, x1=3.84, y1=2.30, x2=shot_x + 0.52, y2=shot_y + 0.56)
    _draw_system_link(slide, x1=3.84, y1=3.42, x2=shot_x + 0.52, y2=shot_y + 1.42)


def _build_slide_4_bsky_feed_timeline(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.28, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.32, w=12.1, h=5.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="UI card texture")

    _add_textbox(
        slide,
        x=0.72,
        y=0.70,
        w=10.6,
        h=0.82,
        lines=["Bluesky UI: what a feed looks like"],
        size_pt=33,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.52,
        w=10.8,
        h=0.42,
        lines=["Takeaway: a feed page renders an algorithm into a timeline of posts."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    shot_x, shot_y, shot_w = 4.05, 1.86, 8.58
    shot_h = shot_w * 9.0 / 16.0
    _add_picture(slide, image_path=assets["B2"].derived_path, x=shot_x, y=shot_y, w=shot_w, h=shot_h)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["B2"], usage="Bluesky feed timeline screenshot (sanitized)")
    _add_picture(slide, image_path=assets["S1"].derived_path, x=11.28, y=0.52, w=1.2, h=1.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S1"], usage="UI slide icon")

    chips = [
        ("A. Feed name", COLORS["accent_cyan"]),
        ("Demand signals", COLORS["accent_cyan"]),
        ("B. Timeline output", COLORS["accent_amber"]),
        ("Exposure decisions", COLORS["accent_amber"]),
        ("C. Post actions", COLORS["accent_amber"]),
        ("Why we measure", COLORS["accent_coral"]),
        ("Manipulation risk", COLORS["accent_coral"]),
    ]
    y = 2.06
    for text, accent in chips:
        _add_chip(slide, x=0.72, y=y, w=3.12, h=0.50, text=text, accent_hex=accent)
        y += 0.56

    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 0.12, d=0.34, label="A", fill_hex=COLORS["accent_cyan"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.50,
        y=shot_y + 0.10,
        w=shot_w * 0.70,
        h=0.62,
        accent_hex=COLORS["accent_cyan"],
        fill_alpha=0.90,
    )
    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 1.05, d=0.34, label="B", fill_hex=COLORS["accent_amber"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.50,
        y=shot_y + 1.02,
        w=shot_w * 0.70,
        h=shot_h * 0.62,
        accent_hex=COLORS["accent_amber"],
        fill_alpha=0.92,
    )
    _add_marker_badge(
        slide,
        x=shot_x + 0.18,
        y=shot_y + (shot_h * 0.62) + 0.34,
        d=0.34,
        label="C",
        fill_hex=COLORS["accent_amber"],
        text_hex=COLORS["bg"],
    )
    _add_highlight_rect(
        slide,
        x=shot_x + 0.70,
        y=shot_y + (shot_h * 0.62) + 0.32,
        w=shot_w * 0.36,
        h=0.58,
        accent_hex=COLORS["accent_amber"],
        fill_alpha=0.94,
    )

    _draw_system_link(slide, x1=3.84, y1=2.30, x2=shot_x + 0.50, y2=shot_y + 0.42)
    _draw_system_link(slide, x1=3.84, y1=3.42, x2=shot_x + 0.50, y2=shot_y + 1.34)


def _build_slide_5_bsky_starterpack(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.28, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.32, w=12.1, h=5.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="UI card texture")

    _add_textbox(
        slide,
        x=0.72,
        y=0.70,
        w=10.9,
        h=0.82,
        lines=["Bluesky UI: starter packs (onboarding bundles)"],
        size_pt=33,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.52,
        w=11.2,
        h=0.42,
        lines=["Takeaway: starter packs bundle people (and sometimes feeds) for fast onboarding."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    shot_x, shot_y, shot_w = 4.05, 1.86, 8.58
    shot_h = shot_w * 9.0 / 16.0
    _add_picture(slide, image_path=assets["B3"].derived_path, x=shot_x, y=shot_y, w=shot_w, h=shot_h)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["B3"], usage="Bluesky starter pack screenshot (sanitized)")
    _add_picture(slide, image_path=assets["S3"].derived_path, x=11.28, y=0.52, w=1.2, h=1.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S3"], usage="UI slide icon")

    chips = [
        ("A. Pack title", COLORS["accent_cyan"]),
        ("Curated list", COLORS["accent_cyan"]),
        ("B. Join button", COLORS["accent_amber"]),
        ("1-click onboarding", COLORS["accent_amber"]),
        ("C. People cards", COLORS["accent_amber"]),
        ("Curation questions", COLORS["accent_coral"]),
        ("Measure exposure", COLORS["accent_coral"]),
    ]
    y = 2.06
    for text, accent in chips:
        _add_chip(slide, x=0.72, y=y, w=3.12, h=0.50, text=text, accent_hex=accent)
        y += 0.56

    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 0.14, d=0.34, label="A", fill_hex=COLORS["accent_cyan"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.50,
        y=shot_y + 0.10,
        w=shot_w * 0.74,
        h=0.88,
        accent_hex=COLORS["accent_cyan"],
        fill_alpha=0.90,
    )
    _add_marker_badge(slide, x=shot_x + 0.18, y=shot_y + 1.18, d=0.34, label="B", fill_hex=COLORS["accent_amber"], text_hex=COLORS["bg"])
    _add_highlight_rect(
        slide,
        x=shot_x + 0.50,
        y=shot_y + 1.14,
        w=shot_w * 0.74,
        h=0.72,
        accent_hex=COLORS["accent_amber"],
        fill_alpha=0.92,
    )
    _add_marker_badge(
        slide,
        x=shot_x + 0.18,
        y=shot_y + 2.10,
        d=0.34,
        label="C",
        fill_hex=COLORS["accent_amber"],
        text_hex=COLORS["bg"],
    )
    _add_highlight_rect(
        slide,
        x=shot_x + 0.50,
        y=shot_y + 2.06,
        w=shot_w * 0.74,
        h=shot_h * 0.52,
        accent_hex=COLORS["accent_amber"],
        fill_alpha=0.94,
    )

    _draw_system_link(slide, x1=3.84, y1=2.30, x2=shot_x + 0.50, y2=shot_y + 0.54)
    _draw_system_link(slide, x1=3.84, y1=2.86, x2=shot_x + 0.50, y2=shot_y + 1.50)


def _build_slide_3(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.28, w=12.1, h=4.98)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="RQ card texture")
    _add_picture(slide, image_path=assets["A9"].derived_path, x=10.1, y=0.25, w=2.18, h=1.08)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A9"], usage="title chart accent")
    _add_picture(slide, image_path=assets["A3"].derived_path, x=0.72, y=1.96, w=4.0, h=0.24)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A3"], usage="divider accent")

    _add_textbox(
        slide,
        x=0.72,
        y=0.66,
        w=7.8,
        h=0.82,
        lines=["Research questions (R1-R5)"],
        size_pt=36,
        bold_first=True,
    )
    _add_textbox(
        slide,
        x=0.72,
        y=1.56,
        w=10.8,
        h=0.4,
        lines=["Takeaway: security and privacy scope."],
        size_pt=20,
        color_hex=COLORS["muted"],
    )

    rq_boxes = [
        (0.82, 2.05, 5.78, 1.02, "R1 Discovery chokepoints? governance.", "S1"),
        (6.74, 2.05, 5.78, 1.02, "R2 Provider buckets dominate? reliability.", "S2"),
        (0.82, 3.17, 5.78, 1.02, "R3 Same winners across feeds? manipulation.", "S3"),
        (6.74, 3.17, 5.78, 1.02, "R4 Labels vary by viewer? safety.", "S4"),
        (
            0.82,
            4.29,
            11.72,
            1.14,
            "R5 Future experiment: gaming plus privacy attacks. extra guardrails.",
            "S5",
        ),
    ]
    for x, y, w, h, text, icon_id in rq_boxes:
        card = _add_card(slide, x=x, y=y, w=w, h=h)
        _set_shape_text(card, [text], size_pt=20)
        icon_x = x + 0.08
        icon_y = y + 0.08
        _add_picture(slide, image_path=assets[icon_id].derived_path, x=icon_x, y=icon_y, w=0.58, h=0.58)
        _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets[icon_id], usage=f"{icon_id} RQ marker")

    _add_picture(slide, image_path=assets["A8"].derived_path, x=10.86, y=4.22, w=1.46, h=1.46)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A8"], usage="RQ5 donut accent")


def _build_slide_4(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.26, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.58, y=1.24, w=12.2, h=5.84)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="artifact card texture")
    _add_picture(slide, image_path=assets["S6"].derived_path, x=11.22, y=0.5, w=1.2, h=1.2)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S6"], usage="artifact slide icon")
    _add_picture(slide, image_path=assets["A8"].derived_path, x=9.74, y=0.06, w=1.62, h=1.62)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A8"], usage="header donut accent")

    _add_textbox(
        slide,
        x=0.72,
        y=0.64,
        w=10.2,
        h=0.78,
        lines=["Data artifacts collected (no results)"],
        size_pt=34,
        bold_first=True,
    )

    c1 = _add_card(slide, x=0.58, y=1.5, w=4.0, h=3.05)
    _set_shape_text(
        c1,
        [
            "Discovery inputs",
            "feed_generators_index.csv",
            "starterpacks.csv",
            "starterpack_feeds.csv",
            "discovery_feed_inclusions.csv",
            "popular_feeds.csv",
        ],
        size_pt=20,
        bold_first=True,
    )

    c2 = _add_card(slide, x=4.66, y=1.5, w=4.0, h=3.05)
    _set_shape_text(
        c2,
        [
            "Panel snapshots",
            "feed_panel.csv",
            "feed_items.csv.gz",
            "posts.csv.gz",
            "authors.csv.gz",
            "post_labels.csv.gz",
        ],
        size_pt=20,
        bold_first=True,
    )

    c3 = _add_card(slide, x=8.74, y=1.5, w=4.0, h=3.05)
    _set_shape_text(
        c3,
        [
            "Trust and reproducibility",
            "state.db",
            "run_metadata.csv",
            "run_summary.csv",
            "manifest.csv",
            "data_dictionary.csv",
            "validation_report.csv",
        ],
        size_pt=20,
        bold_first=True,
    )

    spine = _add_card(slide, x=0.58, y=4.7, w=12.16, h=2.1)
    _set_shape_text(
        spine,
        [
            "Join spine",
            "feed_panel -> feed_items -> posts/authors",
            "post_labels attaches by feed_uri, viewer_mode, post_uri, post_cid.",
            "Result: joinable, reproducible, auditable artifacts.",
        ],
        size_pt=20,
        bold_first=True,
    )
    _add_picture(slide, image_path=assets["A5"].derived_path, x=9.95, y=4.74, w=2.47, h=1.9)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A5"], usage="artifact supporting illustration")


def _build_slide_5(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A3"].derived_path, x=0.0, y=0.0, w=13.333, h=0.35)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A3"], usage="top divider strip")
    _add_picture(slide, image_path=assets["A7"].derived_path, x=8.46, y=1.46, w=4.2, h=4.84)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A7"], usage="photo overlay accent")
    _add_picture(slide, image_path=assets["A9"].derived_path, x=9.82, y=0.32, w=1.85, h=0.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A9"], usage="method chart accent")
    _add_picture(slide, image_path=assets["S4"].derived_path, x=11.25, y=0.52, w=1.18, h=1.18)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S4"], usage="method icon")

    _add_textbox(
        slide,
        x=0.72,
        y=0.62,
        w=10.2,
        h=0.82,
        lines=["Collection credibility (one slide)"],
        size_pt=34,
        bold_first=True,
    )

    left = _add_card(slide, x=0.72, y=1.54, w=5.2, h=4.9)
    _set_shape_text(
        left,
        [
            "Read-only XRPC summary",
            "Relay + AppView reads.",
            "GET-only in this run.",
            "POST auth only if auth mode is used.",
            "No posting actions.",
        ],
        size_pt=20,
        bold_first=True,
    )

    right = _add_card(slide, x=6.15, y=1.54, w=6.45, h=4.9)
    _set_shape_text(
        right,
        [
            "7-stage pipeline",
            "relay -> index -> packs -> popular -> enrich -> panel -> snapshot",
            "Reproducibility: state.db, manifest.csv, data_dictionary.csv, validation_report.csv.",
        ],
        size_pt=20,
        bold_first=True,
    )


def _build_slide_6(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.26, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A8"].derived_path, x=10.92, y=0.24, w=1.34, h=1.34)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A8"], usage="next donut accent")
    _add_picture(slide, image_path=assets["S5"].derived_path, x=11.25, y=0.52, w=1.18, h=1.18)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S5"], usage="next steps icon")

    _add_textbox(
        slide,
        x=0.72,
        y=0.62,
        w=10.6,
        h=0.82,
        lines=["What this enables next + advisor asks"],
        size_pt=34,
        bold_first=True,
    )

    left = _add_card(slide, x=0.72, y=1.56, w=7.4, h=5.0)
    _set_shape_text(
        left,
        [
            "Ready measurements",
            "RQ1: measure discovery concentration.",
            "RQ2: map provider leverage.",
            "RQ3: test same winner overlap.",
            "RQ4: compare label variability.",
        ],
        size_pt=20,
        bold_first=True,
    )

    right = _add_card(slide, x=8.36, y=1.96, w=4.2, h=3.9)
    _set_shape_text(
        right,
        [
            "What I need from you",
            "Ask 1: priority RQs and venue fit.",
            "Ask 2: recollect longitudinal, auth mode, or experiments?",
        ],
        size_pt=20,
        bold_first=True,
    )
    _add_picture(slide, image_path=assets["A6"].derived_path, x=9.46, y=4.22, w=2.72, h=2.12)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A6"], usage="advisor ask illustration")


def _build_slide_7(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A4"].derived_path, x=0.72, y=0.26, w=4.6, h=0.4)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A4"], usage="section header bar")
    _add_picture(slide, image_path=assets["A9"].derived_path, x=10.72, y=0.26, w=1.78, h=0.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A9"], usage="backup chart accent")
    _add_picture(slide, image_path=assets["S2"].derived_path, x=11.25, y=0.52, w=1.18, h=1.18)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S2"], usage="backup API icon")

    _add_textbox(
        slide,
        x=0.72,
        y=0.62,
        w=10.8,
        h=0.82,
        lines=["Backup: API surfaces (names only)"],
        size_pt=32,
        bold_first=True,
    )

    table = _add_card(slide, x=0.9, y=1.7, w=11.5, h=4.95)
    _set_shape_text(
        table,
        [
            "Surface | Names",
            "Relay | com.atproto.sync.listReposByCollection",
            "AppView | getActorFeeds, getFeedGenerators, getStarterPack, getFeed, getProfiles",
            "PDS/Auth | createSession, refreshSession (auth mode only)",
        ],
        size_pt=20,
        bold_first=True,
    )


def _build_slide_8(
    slide,
    prs: Presentation,
    *,
    slide_num: int,
    assets: dict[str, PreparedAsset],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    _apply_base_slide(slide, prs, slide_num=slide_num, assets=assets, manifest=manifest, manifest_ids=manifest_ids)
    _add_picture(slide, image_path=assets["A3"].derived_path, x=0.0, y=0.0, w=13.333, h=0.33)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A3"], usage="top divider strip")
    _add_picture(slide, image_path=assets["A2"].derived_path, x=0.62, y=1.3, w=12.1, h=5.7)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A2"], usage="backup texture")
    _add_picture(slide, image_path=assets["S3"].derived_path, x=11.25, y=0.52, w=1.18, h=1.18)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["S3"], usage="backup folder icon")

    _add_textbox(
        slide,
        x=0.72,
        y=0.62,
        w=10.8,
        h=0.82,
        lines=["Backup: folder layout + trust stamps"],
        size_pt=32,
        bold_first=True,
    )

    left = _add_card(slide, x=0.82, y=1.72, w=6.0, h=5.2)
    _set_shape_text(
        left,
        [
            "Run directory layout",
            "01_state_db/",
            "02_csv_exports/",
            "03_postprocess_metrics/",
            "04_logs/",
            "05_manifest/",
            "06_figures_preview/",
            "07_archive_zip/",
        ],
        size_pt=20,
        bold_first=True,
    )
    right = _add_card(slide, x=7.02, y=1.72, w=5.5, h=5.2)
    _set_shape_text(
        right,
        [
            "Trust stamps",
            "run_metadata.csv",
            "run_summary.csv",
            "manifest.csv",
            "data_dictionary.csv",
            "validation_report.csv",
            "state.db",
        ],
        size_pt=20,
        bold_first=True,
    )
    _add_picture(slide, image_path=assets["A5"].derived_path, x=9.96, y=4.84, w=2.46, h=1.92)
    _record_asset(manifest, manifest_ids, slide_num=slide_num, asset=assets["A5"], usage="backup illustration")


def _ensure_required_artifacts(run_dir: Path) -> None:
    missing = [rel for rel in REQUIRED_ARTIFACTS if not (run_dir / rel).exists()]
    if missing:
        missing_text = "\n".join(f"- {item}" for item in missing)
        raise RuntimeError(f"Missing required run artifacts under {run_dir}:\n{missing_text}")


def _write_slide_index(slides: list[SlideSpec], out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    for spec in slides:
        bullets = "; ".join(spec.speaker_bullets)
        hidden_tag = " [hidden backup]" if spec.hidden else ""
        lines.append(
            f"Slide {spec.num}{hidden_tag} — {spec.title} — {spec.purpose} — {bullets}"
        )
    out_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def _write_asset_manifest(
    entries: list[ManifestEntry],
    out_path: Path,
    manifest_ids: dict[int, list[str]],
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    lines: list[str] = []
    for entry in entries:
        lines.append(f"Slide {entry.slide} | {entry.title}")
        ids = manifest_ids.get(entry.slide, [])
        for idx, asset in enumerate(entry.assets):
            asset_id = ids[idx] if idx < len(ids) else f"asset_{idx + 1}"
            lines.append(
                "ASSET | "
                f"{asset.role} | {asset_id} | {asset.source_pptx} | {asset.media_path} | "
                f"{asset.derived_path} | {asset.usage}"
            )
        lines.append("")
    out_path.write_text("\n".join(lines), encoding="utf-8")


def _validate_asset_usage(
    slide_specs: list[SlideSpec],
    manifest: dict[int, list[AssetRef]],
    manifest_ids: dict[int, list[str]],
) -> None:
    for spec in slide_specs:
        assets = manifest.get(spec.num, [])
        ids = set(manifest_ids.get(spec.num, []))
        if not assets:
            raise RuntimeError(f"Slide {spec.num} has no recorded assets")
        if not set(spec.asset_ids).issubset(ids):
            missing = sorted(set(spec.asset_ids) - ids)
            raise RuntimeError(f"Slide {spec.num} missing planned assets: {missing}")
        roles = {a.role for a in assets}
        if not {"structural", "semantic"}.issubset(roles):
            raise RuntimeError(f"Slide {spec.num} must include both structural and semantic assets")
        if len(assets) < 2:
            raise RuntimeError(f"Slide {spec.num} must include at least two assets")

    for special in (1, 7):
        if len(manifest.get(special, [])) < 3:
            raise RuntimeError(f"Slide {special} must include at least three assets")


def build_deck(config: BuildConfig) -> None:
    _ensure_required_artifacts(config.run_dir)
    assets = prepare_assets(config.assert_dir, config.assets_cache_dir)
    slides = _slide_specs()
    _write_storyboard(slides, config.assets_cache_dir.parent / "storyboard_internal.md")

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH_IN)
    prs.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank = prs.slide_layouts[6]

    manifest: dict[int, list[AssetRef]] = {}
    manifest_ids: dict[int, list[str]] = {}

    for spec in slides:
        slide = prs.slides.add_slide(blank)
        builders = {
            1: _build_slide_1,
            2: _build_slide_2,
            3: _build_slide_3_bsky_feeds_market,
            4: _build_slide_4_bsky_feed_timeline,
            5: _build_slide_5_bsky_starterpack,
            6: _build_slide_3,
            7: _build_slide_4,
            8: _build_slide_5,
            9: _build_slide_6,
            10: _build_slide_7,
            11: _build_slide_8,
        }
        builders[spec.num](
            slide,
            prs,
            slide_num=spec.num,
            assets=assets,
            manifest=manifest,
            manifest_ids=manifest_ids,
        )
        _set_slide_fade_transition(slide, speed="slow")

    # Mark backup slides hidden in OOXML.
    for slide_num in (10, 11):
        prs.slides[slide_num - 1]._element.set("show", "0")

    config.working_out.parent.mkdir(parents=True, exist_ok=True)
    tmp_working = config.working_out.with_suffix(".tmp.pptx")
    if tmp_working.exists():
        tmp_working.unlink()
    prs.save(str(tmp_working))
    shutil.move(tmp_working, config.working_out)

    # Add per-shape click-to-reveal animations (python-pptx cannot author animations).
    from pptx_click_animations import inject_click_reveals  # noqa: PLC0415

    with tempfile.TemporaryDirectory(prefix="prof_deck_anim_") as tmp_dir:
        tmp_animated = Path(tmp_dir) / config.working_out.name
        inject_click_reveals(pptx_in=config.working_out, pptx_out=tmp_animated, exclude_spids={2, 3})
        config.working_out.write_bytes(tmp_animated.read_bytes())

    entries: list[ManifestEntry] = []
    for spec in slides:
        entries.append(
            ManifestEntry(
                slide=spec.num,
                title=spec.title,
                assets=manifest.get(spec.num, []),
            )
        )

    _validate_asset_usage(slides, manifest, manifest_ids)
    _write_slide_index(slides, config.slide_index_out)
    _write_asset_manifest(entries, config.asset_manifest_out, manifest_ids)

    config.final_out.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(config.working_out, config.final_out)


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(
        description="Build a fresh professor-ready Bluesky deck (11 slides, assert-first, no-results)."
    )
    parser.add_argument(
        "--assert-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slides/assert"),
    )
    parser.add_argument(
        "--run-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slide2/data_run_20260206"),
    )
    parser.add_argument("--working-out", type=Path, required=True)
    parser.add_argument("--final-out", type=Path, required=True)
    parser.add_argument("--slide-index-out", type=Path, required=True)
    parser.add_argument("--asset-manifest-out", type=Path, required=True)
    parser.add_argument(
        "--assets-cache-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slide2/prof_build/assets_cache"),
    )
    return parser


def main() -> None:
    args = _build_parser().parse_args()
    cfg = BuildConfig(
        assert_dir=args.assert_dir.resolve(),
        run_dir=args.run_dir.resolve(),
        working_out=args.working_out.resolve(),
        final_out=args.final_out.resolve(),
        slide_index_out=args.slide_index_out.resolve(),
        asset_manifest_out=args.asset_manifest_out.resolve(),
        assets_cache_dir=args.assets_cache_dir.resolve(),
    )
    build_deck(cfg)
    print(f"OK: wrote {cfg.working_out}")
    print(f"OK: wrote {cfg.final_out}")
    print(f"OK: wrote {cfg.slide_index_out}")
    print(f"OK: wrote {cfg.asset_manifest_out}")


if __name__ == "__main__":
    main()
