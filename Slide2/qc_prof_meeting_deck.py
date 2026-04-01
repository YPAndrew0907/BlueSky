#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation


EMU_PER_INCH = 914400


@dataclass(frozen=True)
class QcResult:
    slide_count_ok: bool
    hidden_ok: bool
    bounds_ok: bool
    asset_ok: bool
    font_ok: bool
    words_ok: bool
    inventory_ok: bool
    watermark_ok: bool
    timing_ok: bool
    clicks_ok: bool
    total_click_effects: int
    errors: list[str]


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


def _parse_hidden_slides(value: str) -> set[int]:
    if not value.strip():
        return set()
    out: set[int] = set()
    for token in value.split(","):
        token = token.strip()
        if not token:
            continue
        out.add(int(token))
    return out


def _parse_min_assets_special(value: str) -> dict[int, int]:
    if not value.strip():
        return {}
    out: dict[int, int] = {}
    for token in value.split(","):
        token = token.strip()
        if not token:
            continue
        parts = token.split(":", maxsplit=1)
        if len(parts) != 2:
            raise ValueError(f"Invalid --min-assets-special token: {token}")
        out[int(parts[0])] = int(parts[1])
    return out


def _check_bounds(pptx_path: Path, margin_in: float) -> list[str]:
    prs = Presentation(str(pptx_path))
    margin = int(margin_in * EMU_PER_INCH)
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide):
            left = int(getattr(shape, "left", 0))
            top = int(getattr(shape, "top", 0))
            width = int(getattr(shape, "width", 0))
            height = int(getattr(shape, "height", 0))
            if width <= 0 or height <= 0:
                continue
            right = left + width
            bottom = top + height
            if left < -margin or top < -margin or right > slide_w + margin or bottom > slide_h + margin:
                issues.append(
                    f"slide {slide_num:02d}: shape_id={shape.shape_id} "
                    f"bbox=({left},{top})-({right},{bottom}) slide=({slide_w},{slide_h})"
                )
    return issues


def _extract_words(text: str) -> list[str]:
    return re.findall(r"[A-Za-z0-9_./-]+", text)


def _looks_like_filename(token: str) -> bool:
    return bool(re.search(r"\.(csv|gz|db|json|txt|md|log|zip|png|jpg)$", token))


def _collect_slide_words(slide) -> int:
    words: list[str] = []
    for shape in _iter_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            tokens = _extract_words(paragraph.text)
            for token in tokens:
                if _looks_like_filename(token):
                    continue
                words.append(token)
    return len(words)


def _check_text_and_words(prs: Presentation, min_body_font_pt: float, max_words_per_slide: int) -> tuple[list[str], list[str]]:
    font_issues: list[str] = []
    word_issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        word_count = _collect_slide_words(slide)
        if word_count > max_words_per_slide:
            word_issues.append(
                f"slide {slide_num:02d}: words={word_count} > max_words_per_slide={max_words_per_slide}"
            )

        for shape in _iter_shapes(slide):
            if not getattr(shape, "has_text_frame", False):
                continue
            for para_idx, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                for run_idx, run in enumerate(paragraph.runs, start=1):
                    text = run.text.strip()
                    if not text:
                        continue
                    if run.font.size is None:
                        font_issues.append(
                            f"slide {slide_num:02d}: shape_id={shape.shape_id} paragraph={para_idx} run={run_idx} "
                            "font size is unset"
                        )
                        continue
                    size_pt = run.font.size.pt
                    if size_pt < min_body_font_pt:
                        font_issues.append(
                            f"slide {slide_num:02d}: shape_id={shape.shape_id} paragraph={para_idx} run={run_idx} "
                            f"font={size_pt:.1f}pt < {min_body_font_pt:.1f}pt"
                        )
    return font_issues, word_issues


def _is_under(child: Path, parent: Path) -> bool:
    child_s = str(child.resolve())
    parent_s = str(parent.resolve())
    if child_s == parent_s:
        return True
    return child_s.startswith(parent_s + "/")


def _parse_asset_manifest_lines(lines: Iterable[str]) -> tuple[dict[int, int], dict[int, set[str]], list[tuple[int, str, Path, Path]]]:
    slide_counts: dict[int, int] = {}
    slide_roles: dict[int, set[str]] = {}
    entries: list[tuple[int, str, Path, Path]] = []
    current_slide: int | None = None
    for raw_line in lines:
        line = raw_line.strip()
        if not line:
            continue
        if line.startswith("Slide "):
            match = re.match(r"Slide\s+(\d+)\s+\|", line)
            if not match:
                raise ValueError(f"Malformed slide header in manifest: {line}")
            current_slide = int(match.group(1))
            slide_counts.setdefault(current_slide, 0)
            slide_roles.setdefault(current_slide, set())
            continue
        if line.startswith("ASSET | "):
            if current_slide is None:
                raise ValueError("Manifest ASSET line appears before a Slide header")
            parts = [p.strip() for p in line.split("|")]
            if len(parts) < 7:
                raise ValueError(f"Malformed ASSET line in manifest: {line}")
            role = parts[1]
            source_pptx = Path(parts[3])
            derived = Path(parts[5])
            slide_counts[current_slide] = slide_counts.get(current_slide, 0) + 1
            slide_roles.setdefault(current_slide, set()).add(role)
            entries.append((current_slide, role, source_pptx, derived))
            continue
    return slide_counts, slide_roles, entries


def _check_assets(
    *,
    asset_manifest_path: Path,
    assert_dir: Path,
    allow_source_roots: tuple[Path, ...],
    expected_slides: int,
    min_assets_per_slide: int,
    min_assets_special: dict[int, int],
) -> list[str]:
    if not asset_manifest_path.exists():
        return [f"Missing asset manifest: {asset_manifest_path}"]

    lines = asset_manifest_path.read_text(encoding="utf-8").splitlines()
    try:
        slide_counts, slide_roles, entries = _parse_asset_manifest_lines(lines)
    except ValueError as err:
        return [str(err)]

    issues: list[str] = []
    assert_dir_resolved = assert_dir.resolve()
    allow_source_roots_resolved = tuple(root.resolve() for root in allow_source_roots)
    for slide_num, role, source_pptx, derived in entries:
        if role not in {"structural", "semantic"}:
            issues.append(f"slide {slide_num:02d}: invalid asset role={role!r}")
        if not _is_under(source_pptx, assert_dir_resolved):
            if not any(_is_under(source_pptx, root) for root in allow_source_roots_resolved):
                issues.append(f"slide {slide_num:02d}: source outside allowed roots: {source_pptx}")
        if not source_pptx.exists():
            issues.append(f"slide {slide_num:02d}: source pptx missing: {source_pptx}")
        if not derived.exists():
            issues.append(f"slide {slide_num:02d}: derived asset missing: {derived}")

    for slide_num in range(1, expected_slides + 1):
        count = slide_counts.get(slide_num, 0)
        if count < min_assets_per_slide:
            issues.append(
                f"slide {slide_num:02d}: assets={count} < min_assets_per_slide={min_assets_per_slide}"
            )
        roles = slide_roles.get(slide_num, set())
        if not {"structural", "semantic"}.issubset(roles):
            issues.append(f"slide {slide_num:02d}: must include both structural and semantic roles")

    for slide_num, required_count in min_assets_special.items():
        count = slide_counts.get(slide_num, 0)
        if count < required_count:
            issues.append(
                f"slide {slide_num:02d}: assets={count} < special minimum={required_count}"
            )
    return issues


def _iter_slide_xml_parts(z: zipfile.ZipFile) -> list[str]:
    parts = [
        name
        for name in z.namelist()
        if name.startswith("ppt/slides/slide") and name.endswith(".xml") and "/_rels/" not in name
    ]
    parts.sort(key=lambda s: int(Path(s).stem.replace("slide", "")))
    return parts


def _count_click_effects(xml: str) -> int:
    return len(re.findall(r'nodeType="clickEffect"', xml))


def _check_animations(*, pptx_path: Path, expected_slides: int, min_click_effects: int) -> tuple[list[str], int]:
    issues: list[str] = []
    total_clicks = 0
    missing_timing: list[str] = []

    with zipfile.ZipFile(pptx_path, "r") as z:
        slide_parts = _iter_slide_xml_parts(z)
        for part in slide_parts:
            xml = z.read(part).decode("utf-8", "ignore")
            total_clicks += _count_click_effects(xml)
            if "<p:timing" not in xml:
                missing_timing.append(part)

        if len(slide_parts) != expected_slides:
            issues.append(f"slide xml count mismatch: {len(slide_parts)} != expected_slides={expected_slides}")

    if missing_timing:
        for part in missing_timing:
            issues.append(f"missing <p:timing>: {part}")

    if total_clicks < min_click_effects:
        issues.append(f"total_click_effects {total_clicks} < min_click_effects {min_click_effects}")

    return issues, total_clicks


def _check_hidden(prs: Presentation, expected_hidden: set[int]) -> list[str]:
    issues: list[str] = []
    actual_hidden = {
        idx
        for idx, slide in enumerate(prs.slides, start=1)
        if slide._element.get("show") == "0"
    }
    if actual_hidden != expected_hidden:
        issues.append(f"hidden mismatch: expected={sorted(expected_hidden)} actual={sorted(actual_hidden)}")
    return issues


def _check_inventory_issues(
    *,
    pptx_path: Path,
    inventory_script: Path,
) -> list[str]:
    if not inventory_script.exists():
        return [f"Missing inventory script for overflow/overlap check: {inventory_script}"]

    with tempfile.TemporaryDirectory(prefix="qc_inventory_") as tmp_dir:
        out_json = Path(tmp_dir) / "issues.json"
        cmd = [
            sys.executable,
            str(inventory_script),
            str(pptx_path),
            str(out_json),
            "--issues-only",
        ]
        proc = subprocess.run(cmd, capture_output=True, text=True)
        if proc.returncode != 0:
            stderr = proc.stderr.strip()
            stdout = proc.stdout.strip()
            return [f"inventory.py failed: {stderr or stdout or 'unknown error'}"]
        payload = json.loads(out_json.read_text(encoding="utf-8"))

    issues: list[str] = []
    for slide_key, shapes in payload.items():
        for shape_key, shape_data in shapes.items():
            if not isinstance(shape_data, dict):
                continue
            if "overflow" in shape_data:
                issues.append(f"{slide_key} {shape_key}: overflow issue")
            if "overlap" in shape_data:
                issues.append(f"{slide_key} {shape_key}: overlap issue")
    return issues


def _parse_csv_terms(value: str) -> list[str]:
    terms: list[str] = []
    for token in value.split(","):
        token = token.strip()
        if token:
            terms.append(token.lower())
    return terms


def _check_watermark_text(
    *,
    pptx_path: Path,
    forbidden_terms: list[str],
    dpi: int,
) -> list[str]:
    if not forbidden_terms:
        return []
    for exe in ("soffice", "pdftoppm", "tesseract"):
        if shutil.which(exe) is None:
            return [f"watermark scan missing dependency: {exe}"]

    with tempfile.TemporaryDirectory(prefix="qc_watermark_") as tmp_dir:
        tmp = Path(tmp_dir)
        # Convert PPTX to PDF.
        proc = subprocess.run(
            ["soffice", "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(pptx_path)],
            capture_output=True,
            text=True,
        )
        if proc.returncode != 0:
            stderr = (proc.stderr or proc.stdout or "").strip()
            return [f"watermark scan: soffice convert-to pdf failed: {stderr or 'unknown error'}"]

        pdfs = sorted(tmp.glob("*.pdf"))
        if not pdfs:
            return ["watermark scan: no PDF produced by soffice"]
        pdf_path = pdfs[0]

        # Render PDF to PNG pages.
        prefix = tmp / "page"
        proc = subprocess.run(
            ["pdftoppm", "-png", "-r", str(int(dpi)), str(pdf_path), str(prefix)],
            capture_output=True,
            text=True,
        )
        if proc.returncode != 0:
            stderr = (proc.stderr or proc.stdout or "").strip()
            return [f"watermark scan: pdftoppm failed: {stderr or 'unknown error'}"]

        def _page_num(path: Path) -> int:
            match = re.search(r"(\d+)", path.stem)
            return int(match.group(1)) if match else 0

        images = sorted(tmp.glob("page-*.png"), key=_page_num)
        if not images:
            return ["watermark scan: pdftoppm produced no images"]

        issues: list[str] = []
        for img in images:
            # OCR the slide. This is intentionally simple: we only look for obvious watermark strings.
            proc = subprocess.run(
                ["tesseract", str(img), "stdout", "--dpi", str(int(dpi))],
                capture_output=True,
                text=True,
            )
            if proc.returncode != 0:
                stderr = (proc.stderr or proc.stdout or "").strip()
                issues.append(f"watermark scan: tesseract failed on {img.name}: {stderr or 'unknown error'}")
                continue

            text = (proc.stdout or "").lower()
            for term in forbidden_terms:
                if term in text:
                    issues.append(f"watermark text detected on {img.name}: term={term!r}")
                    break

        return issues


def run_qc(
    *,
    pptx_path: Path,
    asset_manifest_path: Path,
    expected_slides: int,
    hidden_slides: set[int],
    assert_dir: Path,
    allow_source_roots: tuple[Path, ...],
    min_assets_per_slide: int,
    min_assets_special: dict[int, int],
    min_body_font_pt: float,
    max_words_per_slide: int,
    margin_in: float,
    inventory_script: Path,
    forbidden_watermark_terms: list[str],
    watermark_dpi: int,
    min_click_effects: int,
) -> QcResult:
    errors: list[str] = []
    prs = Presentation(str(pptx_path))

    slide_count_ok = len(prs.slides) == expected_slides
    if not slide_count_ok:
        errors.append(f"slide_count={len(prs.slides)} expected={expected_slides}")

    hidden_issues = _check_hidden(prs, hidden_slides)
    hidden_ok = not hidden_issues
    errors.extend(hidden_issues)

    bound_issues = _check_bounds(pptx_path, margin_in)
    bounds_ok = not bound_issues
    errors.extend(bound_issues)

    asset_issues = _check_assets(
        asset_manifest_path=asset_manifest_path,
        assert_dir=assert_dir,
        allow_source_roots=allow_source_roots,
        expected_slides=expected_slides,
        min_assets_per_slide=min_assets_per_slide,
        min_assets_special=min_assets_special,
    )
    asset_ok = not asset_issues
    errors.extend(asset_issues)

    font_issues, word_issues = _check_text_and_words(prs, min_body_font_pt, max_words_per_slide)
    font_ok = not font_issues
    words_ok = not word_issues
    errors.extend(font_issues)
    errors.extend(word_issues)

    inventory_issues = _check_inventory_issues(pptx_path=pptx_path, inventory_script=inventory_script)
    inventory_ok = not inventory_issues
    errors.extend(inventory_issues)

    watermark_issues = _check_watermark_text(
        pptx_path=pptx_path,
        forbidden_terms=forbidden_watermark_terms,
        dpi=watermark_dpi,
    )
    watermark_ok = not watermark_issues
    errors.extend(watermark_issues)

    anim_issues, total_clicks = _check_animations(
        pptx_path=pptx_path,
        expected_slides=expected_slides,
        min_click_effects=min_click_effects,
    )
    timing_ok = not any("missing <p:timing>" in issue for issue in anim_issues)
    clicks_ok = not any("total_click_effects" in issue for issue in anim_issues)
    errors.extend(anim_issues)

    return QcResult(
        slide_count_ok=slide_count_ok,
        hidden_ok=hidden_ok,
        bounds_ok=bounds_ok,
        asset_ok=asset_ok,
        font_ok=font_ok,
        words_ok=words_ok,
        inventory_ok=inventory_ok,
        watermark_ok=watermark_ok,
        timing_ok=timing_ok,
        clicks_ok=clicks_ok,
        total_click_effects=total_clicks,
        errors=errors,
    )


def _build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="QC for the professor-ready Bluesky deck.")
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--asset-manifest", type=Path, required=True)
    parser.add_argument("--expected-slides", type=int, default=11)
    parser.add_argument("--hidden-slides", type=str, default="10,11")
    parser.add_argument("--assert-dir", type=Path, required=True)
    parser.add_argument(
        "--allow-source-roots",
        type=str,
        default=str(Path(__file__).resolve().parent / "prof_build"),
        help="Comma-separated extra roots allowed for ASSET source paths (e.g., captured screenshots).",
    )
    parser.add_argument("--min-assets-per-slide", type=int, default=2)
    parser.add_argument("--min-assets-special", type=str, default="1:3,7:3")
    parser.add_argument("--min-body-font-pt", type=float, default=20.0)
    parser.add_argument("--max-words-per-slide", type=int, default=45)
    parser.add_argument("--margin-in", type=float, default=0.05)
    parser.add_argument("--min-click-effects", type=int, default=140)
    parser.add_argument(
        "--forbidden-watermark-terms",
        type=str,
        default="slidesgo,codex",
        help="Comma-separated list of substrings to forbid in slide OCR (rendered from PPTX).",
    )
    parser.add_argument("--watermark-dpi", type=int, default=200)
    parser.add_argument(
        "--inventory-script",
        type=Path,
        default=Path("/Users/yipengandrewwang/.codex/skills/pptx/scripts/inventory.py"),
    )
    return parser


def main(argv: list[str] | None = None) -> int:
    args = _build_parser().parse_args(argv if argv is not None else sys.argv[1:])
    if not args.pptx.exists():
        print(f"Missing pptx: {args.pptx}", file=sys.stderr)
        return 2

    hidden_slides = _parse_hidden_slides(args.hidden_slides)
    min_assets_special = _parse_min_assets_special(args.min_assets_special)
    forbidden_terms = _parse_csv_terms(args.forbidden_watermark_terms)
    allow_source_roots = tuple(Path(p.strip()).resolve() for p in args.allow_source_roots.split(",") if p.strip())

    result = run_qc(
        pptx_path=args.pptx,
        asset_manifest_path=args.asset_manifest,
        expected_slides=int(args.expected_slides),
        hidden_slides=hidden_slides,
        assert_dir=args.assert_dir,
        allow_source_roots=allow_source_roots,
        min_assets_per_slide=int(args.min_assets_per_slide),
        min_assets_special=min_assets_special,
        min_body_font_pt=float(args.min_body_font_pt),
        max_words_per_slide=int(args.max_words_per_slide),
        margin_in=float(args.margin_in),
        inventory_script=args.inventory_script,
        forbidden_watermark_terms=forbidden_terms,
        watermark_dpi=int(args.watermark_dpi),
        min_click_effects=int(args.min_click_effects),
    )

    if result.errors:
        print("FAIL: QC checks did not pass", file=sys.stderr)
        for issue in result.errors:
            print(f"- {issue}", file=sys.stderr)
        return 1

    print(
        "OK:",
        f"slides={args.expected_slides}",
        f"hidden={sorted(hidden_slides)}",
        "bounds_ok",
        "assets_ok",
        "fonts_ok",
        "words_ok",
        "overflow_overlap_ok",
        "watermark_ok",
        "timing_ok",
        f"clicks_ok(total={result.total_click_effects})",
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
