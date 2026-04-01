#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import gzip
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterator, Sequence

from lxml import etree
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

# Reuse the receipt-card rendering style where helpful.
from build_run_receipts import COLORS as RECEIPT_COLORS  # noqa: E402
from build_run_receipts import FONT_MONO, FONT_SANS  # noqa: E402
from build_run_receipts import render_list_card  # noqa: E402


EMU_PER_INCH = 914400

PPTX_MIME_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

NSMAP = {"p": PML_NS, "r": R_NS}
REL_NSMAP = {"rel": REL_NS}
CT_NSMAP = {"ct": CT_NS}


@dataclass(frozen=True)
class BboxPx:
    x0: int
    y0: int
    x1: int
    y1: int

    @property
    def w(self) -> int:
        return self.x1 - self.x0

    @property
    def h(self) -> int:
        return self.y1 - self.y0


@dataclass(frozen=True)
class BboxIn:
    x: float
    y: float
    w: float
    h: float


def _emu(inches: float) -> int:
    return int(round(inches * EMU_PER_INCH))


def _emu_to_in(emu: int) -> float:
    return emu / EMU_PER_INCH


def _read_zip_xml(z: zipfile.ZipFile, name: str) -> etree._Element:
    return etree.fromstring(z.read(name))


def _write_zip_xml(z: zipfile.ZipFile, name: str, root: etree._Element) -> None:
    z.writestr(name, etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True))


def _max_slide_number(z: zipfile.ZipFile) -> int:
    max_n = 0
    for name in z.namelist():
        if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
            continue
        stem = Path(name).stem
        try:
            n = int(stem.replace("slide", ""))
        except ValueError:
            continue
        max_n = max(max_n, n)
    return max_n


def _next_rid(pres_rels_root: etree._Element) -> str:
    rids: list[int] = []
    for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                rids.append(int(rid[3:]))
            except ValueError:
                continue
    nxt = (max(rids) + 1) if rids else 1
    return f"rId{nxt}"


def _next_slide_id(pres_root: etree._Element) -> int:
    sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
    if sld_id_lst is None:
        raise RuntimeError("presentation.xml missing p:sldIdLst")
    ids: list[int] = []
    for sld_id in sld_id_lst.findall("p:sldId", namespaces=NSMAP):
        try:
            ids.append(int(sld_id.get("id", "0")))
        except ValueError:
            continue
    return (max(ids) + 1) if ids else 256


def duplicate_slide_parts(
    *,
    pptx_in: Path,
    pptx_out: Path,
    base_slide_num: int,
    copies: int,
    insert_after_slide_num: int,
) -> list[int]:
    """
    Duplicate a slide part (slide{base_slide_num}.xml + rels) N times.

    Returns the new slide part numbers (e.g., [14, 15]).
    """
    if copies <= 0:
        raise ValueError("copies must be > 0")

    with zipfile.ZipFile(pptx_in, "r") as zin:
        pres_root = _read_zip_xml(zin, "ppt/presentation.xml")
        pres_rels_root = _read_zip_xml(zin, "ppt/_rels/presentation.xml.rels")
        ct_root = _read_zip_xml(zin, "[Content_Types].xml")

        sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
        if sld_id_lst is None:
            raise RuntimeError("presentation.xml missing p:sldIdLst")

        max_slide = _max_slide_number(zin)
        next_slide_id = _next_slide_id(pres_root)

        slide_src = f"ppt/slides/slide{base_slide_num}.xml"
        slide_rels_src = f"ppt/slides/_rels/slide{base_slide_num}.xml.rels"
        slide_src_bytes = zin.read(slide_src)
        slide_rels_src_bytes = zin.read(slide_rels_src)

        new_nums: list[int] = []
        for _ in range(copies):
            n = max_slide + 1
            max_slide = n
            new_nums.append(n)

            rid = _next_rid(pres_rels_root)
            rel_el = etree.Element(f"{{{REL_NS}}}Relationship")
            rel_el.set("Id", rid)
            rel_el.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel_el.set("Target", f"slides/slide{n}.xml")
            pres_rels_root.append(rel_el)

            sld_id_el = etree.Element(f"{{{PML_NS}}}sldId")
            sld_id_el.set("id", str(next_slide_id))
            sld_id_el.set(f"{{{R_NS}}}id", rid)
            sld_id_lst.append(sld_id_el)
            next_slide_id += 1

            override = etree.Element(f"{{{CT_NS}}}Override")
            override.set("PartName", f"/ppt/slides/slide{n}.xml")
            override.set("ContentType", PPTX_MIME_SLIDE)
            ct_root.append(override)

        # Reorder slides to insert new slides right after insert_after_slide_num.
        rid_to_target: dict[str, str] = {}
        for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
            if rel.get("Type") != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
                continue
            rid = rel.get("Id")
            tgt = rel.get("Target")
            if rid and tgt:
                rid_to_target[rid] = tgt

        target_to_sld_id: dict[str, etree._Element] = {}
        targets_in_order: list[str] = []
        for sld_id in list(sld_id_lst.findall("p:sldId", namespaces=NSMAP)):
            rid = sld_id.get(f"{{{R_NS}}}id")
            tgt = rid_to_target.get(rid or "")
            if not tgt:
                continue
            target_to_sld_id[tgt] = sld_id
            targets_in_order.append(tgt)

        new_targets = [f"slides/slide{n}.xml" for n in new_nums]
        base_target = f"slides/slide{insert_after_slide_num}.xml"
        if base_target not in targets_in_order:
            raise RuntimeError(f"Could not find slide target in order: {base_target}")

        # Remove new targets from current order (they are appended) and re-insert after base_target.
        targets_wo_new = [t for t in targets_in_order if t not in set(new_targets)]
        insert_at = targets_wo_new.index(base_target) + 1
        final_targets = targets_wo_new[:insert_at] + new_targets + targets_wo_new[insert_at:]

        for child in list(sld_id_lst):
            sld_id_lst.remove(child)
        for tgt in final_targets:
            el = target_to_sld_id.get(tgt)
            if el is None:
                raise RuntimeError(f"Missing sldId element for target: {tgt}")
            sld_id_lst.append(el)

        tmp_path = pptx_out.with_suffix(".tmp.pptx")
        if tmp_path.exists():
            tmp_path.unlink()
        pptx_out.parent.mkdir(parents=True, exist_ok=True)

        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}:
                    continue
                zout.writestr(item, zin.read(item.filename))

            # Write new slide parts.
            for n in new_nums:
                zout.writestr(f"ppt/slides/slide{n}.xml", slide_src_bytes)
                zout.writestr(f"ppt/slides/_rels/slide{n}.xml.rels", slide_rels_src_bytes)

            _write_zip_xml(zout, "ppt/presentation.xml", pres_root)
            _write_zip_xml(zout, "ppt/_rels/presentation.xml.rels", pres_rels_root)
            _write_zip_xml(zout, "[Content_Types].xml", ct_root)

        tmp_path.replace(pptx_out)

    return new_nums


def _compact_did(did: str) -> str:
    did = did.strip()
    if not did.startswith("did:"):
        return did[:24]
    if len(did) <= 24:
        return did
    return f"{did[:12]}...{did[-6:]}"


def _iter_gz_csv_rows(path: Path) -> Iterator[dict[str, str]]:
    with gzip.open(path, "rt", encoding="utf-8") as f:
        yield from csv.DictReader(f)


def pick_label_example_lines(*, labels_gz: Path, wanted_vals: Sequence[str], limit: int) -> list[str]:
    wanted = list(dict.fromkeys(wanted_vals))  # stable unique
    chosen: dict[str, tuple[str, str]] = {}
    for row in _iter_gz_csv_rows(labels_gz):
        val = (row.get("label_val") or "").strip()
        if val not in wanted:
            continue
        if val in chosen:
            continue
        src = (row.get("label_src") or "").strip()
        chosen[val] = (src, val)
        if len(chosen) >= min(limit, len(wanted)):
            break

    lines: list[str] = []
    for val in wanted:
        if val not in chosen:
            continue
        src, vv = chosen[val]
        lines.append(f"- {_compact_did(src)}:{vv}")
        if len(lines) >= limit:
            break
    return lines


def build_rich_risk_labels_stack(
    *,
    out_path: Path,
    top_card: Path,
    mid_card: Path,
    labels_gz: Path,
    wanted_vals: Sequence[str],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    # Render a richer labels excerpt (keep card size consistent with existing 1400x420 assets).
    label_lines = pick_label_example_lines(labels_gz=labels_gz, wanted_vals=wanted_vals, limit=11)
    if not label_lines:
        label_lines = ["(no labels found)"]

    rich_labels_card = out_path.with_name("post_labels_row_rich.png")
    render_list_card(
        out_path=rich_labels_card,
        title="post_labels.csv.gz (labels excerpt; Feb 1 run)",
        accent_hex=RECEIPT_COLORS["purple"],
        lines=label_lines,
        height_px=420,
        max_lines=11,
    )

    a = Image.open(top_card).convert("RGBA")
    b = Image.open(mid_card).convert("RGBA")
    c = Image.open(rich_labels_card).convert("RGBA")

    gap = 18
    width = max(a.width, b.width, c.width)
    if not (a.width == b.width == c.width == width):
        raise RuntimeError("Card widths must match to stack cleanly.")

    height = a.height + gap + b.height + gap + c.height
    stack = Image.new("RGBA", (width, height), (0, 0, 0, 0))
    y = 0
    stack.paste(a, (0, y), a)
    y += a.height + gap
    stack.paste(b, (0, y), b)
    y += b.height + gap
    stack.paste(c, (0, y), c)

    stack.save(out_path)
    return out_path


def _render_screenshot_card(
    *,
    out_path: Path,
    title: str,
    src_img: Path,
    crop_box: tuple[int, int, int, int],
    accent_hex: str,
    width_px: int = 1600,
    height_px: int = 288,
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    img = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))

    pad = 34
    radius = 32
    border = 4

    # Shadow
    shadow = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle(
        (12, 14, width_px - 4, height_px - 4),
        radius=radius,
        fill="#00000022",
        outline=None,
    )
    img = Image.alpha_composite(img, shadow)

    draw = ImageDraw.Draw(img)
    draw.rounded_rectangle(
        (0, 0, width_px - 1, height_px - 1),
        radius=radius,
        fill="#FFFFFF",
        outline=f"#{accent_hex}",
        width=border,
    )

    title_font = ImageFont.truetype(str(FONT_SANS), 26)
    draw.text((pad, 18), title, fill="#0B0B0F", font=title_font)

    # Content region
    content_x0 = pad
    content_y0 = 62
    content_x1 = width_px - pad
    content_y1 = height_px - 18
    content_w = content_x1 - content_x0
    content_h = content_y1 - content_y0

    src = Image.open(src_img).convert("RGB")
    crop = src.crop(crop_box)
    # Fit preserves aspect ratio with center crop if necessary.
    crop_fit = ImageOps.fit(crop, (content_w, content_h), method=Image.LANCZOS)
    img.paste(crop_fit, (content_x0, content_y0))

    img.save(out_path)
    return out_path


def _render_list_card_small(
    *,
    out_path: Path,
    title: str,
    lines: Sequence[str],
    accent_hex: str,
    width_px: int = 1600,
    height_px: int = 288,
    max_lines: int = 7,
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    img = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))

    pad = 34
    radius = 32
    border = 4

    shadow = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle(
        (12, 14, width_px - 4, height_px - 4),
        radius=radius,
        fill="#00000022",
        outline=None,
    )
    img = Image.alpha_composite(img, shadow)
    draw = ImageDraw.Draw(img)

    draw.rounded_rectangle(
        (0, 0, width_px - 1, height_px - 1),
        radius=radius,
        fill="#FFFFFF",
        outline=f"#{accent_hex}",
        width=border,
    )

    title_font = ImageFont.truetype(str(FONT_SANS), 26)
    mono = ImageFont.truetype(str(FONT_MONO), 20)
    draw.text((pad, 18), title, fill="#0B0B0F", font=title_font)

    y = 76
    line_h = 26
    for line in lines[:max_lines]:
        draw.text((pad, y), line, fill="#0B0B0F", font=mono)
        y += line_h
        if y > height_px - 30:
            break

    img.save(out_path)
    return out_path


def _compose_three_card_stack(
    *,
    out_path: Path,
    card_a: Path,
    card_b: Path,
    card_c: Path,
    gap_px: int = 18,
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)

    a = Image.open(card_a).convert("RGBA")
    b = Image.open(card_b).convert("RGBA")
    c = Image.open(card_c).convert("RGBA")

    width = max(a.width, b.width, c.width)
    if not (a.width == b.width == c.width == width):
        raise RuntimeError("Card widths must match.")
    if not (a.height == b.height == c.height):
        raise RuntimeError("Card heights must match for fixed 16:9 stacking.")

    height = a.height + gap_px + b.height + gap_px + c.height
    canvas = Image.new("RGBA", (width, height), (0, 0, 0, 0))

    y = 0
    canvas.paste(a, (0, y), a)
    y += a.height + gap_px
    canvas.paste(b, (0, y), b)
    y += b.height + gap_px
    canvas.paste(c, (0, y), c)

    canvas.save(out_path)
    return out_path


def build_labels_docs_stack(*, out_path: Path, webshots_dir: Path) -> Path:
    tmp_dir = out_path.parent
    tmp_dir.mkdir(parents=True, exist_ok=True)

    card_h = 288

    a = _render_screenshot_card(
        out_path=tmp_dir / "labels_docs_card_a.png",
        title="ATProto spec: label object fields",
        src_img=webshots_dir / "atproto_label_slice1.png",
        # Focus on the field list section (keep a wide, short slice for legibility).
        crop_box=(260, 880, 1570, 1060),
        accent_hex=RECEIPT_COLORS["cyan"],
        height_px=card_h,
    )
    b = _render_screenshot_card(
        out_path=tmp_dir / "labels_docs_card_b.png",
        title="Bluesky docs: global label values (excerpt)",
        src_img=webshots_dir / "moderation_slice1.png",
        # Global label values list (include a few of the most commonly referenced values).
        crop_box=(240, 1500, 1570, 1800),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )
    c = _render_screenshot_card(
        out_path=tmp_dir / "labels_docs_card_c.png",
        title="Bluesky docs: viewer preferences (hide/warn/ignore)",
        src_img=webshots_dir / "moderation_slice4.png",
        # Center on the label preference mapping in the config example.
        crop_box=(240, 980, 1570, 1160),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )

    return _compose_three_card_stack(out_path=out_path, card_a=a, card_b=b, card_c=c, gap_px=18)


def build_labels_labelers_stack(*, out_path: Path, webshots_dir: Path) -> Path:
    """
    Three-card stack grounding *where labels come from*:
      A) labeler declaration record
      B) labeler subscriptions (viewer-selected)
      C) self-labels (built-in values)
    """

    tmp_dir = out_path.parent
    tmp_dir.mkdir(parents=True, exist_ok=True)

    card_h = 288

    a = _render_screenshot_card(
        out_path=tmp_dir / "labels_labelers_card_a.png",
        title="Bluesky docs: labeler declaration (record excerpt)",
        src_img=webshots_dir / "moderation_slice3.png",
        # Center on the JSON example showing label values and policies.
        crop_box=(240, 520, 1570, 700),
        accent_hex=RECEIPT_COLORS["cyan"],
        height_px=card_h,
    )
    b = _render_screenshot_card(
        out_path=tmp_dir / "labels_labelers_card_b.png",
        title="Bluesky docs: labeler subscriptions (viewer-selected)",
        src_img=webshots_dir / "moderation_slice3.png",
        # Subscription mechanics: atproto-accept-labelers header + explanation.
        crop_box=(240, 1150, 1570, 1500),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )
    c = _render_screenshot_card(
        out_path=tmp_dir / "labels_labelers_card_c.png",
        title="Bluesky docs: self-labels (supported values)",
        src_img=webshots_dir / "moderation_slice3.png",
        # Bullet list of supported self-label values (!no-unauthenticated/porn/sexual).
        crop_box=(240, 1620, 1570, 1800),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )

    return _compose_three_card_stack(out_path=out_path, card_a=a, card_b=b, card_c=c, gap_px=18)


def build_labels_moderation_stack(*, out_path: Path, webshots_dir: Path) -> Path:
    """
    Three-card stack grounding *how labels become UI actions*:
      A) moderatePost() input
      B) output flags (filter/blur/alert/inform)
      C) UI contexts (contentList/contentView/etc)
    """

    tmp_dir = out_path.parent
    tmp_dir.mkdir(parents=True, exist_ok=True)

    card_h = 288

    a = _render_screenshot_card(
        out_path=tmp_dir / "labels_moderation_card_a.png",
        title="Bluesky docs: moderatePost() (input → decision)",
        src_img=webshots_dir / "moderation_slice6.png",
        # Focus on the moderatePost call + signature context.
        crop_box=(240, 160, 1570, 340),
        accent_hex=RECEIPT_COLORS["cyan"],
        height_px=card_h,
    )
    b = _render_screenshot_card(
        out_path=tmp_dir / "labels_moderation_card_b.png",
        title="Bluesky docs: decision flags (filter / blur / alert / inform)",
        src_img=webshots_dir / "moderation_slice6.png",
        # Output structure: filter/blur/alert/inform/noOverride.
        crop_box=(240, 340, 1570, 520),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )
    c = _render_screenshot_card(
        out_path=tmp_dir / "labels_moderation_card_c.png",
        title="Bluesky docs: UI contexts (contentList vs contentView vs media)",
        src_img=webshots_dir / "moderation_slice6.png",
        # List of contexts where the same label may render differently.
        crop_box=(240, 960, 1570, 1140),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )

    return _compose_three_card_stack(out_path=out_path, card_a=a, card_b=b, card_c=c, gap_px=18)


def build_labels_custom_stack(*, out_path: Path, webshots_dir: Path, labels_gz: Path) -> Path:
    tmp_dir = out_path.parent
    tmp_dir.mkdir(parents=True, exist_ok=True)

    card_h = 288

    a = _render_screenshot_card(
        out_path=tmp_dir / "labels_custom_card_a.png",
        title="Bluesky docs: custom label values (table excerpt)",
        src_img=webshots_dir / "moderation_slice2.png",
        # Show table header + a few rows for blur/severity semantics.
        crop_box=(240, 340, 1570, 520),
        accent_hex=RECEIPT_COLORS["cyan"],
        height_px=card_h,
    )
    b = _render_screenshot_card(
        out_path=tmp_dir / "labels_custom_card_b.png",
        title="Bluesky docs: label targets (account / profile / content)",
        src_img=webshots_dir / "moderation_slice2.png",
        # Label targets bullet list.
        crop_box=(240, 1200, 1570, 1380),
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
    )

    # Example custom values observed in the Feb 1 run (avoid implying prevalence; just show diversity).
    custom_vals: list[str] = []
    for row in _iter_gz_csv_rows(labels_gz):
        val = (row.get("label_val") or "").strip()
        if not val:
            continue
        if ":" not in val and val.lower() == val:
            continue
        if val not in custom_vals:
            custom_vals.append(val)
        if len(custom_vals) >= 7:
            break

    if not custom_vals:
        custom_vals = ["(no custom values found)"]

    custom_lines = [f"- {v}" for v in custom_vals]
    c = _render_list_card_small(
        out_path=tmp_dir / "labels_custom_card_c.png",
        title="Feb 1 run: custom label_val strings (examples)",
        lines=custom_lines,
        accent_hex=RECEIPT_COLORS["amber"],
        height_px=card_h,
        max_lines=7,
    )

    return _compose_three_card_stack(out_path=out_path, card_a=a, card_b=b, card_c=c, gap_px=18)


def _picture_shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh
    raise KeyError(f"picture shape_id={shape_id} not found")


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found")


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001
    blip.set(qn("r:embed"), rid)


def _open_image_from_blob(blob: bytes) -> Image.Image:
    # PIL can open from a file-like object, but Image.open accepts a BytesIO.
    import io

    return Image.open(io.BytesIO(blob))


def _compute_card_bboxes_px_from_picture(pic_shape) -> list[BboxPx]:
    img = _open_image_from_blob(pic_shape.image.blob).convert("RGBA")
    alpha = img.getchannel("A")
    w, h = alpha.size

    px = alpha.load()
    row_has = [False] * h
    for y in range(h):
        for x in range(w):
            if px[x, y] != 0:
                row_has[y] = True
                break

    segments: list[tuple[int, int]] = []
    y = 0
    while y < h:
        if not row_has[y]:
            y += 1
            continue
        y0 = y
        while y < h and row_has[y]:
            y += 1
        y1 = y
        segments.append((y0, y1))

    bboxes: list[BboxPx] = []
    for seg_y0, seg_y1 in segments:
        bbox = alpha.crop((0, seg_y0, w, seg_y1)).getbbox()
        if bbox is None:
            continue
        x0, y0, x1, y1 = bbox
        bboxes.append(BboxPx(x0=x0, y0=seg_y0 + y0, x1=x1, y1=seg_y0 + y1))

    return bboxes


def _map_bbox_to_slide_in(*, bbox: BboxPx, pic_shape) -> BboxIn:
    img_w, img_h = pic_shape.image.size
    pic_left = _emu_to_in(int(pic_shape.left))
    pic_top = _emu_to_in(int(pic_shape.top))
    pic_w = _emu_to_in(int(pic_shape.width))
    pic_h = _emu_to_in(int(pic_shape.height))

    sx = pic_w / img_w
    sy = pic_h / img_h
    x = pic_left + bbox.x0 * sx
    y = pic_top + bbox.y0 * sy
    w = bbox.w * sx
    h = bbox.h * sy
    return BboxIn(x=x, y=y, w=w, h=h)


def align_callouts_to_cards(
    *,
    slide,
    picture_shape_id: int,
    rect_shape_ids: Sequence[int],
    oval_shape_ids: Sequence[int],
    inset_x_in: float,
    inset_y_in: float,
    oval_overlap_in: float,
    oval_top_offset_in: float,
) -> None:
    if len(rect_shape_ids) != len(oval_shape_ids):
        raise ValueError("rect_shape_ids and oval_shape_ids must have same length")

    pic = _picture_shape(slide, picture_shape_id)
    bboxes_px = _compute_card_bboxes_px_from_picture(pic)
    if len(bboxes_px) < len(rect_shape_ids):
        raise RuntimeError(f"expected >= {len(rect_shape_ids)} card bboxes, got {len(bboxes_px)}")

    # Keep only the first N bboxes (top->bottom) for A/B/C.
    for i, (rect_id, oval_id) in enumerate(zip(rect_shape_ids, oval_shape_ids, strict=True)):
        bbox_in = _map_bbox_to_slide_in(bbox=bboxes_px[i], pic_shape=pic)
        x = bbox_in.x + inset_x_in
        y = bbox_in.y + inset_y_in
        w = bbox_in.w - 2 * inset_x_in
        h = bbox_in.h - 2 * inset_y_in
        if w <= 0 or h <= 0:
            raise RuntimeError(f"invalid bbox after inset for card {i}: {bbox_in}")

        rect = _shape(slide, rect_id)
        rect.left = _emu(x)
        rect.top = _emu(y)
        rect.width = _emu(w)
        rect.height = _emu(h)

        oval = _shape(slide, oval_id)
        oval_w = _emu_to_in(int(oval.width))
        oval_h = _emu_to_in(int(oval.height))
        oval.left = _emu(x - oval_w + oval_overlap_in)
        oval.top = _emu(y + oval_top_offset_in)
        oval.width = _emu(oval_w)
        oval.height = _emu(oval_h)


def align_ovals_to_rectangles(
    *,
    slide,
    rect_shape_ids: Sequence[int],
    oval_shape_ids: Sequence[int],
    oval_overlap_in: float,
    oval_top_offset_in: float,
) -> None:
    if len(rect_shape_ids) != len(oval_shape_ids):
        raise ValueError("rect_shape_ids and oval_shape_ids must have same length")

    for rect_id, oval_id in zip(rect_shape_ids, oval_shape_ids, strict=True):
        rect = _shape(slide, rect_id)
        oval = _shape(slide, oval_id)

        rect_left = _emu_to_in(int(rect.left))
        rect_top = _emu_to_in(int(rect.top))
        oval_w = _emu_to_in(int(oval.width))
        oval_h = _emu_to_in(int(oval.height))

        oval.left = _emu(rect_left - oval_w + oval_overlap_in)
        oval.top = _emu(rect_top + oval_top_offset_in)
        oval.width = _emu(oval_w)
        oval.height = _emu(oval_h)


def _set_shape_text(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise TypeError("shape has no text frame")
    tf = shape.text_frame
    # Preserve any existing run properties; replace paragraph text only.
    if not tf.paragraphs:
        p = tf.add_paragraph()
        p.text = text
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for run in p0.runs[1:]:
            run.text = ""
    else:
        p0.text = text
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def update_more_examples_deck(*, pptx_in: Path, pptx_out: Path, keep_staged: bool = False) -> None:
    staged_dir = pptx_out.parent / "_staged"
    staged_dir.mkdir(parents=True, exist_ok=True)

    work1 = staged_dir / f"{pptx_out.stem}.staged.1.pptx"
    # Insert 2 new slides after the ranked exposure receipt slide by duplicating its slide part.
    # Note: in the source deck, the receipt slide is slide part 13 (slide13.xml) even though it is
    # slide order index 9.
    duplicate_slide_parts(
        pptx_in=pptx_in,
        pptx_out=work1,
        base_slide_num=13,
        copies=2,
        insert_after_slide_num=13,
    )

    work2 = staged_dir / f"{pptx_out.stem}.staged.2.pptx"
    # Insert 2 more new slides after the 2nd duplicate (slide part 15).
    duplicate_slide_parts(
        pptx_in=work1,
        pptx_out=work2,
        base_slide_num=13,
        copies=2,
        insert_after_slide_num=15,
    )

    prs = Presentation(str(work2))

    # --- Asset generation (images will be embedded into the PPTX) ---
    extra_dir = Path("/Users/yipengandrewwang/BlueSky/Slide2/prof_build/extra_assets")
    receipts_dir = extra_dir / "receipts"
    webshots_dir = extra_dir / "webshots"

    feb1_labels_gz = Path(
        "/Users/yipengandrewwang/BlueSky/bsky_fair_collect/out/bsky_fair_run_20260201T090454Z/csv/post_labels.csv.gz"
    )

    # Richer labels stack for slide 9 (keep the same “receipt stack” format, but with more label examples).
    risk_stack_rich = build_rich_risk_labels_stack(
        out_path=receipts_dir / "receipts_stack_risk_labels_rich.png",
        top_card=Path(
            "/Users/yipengandrewwang/BlueSky/Slide2/data_run_20260206/06_figures_preview/receipts/impression_row.png"
        ),
        mid_card=Path(
            "/Users/yipengandrewwang/BlueSky/Slide2/data_run_20260206/06_figures_preview/receipts/post_row.png"
        ),
        labels_gz=feb1_labels_gz,
        wanted_vals=(
            "porn",
            "sexual",
            "nudity",
            "graphic-media",
            "self-harm",
            "intolerant",
            "rude",
            "!no-unauthenticated",
            "comment_count:0",
            "reaction_count:5",
            "Adult Content",
        ),
    )

    labels_docs_stack = build_labels_docs_stack(out_path=receipts_dir / "labels_docs_stack_16x9.png", webshots_dir=webshots_dir)
    labels_custom_stack = build_labels_custom_stack(
        out_path=receipts_dir / "labels_custom_stack_16x9.png",
        webshots_dir=webshots_dir,
        labels_gz=feb1_labels_gz,
    )

    labels_labelers_stack = build_labels_labelers_stack(
        out_path=receipts_dir / "labels_labelers_stack_16x9.png",
        webshots_dir=webshots_dir,
    )
    labels_moderation_stack = build_labels_moderation_stack(
        out_path=receipts_dir / "labels_moderation_stack_16x9.png",
        webshots_dir=webshots_dir,
    )

    # Slides after insertion:
    #  8 = discovery receipts
    #  9 = ranked exposure + risk labels
    # 10 = labels docs
    # 11 = labels custom + dataset
    # 12 = labelers: declarations/subscriptions/self-labels
    # 13 = moderation: decision flags + contexts
    slide3 = prs.slides[2]
    slide4 = prs.slides[3]
    slide5 = prs.slides[4]
    slide8 = prs.slides[7]
    slide9 = prs.slides[8]
    slide10 = prs.slides[9]
    slide11 = prs.slides[10]
    slide12 = prs.slides[11]
    slide13 = prs.slides[12]

    # Fix circle markers on the UI screenshot example slides by snapping ovals to their rectangles.
    for slide in (slide3, slide4, slide5):
        align_ovals_to_rectangles(
            slide=slide,
            rect_shape_ids=(18, 20, 22),
            oval_shape_ids=(17, 19, 21),
            oval_overlap_in=0.05,
            oval_top_offset_in=0.04,
        )

    # Replace the opaque screenshots on the existing receipt slides with the transparent card stacks,
    # then align baskets/circles to the detected card bboxes.
    _replace_picture_image(
        _picture_shape(slide8, 8),
        slide8,
        Path("/Users/yipengandrewwang/BlueSky/Slide2/data_run_20260206/06_figures_preview/receipts/receipts_stack_discovery.png"),
    )
    _replace_picture_image(_picture_shape(slide9, 8), slide9, risk_stack_rich)

    for slide in (slide8, slide9):
        align_callouts_to_cards(
            slide=slide,
            picture_shape_id=8,
            rect_shape_ids=(18, 20, 22),
            oval_shape_ids=(17, 19, 21),
            inset_x_in=0.08,
            inset_y_in=0.06,
            oval_overlap_in=0.05,
            oval_top_offset_in=0.04,
        )

    # --- Slide 10: labels docs grounding ---
    _set_shape_text(_shape(slide10, 6), "Data receipt: labels (schema + global values + prefs)")
    _set_shape_text(_shape(slide10, 7), "Takeaway: labels are structured safety signals; behavior depends on viewer prefs.")

    _set_shape_text(_shape(slide10, 10), "A. Label schema")
    _set_shape_text(_shape(slide10, 11), "(src, uri/cid, val, neg, cts)")
    _set_shape_text(_shape(slide10, 12), "B. Global values")
    _set_shape_text(_shape(slide10, 13), "protocol-defined strings")
    _set_shape_text(_shape(slide10, 14), "C. Viewer prefs")
    _set_shape_text(_shape(slide10, 15), "hide / warn / ignore")
    _set_shape_text(_shape(slide10, 16), "Docs grounding")

    _replace_picture_image(_picture_shape(slide10, 8), slide10, labels_docs_stack)
    align_callouts_to_cards(
        slide=slide10,
        picture_shape_id=8,
        rect_shape_ids=(18, 20, 22),
        oval_shape_ids=(17, 19, 21),
        inset_x_in=0.08,
        inset_y_in=0.06,
        oval_overlap_in=0.05,
        oval_top_offset_in=0.04,
    )

    # --- Slide 11: custom labels + targets + observed values ---
    _set_shape_text(_shape(slide11, 6), "Data receipt: custom labels (definitions + targets)")
    _set_shape_text(_shape(slide11, 7), "Takeaway: labelers can define custom taxonomies; we preserve raw label_val strings.")

    _set_shape_text(_shape(slide11, 10), "A. Custom definitions")
    _set_shape_text(_shape(slide11, 11), "blur + severity + defaults")
    _set_shape_text(_shape(slide11, 12), "B. Label targets")
    _set_shape_text(_shape(slide11, 13), "account vs profile vs content")
    _set_shape_text(_shape(slide11, 14), "C. Observed values")
    _set_shape_text(_shape(slide11, 15), "dataset examples")
    _set_shape_text(_shape(slide11, 16), "Why it matters")

    _replace_picture_image(_picture_shape(slide11, 8), slide11, labels_custom_stack)
    align_callouts_to_cards(
        slide=slide11,
        picture_shape_id=8,
        rect_shape_ids=(18, 20, 22),
        oval_shape_ids=(17, 19, 21),
        inset_x_in=0.08,
        inset_y_in=0.06,
        oval_overlap_in=0.05,
        oval_top_offset_in=0.04,
    )

    # --- Slide 12: where labels come from (labelers + subscriptions + self-labels) ---
    _set_shape_text(_shape(slide12, 6), "Data receipt: labelers (declarations + subscriptions + self‑labels)")
    _set_shape_text(_shape(slide12, 7), "Takeaway: labels come from services; viewers choose which labelers to trust.")

    _set_shape_text(_shape(slide12, 10), "A. Declaration")
    _set_shape_text(_shape(slide12, 11), "values + policies")
    _set_shape_text(_shape(slide12, 12), "B. Subscriptions")
    _set_shape_text(_shape(slide12, 13), "viewer-selected")
    _set_shape_text(_shape(slide12, 14), "C. Self-labels")
    _set_shape_text(_shape(slide12, 15), "built-in values")
    _set_shape_text(_shape(slide12, 16), "Why it matters")

    _replace_picture_image(_picture_shape(slide12, 8), slide12, labels_labelers_stack)
    align_callouts_to_cards(
        slide=slide12,
        picture_shape_id=8,
        rect_shape_ids=(18, 20, 22),
        oval_shape_ids=(17, 19, 21),
        inset_x_in=0.08,
        inset_y_in=0.06,
        oval_overlap_in=0.05,
        oval_top_offset_in=0.04,
    )

    # --- Slide 13: how labels become UI actions (moderation decisions + contexts) ---
    _set_shape_text(_shape(slide13, 6), "Data receipt: moderation decisions (flags + UI contexts)")
    _set_shape_text(_shape(slide13, 7), "Takeaway: the same label can blur/warn/filter depending on context + prefs.")

    _set_shape_text(_shape(slide13, 10), "A. moderatePost()")
    _set_shape_text(_shape(slide13, 11), "post + prefs → decision")
    _set_shape_text(_shape(slide13, 12), "B. Output flags")
    _set_shape_text(_shape(slide13, 13), "filter / blur / alert / inform")
    _set_shape_text(_shape(slide13, 14), "C. UI contexts")
    _set_shape_text(_shape(slide13, 15), "contentList vs view vs media")
    _set_shape_text(_shape(slide13, 16), "Docs grounding")

    _replace_picture_image(_picture_shape(slide13, 8), slide13, labels_moderation_stack)
    align_callouts_to_cards(
        slide=slide13,
        picture_shape_id=8,
        rect_shape_ids=(18, 20, 22),
        oval_shape_ids=(17, 19, 21),
        inset_x_in=0.08,
        inset_y_in=0.06,
        oval_overlap_in=0.05,
        oval_top_offset_in=0.04,
    )

    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(pptx_out))

    if not keep_staged:
        for path in (work1, work2):
            try:
                path.unlink()
            except FileNotFoundError:
                pass


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Improve the professor deck: align callouts, add labels example slides.")
    p.add_argument("--in", dest="pptx_in", type=Path, required=True)
    p.add_argument("--out", dest="pptx_out", type=Path, required=True)
    p.add_argument("--keep-staged", action="store_true", help="Keep intermediate staged PPTX files for debugging.")
    return p


def main() -> int:
    args = _build_parser().parse_args()
    pptx_in = args.pptx_in.resolve()
    pptx_out = args.pptx_out.resolve()
    if not pptx_in.exists():
        raise SystemExit(f"Missing input PPTX: {pptx_in}")
    update_more_examples_deck(pptx_in=pptx_in, pptx_out=pptx_out, keep_staged=bool(args.keep_staged))
    print(f"OK: wrote {pptx_out}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
