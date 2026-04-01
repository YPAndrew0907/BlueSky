#!/usr/bin/env python3
from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path
from typing import Final, Iterable, Literal

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR, MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


EMU_PER_INCH: Final[int] = 914_400

SlideKind = Literal[
    "COVER",
    "SYSTEM",
    "UI_SCREENSHOT",
    "RQS",
    "ARTIFACTS",
    "RECEIPT",
    "CREDIBILITY",
    "NEXT",
    "APPENDIX",
    "UNKNOWN",
]


@dataclass(frozen=True)
class RectRatio:
    x: float
    y: float
    w: float
    h: float


@dataclass(frozen=True)
class UiReceiptSource:
    title: str
    takeaway: str
    chips: tuple[str, str, str, str, str, str, str]
    rects: dict[str, RectRatio]  # keys: A/B/C
    has_connectors: bool


@dataclass(frozen=True)
class CoverSource:
    title: str
    subtitle: str
    thesis: str
    badge: str


@dataclass(frozen=True)
class SystemSource:
    title: str
    takeaway: str
    nodes: tuple[str, str, str, str, str, str]  # viewer, discovery, feeds, hosting, outcomes, labelers
    bottom: str


@dataclass(frozen=True)
class RqsSource:
    title: str
    takeaway: str
    items: tuple[str, str, str, str, str]


@dataclass(frozen=True)
class ArtifactsSource:
    title: str
    takeaway: str
    groups: tuple[tuple[str, str], tuple[str, str], tuple[str, str], tuple[str, str]]


@dataclass(frozen=True)
class TwoCardSource:
    title: str
    left_card: tuple[str, str]
    right_card: tuple[str, str]


@dataclass(frozen=True)
class AppendixSource:
    title: str
    body_lines: tuple[str, ...]


@dataclass(frozen=True)
class BuildConfig:
    template_pptx: Path
    source_pptx: Path
    out_pptx: Path
    bsky_raw_dir: Path
    receipt_dir: Path
    official_dir: Path
    keep_project_label: bool


class Colors:
    # From Data Visual template palette screenshot.
    INK = RGBColor(0x26, 0x26, 0x26)
    MUTED = RGBColor(0x4B, 0x55, 0x63)
    LIGHT_LINE = RGBColor(0xCB, 0xD5, 0xE1)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)

    ACCENT = RGBColor(0x67, 0x6F, 0xFE)  # #676FFE
    ACCENT_2 = RGBColor(0x92, 0x9C, 0xF4)  # #929CF4
    GREEN = RGBColor(0x1D, 0xD3, 0x8B)  # #1DD38B

    # A/B/C scheme (distinct + template-consistent).
    A = ACCENT
    B = GREEN
    C = ACCENT_2

    A_TINT = RGBColor(0xE1, 0xEA, 0xFF)  # #E1EAFF
    B_TINT = RGBColor(0xE6, 0xFB, 0xF3)
    C_TINT = RGBColor(0xF0, 0xF1, 0xFF)

    BADGE_DARK = INK


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _iter_shapes(slide) -> Iterable:
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


def _remove_shape(slide, shape) -> None:
    slide.shapes._spTree.remove(shape._element)  # noqa: SLF001


def _set_text(
    shape,
    text: str,
    *,
    size_pt: float,
    bold: bool = False,
    color: RGBColor = Colors.INK,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    word_wrap: bool = True,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bold
        run.font.color.rgb = color


def _set_lines(
    shape,
    lines: list[str],
    *,
    size_pt: float,
    color: RGBColor = Colors.INK,
    bold_first: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
    for run in p.runs:
        run.font.size = Pt(size_pt)
        run.font.bold = bool(bold_first and idx == 0)
        run.font.color.rgb = color


def _clear_project_month(slide, *, keep_project_label: bool, keep_date_label: bool = True) -> None:
    for sh in list(slide.shapes):
        if not sh.is_placeholder:
            continue
        if sh.placeholder_format.type != 4:  # SUBTITLE in this template
            continue
        raw = (sh.text_frame.text or "").strip()
        if not raw:
            continue
        up = raw.upper()
        if up.startswith("PROJECT"):
            sh.name = "PROJECT_LABEL"
            if keep_project_label:
                _set_text(sh, "BLUESKY\nRQs + DATA", size_pt=10, bold=True, color=Colors.MUTED)
                sh.top = Inches(0.12)
                sh.height = Inches(0.28)
            else:
                sh.text_frame.clear()
            continue
        if up.startswith("MONTH"):
            sh.name = "DATE_LABEL"
            if keep_date_label:
                _set_text(sh, "FEB 2026", size_pt=11, bold=True, color=Colors.MUTED)
            else:
                sh.text_frame.clear()
            continue

        # Any other subtitle placeholders: clear.
        sh.text_frame.clear()


def _title_size_pt(title: str) -> float:
    # Heuristic: keep long "Data receipt: ..." / "System framing: ..." titles from colliding
    # with the takeaway line by shrinking font size when wrapping is likely.
    n = len(title.strip())
    if n >= 70:
        return 26.0
    if n >= 56:
        return 28.0
    return 30.0


@dataclass(frozen=True)
class HeaderLayout:
    title_left: float
    title_top: float
    title_w: float
    title_h: float
    takeaway_left: float
    takeaway_top: float
    takeaway_w: float
    takeaway_h: float


HEADER_STD: Final[HeaderLayout] = HeaderLayout(
    title_left=0.55,
    title_top=0.45,
    title_w=9.3,
    title_h=0.85,
    takeaway_left=0.55,
    takeaway_top=1.28,
    takeaway_w=9.3,
    takeaway_h=0.32,
)


def _delete_template_split_scaffold(slide) -> None:
    # Remove right-side picture placeholder and left-body placeholder, plus divider lines.
    for sh in list(slide.shapes):
        if sh.is_placeholder and sh.placeholder_format.type in {2, 18}:  # BODY / PICTURE
            _remove_shape(slide, sh)
            continue
        if sh.shape_type == MSO_SHAPE_TYPE.LINE:
            _remove_shape(slide, sh)


def _classify_v16_slide(slide) -> SlideKind:
    text = " ".join(
        p.text.strip()
        for shape in _iter_shapes(slide)
        if getattr(shape, "has_text_frame", False)
        for p in shape.text_frame.paragraphs
        if p.text.strip()
    )
    if "Algorithmic choice relocates power" in text:
        return "COVER"
    if "System framing:" in text:
        return "SYSTEM"
    if "Research questions" in text:
        return "RQS"
    if "Data artifacts collected" in text:
        return "ARTIFACTS"
    if "Collection credibility" in text:
        return "CREDIBILITY"
    if "What this enables next" in text:
        return "NEXT"
    if "Bluesky UI:" in text:
        return "UI_SCREENSHOT"
    if "Data receipt:" in text:
        return "RECEIPT"
    return "UNKNOWN"


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found")


def _picture(slide, shape_id: int):
    sh = _shape(slide, shape_id)
    if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
        raise TypeError(f"shape_id={shape_id} is not a picture (got {sh.shape_type})")
    return sh


def _rect_ratios_from_v16(slide) -> dict[str, RectRatio]:
    pic = _picture(slide, 8)
    pic_left = int(pic.left)
    pic_top = int(pic.top)
    pic_w = int(pic.width)
    pic_h = int(pic.height)

    out: dict[str, RectRatio] = {}
    for key, rect_id in (("A", 18), ("B", 20), ("C", 22)):
        rect = _shape(slide, rect_id)
        out[key] = RectRatio(
            x=(int(rect.left) - pic_left) / pic_w,
            y=(int(rect.top) - pic_top) / pic_h,
            w=int(rect.width) / pic_w,
            h=int(rect.height) / pic_h,
        )
    return out


def _chips_from_v16(slide) -> tuple[str, str, str, str, str, str, str]:
    chips: list[str] = []
    for sid in range(10, 17):
        sh = _shape(slide, sid)
        chips.append(sh.text_frame.text.strip())
    if len(chips) != 7:
        raise RuntimeError("expected 7 chips")
    return tuple(chips)  # type: ignore[return-value]


def _extract_v16_sources(v16: Presentation) -> dict[int, object]:
    sources: dict[int, object] = {}

    # Slide 0: cover
    s0 = v16.slides[0]
    sources[0] = CoverSource(
        title=_shape(s0, 7).text_frame.text.strip(),
        subtitle=_shape(s0, 8).text_frame.text.strip(),
        thesis=_shape(s0, 11).text_frame.text.strip(),
        badge=_shape(s0, 12).text_frame.text.strip(),
    )

    # Slide 1: system
    s1 = v16.slides[1]
    sources[1] = SystemSource(
        title=_shape(s1, 7).text_frame.text.strip(),
        takeaway=_shape(s1, 8).text_frame.text.strip(),
        nodes=(
            _shape(s1, 9).text_frame.text.strip(),
            _shape(s1, 10).text_frame.text.strip(),
            _shape(s1, 11).text_frame.text.strip(),
            _shape(s1, 12).text_frame.text.strip(),
            _shape(s1, 13).text_frame.text.strip(),
            _shape(s1, 14).text_frame.text.strip(),
        ),
        bottom=_shape(s1, 21).text_frame.text.strip(),
    )

    # UI slides 2-4
    for idx in (2, 3, 4):
        s = v16.slides[idx]
        sources[idx] = UiReceiptSource(
            title=_shape(s, 6).text_frame.text.strip(),
            takeaway=_shape(s, 7).text_frame.text.strip(),
            chips=_chips_from_v16(s),
            rects=_rect_ratios_from_v16(s),
            has_connectors=True,
        )

    # Slide 5: RQs
    s5 = v16.slides[5]
    items: list[str] = []
    for sid in (9, 11, 13, 15, 17):
        items.append(_shape(s5, sid).text_frame.text.strip())
    sources[5] = RqsSource(
        title=_shape(s5, 7).text_frame.text.strip(),
        takeaway=_shape(s5, 8).text_frame.text.strip(),
        items=tuple(items),  # type: ignore[assignment]
    )

    # Slide 6: artifacts collected
    s6 = v16.slides[6]
    groups: list[tuple[str, str]] = []
    for sid in (9, 10, 11, 12):
        lines = _shape(s6, sid).text_frame.text.strip().splitlines()
        head = lines[0].strip() if lines else ""
        body = "\n".join(line.strip() for line in lines[1:] if line.strip())
        groups.append((head, body))
    sources[6] = ArtifactsSource(
        title=_shape(s6, 8).text_frame.text.strip(),
        takeaway="Takeaway: artifacts are joinable, reproducible, and auditable.",
        groups=tuple(groups),  # type: ignore[assignment]
    )

    # Receipt slides 7-12
    for idx in range(7, 13):
        s = v16.slides[idx]
        sources[idx] = UiReceiptSource(
            title=_shape(s, 6).text_frame.text.strip(),
            takeaway=_shape(s, 7).text_frame.text.strip(),
            chips=_chips_from_v16(s),
            rects=_rect_ratios_from_v16(s),
            has_connectors=False,
        )

    # Slide 13: credibility (two cards)
    s13 = v16.slides[13]
    left_lines = _shape(s13, 9).text_frame.text.strip().splitlines()
    right_lines = _shape(s13, 10).text_frame.text.strip().splitlines()
    sources[13] = TwoCardSource(
        title=_shape(s13, 8).text_frame.text.strip(),
        left_card=(left_lines[0], "\n".join(left_lines[1:])),
        right_card=(right_lines[0], "\n".join(right_lines[1:])),
    )

    # Slide 14: next (two cards)
    s14 = v16.slides[14]
    left_lines = _shape(s14, 8).text_frame.text.strip().splitlines()
    right_lines = _shape(s14, 9).text_frame.text.strip().splitlines()
    sources[14] = TwoCardSource(
        title=_shape(s14, 7).text_frame.text.strip(),
        left_card=(left_lines[0], "\n".join(left_lines[1:])),
        right_card=(right_lines[0], "\n".join(right_lines[1:])),
    )

    # Official appendix slides 17-23
    for idx in range(17, 24):
        s = v16.slides[idx]
        sources[idx] = UiReceiptSource(
            title=_shape(s, 6).text_frame.text.strip(),
            takeaway=_shape(s, 7).text_frame.text.strip(),
            chips=_chips_from_v16(s),
            rects=_rect_ratios_from_v16(s),
            has_connectors=True,
        )

    return sources


def _add_divider_line(slide, *, x: float, y0: float, y1: float, width_pt: float = 1.25) -> None:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(y0), Inches(x), Inches(y1))
    line.name = "DIVIDER"
    line.line.width = Pt(width_pt)
    line.line.color.rgb = Colors.LIGHT_LINE


def _add_badge(slide, *, letter: str, x: float, y: float, d: float, color: RGBColor, name: str) -> None:
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    badge.name = name
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.color.rgb = color
    tf = badge.text_frame
    tf.clear()
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = letter
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = Colors.WHITE


def _add_highlight_rect(slide, *, x: float, y: float, w: float, h: float, color: RGBColor, name: str) -> None:
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    rect.name = name
    rect.fill.background()
    rect.line.color.rgb = color
    rect.line.width = Pt(2.5)


def _add_image_frame(slide, *, x: float, y: float, w: float, h: float) -> None:
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    frame.name = "IMAGE_FRAME"
    frame.fill.background()
    frame.line.color.rgb = Colors.LIGHT_LINE
    frame.line.width = Pt(1.25)


def _add_chip(
    slide,
    *,
    name: str,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    group: Literal["A", "B", "C"],
    is_heading: bool,
) -> None:
    color = {"A": Colors.A, "B": Colors.B, "C": Colors.C}[group]
    tint = {"A": Colors.A_TINT, "B": Colors.B_TINT, "C": Colors.C_TINT}[group]
    chip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    chip.name = name
    chip.fill.solid()
    chip.fill.fore_color.rgb = tint if is_heading else Colors.WHITE
    chip.line.color.rgb = color
    chip.line.width = Pt(2.0)

    tf = chip.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    for run in p.runs:
        run.font.size = Pt(13.25)
        run.font.color.rgb = Colors.INK
        if is_heading:
            run.font.bold = True


def _add_connector(slide, *, x1: float, y1: float, x2: float, y2: float, color: RGBColor, name: str) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.name = name
    conn.line.color.rgb = color
    conn.line.width = Pt(2.0)


def _build_screenshot_like_slide(
    slide,
    *,
    src: UiReceiptSource,
    image_path: Path,
    keep_project_label: bool,
) -> None:
    _clear_project_month(slide, keep_project_label=keep_project_label)
    _delete_template_split_scaffold(slide)

    # Title placeholder if present.
    title_ph = next(
        (sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 1),
        None,
    )
    if title_ph is None:
        title_ph = slide.shapes.add_textbox(
            Inches(HEADER_STD.title_left),
            Inches(HEADER_STD.title_top),
            Inches(HEADER_STD.title_w),
            Inches(HEADER_STD.title_h),
        )
    title_ph.name = "TITLE"
    _set_text(title_ph, src.title, size_pt=_title_size_pt(src.title), bold=False, color=Colors.INK)
    title_ph.left = Inches(HEADER_STD.title_left)
    title_ph.top = Inches(HEADER_STD.title_top)
    title_ph.width = Inches(HEADER_STD.title_w)
    title_ph.height = Inches(HEADER_STD.title_h)

    takeaway = slide.shapes.add_textbox(
        Inches(HEADER_STD.takeaway_left),
        Inches(HEADER_STD.takeaway_top),
        Inches(HEADER_STD.takeaway_w),
        Inches(HEADER_STD.takeaway_h),
    )
    takeaway.name = "TAKEAWAY"
    _set_text(takeaway, src.takeaway, size_pt=13.75, color=Colors.MUTED)

    # Divider between chips and screenshot.
    _add_divider_line(slide, x=3.65, y0=1.6, y1=5.15)

    # Main image (16:9).
    img_x = 3.8
    img_y = 1.65
    img_w = 6.05
    img_h = img_w * 9 / 16
    _add_image_frame(slide, x=img_x, y=img_y, w=img_w, h=img_h)
    main_pic = slide.shapes.add_picture(str(image_path), Inches(img_x), Inches(img_y), Inches(img_w), Inches(img_h))
    main_pic.name = "MAIN_IMAGE"

    # Chips.
    chip_x = 0.55
    chip_y0 = 1.65
    chip_w = 3.0
    chip_h = 0.42
    chip_gap = 0.08
    chip_groups: list[Literal["A", "B", "C"]] = ["A", "A", "B", "B", "C", "C", "C"]
    for idx, (text, group) in enumerate(zip(src.chips, chip_groups, strict=True), start=1):
        is_heading = text.strip().startswith(("A.", "B.", "C.")) or text.strip().lower().startswith("why")
        _add_chip(
            slide,
            name=f"CHIP_{idx:02d}",
            text=text,
            x=chip_x,
            y=chip_y0 + (idx - 1) * (chip_h + chip_gap),
            w=chip_w,
            h=chip_h,
            group=group,
            is_heading=is_heading,
        )

    # Highlight rectangles (ratio-mapped to new image box).
    def map_rect(r: RectRatio) -> tuple[float, float, float, float]:
        return (
            img_x + r.x * img_w,
            img_y + r.y * img_h,
            r.w * img_w,
            r.h * img_h,
        )

    a_x, a_y, a_w, a_h = map_rect(src.rects["A"])
    b_x, b_y, b_w, b_h = map_rect(src.rects["B"])
    c_x, c_y, c_w, c_h = map_rect(src.rects["C"])

    _add_highlight_rect(slide, x=a_x, y=a_y, w=a_w, h=a_h, color=Colors.A, name="HILITE_A")
    _add_highlight_rect(slide, x=b_x, y=b_y, w=b_w, h=b_h, color=Colors.B, name="HILITE_B")
    _add_highlight_rect(slide, x=c_x, y=c_y, w=c_w, h=c_h, color=Colors.C, name="HILITE_C")

    # Badges aligned to left edge of each highlight.
    badge_d = 0.30
    badge_overlap = 0.06
    _add_badge(slide, letter="A", x=a_x - badge_d + badge_overlap, y=a_y + 0.04, d=badge_d, color=Colors.A, name="BADGE_A")
    _add_badge(slide, letter="B", x=b_x - badge_d + badge_overlap, y=b_y + 0.04, d=badge_d, color=Colors.B, name="BADGE_B")
    _add_badge(slide, letter="C", x=c_x - badge_d + badge_overlap, y=c_y + 0.04, d=badge_d, color=Colors.C, name="BADGE_C")

    # Optional connectors (UI slides only): connect chip A and B groups to their highlights.
    if src.has_connectors:
        chip_a_top = chip_y0
        chip_a_bottom = chip_y0 + 2 * chip_h + chip_gap
        chip_b_top = chip_y0 + 2 * (chip_h + chip_gap)
        chip_b_bottom = chip_y0 + 4 * chip_h + 3 * chip_gap

        _add_connector(
            slide,
            x1=chip_x + chip_w,
            y1=(chip_a_top + chip_a_bottom) / 2,
            x2=a_x,
            y2=a_y + a_h / 2,
            color=Colors.A,
            name="CONN_A",
        )
        _add_connector(
            slide,
            x1=chip_x + chip_w,
            y1=(chip_b_top + chip_b_bottom) / 2,
            x2=b_x,
            y2=b_y + b_h / 2,
            color=Colors.B,
            name="CONN_B",
        )


def _build_cover_slide(slide, *, src: CoverSource, keep_project_label: bool) -> None:
    # Cover slide has 4 placeholders (see debug snippet in handoff).
    title_ph = _shape(slide, 118)
    title_ph.name = "TITLE"
    _set_lines(title_ph, src.title.splitlines() or [src.title], size_pt=54, color=Colors.WHITE, bold_first=False, align=PP_ALIGN.LEFT)

    subtitle_ph = _shape(slide, 119)
    subtitle_ph.name = "SUBTITLE"
    _set_text(subtitle_ph, src.subtitle, size_pt=18, color=Colors.WHITE, align=PP_ALIGN.LEFT)

    proj = _shape(slide, 120)
    proj.name = "PROJECT_LABEL"
    date = _shape(slide, 121)
    date.name = "DATE_LABEL"
    _set_text(proj, "BLUESKY RQs + DATA", size_pt=11, bold=True, color=Colors.WHITE)
    _set_text(date, "FEB 2026", size_pt=11, bold=True, color=Colors.WHITE)

    # Badge in the top-right.
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(7.35), Inches(0.55), Inches(2.1), Inches(0.55))
    badge.name = "BADGE_NO_RESULTS"
    badge.fill.solid()
    badge.fill.fore_color.rgb = Colors.BADGE_DARK
    badge.line.color.rgb = Colors.BADGE_DARK
    tf = badge.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = src.badge
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.bold = True
        run.font.size = Pt(14)
        run.font.color.rgb = Colors.WHITE


def _build_system_slide(slide, *, src: SystemSource, keep_project_label: bool) -> None:
    _clear_project_month(slide, keep_project_label=keep_project_label)
    _delete_template_split_scaffold(slide)

    title_ph = next(
        (sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 1),
        None,
    )
    if title_ph is None:
        title_ph = slide.shapes.add_textbox(
            Inches(HEADER_STD.title_left),
            Inches(HEADER_STD.title_top),
            Inches(HEADER_STD.title_w),
            Inches(HEADER_STD.title_h),
        )
    title_ph.name = "TITLE"
    _set_text(title_ph, src.title, size_pt=_title_size_pt(src.title), color=Colors.INK)
    title_ph.left = Inches(HEADER_STD.title_left)
    title_ph.top = Inches(HEADER_STD.title_top)
    title_ph.width = Inches(HEADER_STD.title_w)
    title_ph.height = Inches(HEADER_STD.title_h)

    takeaway = slide.shapes.add_textbox(
        Inches(HEADER_STD.takeaway_left),
        Inches(HEADER_STD.takeaway_top),
        Inches(HEADER_STD.takeaway_w),
        Inches(HEADER_STD.takeaway_h),
    )
    takeaway.name = "TAKEAWAY"
    _set_text(takeaway, src.takeaway, size_pt=13.75, color=Colors.MUTED)

    viewer, discovery, feeds, hosting, outcomes, labelers = src.nodes

    y = 2.15
    h = 0.68
    gap = 0.25
    x0 = 0.65
    widths = [1.45, 1.85, 1.55, 1.65, 1.55]
    labels = [viewer, discovery, feeds, hosting, outcomes]
    names = ["NODE_VIEWER", "NODE_DISCOVERY", "NODE_FEEDS", "NODE_HOSTING", "NODE_OUTCOMES"]

    boxes: list[tuple[float, float, object]] = []
    x = x0
    for w, text, name in zip(widths, labels, names, strict=True):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        box.name = name
        box.fill.solid()
        box.fill.fore_color.rgb = Colors.WHITE
        box.line.color.rgb = Colors.ACCENT
        box.line.width = Pt(2.0)
        tf = box.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        for run in p.runs:
            run.font.size = Pt(14)
            run.font.bold = True
            run.font.color.rgb = Colors.INK
        boxes.append((x, w, box))
        x += w + gap

    # Labelers node below feeds+hosting midpoint.
    feeds_x, feeds_w, _ = boxes[2]
    hosting_x, hosting_w, _ = boxes[3]
    mid_x = (feeds_x + feeds_w / 2 + hosting_x + hosting_w / 2) / 2
    label_w = 2.15
    label_x = mid_x - label_w / 2
    label_y = 3.2
    lab = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(label_x), Inches(label_y), Inches(label_w), Inches(h))
    lab.name = "NODE_LABELERS"
    lab.fill.solid()
    lab.fill.fore_color.rgb = Colors.WHITE
    lab.line.color.rgb = Colors.ACCENT
    lab.line.width = Pt(2.0)
    tf = lab.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = labelers
    p.alignment = PP_ALIGN.CENTER
    for run in p.runs:
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = Colors.INK

    # Connectors along the top row.
    for i in range(len(boxes) - 1):
        x_left, w_left, _ = boxes[i]
        x_right, _w_right, _ = boxes[i + 1]
        _add_connector(
            slide,
            x1=x_left + w_left,
            y1=y + h / 2,
            x2=x_right,
            y2=y + h / 2,
            color=Colors.ACCENT,
            name=f"LINK_{i+1}",
        )

    # Labelers connector (down from feeds).
    _add_connector(
        slide,
        x1=feeds_x + feeds_w / 2,
        y1=y + h,
        x2=label_x + label_w / 2,
        y2=label_y,
        color=Colors.ACCENT,
        name="LINK_LABELERS",
    )

    bottom = slide.shapes.add_textbox(Inches(0.55), Inches(4.65), Inches(9.3), Inches(0.5))
    bottom.name = "BOTTOM"
    _set_text(bottom, src.bottom, size_pt=14.25, color=Colors.MUTED)


def _build_rqs_slide(slide, *, src: RqsSource) -> None:
    # Use the table-of-contents template: 8 slots -> fill first 5.
    _clear_project_month(slide, keep_project_label=False)

    title_ph = next((sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 1), None)
    if title_ph is not None:
        title_ph.name = "TITLE"
        _set_lines(title_ph, [src.title], size_pt=40, color=Colors.INK, align=PP_ALIGN.LEFT)

    takeaway = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(9.3), Inches(0.4))
    takeaway.name = "TAKEAWAY"
    _set_text(takeaway, src.takeaway, size_pt=14.0, color=Colors.MUTED)

    # Collect slot shapes.
    section_shapes = [sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 4]
    number_shapes = [
        sh
        for sh in slide.shapes
        if sh.is_placeholder and sh.placeholder_format.type == 1 and (sh.text_frame.text or "").strip().isdigit()
    ]

    def _split_x(shapes: list[object]) -> int:
        lefts = [int(getattr(s, "left", 0)) for s in shapes]
        if not lefts:
            return 0
        return (min(lefts) + max(lefts)) // 2

    split_sections = _split_x(section_shapes)
    split_nums = _split_x(number_shapes)

    left_sections = sorted([s for s in section_shapes if int(s.left) < split_sections], key=lambda s: int(s.top))
    right_sections = sorted([s for s in section_shapes if int(s.left) >= split_sections], key=lambda s: int(s.top))
    left_nums = sorted([n for n in number_shapes if int(n.left) < split_nums], key=lambda s: int(s.top))
    right_nums = sorted([n for n in number_shapes if int(n.left) >= split_nums], key=lambda s: int(s.top))

    # Left column: R1-R4; Right column: R5 then blanks.
    items = list(src.items)

    def fill(col_sections, col_nums, start_idx: int, count: int) -> None:
        usable = min(count, len(col_sections), len(col_nums))
        for i in range(usable):
            slot_idx = start_idx + i
            sec = col_sections[i]
            num = col_nums[i]
            sec.name = f"RQ_SLOT_{slot_idx:02d}"
            num.name = f"RQ_NUM_{slot_idx:02d}"
            if slot_idx <= len(items):
                _set_text(sec, items[slot_idx - 1], size_pt=18, color=Colors.INK)
                _set_text(num, f"{slot_idx:02d}", size_pt=30, bold=True, color=Colors.ACCENT)
            else:
                sec.text_frame.clear()
                num.text_frame.clear()

    fill(left_sections, left_nums, start_idx=1, count=4)
    fill(right_sections, right_nums, start_idx=5, count=4)


def _summarize_group_body(body: str, *, max_items: int = 3) -> str:
    lines = [ln.strip() for ln in body.splitlines() if ln.strip()]
    if not lines:
        return ""
    # Prefer a compact, readable summary: join first few tokens.
    if len(lines) <= max_items:
        return " • ".join(lines)
    return " • ".join(lines[:max_items]) + " • …"


def _build_artifacts_slide(slide, *, src: ArtifactsSource, keep_project_label: bool) -> None:
    _clear_project_month(slide, keep_project_label=keep_project_label)

    title_ph = _shape(slide, 219)
    title_ph.name = "TITLE"
    _set_lines(title_ph, [src.title], size_pt=34, color=Colors.INK, align=PP_ALIGN.LEFT)

    takeaway = slide.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(9.3), Inches(0.35))
    takeaway.name = "TAKEAWAY"
    _set_text(takeaway, src.takeaway, size_pt=13.5, color=Colors.MUTED)

    num_ids = [228, 229, 230, 231]
    head_ids = [220, 222, 224, 226]
    body_ids = [221, 223, 225, 227]

    for idx, ((head, body), num_id, head_id, body_id) in enumerate(
        zip(src.groups, num_ids, head_ids, body_ids, strict=True), start=1
    ):
        num = _shape(slide, num_id)
        num.name = f"GROUP_{idx:02d}_NUM"
        _set_text(num, f"{idx:02d}", size_pt=26, bold=True, color=Colors.ACCENT)

        hsh = _shape(slide, head_id)
        hsh.name = f"GROUP_{idx:02d}_HEAD"
        _set_text(hsh, head, size_pt=18, bold=True, color=Colors.ACCENT)

        bsh = _shape(slide, body_id)
        bsh.name = f"GROUP_{idx:02d}_BODY"
        _set_text(bsh, _summarize_group_body(body, max_items=3), size_pt=13.0, color=Colors.INK)


def _build_two_card_slide(slide, *, src: TwoCardSource, keep_project_label: bool) -> None:
    _clear_project_month(slide, keep_project_label=keep_project_label)
    _delete_template_split_scaffold(slide)

    title_ph = next(
        (sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 1),
        None,
    )
    if title_ph is None:
        title_ph = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(9.3), Inches(0.7))
    title_ph.name = "TITLE"
    _set_text(title_ph, src.title, size_pt=30, color=Colors.INK)
    title_ph.left = Inches(0.55)
    title_ph.top = Inches(0.35)
    title_ph.width = Inches(9.3)
    title_ph.height = Inches(0.7)

    card_y = 1.55
    card_h = 3.75
    card_w = 4.45
    gap = 0.35
    left_x = 0.55
    right_x = left_x + card_w + gap

    def add_card(x: float, name: str, header: str, body: str, border: RGBColor) -> None:
        rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(card_y), Inches(card_w), Inches(card_h))
        rect.name = name
        rect.fill.solid()
        rect.fill.fore_color.rgb = Colors.WHITE
        rect.line.color.rgb = border
        rect.line.width = Pt(2.0)
        tf = rect.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_right = Inches(0.25)
        tf.margin_top = Inches(0.18)
        tf.margin_bottom = Inches(0.18)

        p0 = tf.paragraphs[0]
        p0.text = header
        p0.alignment = PP_ALIGN.LEFT
        for run in p0.runs:
            run.font.size = Pt(18)
            run.font.bold = True
            run.font.color.rgb = border

        for line in [ln.strip() for ln in body.splitlines() if ln.strip()]:
            p = tf.add_paragraph()
            p.text = f"• {line}"
            p.alignment = PP_ALIGN.LEFT
            for run in p.runs:
                run.font.size = Pt(14)
                run.font.color.rgb = Colors.INK

    add_card(left_x, "CARD_LEFT", src.left_card[0], src.left_card[1], border=Colors.ACCENT)
    add_card(right_x, "CARD_RIGHT", src.right_card[0], src.right_card[1], border=Colors.ACCENT_2)


def _build_appendix_slide(slide, *, src: AppendixSource, keep_project_label: bool) -> None:
    _clear_project_month(slide, keep_project_label=keep_project_label)
    title_ph = next((sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 1), None)
    if title_ph is None:
        title_ph = slide.shapes.add_textbox(Inches(0.55), Inches(0.55), Inches(9.3), Inches(1.2))
    title_ph.name = "TITLE"
    _set_lines(title_ph, [src.title], size_pt=44, color=Colors.INK, align=PP_ALIGN.LEFT)

    body_ph = next((sh for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format.type == 4), None)
    if body_ph is None:
        body_ph = slide.shapes.add_textbox(Inches(0.55), Inches(3.2), Inches(9.3), Inches(2.2))
    body_ph.name = "BODY"
    _set_lines(body_ph, [f"• {ln}" for ln in src.body_lines], size_pt=18, color=Colors.MUTED, align=PP_ALIGN.LEFT)


def build_deck(cfg: BuildConfig) -> None:
    template = Presentation(str(cfg.template_pptx))
    v16 = Presentation(str(cfg.source_pptx))

    if len(template.slides) != 24:
        raise ValueError(f"expected 24 slides in template, got {len(template.slides)}")
    if len(v16.slides) < 24:
        raise ValueError(f"expected >=24 slides in v16 source, got {len(v16.slides)}")

    sources = _extract_v16_sources(v16)

    # Fixed image mapping for the main UI slides.
    ui_images = {
        2: cfg.bsky_raw_dir / "bsky_raw_feeds.png",
        3: cfg.bsky_raw_dir / "bsky_raw_feed_for_science.png",
        4: cfg.bsky_raw_dir / "bsky_raw_starterpack.png",
    }

    for idx, slide in enumerate(template.slides):
        kind = _classify_v16_slide(v16.slides[idx])

        if idx == 0:
            _build_cover_slide(slide, src=sources[0], keep_project_label=cfg.keep_project_label)  # type: ignore[arg-type]
            continue
        if idx == 1:
            _build_system_slide(slide, src=sources[1], keep_project_label=cfg.keep_project_label)  # type: ignore[arg-type]
            continue
        # Slides 17-23: official composites (also classified as UI_SCREENSHOT).
        if 17 <= idx <= 23:
            official_img = cfg.official_dir / f"official_ui_slide{idx+1:02d}.png"
            _build_screenshot_like_slide(
                slide,
                src=sources[idx],  # type: ignore[arg-type]
                image_path=official_img,
                keep_project_label=cfg.keep_project_label,
            )
            continue

        if kind == "UI_SCREENSHOT":
            image = ui_images.get(idx)
            if image is None:
                raise KeyError(f"missing UI image mapping for slide index {idx}")
            _build_screenshot_like_slide(
                slide,
                src=sources[idx],  # type: ignore[arg-type]
                image_path=image,
                keep_project_label=cfg.keep_project_label,
            )
            continue
        if kind == "RQS":
            _build_rqs_slide(slide, src=sources[idx])  # type: ignore[arg-type]
            continue
        if kind == "ARTIFACTS":
            _build_artifacts_slide(slide, src=sources[idx], keep_project_label=cfg.keep_project_label)  # type: ignore[arg-type]
            continue
        if kind == "RECEIPT":
            receipt_img = cfg.receipt_dir / f"receipt_hd_slide{idx+1:02d}.png"
            _build_screenshot_like_slide(
                slide,
                src=sources[idx],  # type: ignore[arg-type]
                image_path=receipt_img,
                keep_project_label=cfg.keep_project_label,
            )
            continue
        if kind in {"CREDIBILITY", "NEXT"}:
            _build_two_card_slide(slide, src=sources[idx], keep_project_label=cfg.keep_project_label)  # type: ignore[arg-type]
            continue

        # Slide 15-16: turn blanks into visible appendix intro slides.
        if idx == 15:
            _build_appendix_slide(
                slide,
                src=AppendixSource(
                    title="Appendix: official Bluesky UI screens",
                    body_lines=(
                        "Screenshots are official (for backup/citation).",
                        "Each slide uses A/B/C chips + callouts.",
                        "Animations: chips cascade, then A → B → C.",
                    ),
                ),
                keep_project_label=cfg.keep_project_label,
            )
            continue
        if idx == 16:
            _build_appendix_slide(
                slide,
                src=AppendixSource(
                    title="Appendix: how to read the callouts",
                    body_lines=(
                        "A/B/C badges map to highlighted regions in the screenshot.",
                        "Chip text summarizes the surface + why it matters.",
                        "Use these as optional backup in the meeting.",
                    ),
                ),
                keep_project_label=cfg.keep_project_label,
            )
            continue

        # Fallback: clear project/month and leave template slide minimal.
        _clear_project_month(slide, keep_project_label=cfg.keep_project_label)

    cfg.out_pptx.parent.mkdir(parents=True, exist_ok=True)
    template.save(str(cfg.out_pptx))


def _build_arg_parser() -> argparse.ArgumentParser:
    here = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(description="Re-template the professor meeting deck to Slidesgo Data Visual (v17).")
    parser.add_argument(
        "--template",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v17.base.pptx",
        help="Template-based skeleton PPTX (v17.base).",
    )
    parser.add_argument(
        "--source",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v16.pptx",
        help="Source-of-truth deck for text + highlight ratios (v16).",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v17.preanim.pptx",
        help="Output PPTX before animation injection.",
    )
    parser.add_argument("--bsky-raw-dir", type=Path, default=here / "prof_build/bsky/raw")
    parser.add_argument("--receipt-dir", type=Path, default=here / "prof_build/assets_cache/polish_v15")
    parser.add_argument("--official-dir", type=Path, default=here / "prof_build/assets_cache/official_ui")
    parser.add_argument("--keep-project-label", action="store_true", help="Keep small PROJECT/MONTH labels on slides.")
    return parser


def main() -> None:
    args = _build_arg_parser().parse_args()

    cfg = BuildConfig(
        template_pptx=args.template.resolve(),
        source_pptx=args.source.resolve(),
        out_pptx=args.out.resolve(),
        bsky_raw_dir=args.bsky_raw_dir.resolve(),
        receipt_dir=args.receipt_dir.resolve(),
        official_dir=args.official_dir.resolve(),
        keep_project_label=bool(args.keep_project_label),
    )

    for p in (cfg.template_pptx, cfg.source_pptx):
        if not p.exists():
            raise FileNotFoundError(p)

    build_deck(cfg)


if __name__ == "__main__":
    main()
