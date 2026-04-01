#!/usr/bin/env python3
from __future__ import annotations

import argparse
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Literal, Sequence

from lxml import etree
from PIL import Image, ImageDraw, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn


EMU_PER_INCH = 914400

PPTX_MIME_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

NSMAP = {"p": PML_NS, "r": R_NS}
REL_NSMAP = {"rel": REL_NS}
CT_NSMAP = {"ct": CT_NS}


@dataclass(frozen=True)
class BboxPx:
    x0: int
    y0: int
    x1: int
    y1: int

    @property
    def w(self) -> int:
        return self.x1 - self.x0

    @property
    def h(self) -> int:
        return self.y1 - self.y0


@dataclass(frozen=True)
class BboxIn:
    x: float
    y: float
    w: float
    h: float


LayoutId = Literal["L1", "L2", "L3", "L4"]


@dataclass(frozen=True)
class AppendixSlideSpec:
    title: str
    takeaway: str
    chips: tuple[str, str, str, str, str, str, str]
    layout: LayoutId
    img_a: str
    img_b: str
    img_c: str


OFFICIAL_SPECS: tuple[AppendixSlideSpec, ...] = (
    AppendixSlideSpec(
        title="Bluesky UI: labeling surfaces (official)",
        takeaway="Takeaway: labels are user-facing controls that shape what gets shown or filtered.",
        chips=(
            "A. Self-label button",
            "Composer surface",
            "B. Content warning picker",
            "Adult Content options",
            "C. Labeled post banner",
            "Attribution to labeler",
            "Why it matters: filtering + exposure",
        ),
        layout="L1",
        img_a="01_self-label_button_crop.jpg",
        img_b="02_self-label_picker_crop.jpg",
        img_c="05_labeled_post.png",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: reply controls (official)",
        takeaway="Takeaway: reply controls gate interaction, changing who can respond and what feedback loops form.",
        chips=(
            "A. Reply controls button",
            "Composer setting",
            "B. Who can reply picker",
            "Audience gating",
            "C. Post notice",
            "Visible constraint",
            "Why it matters: interaction graphs",
        ),
        layout="L1",
        img_a="06_reply_controls_button.jpg",
        img_b="07_reply_controls_picker_crop.jpg",
        img_c="08_reply_controls_notice_crop.jpg",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: labelers + moderation tools (official)",
        takeaway="Takeaway: moderation is modular—users can subscribe to labelers and tune label handling.",
        chips=(
            "A. Subscribe to labeler",
            "Third-party policies",
            "B. Per-label controls",
            "Off / Warn / Hide",
            "C. Ozone moderation UI",
            "Operational tooling",
            "Why it matters: power shifts",
        ),
        layout="L3",
        img_a="03_labeler_subscribe_crop.jpg",
        img_b="04_label_options_crop.jpg",
        img_c="07_ozone_moderation_tool_crop.png",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: discovery surfaces (official)",
        takeaway="Takeaway: discovery is shaped by navigation, defaults, and safety settings surfaced in the UI.",
        chips=(
            "A. Custom feeds tab",
            "Algorithm marketplace",
            "B. Moderation menu",
            "Filtering settings",
            "C. Starter packs tab",
            "Onboarding surface",
            "Why it matters: defaults steer",
        ),
        layout="L2",
        img_a="01_custom_feeds_tab_crop.jpg",
        img_b="06_moderation_menu_crop.jpg",
        img_c="02_starter_packs_tab_crop.jpg",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: starter packs flow (official)",
        takeaway="Takeaway: starter packs bundle people (and context) into a single shareable onboarding object.",
        chips=(
            "A. Create starter pack",
            "Name + description",
            "B. Choose people",
            "Curation choices",
            "C. Share pack",
            "Link / QR / image",
            "Why it matters: bundled exposure",
        ),
        layout="L4",
        img_a="03_starter_pack_create_name_crop.jpg",
        img_b="04_starter_pack_choose_people_crop.jpg",
        img_c="05_starter_pack_share_crop.jpg",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: domain handles (surface → settings) (official)",
        takeaway="Takeaway: domain handles connect identity and credibility to a verifiable control channel (DNS).",
        chips=(
            "A. Domain handle profile",
            "Trust signal",
            "B. Settings: handle",
            "Entry point",
            "C. Use own domain",
            "Verification flow",
            "Why it matters: credibility",
        ),
        layout="L1",
        img_a="09_domain_handle_profile_crop.png",
        img_b="10_account_settings_handle_crop.jpg",
        img_c="11_change_handle_domain_crop.jpg",
    ),
    AppendixSlideSpec(
        title="Bluesky UI: domain verification + examples (official)",
        takeaway="Takeaway: verification relies on external DNS tooling, creating measurable and potentially gameable steps.",
        chips=(
            "A. DNS TXT record",
            "Proof-of-control",
            "B. Domain handle example",
            "@npr.org",
            "C. DNS management UI",
            "External dependency",
            "Why it matters: measurable steps",
        ),
        layout="L2",
        img_a="12_domain_dns_record_crop.jpg",
        img_b="13_domain_handle_npr_crop.png",
        img_c="14_dns_management_example.png",
    ),
)


LAYOUTS_PX: dict[LayoutId, dict[str, BboxPx]] = {
    "L1": {
        "A": BboxPx(60, 60, 1860, 340),
        "B": BboxPx(60, 376, 942, 1020),
        "C": BboxPx(978, 376, 1860, 1020),
    },
    "L2": {
        "A": BboxPx(60, 60, 942, 724),
        "B": BboxPx(978, 60, 1860, 724),
        "C": BboxPx(60, 760, 1860, 1020),
    },
    "L3": {
        "A": BboxPx(60, 60, 1860, 356),
        "B": BboxPx(60, 392, 1860, 688),
        "C": BboxPx(60, 724, 1860, 1020),
    },
    "L4": {
        "A": BboxPx(60, 60, 636, 1020),
        "B": BboxPx(672, 60, 1248, 1020),
        "C": BboxPx(1284, 60, 1860, 1020),
    },
}


def _emu(inches: float) -> int:
    return int(round(inches * EMU_PER_INCH))


def _emu_to_in(emu: int) -> float:
    return emu / EMU_PER_INCH


def _read_zip_xml(z: zipfile.ZipFile, name: str) -> etree._Element:
    return etree.fromstring(z.read(name))


def _write_zip_xml(z: zipfile.ZipFile, name: str, root: etree._Element) -> None:
    z.writestr(name, etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True))


def _max_slide_number(z: zipfile.ZipFile) -> int:
    max_n = 0
    for name in z.namelist():
        if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
            continue
        stem = Path(name).stem
        try:
            n = int(stem.replace("slide", ""))
        except ValueError:
            continue
        max_n = max(max_n, n)
    return max_n


def _next_rid(pres_rels_root: etree._Element) -> str:
    rids: list[int] = []
    for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                rids.append(int(rid[3:]))
            except ValueError:
                continue
    nxt = (max(rids) + 1) if rids else 1
    return f"rId{nxt}"


def _next_slide_id(pres_root: etree._Element) -> int:
    sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
    if sld_id_lst is None:
        raise RuntimeError("presentation.xml missing p:sldIdLst")
    ids: list[int] = []
    for sld_id in sld_id_lst.findall("p:sldId", namespaces=NSMAP):
        try:
            ids.append(int(sld_id.get("id", "0")))
        except ValueError:
            continue
    return (max(ids) + 1) if ids else 256


def duplicate_slide_parts(
    *,
    pptx_in: Path,
    pptx_out: Path,
    base_slide_num: int,
    copies: int,
    insert_after_slide_num: int,
) -> list[int]:
    """
    Duplicate a slide part (slide{base_slide_num}.xml + rels) N times.

    Returns the new slide part numbers (e.g., [18, 19]).
    """
    if copies <= 0:
        raise ValueError("copies must be > 0")

    with zipfile.ZipFile(pptx_in, "r") as zin:
        pres_root = _read_zip_xml(zin, "ppt/presentation.xml")
        pres_rels_root = _read_zip_xml(zin, "ppt/_rels/presentation.xml.rels")
        ct_root = _read_zip_xml(zin, "[Content_Types].xml")

        sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
        if sld_id_lst is None:
            raise RuntimeError("presentation.xml missing p:sldIdLst")

        max_slide = _max_slide_number(zin)
        next_slide_id = _next_slide_id(pres_root)

        slide_src = f"ppt/slides/slide{base_slide_num}.xml"
        slide_rels_src = f"ppt/slides/_rels/slide{base_slide_num}.xml.rels"
        slide_src_bytes = zin.read(slide_src)
        slide_rels_src_bytes = zin.read(slide_rels_src)

        new_nums: list[int] = []
        for _ in range(copies):
            n = max_slide + 1
            max_slide = n
            new_nums.append(n)

            rid = _next_rid(pres_rels_root)
            rel_el = etree.Element(f"{{{REL_NS}}}Relationship")
            rel_el.set("Id", rid)
            rel_el.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel_el.set("Target", f"slides/slide{n}.xml")
            pres_rels_root.append(rel_el)

            sld_id_el = etree.Element(f"{{{PML_NS}}}sldId")
            sld_id_el.set("id", str(next_slide_id))
            sld_id_el.set(f"{{{R_NS}}}id", rid)
            sld_id_lst.append(sld_id_el)
            next_slide_id += 1

            override = etree.Element(f"{{{CT_NS}}}Override")
            override.set("PartName", f"/ppt/slides/slide{n}.xml")
            override.set("ContentType", PPTX_MIME_SLIDE)
            ct_root.append(override)

        # Reorder slides to insert new slides right after insert_after_slide_num.
        rid_to_target: dict[str, str] = {}
        for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
            if rel.get("Type") != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
                continue
            rid = rel.get("Id")
            tgt = rel.get("Target")
            if rid and tgt:
                rid_to_target[rid] = tgt

        target_to_sld_id: dict[str, etree._Element] = {}
        targets_in_order: list[str] = []
        for sld_id in list(sld_id_lst.findall("p:sldId", namespaces=NSMAP)):
            rid = sld_id.get(f"{{{R_NS}}}id")
            tgt = rid_to_target.get(rid or "")
            if not tgt:
                continue
            target_to_sld_id[tgt] = sld_id
            targets_in_order.append(tgt)

        new_targets = [f"slides/slide{n}.xml" for n in new_nums]
        base_target = f"slides/slide{insert_after_slide_num}.xml"
        if base_target not in targets_in_order:
            raise RuntimeError(f"Could not find slide target in order: {base_target}")

        targets_wo_new = [t for t in targets_in_order if t not in set(new_targets)]
        insert_at = targets_wo_new.index(base_target) + 1
        final_targets = targets_wo_new[:insert_at] + new_targets + targets_wo_new[insert_at:]

        for child in list(sld_id_lst):
            sld_id_lst.remove(child)
        for tgt in final_targets:
            el = target_to_sld_id.get(tgt)
            if el is None:
                raise RuntimeError(f"Missing sldId element for target: {tgt}")
            sld_id_lst.append(el)

        tmp_path = pptx_out.with_suffix(".tmp.pptx")
        if tmp_path.exists():
            tmp_path.unlink()
        pptx_out.parent.mkdir(parents=True, exist_ok=True)

        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}:
                    continue
                zout.writestr(item, zin.read(item.filename))

            _write_zip_xml(zout, "ppt/presentation.xml", pres_root)
            _write_zip_xml(zout, "ppt/_rels/presentation.xml.rels", pres_rels_root)
            _write_zip_xml(zout, "[Content_Types].xml", ct_root)

            for n in new_nums:
                zout.writestr(f"ppt/slides/slide{n}.xml", slide_src_bytes)
                zout.writestr(f"ppt/slides/_rels/slide{n}.xml.rels", slide_rels_src_bytes)

        if pptx_out.exists():
            pptx_out.unlink()
        tmp_path.rename(pptx_out)

    return new_nums


def _resolve_image_path(*, official_screens_dir: Path, rel_name: str) -> Path:
    matches = [p for p in official_screens_dir.rglob(rel_name) if p.is_file() and p.name == rel_name]
    if not matches:
        raise FileNotFoundError(f"missing official screenshot: {rel_name} (root={official_screens_dir})")
    if len(matches) > 1:
        pretty = "\n".join(f"- {m}" for m in sorted(matches))
        raise RuntimeError(f"ambiguous official screenshot: {rel_name}\n{pretty}")
    return matches[0]


def _render_three_panel_composite(
    *,
    out_path: Path,
    layout: LayoutId,
    img_a: Path,
    img_b: Path,
    img_c: Path,
    pad_px: int = 60,
    border_px: int = 4,
    inset_px: int = 0,
) -> tuple[Path, dict[str, BboxPx]]:
    """
    Render a 1920x1080 composite with three panels (A/B/C) and return the panel bboxes in px.

    NOTE: `inset_px` is for display-only in the composite; highlight bboxes should be computed separately.
    """
    out_path.parent.mkdir(parents=True, exist_ok=True)

    canvas_w, canvas_h = 1920, 1080
    bg = Image.new("RGBA", (canvas_w, canvas_h), "#0B1320")
    draw = ImageDraw.Draw(bg)

    panels = LAYOUTS_PX[layout]
    images = {"A": img_a, "B": img_b, "C": img_c}

    for key in ("A", "B", "C"):
        bbox = panels[key]
        x0, y0, x1, y1 = bbox.x0, bbox.y0, bbox.x1, bbox.y1

        draw.rectangle([x0, y0, x1, y1], fill="#111C2D", outline="#2B3A55", width=border_px)

        content = BboxPx(
            x0=x0 + pad_px + inset_px,
            y0=y0 + pad_px + inset_px,
            x1=x1 - pad_px - inset_px,
            y1=y1 - pad_px - inset_px,
        )
        if content.w <= 0 or content.h <= 0:
            raise RuntimeError(f"invalid content bbox for panel {key}: {content}")

        img = ImageOps.exif_transpose(Image.open(images[key])).convert("RGBA")
        resampling = getattr(Image, "Resampling", Image)
        fitted = ImageOps.contain(img, (content.w, content.h), method=resampling.LANCZOS)

        ox = content.x0 + (content.w - fitted.size[0]) // 2
        oy = content.y0 + (content.h - fitted.size[1]) // 2
        bg.alpha_composite(fitted, (ox, oy))

    bg.save(out_path, format="PNG")
    return out_path, panels


def _picture_shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh
    raise KeyError(f"picture shape_id={shape_id} not found")


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found")


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001
    blip.set(qn("r:embed"), rid)


def _set_shape_text(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        raise TypeError("shape has no text frame")
    tf = shape.text_frame
    if not tf.paragraphs:
        p = tf.add_paragraph()
        p.text = text
        return
    p0 = tf.paragraphs[0]
    if p0.runs:
        p0.runs[0].text = text
        for run in p0.runs[1:]:
            run.text = ""
    else:
        p0.text = text
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _map_bbox_to_slide_in(*, bbox: BboxPx, pic_shape, img_w: int, img_h: int) -> BboxIn:
    pic_left = _emu_to_in(int(pic_shape.left))
    pic_top = _emu_to_in(int(pic_shape.top))
    pic_w = _emu_to_in(int(pic_shape.width))
    pic_h = _emu_to_in(int(pic_shape.height))

    sx = pic_w / img_w
    sy = pic_h / img_h
    x = pic_left + bbox.x0 * sx
    y = pic_top + bbox.y0 * sy
    w = bbox.w * sx
    h = bbox.h * sy
    return BboxIn(x=x, y=y, w=w, h=h)


def _apply_panel_highlights(
    *,
    slide,
    pic_shape_id: int,
    panel_bboxes_px: dict[str, BboxPx],
    inset_px: int,
) -> None:
    pic = _picture_shape(slide, pic_shape_id)
    img_w, img_h = pic.image.size

    mapping = {
        18: panel_bboxes_px["A"],
        20: panel_bboxes_px["B"],
        22: panel_bboxes_px["C"],
    }
    for rect_id, bbox in mapping.items():
        bbox_inset = BboxPx(
            x0=bbox.x0 + inset_px,
            y0=bbox.y0 + inset_px,
            x1=bbox.x1 - inset_px,
            y1=bbox.y1 - inset_px,
        )
        bbox_in = _map_bbox_to_slide_in(bbox=bbox_inset, pic_shape=pic, img_w=img_w, img_h=img_h)
        rect = _shape(slide, rect_id)
        rect.left = _emu(bbox_in.x)
        rect.top = _emu(bbox_in.y)
        rect.width = _emu(bbox_in.w)
        rect.height = _emu(bbox_in.h)


def _align_ovals_to_rectangles(
    *,
    slide,
    rect_shape_ids: Sequence[int],
    oval_shape_ids: Sequence[int],
    oval_overlap_in: float,
    oval_top_offset_in: float,
) -> None:
    if len(rect_shape_ids) != len(oval_shape_ids):
        raise ValueError("rect_shape_ids and oval_shape_ids must have same length")

    for rect_id, oval_id in zip(rect_shape_ids, oval_shape_ids, strict=True):
        rect = _shape(slide, rect_id)
        oval = _shape(slide, oval_id)

        rect_left = _emu_to_in(int(rect.left))
        rect_top = _emu_to_in(int(rect.top))
        oval_w = _emu_to_in(int(oval.width))
        oval_h = _emu_to_in(int(oval.height))

        oval.left = _emu(rect_left - oval_w + oval_overlap_in)
        oval.top = _emu(rect_top + oval_top_offset_in)
        oval.width = _emu(oval_w)
        oval.height = _emu(oval_h)


def _shape_bbox(shape) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0))
    top = int(getattr(shape, "top", 0))
    width = int(getattr(shape, "width", 0))
    height = int(getattr(shape, "height", 0))
    return left, top, left + width, top + height


def _center_y(top: int, bottom: int) -> int:
    return int(round((top + bottom) / 2))


def _set_connectors_to_highlights(slide) -> None:
    # Connector 23: A chips -> rect 18
    chip_a0 = _shape(slide, 10)
    chip_a1 = _shape(slide, 11)
    rect_a = _shape(slide, 18)
    conn_a = _shape(slide, 23)

    a_l0, a_t0, a_r0, a_b0 = _shape_bbox(chip_a0)
    a_l1, a_t1, a_r1, a_b1 = _shape_bbox(chip_a1)
    a_right = max(a_r0, a_r1)
    a_top = min(a_t0, a_t1)
    a_bottom = max(a_b0, a_b1)

    ra_l, ra_t, ra_r, ra_b = _shape_bbox(rect_a)
    conn_a.begin_x = a_right
    conn_a.begin_y = _center_y(a_top, a_bottom)
    conn_a.end_x = ra_l
    conn_a.end_y = _center_y(ra_t, ra_b)

    # Connector 24: B chips -> rect 20
    chip_b0 = _shape(slide, 12)
    chip_b1 = _shape(slide, 13)
    rect_b = _shape(slide, 20)
    conn_b = _shape(slide, 24)

    b_l0, b_t0, b_r0, b_b0 = _shape_bbox(chip_b0)
    b_l1, b_t1, b_r1, b_b1 = _shape_bbox(chip_b1)
    b_right = max(b_r0, b_r1)
    b_top = min(b_t0, b_t1)
    b_bottom = max(b_b0, b_b1)

    rb_l, rb_t, rb_r, rb_b = _shape_bbox(rect_b)
    conn_b.begin_x = b_right
    conn_b.begin_y = _center_y(b_top, b_bottom)
    conn_b.end_x = rb_l
    conn_b.end_y = _center_y(rb_t, rb_b)


def append_official_ui_screens(
    *,
    pptx_in: Path,
    pptx_out: Path,
    official_screens_dir: Path,
    cache_dir: Path,
    template_slide_part: int,
    insert_after_slide_part: int,
    hide_new_slides: bool,
) -> None:
    import tempfile

    pptx_in = pptx_in.resolve()
    pptx_out = pptx_out.resolve()
    official_screens_dir = official_screens_dir.resolve()
    cache_dir = cache_dir.resolve()

    if not pptx_in.exists():
        raise FileNotFoundError(f"missing --in: {pptx_in}")
    if not official_screens_dir.exists():
        raise FileNotFoundError(f"missing --official-screens-dir: {official_screens_dir}")

    with tempfile.TemporaryDirectory(prefix="pptx_official_ui_") as tmp_dir:
        tmp_duped = Path(tmp_dir) / "duped.pptx"

        new_parts = duplicate_slide_parts(
            pptx_in=pptx_in,
            pptx_out=tmp_duped,
            base_slide_num=int(template_slide_part),
            copies=len(OFFICIAL_SPECS),
            insert_after_slide_num=int(insert_after_slide_part),
        )

        prs = Presentation(str(tmp_duped))

    slides_by_part: dict[int, object] = {}
    order_by_part: dict[int, int] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        partname = str(slide.part.partname)
        try:
            num = int(Path(partname).stem.replace("slide", ""))
        except ValueError:
            continue
        slides_by_part[num] = slide
        order_by_part[num] = idx

    new_parts_sorted = sorted(new_parts, key=lambda p: order_by_part.get(p, 10**9))
    if len(new_parts_sorted) != len(OFFICIAL_SPECS):
        raise RuntimeError("internal error: new slide part count mismatch")

    for part_num, spec in zip(new_parts_sorted, OFFICIAL_SPECS, strict=True):
        slide = slides_by_part.get(part_num)
        if slide is None:
            raise RuntimeError(f"missing duplicated slide part in python-pptx view: slide{part_num}.xml")
        slide_num = order_by_part.get(part_num, part_num)

        img_a = _resolve_image_path(official_screens_dir=official_screens_dir, rel_name=spec.img_a)
        img_b = _resolve_image_path(official_screens_dir=official_screens_dir, rel_name=spec.img_b)
        img_c = _resolve_image_path(official_screens_dir=official_screens_dir, rel_name=spec.img_c)

        composite_path = cache_dir / f"official_ui_slide{slide_num:02d}.png"
        composite_path, panels_px = _render_three_panel_composite(
            out_path=composite_path,
            layout=spec.layout,
            img_a=img_a,
            img_b=img_b,
            img_c=img_c,
        )

        _set_shape_text(_shape(slide, 6), spec.title)
        _set_shape_text(_shape(slide, 7), spec.takeaway)
        for shape_id, line in zip(range(10, 17), spec.chips, strict=True):
            _set_shape_text(_shape(slide, shape_id), line)

        _replace_picture_image(_picture_shape(slide, 8), slide, composite_path)

        _apply_panel_highlights(
            slide=slide,
            pic_shape_id=8,
            panel_bboxes_px=panels_px,
            inset_px=18,
        )
        _align_ovals_to_rectangles(
            slide=slide,
            rect_shape_ids=(18, 20, 22),
            oval_shape_ids=(17, 19, 21),
            oval_overlap_in=0.05,
            oval_top_offset_in=0.04,
        )
        _set_connectors_to_highlights(slide)

        if hide_new_slides:
            slide._element.set("show", "0")  # noqa: SLF001

        pptx_out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(pptx_out))


def main() -> None:
    here = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(description="Append official Bluesky UI screenshots as hidden backup slides.")
    parser.add_argument(
        "--in",
        dest="pptx_in",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v13.base.pptx",
    )
    parser.add_argument(
        "--out",
        dest="pptx_out",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v13.preanim.pptx",
    )
    parser.add_argument(
        "--official-screens-dir",
        type=Path,
        default=Path(
            "/Users/yipengandrewwang/Linedata-main/fabrics/effa7ec1-4d1c-4100-ad5a-29cfd890f999/"
            "SynapseNotebook/589f3786-8338-47c2-8c8f-7dd174f8a56d/Notebook_1/bluesky_official_screens"
        ),
    )
    parser.add_argument(
        "--cache-dir",
        type=Path,
        default=Path("/Users/yipengandrewwang/BlueSky/Slide2/prof_build/assets_cache/official_ui"),
    )
    parser.add_argument("--template-slide-part", type=int, default=3)
    parser.add_argument("--insert-after-slide-part", type=int, default=17)
    parser.add_argument(
        "--hide-new-slides",
        action=argparse.BooleanOptionalAction,
        default=True,
        help="Hide newly appended slides (p:sld @show=\"0\").",
    )

    args = parser.parse_args()

    append_official_ui_screens(
        pptx_in=args.pptx_in,
        pptx_out=args.pptx_out,
        official_screens_dir=args.official_screens_dir,
        cache_dir=args.cache_dir,
        template_slide_part=int(args.template_slide_part),
        insert_after_slide_part=int(args.insert_after_slide_part),
        hide_new_slides=bool(args.hide_new_slides),
    )
    print(f"OK: wrote {args.pptx_out.resolve()}")


if __name__ == "__main__":
    main()
