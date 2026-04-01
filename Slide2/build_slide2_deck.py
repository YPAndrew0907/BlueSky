#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import shutil
import zipfile
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Literal, Sequence

from lxml import etree
from PIL import Image, ImageChops, ImageDraw, ImageEnhance, ImageFilter, ImageFont, ImageOps
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


EMU_PER_INCH = 914400

PPTX_MIME_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

NSMAP = {"p": PML_NS, "r": R_NS}
CT_NSMAP = {"ct": CT_NS}
REL_NSMAP = {"rel": REL_NS}


FONT_SANS = Path("/System/Library/Fonts/SFNS.ttf")
FONT_MONO = Path("/System/Library/Fonts/SFNSMono.ttf")


_DASH_TRANSLATION = str.maketrans(
    {
        "\u2010": "-",  # hyphen
        "\u2011": "-",  # non-breaking hyphen
        "\u2012": "-",  # figure dash
        "\u2013": "-",  # en dash
        "\u2014": "-",  # em dash
        "\u2212": "-",  # minus sign
        "\u00A0": " ",  # non-breaking space
    }
)


def _norm_text(text: str) -> str:
    return text.translate(_DASH_TRANSLATION)


AssertAssetRole = Literal["structural", "semantic"]


@dataclass(frozen=True)
class AssertAssetRef:
    source_pptx: Path
    media_path: str
    role: AssertAssetRole
    derived_path: Path


@dataclass(frozen=True)
class PreparedAssets:
    # index 1..7 (index 0 is unused but filled with chapter 1 for convenience)
    bg_full_by_chapter: tuple[AssertAssetRef, ...]
    bg_overlay_by_chapter: tuple[AssertAssetRef, ...]
    card_by_chapter: tuple[AssertAssetRef, ...]
    wm_by_chapter: tuple[AssertAssetRef, ...]
    stamp_by_chapter: tuple[AssertAssetRef, ...]


@dataclass(frozen=True)
class SlideManifestEntry:
    slide_num: int
    title: str
    assets: list[AssertAssetRef]


@dataclass(frozen=True)
class BuildPaths:
    template_pptx: Path
    working_pptx: Path
    out_pptx: Path
    assets_dir: Path
    assert_dir: Path


def _emu(inches: float) -> int:
    return int(inches * EMU_PER_INCH)


def _read_zip_xml(z: zipfile.ZipFile, name: str) -> etree._Element:
    return etree.fromstring(z.read(name))


def _write_zip_xml(z: zipfile.ZipFile, name: str, root: etree._Element) -> None:
    z.writestr(name, etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True))


def _max_slide_number(z: zipfile.ZipFile) -> int:
    max_n = 0
    for name in z.namelist():
        if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
            continue
        stem = Path(name).stem
        try:
            n = int(stem.replace("slide", ""))
        except ValueError:
            continue
        max_n = max(max_n, n)
    return max_n


def _next_rid(pres_rels_root: etree._Element) -> str:
    rids: list[int] = []
    for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                rids.append(int(rid[3:]))
            except ValueError:
                continue
    nxt = (max(rids) + 1) if rids else 1
    return f"rId{nxt}"


def _next_slide_id(pres_root: etree._Element) -> int:
    sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
    if sld_id_lst is None:
        raise RuntimeError("presentation.xml missing p:sldIdLst")
    ids: list[int] = []
    for sld_id in sld_id_lst.findall("p:sldId", namespaces=NSMAP):
        try:
            ids.append(int(sld_id.get("id", "0")))
        except ValueError:
            continue
    return (max(ids) + 1) if ids else 256


def _remove_notes_relationship(slide_rels_root: etree._Element) -> None:
    for rel in list(slide_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP)):
        if rel.get("Type") == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide":
            slide_rels_root.remove(rel)


def _duplicate_slides(
    *,
    zin: zipfile.ZipFile,
    pres_root: etree._Element,
    pres_rels_root: etree._Element,
    ct_root: etree._Element,
    duplicates: dict[int, int],
) -> dict[int, list[int]]:
    """
    Duplicate template slides by copying their slide XML + rels into new slide part numbers.

    Returns: base_slide_num -> list[new_slide_nums]
    """
    max_slide = _max_slide_number(zin)
    sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
    if sld_id_lst is None:
        raise RuntimeError("presentation.xml missing p:sldIdLst")

    next_slide_id = _next_slide_id(pres_root)
    next_new = max_slide + 1

    created: dict[int, list[int]] = {k: [] for k in duplicates}
    for base_slide, count in duplicates.items():
        if count <= 0:
            continue
        for _ in range(count):
            n = next_new
            next_new += 1

            rid = _next_rid(pres_rels_root)
            rel_el = etree.Element(f"{{{REL_NS}}}Relationship")
            rel_el.set("Id", rid)
            rel_el.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel_el.set("Target", f"slides/slide{n}.xml")
            pres_rels_root.append(rel_el)

            sld_id_el = etree.Element(f"{{{PML_NS}}}sldId")
            sld_id_el.set("id", str(next_slide_id))
            sld_id_el.set(f"{{{R_NS}}}id", rid)
            sld_id_lst.append(sld_id_el)
            next_slide_id += 1

            override = etree.Element(f"{{{CT_NS}}}Override")
            override.set("PartName", f"/ppt/slides/slide{n}.xml")
            override.set("ContentType", PPTX_MIME_SLIDE)
            ct_root.append(override)

            created[base_slide].append(n)

    return created


def build_working_structure(*, paths: BuildPaths) -> list[int]:
    """
    Create Slide2_Working.pptx by duplicating high-animation template slides and reordering to a 54-slide recipe.

    Returns: recipe_base_slides (len == slide_count) for downstream content edits.
    """
    paths.working_pptx.parent.mkdir(parents=True, exist_ok=True)

    recipe: list[int] = [
        # Chapter 1 — Hook (1–6)
        1,
        2,
        16,
        16,
        16,
        16,
        # Chapter 2 — Research Questions (7–18)
        18,
        16,
        13,
        16,
        13,
        16,
        13,
        16,
        13,
        16,
        16,
        2,
        # Chapter 3 — Credibility / Method (19–26)
        5,
        4,
        3,
        13,
        16,
        13,
        16,
        2,
        # Chapter 4 — Discovery surfaces & catalog (27–38)
        6,
        7,
        13,
        8,
        13,
        13,
        9,
        13,
        10,
        13,
        14,
        2,
        # Chapter 5 — Panel & snapshots (39–49)
        11,
        13,
        16,
        12,
        13,
        13,
        13,
        16,
        13,
        15,
        15,
        # Chapter 6 — Reproducibility (50–51)
        16,
        17,
        # Chapter 7 — What this enables + appendix (52–54)
        14,
        16,
        16,
    ]
    if len(recipe) != 54:
        raise RuntimeError(f"internal error: expected recipe len=54, got {len(recipe)}")

    # Slide duplications needed to satisfy recipe counts.
    duplicates = {
        2: 3,   # total 4
        13: 15,  # total 16
        14: 1,  # total 2
        15: 1,  # total 2
        16: 16,  # total 17
    }

    with zipfile.ZipFile(paths.template_pptx, "r") as zin:
        pres_root = _read_zip_xml(zin, "ppt/presentation.xml")
        pres_rels_root = _read_zip_xml(zin, "ppt/_rels/presentation.xml.rels")
        ct_root = _read_zip_xml(zin, "[Content_Types].xml")

        created = _duplicate_slides(
            zin=zin,
            pres_root=pres_root,
            pres_rels_root=pres_rels_root,
            ct_root=ct_root,
            duplicates=duplicates,
        )

        # Copy slide XML + rels for all new slides.
        new_slide_bytes: dict[int, bytes] = {}
        new_slide_rels_bytes: dict[int, bytes] = {}
        for base_slide, new_nums in created.items():
            if not new_nums:
                continue
            slide_src = f"ppt/slides/slide{base_slide}.xml"
            slide_rels_src = f"ppt/slides/_rels/slide{base_slide}.xml.rels"
            slide_src_bytes = zin.read(slide_src)
            slide_rels_root = _read_zip_xml(zin, slide_rels_src)
            _remove_notes_relationship(slide_rels_root)
            slide_rels_bytes = etree.tostring(slide_rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)

            for n in new_nums:
                new_slide_bytes[n] = slide_src_bytes
                new_slide_rels_bytes[n] = slide_rels_bytes

        # Build a pool of available slide part numbers for each base slide in the recipe.
        pool: dict[int, list[int]] = {}
        for base in sorted(set(recipe)):
            pool[base] = [base] + list(created.get(base, []))

        desired_part_nums: list[int] = []
        for base in recipe:
            avail = pool.get(base)
            if not avail:
                raise RuntimeError(f"missing slide pool for base slide {base}")
            desired_part_nums.append(avail.pop(0))

        # Reorder sldIdLst to match desired_part_nums.
        sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
        if sld_id_lst is None:
            raise RuntimeError("presentation.xml missing p:sldIdLst")

        rid_to_target: dict[str, str] = {}
        for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
            if rel.get("Type") != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
                continue
            rid = rel.get("Id")
            tgt = rel.get("Target")
            if rid and tgt:
                rid_to_target[rid] = tgt

        target_to_sld_id: dict[str, etree._Element] = {}
        for sld_id in list(sld_id_lst.findall("p:sldId", namespaces=NSMAP)):
            rid = sld_id.get(f"{{{R_NS}}}id")
            tgt = rid_to_target.get(rid or "")
            if tgt:
                target_to_sld_id[tgt] = sld_id

        for child in list(sld_id_lst):
            sld_id_lst.remove(child)

        for n in desired_part_nums:
            target = f"slides/slide{n}.xml"
            el = target_to_sld_id.get(target)
            if el is None:
                raise RuntimeError(f"Missing slide target in presentation: {target}")
            sld_id_lst.append(deepcopy(el))

        tmp_path = paths.working_pptx.with_suffix(".tmp.pptx")
        if tmp_path.exists():
            tmp_path.unlink()

        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}:
                    continue
                zout.writestr(item, zin.read(item.filename))

            for n, b in new_slide_bytes.items():
                zout.writestr(f"ppt/slides/slide{n}.xml", b)
                zout.writestr(f"ppt/slides/_rels/slide{n}.xml.rels", new_slide_rels_bytes[n])

            _write_zip_xml(zout, "ppt/presentation.xml", pres_root)
            _write_zip_xml(zout, "ppt/_rels/presentation.xml.rels", pres_rels_root)
            _write_zip_xml(zout, "[Content_Types].xml", ct_root)

        shutil.move(tmp_path, paths.working_pptx)

    return recipe


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found on slide")


def _clear_text(shape) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""


def _set_run_text(shape, text: str) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    text = _norm_text(text)
    tf = shape.text_frame
    if not tf.paragraphs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]

    if p.runs:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""
    else:
        r = p.add_run()
        r.text = text

    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001
    blip.set(f"{{{R_NS}}}embed", rid)


def _fill_shape_with_picture(shape, slide, image_path: Path) -> None:
    # Keep the shape (and animations) but change its fill to a:blipFill.
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    sp = shape._element  # noqa: SLF001
    sp_pr = sp.xpath("./p:spPr")[0]

    fill_tags = {"noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}
    for child in list(sp_pr):
        local = etree.QName(child).localname
        if local in fill_tags:
            sp_pr.remove(child)

    blip_fill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), rid)
    blip_fill.append(blip)
    stretch = OxmlElement("a:stretch")
    stretch.append(OxmlElement("a:fillRect"))
    blip_fill.append(stretch)

    ln_els = sp_pr.xpath("./a:ln")
    insert_at = sp_pr.index(ln_els[0]) if ln_els else len(sp_pr)
    sp_pr.insert(insert_at, blip_fill)


def _move_shape_to_back(slide, shape, *, after_background: bool) -> None:
    sp_tree = slide.shapes._spTree  # noqa: SLF001
    insert_at = 3 if after_background else 2
    sp_tree.insert(insert_at, shape._element)  # noqa: SLF001


def _add_fullslide_picture(slide, prs: Presentation, image_path: Path) -> None:
    pic = slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    _move_shape_to_back(slide, pic, after_background=False)


def _add_card_panel(slide, *, image_path: Path, left_in: float, top_in: float, w_in: float, h_in: float) -> None:
    pic = slide.shapes.add_picture(
        str(image_path),
        _emu(left_in),
        _emu(top_in),
        width=_emu(w_in),
        height=_emu(h_in),
    )
    _move_shape_to_back(slide, pic, after_background=True)


def _add_icon(
    *,
    slide,
    image_path: Path,
    left_in: float,
    top_in: float,
    size_in: float,
    send_to_back: bool,
) -> None:
    pic = slide.shapes.add_picture(
        str(image_path),
        _emu(left_in),
        _emu(top_in),
        width=_emu(size_in),
        height=_emu(size_in),
    )
    if send_to_back:
        _move_shape_to_back(slide, pic, after_background=True)


def _extract_media(*, source_pptx: Path, media_path: str, out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    with zipfile.ZipFile(source_pptx, "r") as z:
        out_path.write_bytes(z.read(media_path))
    return out_path


def _apply_alpha_multiplier(img: Image.Image, alpha_mult: float) -> Image.Image:
    if alpha_mult >= 0.999:
        return img
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    r, g, b, a = img.split()
    a = a.point(lambda v: int(v * alpha_mult))
    return Image.merge("RGBA", (r, g, b, a))


def _prep_blueprint_overlay(*, src_path: Path, out_path: Path, tint_hex: str, alpha_max: float, size_px: tuple[int, int]) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    img = Image.open(src_path).convert("RGB")
    g = ImageOps.grayscale(img)
    g = ImageOps.invert(g)
    g = ImageOps.autocontrast(g, cutoff=2)
    g = ImageEnhance.Contrast(g).enhance(1.6)
    g = ImageEnhance.Brightness(g).enhance(0.95)
    g = g.resize(size_px, resample=Image.LANCZOS)

    tint_rgb = tuple(int(tint_hex[i : i + 2], 16) for i in (0, 2, 4))
    alpha_scale = max(0.0, min(1.0, alpha_max))
    a = g.point(lambda v: int(v * alpha_scale))
    overlay = Image.new("RGBA", size_px, tint_rgb + (0,))
    overlay.putalpha(a)
    overlay.save(out_path)
    return out_path


def _prep_blueprint_card_texture(
    *,
    src_path: Path,
    out_path: Path,
    base_hex: str,
    tint_hex: str,
    base_alpha: float,
    overlay_alpha: float,
    size_px: tuple[int, int],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path

    # Base: semi-transparent light card.
    base_rgb = tuple(int(base_hex[i : i + 2], 16) for i in (0, 2, 4))
    base_a = int(max(0.0, min(1.0, base_alpha)) * 255)
    base = Image.new("RGBA", size_px, base_rgb + (base_a,))

    overlay_tmp = out_path.with_name(out_path.stem + "_overlay_tmp.png")
    _prep_blueprint_overlay(src_path=src_path, out_path=overlay_tmp, tint_hex=tint_hex, alpha_max=overlay_alpha, size_px=size_px)
    overlay = Image.open(overlay_tmp).convert("RGBA")
    comp = Image.alpha_composite(base, overlay)

    # Add rounded corners + subtle drop shadow for depth on the dark board.
    w, h = size_px
    radius = max(22, int(min(w, h) * 0.04))

    mask = Image.new("L", size_px, 0)
    md = ImageDraw.Draw(mask)
    md.rounded_rectangle((0, 0, w - 1, h - 1), radius=radius, fill=255)

    r, g, b, a = comp.split()
    a = ImageChops.multiply(a, mask)
    comp = Image.merge("RGBA", (r, g, b, a))

    shadow = Image.new("RGBA", size_px, (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle((12, 14, w - 6, h - 6), radius=radius, fill=(0, 0, 0, 44))
    shadow = shadow.filter(ImageFilter.GaussianBlur(18))

    out = Image.alpha_composite(shadow, comp)
    out.save(out_path)
    try:
        overlay_tmp.unlink()
    except OSError:
        pass
    return out_path


def _prep_icon(*, src_path: Path, out_path: Path, size_px: int, alpha_mult: float) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    img = Image.open(src_path).convert("RGBA")
    bbox = img.split()[-1].getbbox()
    if bbox:
        img = img.crop(bbox)
    img = _apply_alpha_multiplier(img, alpha_mult=alpha_mult)

    w, h = img.size
    if w <= 0 or h <= 0:
        raise ValueError(f"empty icon image: {src_path}")
    scale = min(size_px / w, size_px / h)
    new_w = max(1, int(w * scale))
    new_h = max(1, int(h * scale))
    img = img.resize((new_w, new_h), resample=Image.LANCZOS)

    canvas = Image.new("RGBA", (size_px, size_px), (0, 0, 0, 0))
    canvas.paste(img, ((size_px - new_w) // 2, (size_px - new_h) // 2), img)
    canvas.save(out_path)
    return out_path


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    h = hex_color.strip().lstrip("#")
    if len(h) != 6:
        raise ValueError(f"Expected 6-char hex color, got: {hex_color}")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def _prep_space_base_background(
    *,
    src_path: Path,
    out_path: Path,
    tint_hex: str,
    size_px: tuple[int, int],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path

    base = Image.open(src_path).convert("RGB")
    base = ImageOps.fit(base, size_px, method=Image.LANCZOS)
    base = ImageEnhance.Brightness(base).enhance(0.72)
    base = ImageEnhance.Contrast(base).enhance(1.15)
    base = base.filter(ImageFilter.GaussianBlur(radius=1.2))

    tint = Image.new("RGBA", size_px, _hex_to_rgb(tint_hex) + (36,))
    out = Image.alpha_composite(base.convert("RGBA"), tint)
    out.save(out_path)
    return out_path


def _prep_space_overlay_layer(
    *,
    src_path: Path,
    out_path: Path,
    tint_hex: str,
    alpha_mult: float,
    size_px: tuple[int, int],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path

    overlay = Image.open(src_path).convert("RGBA")
    overlay = ImageOps.fit(overlay, size_px, method=Image.LANCZOS)

    rgb = overlay.convert("RGB")
    gray = ImageOps.grayscale(rgb)
    tint_img = ImageOps.colorize(gray, black="#0B1020", white=f"#{tint_hex}")
    alpha = overlay.split()[-1]
    tinted = Image.merge("RGBA", (*tint_img.split(), alpha))
    tinted = _apply_alpha_multiplier(tinted, alpha_mult=alpha_mult)
    tinted.save(out_path)
    return out_path


def _prep_space_card_texture(
    *,
    src_path: Path,
    out_path: Path,
    accent_hex: str,
    size_px: tuple[int, int],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path

    tex = Image.open(src_path).convert("RGB")
    tex = ImageOps.fit(tex, size_px, method=Image.LANCZOS)
    tex = ImageEnhance.Brightness(tex).enhance(1.08)
    tex = ImageEnhance.Contrast(tex).enhance(0.88)
    tex_rgba = _apply_alpha_multiplier(tex.convert("RGBA"), alpha_mult=0.24)

    base = Image.new("RGBA", size_px, (255, 255, 255, 232))
    comp = Image.alpha_composite(base, tex_rgba)
    comp = Image.alpha_composite(comp, Image.new("RGBA", size_px, _hex_to_rgb(accent_hex) + (14,)))

    w, h = size_px
    radius = max(20, int(min(w, h) * 0.035))
    mask = Image.new("L", size_px, 0)
    md = ImageDraw.Draw(mask)
    md.rounded_rectangle((0, 0, w - 1, h - 1), radius=radius, fill=255)

    r, g, b, a = comp.split()
    a = ImageChops.multiply(a, mask)
    comp = Image.merge("RGBA", (r, g, b, a))

    shadow = Image.new("RGBA", size_px, (0, 0, 0, 0))
    sd = ImageDraw.Draw(shadow)
    sd.rounded_rectangle((12, 14, w - 8, h - 8), radius=radius, fill=(0, 0, 0, 56))
    shadow = shadow.filter(ImageFilter.GaussianBlur(20))
    out = Image.alpha_composite(shadow, comp)

    od = ImageDraw.Draw(out)
    od.rounded_rectangle(
        (3, 3, w - 4, h - 4),
        radius=radius,
        outline=_hex_to_rgb(accent_hex) + (142,),
        width=6,
    )
    out.save(out_path)
    return out_path


def _prep_stamp_from_gif(
    *,
    src_path: Path,
    out_path: Path,
    tint_hex: str,
    size_px: int,
    alpha_mult: float,
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path

    gif = Image.open(src_path)
    try:
        gif.seek(0)
    except EOFError:
        pass
    frame = gif.convert("RGBA")

    bbox = frame.split()[-1].getbbox()
    if bbox:
        frame = frame.crop(bbox)

    frame = ImageOps.fit(frame, (size_px, size_px), method=Image.LANCZOS)
    gray = ImageOps.grayscale(frame.convert("RGB"))
    tint_img = ImageOps.colorize(gray, black="#111827", white=f"#{tint_hex}")
    alpha = frame.split()[-1]
    tinted = Image.merge("RGBA", (*tint_img.split(), alpha))
    tinted = _apply_alpha_multiplier(tinted, alpha_mult=alpha_mult)
    tinted.save(out_path)
    return out_path


def prepare_assert_assets(*, assert_dir: Path, cache_dir: Path) -> PreparedAssets:
    space_pptx = assert_dir / "Space Exploration Mission Pitch Deck by Slidesgo.pptx"
    data_pptx = assert_dir / "Data Visual by Slidesgo.pptx"
    multispace_pptx = assert_dir / "Multi-Space Theme for May _ by Slidesgo.pptx"

    base_bg_media = "ppt/media/image9.jpg"
    base_bg_src = _extract_media(
        source_pptx=space_pptx,
        media_path=base_bg_media,
        out_path=cache_dir / "raw_space_bg_image9.jpg",
    )

    chapter_accent: dict[int, str] = {
        1: "22D3EE",
        2: "A78BFA",
        3: "34D399",
        4: "FBBF24",
        5: "FB7185",
        6: "60A5FA",
        7: "A3E635",
    }
    chapter_overlay_media: dict[int, str] = {
        1: "ppt/media/image14.png",
        2: "ppt/media/image3.png",
        3: "ppt/media/image4.png",
        4: "ppt/media/image2.png",
        5: "ppt/media/image7.png",
        6: "ppt/media/image1.png",
        7: "ppt/media/image5.png",
    }
    chapter_card_media: dict[int, str] = {
        1: "ppt/media/image6.jpg",
        2: "ppt/media/image13.jpg",
        3: "ppt/media/image15.jpg",
        4: "ppt/media/image2.jpg",
        5: "ppt/media/image3.jpg",
        6: "ppt/media/image10.jpg",
        7: "ppt/media/image14.jpg",
    }
    chapter_wm_media: dict[int, str] = {
        1: "ppt/media/image14.gif",
        2: "ppt/media/image8.gif",
        3: "ppt/media/image10.gif",
        4: "ppt/media/image13.gif",
        5: "ppt/media/image9.gif",
        6: "ppt/media/image14.gif",
        7: "ppt/media/image8.gif",
    }

    bg_full: dict[int, AssertAssetRef] = {}
    bg_overlay: dict[int, AssertAssetRef] = {}
    cards: dict[int, AssertAssetRef] = {}
    watermarks: dict[int, AssertAssetRef] = {}
    stamps: dict[int, AssertAssetRef] = {}

    for ch in range(1, 8):
        accent = chapter_accent[ch]
        overlay_media = chapter_overlay_media[ch]
        card_media = chapter_card_media[ch]
        wm_media = chapter_wm_media[ch]

        overlay_src = _extract_media(
            source_pptx=space_pptx,
            media_path=overlay_media,
            out_path=cache_dir / f"raw_space_overlay_ch{ch}{Path(overlay_media).suffix}",
        )
        card_src = _extract_media(
            source_pptx=data_pptx,
            media_path=card_media,
            out_path=cache_dir / f"raw_data_card_ch{ch}{Path(card_media).suffix}",
        )
        wm_src = _extract_media(
            source_pptx=multispace_pptx,
            media_path=wm_media,
            out_path=cache_dir / f"wm_space_ch{ch}.gif",
        )

        bg_full_path = cache_dir / f"bg_space_full_ch{ch}.png"
        bg_overlay_path = cache_dir / f"bg_space_overlay_ch{ch}.png"
        card_path = cache_dir / f"card_space_glass_ch{ch}.png"
        stamp_path = cache_dir / f"stamp_space_ch{ch}.png"

        _prep_space_base_background(
            src_path=base_bg_src,
            out_path=bg_full_path,
            tint_hex=accent,
            size_px=(1920, 1080),
        )
        _prep_space_overlay_layer(
            src_path=overlay_src,
            out_path=bg_overlay_path,
            tint_hex=accent,
            alpha_mult=0.16,
            size_px=(1920, 1080),
        )
        _prep_space_card_texture(
            src_path=card_src,
            out_path=card_path,
            accent_hex=accent,
            size_px=(1900, 980),
        )
        _prep_stamp_from_gif(
            src_path=wm_src,
            out_path=stamp_path,
            tint_hex=accent,
            size_px=768,
            alpha_mult=0.42,
        )

        bg_full[ch] = AssertAssetRef(
            source_pptx=space_pptx,
            media_path=base_bg_media,
            role="structural",
            derived_path=bg_full_path,
        )
        bg_overlay[ch] = AssertAssetRef(
            source_pptx=space_pptx,
            media_path=overlay_media,
            role="structural",
            derived_path=bg_overlay_path,
        )
        cards[ch] = AssertAssetRef(
            source_pptx=data_pptx,
            media_path=card_media,
            role="structural",
            derived_path=card_path,
        )
        watermarks[ch] = AssertAssetRef(
            source_pptx=multispace_pptx,
            media_path=wm_media,
            role="semantic",
            derived_path=wm_src,
        )
        stamps[ch] = AssertAssetRef(
            source_pptx=multispace_pptx,
            media_path=wm_media,
            role="semantic",
            derived_path=stamp_path,
        )

    # index 0 is unused; fill with chapter 1 as safe fallback
    bg_full_by_chapter: list[AssertAssetRef] = [bg_full[1]] * 8
    bg_overlay_by_chapter: list[AssertAssetRef] = [bg_overlay[1]] * 8
    card_by_chapter: list[AssertAssetRef] = [cards[1]] * 8
    wm_by_chapter: list[AssertAssetRef] = [watermarks[1]] * 8
    stamp_by_chapter: list[AssertAssetRef] = [stamps[1]] * 8
    for ch in range(1, 8):
        bg_full_by_chapter[ch] = bg_full[ch]
        bg_overlay_by_chapter[ch] = bg_overlay[ch]
        card_by_chapter[ch] = cards[ch]
        wm_by_chapter[ch] = watermarks[ch]
        stamp_by_chapter[ch] = stamps[ch]

    return PreparedAssets(
        bg_full_by_chapter=tuple(bg_full_by_chapter),
        bg_overlay_by_chapter=tuple(bg_overlay_by_chapter),
        card_by_chapter=tuple(card_by_chapter),
        wm_by_chapter=tuple(wm_by_chapter),
        stamp_by_chapter=tuple(stamp_by_chapter),
    )


def _record_asset(manifest: dict[int, list[AssertAssetRef]], slide_num: int, asset: AssertAssetRef) -> None:
    items = manifest.setdefault(slide_num, [])
    key = (str(asset.source_pptx), asset.media_path, asset.role, str(asset.derived_path))
    for a in items:
        if (str(a.source_pptx), a.media_path, a.role, str(a.derived_path)) == key:
            return
    items.append(asset)


def _chapter_backgrounds(assets: PreparedAssets, chapter: int) -> tuple[AssertAssetRef, AssertAssetRef, AssertAssetRef]:
    ch = chapter
    if ch < 1 or ch > 7:
        ch = 1
    full = assets.bg_full_by_chapter[ch]
    overlay = assets.bg_overlay_by_chapter[ch]
    card = assets.card_by_chapter[ch]
    return full, overlay, card


def _chapter_icons(assets: PreparedAssets, chapter: int) -> tuple[AssertAssetRef, AssertAssetRef]:
    ch = chapter
    if ch < 1 or ch > 7:
        ch = 1
    return assets.wm_by_chapter[ch], assets.stamp_by_chapter[ch]


def _add_common_visuals(
    *,
    slide,
    prs: Presentation,
    assets: PreparedAssets,
    slide_num: int,
    chapter: int,
    asset_manifest: dict[int, list[AssertAssetRef]],
) -> None:
    full, overlay, card = _chapter_backgrounds(assets, chapter)

    # Full-slide space base background (structural).
    _add_fullslide_picture(slide, prs, full.derived_path)
    _record_asset(asset_manifest, slide_num, full)

    # Chapter overlay on top of base background (structural).
    _add_fullslide_picture(slide, prs, overlay.derived_path)
    _record_asset(asset_manifest, slide_num, overlay)

    # Glass card texture behind content (structural).
    _add_card_panel(slide, image_path=card.derived_path, left_in=0.65, top_in=0.85, w_in=12.15, h_in=6.00)
    _record_asset(asset_manifest, slide_num, card)

    wm, stamp = _chapter_icons(assets, chapter)
    # Animated (or static fallback) watermark bottom-right (semantic).
    _add_icon(slide=slide, image_path=wm.derived_path, left_in=11.35, top_in=5.25, size_in=1.70, send_to_back=True)
    _record_asset(asset_manifest, slide_num, wm)

    # Small chapter stamp near title (semantic).
    _add_icon(slide=slide, image_path=stamp.derived_path, left_in=12.30, top_in=0.32, size_in=0.50, send_to_back=False)
    _record_asset(asset_manifest, slide_num, stamp)


@dataclass(frozen=True)
class SlideEdits:
    base: int
    chapter: int
    text: dict[int, str]
    # picture shape_id -> image path
    pictures: dict[int, Path]
    # auto-shape shape_id -> image fill
    fills: dict[int, Path]
    # optional extra note text boxes
    extra_notes: list[str]


def _receipt(paths: dict[str, Path], key: str) -> Path:
    p = paths.get(key)
    if p is None:
        raise KeyError(f"missing receipt: {key}")
    return p


def _build_slide_edits(*, run_date: str, run_id: str, receipts: dict[str, Path]) -> list[SlideEdits]:
    # Chapters: 1–7
    r = receipts

    def rec(k: str) -> Path:
        return _receipt(r, k)

    slide_edits: list[SlideEdits] = []

    # 01 Title (base 1)
    slide_edits.append(
        SlideEdits(
            base=1,
            chapter=1,
            text={
                2: "Bluesky feed marketplace",
                3: "Research questions + reproducible data collection (no results yet)",
                4: f"Fresh run: {run_date} (UTC) • receipts + manifest included",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 02 Thesis (base 2)
    slide_edits.append(
        SlideEdits(
            base=2,
            chapter=1,
            text={
                2: "Algorithmic choice ≠ decentralized power",
                3: "• Feeds + labelers turn ranking/moderation into a marketplace.",
                4: "• Discovery + hosting create chokepoints (and targets).",
                5: "• We built a reproducible, read-only measurement spine.",
                6: "Thesis: pluralism shifts power to discovery + infrastructure.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 03–06 Hook writing slides (base 16 x4)
    slide_edits += [
        SlideEdits(
            base=16,
            chapter=1,
            text={
                2: "Feeds + labelers = marketplace services",
                3: "Ranking and moderation become modular (third-party) services.",
                5: "Claim: ranking/moderation become services you can choose.",
                6: "So power shifts to discovery + hosting layers.",
                8: "Ranking services",
                10: "Moderation services",
                12: "Distribution services",
                14: "Why this matters",
                15: "Marketplace dynamics create new chokepoints and targets.",
                17: "Concrete artifacts (exports)",
                18: "• Feeds: algorithm + curation objective",
                19: "• Labelers: safety schema + enforcement signal",
                20: "• Starter packs: discovery bundling",
                21: "• Hosting/provider: availability + control",
                22: "• AppView: default discovery + UX",
                23: "• Relay: enumeration surface",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
        SlideEdits(
            base=16,
            chapter=1,
            text={
                2: "New choke points, new attack surfaces",
                3: "Pluralism ≠ no power; it relocates.",
                5: "Chokepoint: discovery bundling (starter packs + defaults).",
                6: "Attack surface: hosting/provider + ranking supply chain.",
                8: "Concentration",
                10: "Gaming",
                12: "Reliability",
                14: "Questions that follow",
                15: "We define RQs first, then show what data exists (no results).",
                17: "Examples (no claims)",
                18: "• Do a few feeds dominate discovery?",
                19: "• Can actors manipulate inclusion/visibility?",
                20: "• Which providers host the default experience?",
                21: "• What fails under rate limits / 5xx?",
                22: "• Where do safety signals differ by feed?",
                23: "• What privacy risks arise at scale?",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
        SlideEdits(
            base=16,
            chapter=1,
            text={
                2: "Guardrails for this deck",
                3: "Only RQs, method, and file receipts.",
                5: "No results today: no H1–H6 values, no charts, no concentration numbers.",
                6: "We only show parameters + existence of artifacts (sample rows).",
                8: "Allowed",
                10: "Not allowed",
                12: "Why",
                14: "Credibility first",
                15: "We want confidence the pipeline is real before any claims.",
                17: "Deck content policy",
                18: "• Run parameters (auth_mode, targets)",
                19: "• Run window + run_id",
                20: "• Sample rows (no post text)",
                21: "• Manifest + data dictionary",
                22: "• Derived metrics exist but hidden",
                23: "• Goal: auditability before analysis",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
        SlideEdits(
            base=16,
            chapter=1,
            text={
                2: "Story today (micro-steps)",
                3: "Many small reveals; smooth click-to-reveal.",
                5: "We walk from questions → method → artifacts → joins.",
                6: "Every slide is backed by a file receipt from the run folder.",
                8: "Arc",
                10: "Proof",
                12: "Keys",
                14: "End state",
                15: "A joinable snapshot: panel → impressions → posts → authors.",
                17: "Chapters",
                18: "• RQs (what we’d ask next)",
                19: "• Read-only collection (credibility)",
                20: "• What artifacts exist (receipts)",
                21: "• Panel + snapshots (ranked spine)",
                22: "• Join keys line up",
                23: "• Reproducibility (manifest + zip)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
    ]

    # 07 RQ overview (base 18)
    slide_edits.append(
        SlideEdits(
            base=18,
            chapter=2,
            text={
                2: "Research questions (next steps)",
                3: "S&P framing: marketplace dynamics + new attack surfaces.",
                5: "RQ1",
                6: "Discovery chokepoints via starter packs (bundling leverage).",
                8: "RQ2",
                9: "Provider leverage / feed hosting supply chain.",
                11: "RQ3",
                12: "Exposure concentration + overlap (“same winners”).",
                14: "RQ4",
                15: "Safety risk variability by feed (labels).",
                17: "RQ5 (future)",
                18: "Malicious/growth-gaming feeds + privacy risk.",
                19: "Next: how we collected.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 08 RQ1 statement (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ1 — Discovery chokepoints",
                3: "Starter-pack inclusion as discovery leverage.",
                5: "Do a small number of feeds dominate starter-pack inclusions?",
                6: "How stable is inclusion across time/runs (future)?",
                8: "Threat model",
                10: "What we’d measure",
                12: "What data exists",
                14: "No claims",
                15: "We’re only defining questions + showing the artifacts.",
                17: "Files (names only)",
                18: "• starterpacks.csv",
                19: "• starterpack_feeds.csv",
                20: "• discovery_feed_inclusions.csv",
                21: "• popular_feeds.csv",
                22: "• feed_panel.csv",
                23: "• manifest.csv (integrity)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 09 RQ1 receipt (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=2,
            text={
                2: "Receipt: starterpacks.csv",
                4: "What this proves",
                5: "Starter packs exist as concrete records.",
                6: "We store pack metadata (no interpretation).",
                7: "This slide shows one sanitized row excerpt.",
                9: "Where it lives",
                10: "02_csv_exports/starterpacks.csv",
                11: "02_csv_exports/starterpack_feeds.csv",
                13: "Related (discovery)",
                14: "• discovery_feed_inclusions.csv",
                15: "• feed_panel.csv",
                16: "• popular_feeds.csv",
                17: "• manifest.csv",
                18: "• data_dictionary.csv",
                19: "• validation_report.csv",
                21: "(sample row; truncated; no post text)",
            },
            pictures={20: rec("starterpack_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 10 RQ2 statement (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ2 — Provider leverage",
                3: "Feed hosting supply chain as hidden power.",
                5: "Which provider buckets host the most visible feeds?",
                6: "Do defaults depend on a small set of domains/providers?",
                8: "Where power sits",
                10: "What we’d compute",
                12: "What data exists",
                14: "No results",
                15: "We only show that provider_bucket is recorded and joinable.",
                17: "Files (names only)",
                18: "• feed_generators_index.csv (provider_bucket)",
                19: "• feed_panel.csv (group + provider_bucket)",
                20: "• feed_items.csv.gz (exposure spine)",
                21: "• posts.csv.gz (metadata)",
                22: "• authors.csv.gz (profiles)",
                23: "• postprocess/h2_provider_leverage.csv (stored; not shown)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 11 RQ2 receipt (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=2,
            text={
                2: "Receipt: feed_generators_index.csv",
                4: "What this proves",
                5: "Feed generator metadata is recorded (catalog).",
                6: "provider_bucket is captured for hosting analysis (later).",
                7: "This slide shows one row excerpt (truncated).",
                9: "Where it lives",
                10: "02_csv_exports/feed_generators_index.csv",
                11: "05_manifest/data_dictionary.csv",
                13: "Downstream joins",
                14: "• feed_panel.csv includes feed_uri",
                15: "• feed_items.csv.gz includes feed_uri + viewer_mode",
                16: "• posts.csv.gz includes post_uri + post_cid",
                17: "• authors.csv.gz includes author_did",
                18: "• manifest.csv (sha256)",
                19: "• validation_report.csv (PASS gate)",
                21: "(example row; provider_bucket visible; no claims)",
            },
            pictures={20: rec("feed_index_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 12 RQ3 statement (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ3 — Exposure concentration + overlap",
                3: "Do the same winners appear everywhere?",
                5: "How concentrated are ranked impressions across the panel?",
                6: "How similar are feeds in what they show?",
                8: "Spine",
                10: "Joins",
                12: "Metrics later",
                14: "No results",
                15: "We only show the ranked impression spine + join keys.",
                17: "Files (names only)",
                18: "• feed_items.csv.gz (ranked impressions)",
                19: "• posts.csv.gz (metadata; no text)",
                20: "• authors.csv.gz (profiles)",
                21: "• feed_panel.csv (selection)",
                22: "• postprocess/h3_* + h4_* (stored; not shown)",
                23: "• state.db (resumable truth)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 13 RQ3 receipt (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=2,
            text={
                2: "Receipt: feed_items.csv.gz",
                4: "What this proves",
                5: "We captured a ranked impression spine (not screenshots).",
                6: "Each row is keyed and joinable (feed_uri + post_uri + cid).",
                7: "This slide shows rank=1 example row.",
                9: "Where it lives",
                10: "02_csv_exports/feed_items.csv.gz",
                11: "02_csv_exports/feed_snapshot_status.csv",
                13: "Join keys",
                14: "• feed_uri + viewer_mode + rank",
                15: "• post_uri + post_cid",
                16: "• author_did",
                17: "• post_labels attach on same keys",
                18: "• posts.csv.gz matches on (post_uri, post_cid)",
                19: "• authors.csv.gz matches on author_did",
                21: "(ranked impression excerpt; no analysis)",
            },
            pictures={20: rec("impression_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 14 RQ4 statement (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ4 — Safety signal variability",
                3: "Labels differ by feed and viewer context.",
                5: "Do safety labels vary across feeds for the same post?",
                6: "How do label sources/values map to risks (future)?",
                8: "Normalized labels",
                10: "Attachment",
                12: "Limits",
                14: "No claims",
                15: "We only show normalized label rows (excerpt).",
                17: "Files (names only)",
                18: "• post_labels.csv.gz",
                19: "• feed_items.csv.gz (spine)",
                20: "• feed_uri + viewer_mode + post_uri + post_cid",
                21: "• labels are shown as short list only",
                22: "• no prevalence / risk claims today",
                23: "• no post text shown",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 15 RQ4 receipt (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=2,
            text={
                2: "Receipt: post_labels.csv.gz",
                4: "What this proves",
                5: "Labels are stored as normalized rows (not anecdotes).",
                6: "They attach to impressions by explicit join keys.",
                7: "This slide shows a short labels excerpt (truncated).",
                9: "Where it lives",
                10: "02_csv_exports/post_labels.csv.gz",
                11: "05_manifest/data_dictionary.csv",
                13: "Attachment keys",
                14: "• feed_uri",
                15: "• viewer_mode",
                16: "• post_uri",
                17: "• post_cid",
                18: "• collected_at_utc",
                19: "• label_src + label_val (+ label_neg)",
                21: "(labels excerpt only; no prevalence numbers)",
            },
            pictures={20: rec("post_labels_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 16–17 RQ5 future (base 16 x2)
    slide_edits += [
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ5 (future) — Malicious / growth-gaming feeds",
                3: "Adversarial measurement needs guardrails.",
                5: "How do we detect feeds that manipulate exposure or harvest data?",
                6: "What privacy risks emerge from snapshotting at scale?",
                8: "Signals we’d need",
                10: "Mitigations",
                12: "Ethics",
                14: "Future work",
                15: "We are not running these experiments in this deck.",
                17: "Examples",
                18: "• repeated runs (stability over time)",
                19: "• anomalies in inclusion / rank shifts",
                20: "• provider / domain changes",
                21: "• opt-outs / rate limits / safe release",
                22: "• de-identification and aggregation",
                23: "• avoid publishing user content",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
        SlideEdits(
            base=16,
            chapter=2,
            text={
                2: "RQ5 — Inputs required (future)",
                3: "Not doing it today; defining inputs only.",
                5: "We need repeated runs + provenance to detect gaming.",
                6: "We need privacy-preserving aggregation before publishing.",
                8: "Data",
                10: "Process",
                12: "Outputs",
                14: "Why it’s hard",
                15: "Adversaries adapt; measurement must be cautious.",
                17: "Run receipts enable",
                18: "• run_metadata.csv across dates",
                19: "• manifest.csv (integrity)",
                20: "• state.db (resumability)",
                21: "• snapshot deltas (not shown)",
                22: "• controlled release policies",
                23: "• ethics review for privacy",
            },
            pictures={},
            fills={},
            extra_notes=[],
        ),
    ]

    # 18 Transition (base 2)
    slide_edits.append(
        SlideEdits(
            base=2,
            chapter=2,
            text={
                2: "How we collected (read-only + reproducible)",
                3: "• Relay + AppView reads only (GET).",
                4: "• One state DB: resumable + auditable.",
                5: "• Exports are hashed + documented.",
                6: "Next: method anatomy, then receipts from the run folder.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 19 Read-only XRPC (base 5)
    slide_edits.append(
        SlideEdits(
            base=5,
            chapter=3,
            text={
                2: "Read-only XRPC collection (credibility)",
                3: "Relay + AppView reads only; no posting/likes/follows/labels.",
                6: "GET /xrpc/app.bsky.feed.getFeed?feed=at://...&limit=...&cursor=...",
                9: "“Which service am I asking?”\n• AppView (public): feeds, posts, profiles\n• Relay: enumeration index",
                12: "A stable API name describing what you fetch (names only on slides).",
                13: "Parameters control paging + filtering (e.g., feed URI, limit, cursor). Auth-only POSTs exist only for sessions (not used here).",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 20 Hosts + endpoints (base 4)
    slide_edits.append(
        SlideEdits(
            base=4,
            chapter=3,
            text={
                2: "Where requests go (hosts)",
                3: "Relay for enumeration; AppView for reads; PDS only for auth context.",
                8: "Bulk enumeration\n(actor DIDs)",
                11: "Read APIs\n(feeds, profiles)",
                14: "Session APIs\n(auth context)",
                18: "Footnote: endpoint names only (see README_RUN.md for the full list).",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 21 Pipeline overview (base 3)
    slide_edits.append(
        SlideEdits(
            base=3,
            chapter=3,
            text={
                2: "7-stage pipeline overview",
                3: "One run turns public APIs into analysis-ready datasets.",
                38: "Next: we show the artifacts this run produced (as receipts), and how keys join.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 22 Receipt: run_metadata (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=3,
            text={
                2: "Receipt: run_metadata.csv",
                4: "What this proves",
                5: f"Fresh run_id is recorded: {run_id[:8]}...{run_id[-6:]}",
                6: "Run window + parameters are captured in a single row.",
                7: "This is method metadata (not results).",
                9: "Where it lives",
                10: "05_manifest/run_metadata.csv",
                11: "05_manifest/run_summary.csv",
                13: "Parameters shown",
                14: "• auth_mode",
                15: "• posts_per_feed",
                16: "• n_discovery / n_popular / n_less_known",
                17: "• appview_host / relay_host",
                18: "• rps / max_retries",
                19: "• started_at_utc / finished_at_utc",
                21: "(receipt card; truncated)",
            },
            pictures={20: rec("run_metadata_receipt")},
            fills={},
            extra_notes=[],
        )
    )

    # 23 API surface (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=3,
            text={
                2: "API surface (names only)",
                3: "What this collector calls (no URLs on slides).",
                5: "Relay: com.atproto.sync.listReposByCollection",
                6: "AppView: getActorFeeds, getFeedGenerators, starterpacks, popular, getFeed, getProfiles",
                8: "Discovery",
                10: "Ranking",
                12: "Hydration",
                14: "Auth-only",
                15: "Session create/refresh exists only if auth_mode enabled (not used).",
                17: "AppView endpoints",
                18: "• app.bsky.feed.getActorFeeds",
                19: "• app.bsky.feed.getFeedGenerators",
                20: "• app.bsky.graph.getActorStarterPacks",
                21: "• app.bsky.graph.getStarterPack",
                22: "• app.bsky.unspecced.getPopularFeedGenerators",
                23: "• app.bsky.feed.getFeed + app.bsky.actor.getProfiles",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 24 Receipt: validation_report (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=3,
            text={
                2: "Receipt: validation_report.csv",
                4: "What this proves",
                5: "We gate the run with explicit PASS/FAIL checks.",
                6: "This run’s checks are all PASS (no results shown).",
                7: "This slide shows an excerpt of the report.",
                9: "Where it lives",
                10: "05_manifest/validation_report.csv",
                11: "05_manifest/run_summary.csv",
                13: "What gets checked",
                14: "• discovery/popular coverage",
                15: "• panel coverage + diversity signal",
                16: "• snapshot success rate threshold",
                17: "• author hydration rate threshold",
                18: "• feed index non-empty",
                19: "• starterpack non-empty",
                21: "(PASS excerpt; no metric values on slides)",
            },
            pictures={20: rec("validation_report_receipt")},
            fills={},
            extra_notes=[],
        )
    )

    # 25 Audit trail (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=3,
            text={
                2: "Audit trail",
                3: "Reproducible + auditable outputs (no results).",
                5: "Everything is stored: state.db + logs + hashed exports.",
                6: "We can re-run, diff, and join without screenshots.",
                8: "Integrity",
                10: "Schema",
                12: "Portability",
                14: "Receipts are built from files",
                15: "Every receipt PNG is rendered from the run folder outputs.",
                17: "Artifacts",
                18: "• manifest.csv (sha256 list)",
                19: "• data_dictionary.csv (column meanings)",
                20: "• validation_report.csv (PASS gate)",
                21: "• state.db (single source of truth)",
                22: "• logs/run.log (request trace)",
                23: "• archive zip (portable bundle)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 26 Transition (base 2)
    slide_edits.append(
        SlideEdits(
            base=2,
            chapter=3,
            text={
                2: "Now: discovery → panel → snapshots (with receipts)",
                3: "• Discovery surfaces: starter packs + popular list.",
                4: "• Panel rules: deterministic selection + provider balance.",
                5: "• Snapshot spine: ranked impressions → posts/authors.",
                6: "Next: we show the discovery artifacts first.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 27 Stage 1 (base 6) + fill output with feed_index
    slide_edits.append(
        SlideEdits(
            base=6,
            chapter=4,
            text={
                2: "Stage 1 — Relay scan",
                4: "• Goal: enumerate candidate actors that publish a record type.",
            },
            pictures={},
            fills={12: rec("feed_index_row")},
            extra_notes=[],
        )
    )

    # 28 Stage 2 (base 7) + fill
    slide_edits.append(
        SlideEdits(
            base=7,
            chapter=4,
            text={
                2: "Stage 2 — Build a feed index",
                4: "• Goal: map each feed generator to metadata (creator, name, provider).",
            },
            pictures={},
            fills={12: rec("feed_index_row")},
            extra_notes=[],
        )
    )

    # 29 Receipt: feed index row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=4,
            text={
                2: "Receipt: feed_generators_index.csv",
                4: "What this proves",
                5: "We can point to a concrete catalog row.",
                6: "This is the basis for provider_bucket joins later.",
                7: "This slide shows one row excerpt.",
                9: "Where it lives",
                10: "02_csv_exports/feed_generators_index.csv",
                11: "05_manifest/data_dictionary.csv",
                13: "Used by",
                14: "• feed_panel.csv (selection)",
                15: "• postprocess/feeds_flat.csv (join convenience)",
                16: "• provider bucket analysis (later)",
                17: "• snapshot targeting",
                18: "• integrity checks (manifest.csv)",
                19: "• reproducibility (state.db)",
                21: "(example row; truncated)",
            },
            pictures={20: rec("feed_index_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 30 Stage 3 starter packs (base 8) + fill
    slide_edits.append(
        SlideEdits(
            base=8,
            chapter=4,
            text={
                2: "Stage 3 — Starter packs (discovery)",
                4: "• Goal: observe discovery via starter-pack inclusions.",
            },
            pictures={},
            fills={12: rec("starterpack_inclusion_row")},
            extra_notes=[],
        )
    )

    # 31 Receipt: starterpack inclusion (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=4,
            text={
                2: "Receipt: starterpack_feeds.csv",
                4: "What this proves",
                5: "Starter packs link to specific feed URIs (inclusions).",
                6: "Those inclusions become discovery signals (later).",
                7: "This slide shows one inclusion excerpt.",
                9: "Where it lives",
                10: "02_csv_exports/starterpack_feeds.csv",
                11: "02_csv_exports/discovery_feed_inclusions.csv",
                13: "Downstream",
                14: "• panel selection rules",
                15: "• overlap with popular list",
                16: "• exposure spine (feed_items.csv.gz)",
                17: "• provider buckets (feed index)",
                18: "• integrity (manifest.csv)",
                19: "• schema (data_dictionary.csv)",
                21: "(inclusion excerpt; truncated)",
            },
            pictures={20: rec("starterpack_inclusion_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 32 Receipt: starterpack row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=4,
            text={
                2: "Receipt: starterpacks.csv",
                4: "What this proves",
                5: "Starter packs are concrete records with creator + name + URI.",
                6: "We show a sanitized example row (truncated).",
                7: "No claims; just existence + joinability.",
                9: "Where it lives",
                10: "02_csv_exports/starterpacks.csv",
                11: "02_csv_exports/starterpack_feeds.csv",
                13: "Used by",
                14: "• discovery surfaced candidates",
                15: "• panel selection",
                16: "• later concentration questions",
                17: "• integrity (manifest.csv)",
                18: "• schema (data_dictionary.csv)",
                19: "• validations (validation_report.csv)",
                21: "(sample row; truncated)",
            },
            pictures={20: rec("starterpack_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 33 Stage 4 popular (base 9) + fill
    slide_edits.append(
        SlideEdits(
            base=9,
            chapter=4,
            text={
                2: "Stage 4 — Popular feeds (ranking surface)",
                4: "• Goal: capture the platform’s popular feed generators list.",
            },
            pictures={},
            fills={12: rec("popular_row")},
            extra_notes=[],
        )
    )

    # 34 Receipt: popular row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=4,
            text={
                2: "Receipt: popular_feeds.csv",
                4: "What this proves",
                5: "We captured a ranked “popular feeds” surface (inputs).",
                6: "This slide shows the first row (rank + URI).",
                7: "No claims; just a receipt that it exists.",
                9: "Where it lives",
                10: "02_csv_exports/popular_feeds.csv",
                11: "02_csv_exports/feed_panel.csv",
                13: "Used by",
                14: "• panel selection rules",
                15: "• snapshot targeting",
                16: "• later exposure analyses (not shown)",
                17: "• integrity (manifest.csv)",
                18: "• schema (data_dictionary.csv)",
                19: "• validations (validation_report.csv)",
                21: "(first row excerpt; truncated)",
            },
            pictures={20: rec("popular_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 35 Stage 5 enrich (base 10) + fill with feed index
    slide_edits.append(
        SlideEdits(
            base=10,
            chapter=4,
            text={
                2: "Stage 5 — Enrich missing metadata",
                4: "• Goal: batch hydrate missing feed metadata (provider bucket).",
            },
            pictures={},
            fills={12: rec("feed_index_row")},
            extra_notes=[],
        )
    )

    # 36 Receipt: run_summary (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=4,
            text={
                2: "Receipt: run_summary.csv",
                4: "What this proves",
                5: "We record high-level coverage (dataset size) for the run.",
                6: "Counts shown are “what exists”, not interpretation.",
                7: "This slide shows a coverage excerpt.",
                9: "Where it lives",
                10: "05_manifest/run_summary.csv",
                11: "05_manifest/run_metadata.csv",
                13: "Coverage fields",
                14: "• indexed feeds",
                15: "• starter packs + unique feeds",
                16: "• popular list size",
                17: "• panel size + snapshot success count",
                18: "• feed_items + unique posts/authors",
                19: "• mapping_notes (endpoints)",
                21: "(coverage excerpt; no analysis)",
            },
            pictures={20: rec("run_summary_receipt")},
            fills={},
            extra_notes=[],
        )
    )

    # 37 Montage discovery artifacts (base 14)
    slide_edits.append(
        SlideEdits(
            base=14,
            chapter=4,
            text={
                2: "Discovery leaves artifacts",
                3: "We can point to concrete rows, not anecdotes.",
                5: "Starter packs → included feeds",
                9: "Popular list (ranked)",
                11: "These are inputs to the panel selection rules.",
            },
            pictures={
                6: rec("starterpack_inclusion_row"),
                7: rec("starterpack_row"),
                10: rec("popular_row"),
            },
            fills={},
            extra_notes=[],
        )
    )

    # 38 Transition (base 2)
    slide_edits.append(
        SlideEdits(
            base=2,
            chapter=4,
            text={
                2: "Panel & snapshots: the ranked impressions spine",
                3: "• Panel: deterministic selection rules (targets only).",
                4: "• Snapshot: ranked feed_items are the impression spine.",
                5: "• Joins: panel → impressions → posts → authors.",
                6: "Next: panel receipt, then impression/post/labels/authors receipts.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 39 Stage 6 panel (base 11) + fill
    slide_edits.append(
        SlideEdits(
            base=11,
            chapter=5,
            text={
                2: "Stage 6 — Build the feed panel",
                4: "• Goal: choose a fixed set of feeds we will snapshot.",
            },
            pictures={},
            fills={12: rec("panel_row")},
            extra_notes=[],
        )
    )

    # 40 Receipt: panel row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=5,
            text={
                2: "Receipt: feed_panel.csv",
                4: "What this proves",
                5: "Panel selection is explicit (group + reason + provider_bucket).",
                6: "This slide shows a sample panel row excerpt.",
                7: "No claims about feeds; just the selection artifact.",
                9: "Where it lives",
                10: "02_csv_exports/feed_panel.csv",
                11: "02_csv_exports/feed_snapshot_status.csv",
                13: "Downstream",
                14: "• snapshot feeds via getFeed",
                15: "• join to feed_items on feed_uri",
                16: "• join to feed index on feed_uri",
                17: "• provider balance in less_known group",
                18: "• integrity (manifest.csv)",
                19: "• schema (data_dictionary.csv)",
                21: "(panel row excerpt; truncated)",
            },
            pictures={20: rec("panel_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 41 Panel rules (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=5,
            text={
                2: "Panel rules (selection)",
                3: "How 15 feeds were chosen (targets only).",
                5: "Groups: discovery_surfaced + popular + less_known.",
                6: "less_known is provider-balanced sampling from the remaining index.",
                8: "Discovery surfaced",
                10: "Popular",
                12: "Less-known",
                14: "Deterministic",
                15: "Seeded by run_id (repeatable).",
                17: "Outputs",
                18: "• feed_panel.csv",
                19: "• feed_snapshot_status.csv",
                20: "• feed_items.csv.gz",
                21: "• posts.csv.gz",
                22: "• post_labels.csv.gz",
                23: "• authors.csv.gz",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 42 Stage 7 snapshot (base 12) + fill
    slide_edits.append(
        SlideEdits(
            base=12,
            chapter=5,
            text={
                2: "Stage 7 — Snapshot each feed",
                4: "• Goal: fetch ranked items for each feed (unauth + auth when available).",
            },
            pictures={},
            fills={12: rec("impression_row")},
            extra_notes=[],
        )
    )

    # 43 Receipt: impression row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=5,
            text={
                2: "Receipt: feed_items.csv.gz",
                4: "What this proves",
                5: "Ranked impressions are captured as rows (joinable spine).",
                6: "This slide shows rank=1 excerpt for one feed snapshot.",
                7: "No results; just a receipt of the spine.",
                9: "Where it lives",
                10: "02_csv_exports/feed_items.csv.gz",
                11: "02_csv_exports/posts.csv.gz",
                13: "Keys",
                14: "• feed_uri + viewer_mode + rank",
                15: "• post_uri + post_cid",
                16: "• author_did",
                17: "• collected_at_utc",
                18: "• join to panel on feed_uri",
                19: "• join to labels on (feed_uri, viewer_mode, post_uri, post_cid)",
                21: "(ranked row excerpt; truncated)",
            },
            pictures={20: rec("impression_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 44 Receipt: post row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=5,
            text={
                2: "Receipt: posts.csv.gz",
                4: "What this proves",
                5: "We store post metadata for observed impressions.",
                6: "This slide shows metadata only (no post text).",
                7: "Post identity is (post_uri, post_cid).",
                9: "Where it lives",
                10: "02_csv_exports/posts.csv.gz",
                11: "02_csv_exports/feed_items.csv.gz",
                13: "Notes",
                14: "• text is intentionally not shown here",
                15: "• embeds/domains are stored as metadata",
                16: "• join keys are stable identifiers",
                17: "• no analysis in this deck",
                18: "• schema in data_dictionary.csv",
                19: "• integrity in manifest.csv",
                21: "(post metadata excerpt; truncated)",
            },
            pictures={20: rec("post_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 45 Receipt: post labels (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=5,
            text={
                2: "Receipt: post_labels.csv.gz",
                4: "What this proves",
                5: "Labels are normalized rows keyed to impressions.",
                6: "This slide shows a short labels excerpt for one impression.",
                7: "No claims; just evidence the table exists.",
                9: "Where it lives",
                10: "02_csv_exports/post_labels.csv.gz",
                11: "02_csv_exports/feed_items.csv.gz",
                13: "Attachment keys",
                14: "• feed_uri",
                15: "• viewer_mode",
                16: "• post_uri",
                17: "• post_cid",
                18: "• label_src + label_val",
                19: "• collected_at_utc",
                21: "(labels excerpt; truncated)",
            },
            pictures={20: rec("post_labels_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 46 Authors hydration (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=5,
            text={
                2: "Authors (profile hydration)",
                3: "Turn DIDs into stable profile metadata.",
                5: "Goal: hydrate author profiles for posts observed in impressions.",
                6: "API: app.bsky.actor.getProfiles (batched).",
                8: "Input",
                10: "Output",
                12: "Why",
                14: "No counts",
                15: "We avoid follower-count claims in the deck.",
                17: "Files (names only)",
                18: "• feed_items.csv.gz provides author_did",
                19: "• authors.csv.gz stores handle/display_name/etc",
                20: "• joins enable later analyses",
                21: "• no post text shown",
                22: "• schema in data_dictionary.csv",
                23: "• integrity in manifest.csv",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 47 Receipt: author row (base 13)
    slide_edits.append(
        SlideEdits(
            base=13,
            chapter=5,
            text={
                2: "Receipt: authors.csv.gz",
                4: "What this proves",
                5: "We hydrate author profiles for observed impressions.",
                6: "This slide shows profile fields only (no follower counts).",
                7: "Join key is author_did.",
                9: "Where it lives",
                10: "02_csv_exports/authors.csv.gz",
                11: "02_csv_exports/feed_items.csv.gz",
                13: "Join",
                14: "• posts.csv.gz.author_did → authors.csv.gz.author_did",
                15: "• feed_items.csv.gz.author_did → authors.csv.gz.author_did",
                16: "• viewer_mode does not affect author identity",
                17: "• schema in data_dictionary.csv",
                18: "• integrity in manifest.csv",
                19: "• read-only hydration via getProfiles",
                21: "(author row excerpt; truncated)",
            },
            pictures={20: rec("author_row")},
            fills={},
            extra_notes=[],
        )
    )

    # 48 Join spine (base 15)
    slide_edits.append(
        SlideEdits(
            base=15,
            chapter=5,
            text={
                2: "Join spine (keys line up)",
                3: "Panel → Impression → Post metadata → Author profile",
                8: "Panel entry",
                9: "Ranked impression",
                10: "Post metadata (no text)",
                11: "Author profile",
            },
            pictures={
                12: rec("panel_row"),
                13: rec("impression_row"),
                14: rec("post_row"),
                15: rec("author_row"),
            },
            fills={},
            extra_notes=[],
        )
    )

    # 49 Labels attach (base 15)
    slide_edits.append(
        SlideEdits(
            base=15,
            chapter=5,
            text={
                2: "Labels attach (normalized)",
                3: "post_labels rows attach to impressions via explicit keys.",
                8: "Panel entry",
                9: "Ranked impression",
                10: "Post metadata (no text)",
                11: "Author profile",
            },
            pictures={
                12: rec("panel_row"),
                13: rec("impression_row"),
                14: rec("post_row"),
                15: rec("author_row"),
            },
            fills={},
            extra_notes=[
                "Labels attach via feed_uri + viewer_mode + post_uri + post_cid (see post_labels.csv.gz).",
            ],
        )
    )

    # 50 Run folder map (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=6,
            text={
                2: "Run folder layout (00–07)",
                3: "Everything you need to reproduce receipts.",
                5: "Organized exports + state + logs + postprocess + archive.",
                6: "Deck images are generated from the run folder outputs.",
                8: "Core",
                10: "Receipts",
                12: "Archive",
                14: "Naming",
                15: f"Bluesky_Run_{run_date}_{run_id}.zip",
                17: "Folders",
                18: "• 01_state_db/ + state.db",
                19: "• 02_csv_exports/ (csv + csv.gz tables)",
                20: "• 03_postprocess_metrics/ (joins + H tables; not shown)",
                21: "• 05_manifest/ (manifest + dictionary + validations)",
                22: "• 06_figures_preview/receipts/ (white cards)",
                23: "• 07_archive_zip/ (portable zip)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 51 Trust stamp (base 17)
    slide_edits.append(
        SlideEdits(
            base=17,
            chapter=6,
            text={
                2: "Reproducible outputs",
                3: "A tiny trust stamp: manifest + data dictionary.",
                5: "manifest.csv (hash list)",
                8: "data_dictionary.csv (column meanings)",
            },
            pictures={6: rec("manifest_excerpt"), 9: rec("data_dictionary_excerpt")},
            fills={},
            extra_notes=[],
        )
    )

    # 52 Provenance montage (base 14)
    slide_edits.append(
        SlideEdits(
            base=14,
            chapter=7,
            text={
                2: "Provenance montage",
                3: "Receipts are rendered from the run folder (not hand-made).",
                5: "Folder map + archive + log excerpt",
                9: "Audit trail",
                11: "No results shown; only provenance + file existence.",
            },
            pictures={
                6: rec("folder_tree_receipt"),
                7: rec("zip_receipt"),
                10: rec("log_excerpt_receipt"),
            },
            fills={},
            extra_notes=[],
        )
    )

    # 53 What enables later (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=7,
            text={
                2: "What this enables later (not today)",
                3: "Analyses we can run after Q&A.",
                5: "Compare discovery vs exposure vs safety signals.",
                6: "Test stability across time with repeated runs (future).",
                8: "Measurement",
                10: "Threats",
                12: "Mitigations",
                14: "No claims",
                15: "This deck ends before any results or charts.",
                17: "Examples",
                18: "• H1–H6 tables exist (stored; not shown)",
                19: "• robustness checks across viewer modes (future)",
                20: "• provider bucket risk mapping",
                21: "• feed overlap / same-winners tests",
                22: "• privacy-preserving releases",
                23: "• mitigation experiments (future)",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    # 54 Glossary (base 16)
    slide_edits.append(
        SlideEdits(
            base=16,
            chapter=7,
            text={
                2: "Glossary (for Q&A)",
                3: "Key identifiers and join keys.",
                5: "feed_uri: at://.../app.bsky.feed.generator/<rkey>",
                6: "post_uri: at://.../app.bsky.feed.post/<rkey> (+ post_cid)",
                8: "Identifiers",
                10: "Viewer context",
                12: "Where to look",
                14: "Closing",
                15: "See data_dictionary.csv for the full schema.",
                17: "Terms",
                18: "• did: decentralized identifier for actors",
                19: "• cid: content identifier for a record version",
                20: "• viewer_mode: unauth/auth snapshot context",
                21: "• join spine: panel → feed_items → posts/authors",
                22: "• labels attach on (feed_uri, viewer_mode, post_uri, post_cid)",
                23: "Thanks - questions welcome.",
            },
            pictures={},
            fills={},
            extra_notes=[],
        )
    )

    if len(slide_edits) != 54:
        raise RuntimeError(f"internal error: expected 54 slide edits, got {len(slide_edits)}")
    return slide_edits


def build_deck(*, paths: BuildPaths, run_dir: Path) -> None:
    # 1) Build working structure (OOXML zip edits).
    recipe = build_working_structure(paths=paths)

    # 2) Load run metadata (for date/run_id text).
    run_meta = (run_dir / "05_manifest" / "run_metadata.csv").read_text(encoding="utf-8").splitlines()
    if len(run_meta) < 2:
        raise RuntimeError("run_metadata.csv is empty")
    header = run_meta[0].split(",")
    row = run_meta[1].split(",")
    meta = dict(zip(header, row, strict=False))
    run_id = meta.get("run_id", "").strip()
    run_date = run_dir.name.replace("data_run_", "")

    receipts_dir = run_dir / "06_figures_preview" / "receipts"
    if not receipts_dir.exists():
        raise RuntimeError(f"Missing receipts dir (run build_run_receipts.py first): {receipts_dir}")

    receipts: dict[str, Path] = {
        "run_metadata_receipt": receipts_dir / "run_metadata_receipt.png",
        "run_summary_receipt": receipts_dir / "run_summary_receipt.png",
        "feed_index_row": receipts_dir / "feed_index_row.png",
        "starterpack_row": receipts_dir / "starterpack_row.png",
        "starterpack_inclusion_row": receipts_dir / "starterpack_inclusion_row.png",
        "popular_row": receipts_dir / "popular_row.png",
        "panel_row": receipts_dir / "panel_row.png",
        "impression_row": receipts_dir / "impression_row.png",
        "post_row": receipts_dir / "post_row.png",
        "post_labels_row": receipts_dir / "post_labels_row.png",
        "author_row": receipts_dir / "author_row.png",
        "manifest_excerpt": receipts_dir / "manifest_excerpt.png",
        "data_dictionary_excerpt": receipts_dir / "data_dictionary_excerpt.png",
        "validation_report_receipt": receipts_dir / "validation_report_receipt.png",
        "folder_tree_receipt": receipts_dir / "folder_tree_receipt.png",
        "zip_receipt": receipts_dir / "zip_receipt.png",
        "log_excerpt_receipt": receipts_dir / "log_excerpt_receipt.png",
    }
    missing = [k for k, p in receipts.items() if not p.exists()]
    if missing:
        raise RuntimeError(f"Missing receipt PNGs: {missing}")

    slide_edits = _build_slide_edits(run_date=run_date, run_id=run_id, receipts=receipts)

    # 3) Prepare assets (assert-derived).
    paths.assets_dir.mkdir(parents=True, exist_ok=True)
    assets = prepare_assert_assets(assert_dir=paths.assert_dir, cache_dir=paths.assets_dir)

    # 4) Apply content edits (python-pptx), preserving existing shape IDs and timing.
    prs = Presentation(str(paths.working_pptx))
    if len(prs.slides) != 54:
        raise RuntimeError(f"unexpected slide count in working deck: {len(prs.slides)}")

    asset_manifest: dict[int, list[AssertAssetRef]] = {}
    manifest_entries: list[SlideManifestEntry] = []

    for slide_num, (slide, edits) in enumerate(zip(prs.slides, slide_edits, strict=True), start=1):
        _add_common_visuals(
            slide=slide,
            prs=prs,
            assets=assets,
            slide_num=slide_num,
            chapter=edits.chapter,
            asset_manifest=asset_manifest,
        )

        # Sanity: base slide type should match recipe.
        expected_base = recipe[slide_num - 1]
        if edits.base != expected_base:
            raise RuntimeError(f"slide {slide_num}: edits.base={edits.base} but recipe expects {expected_base}")

        # Text updates.
        for shape_id, text in edits.text.items():
            _set_run_text(_shape(slide, shape_id), text)

        # Picture replacements (existing picture shapes).
        for shape_id, img_path in edits.pictures.items():
            _replace_picture_image(_shape(slide, shape_id), slide, img_path)

        # Fill auto-shapes with pictures (keep shape; preserve animations).
        for shape_id, img_path in edits.fills.items():
            sh = _shape(slide, shape_id)
            _clear_text(sh)
            _fill_shape_with_picture(sh, slide, img_path)

        # Extra note textbox for some slides (adds a new shape; does not alter animations).
        for note in edits.extra_notes:
            tb = slide.shapes.add_textbox(_emu(0.75), _emu(6.85), _emu(11.5), _emu(0.4))
            tf = tb.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = _norm_text(note)
            p.font.name = "Calibri"
            p.font.size = Pt(12)

        # Capture title for the asset manifest.
        title = ""
        try:
            title = _shape(slide, 2).text_frame.text.strip() if getattr(_shape(slide, 2), "has_text_frame", False) else ""
        except KeyError:
            title = ""
        manifest_entries.append(SlideManifestEntry(slide_num=slide_num, title=title, assets=asset_manifest.get(slide_num, [])))

    tmp_out = paths.out_pptx.with_suffix(".tmp.pptx")
    if tmp_out.exists():
        tmp_out.unlink()
    prs.save(str(tmp_out))
    shutil.move(tmp_out, paths.out_pptx)

    # Write asset manifest JSON for audit.
    out_manifest = paths.out_pptx.parent / "Slide2_asset_manifest.json"
    payload: list[dict[str, object]] = []
    for e in manifest_entries:
        payload.append(
            {
                "slide": e.slide_num,
                "title": e.title,
                "assets": [
                    {
                        "role": a.role,
                        "source_pptx": str(a.source_pptx),
                        "media_path": a.media_path,
                        "derived_path": str(a.derived_path),
                    }
                    for a in e.assets
                ],
            }
        )
    out_manifest.write_text(json.dumps(payload, indent=2, sort_keys=True) + "\n", encoding="utf-8")

    # Write slide index for quick auditing.
    out_index = paths.out_pptx.parent / "Slide2_slide_index.json"
    out_index.write_text(
        json.dumps(
            [{"slide": i + 1, "base_slide": recipe[i]} for i in range(len(recipe))],
            indent=2,
            sort_keys=True,
        )
        + "\n",
        encoding="utf-8",
    )


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="Build Slide2 animated blackboard deck from a fresh run folder (54 slides).")
    p.add_argument("--run-dir", type=Path, required=True)
    p.add_argument(
        "--template",
        type=Path,
        default=Path("Slides/PPTXWT2/bluesky-data-collection-pipeline-blackboard-animated-run-20260201.pptx"),
    )
    p.add_argument("--out", type=Path, default=Path("Slide2/Slide2_RQs_and_Data.pptx"))
    p.add_argument("--working", type=Path, default=Path("Slide2/Slide2_Working.pptx"))
    p.add_argument("--assets-dir", type=Path, default=Path("Slide2/_assets"))
    p.add_argument("--assert-dir", type=Path, default=Path("Slides/assert"))
    return p


def main(argv: Sequence[str] | None = None) -> int:
    args = _build_parser().parse_args(list(argv) if argv is not None else None)
    run_dir = args.run_dir.resolve()
    if not run_dir.exists():
        raise SystemExit(f"Missing --run-dir: {run_dir}")

    paths = BuildPaths(
        template_pptx=args.template.resolve(),
        working_pptx=args.working.resolve(),
        out_pptx=args.out.resolve(),
        assets_dir=args.assets_dir.resolve(),
        assert_dir=args.assert_dir.resolve(),
    )

    build_deck(paths=paths, run_dir=run_dir)
    print(f"OK: wrote {paths.out_pptx}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
