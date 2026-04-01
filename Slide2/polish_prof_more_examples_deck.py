#!/usr/bin/env python3
from __future__ import annotations

import argparse
import io
import tempfile
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Sequence

from lxml import etree
from PIL import Image, ImageFilter, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_FILL
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Pt


EMU_PER_INCH = 914400

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

NSMAP = {"p": PML_NS, "a": A_NS, "r": R_NS}


PIL_RGBA = "RGBA"


@dataclass(frozen=True)
class ThemeColors:
    card_fill_old: RGBColor
    card_fill_new: RGBColor
    card_line_old: RGBColor
    card_line_new: RGBColor
    muted_text_old: RGBColor
    muted_text_new: RGBColor
    abc_a_old: RGBColor
    abc_b_old: RGBColor
    abc_c_old: RGBColor
    abc_a_new: RGBColor
    abc_b_new: RGBColor
    abc_c_new: RGBColor
    connector_old: RGBColor
    connector_new: RGBColor


DV_POLISH = ThemeColors(
    card_fill_old=RGBColor(0x16, 0x20, 0x33),
    card_fill_new=RGBColor(0x11, 0x1C, 0x2D),
    card_line_old=RGBColor(0x3A, 0x4D, 0x66),
    card_line_new=RGBColor(0x2B, 0x3A, 0x55),
    muted_text_old=RGBColor(0xC9, 0xD2, 0xE3),
    muted_text_new=RGBColor(0xB7, 0xC2, 0xD6),
    abc_a_old=RGBColor(0x58, 0xB7, 0xE6),
    abc_b_old=RGBColor(0xF4, 0xC5, 0x5D),
    abc_c_old=RGBColor(0xFF, 0x7D, 0x77),
    abc_a_new=RGBColor(0xAA, 0x6E, 0xE7),  # template accent purple
    abc_b_new=RGBColor(0xFF, 0xA5, 0x65),  # template accent orange
    abc_c_new=RGBColor(0xFF, 0x8B, 0x9B),  # template accent pink
    connector_old=RGBColor(0xC9, 0xD2, 0xE3),
    connector_new=RGBColor(0xB7, 0xC2, 0xD6),
)


@dataclass(frozen=True)
class BboxPx:
    x0: int
    y0: int
    x1: int
    y1: int

    @property
    def w(self) -> int:
        return self.x1 - self.x0

    @property
    def h(self) -> int:
        return self.y1 - self.y0


@dataclass(frozen=True)
class BboxIn:
    x: float
    y: float
    w: float
    h: float


def _qn(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _emu(inches: float) -> int:
    return int(round(inches * EMU_PER_INCH))


def _emu_to_in(emu: int) -> float:
    return emu / EMU_PER_INCH


def _slide_text(slide) -> str:
    tokens: list[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        for p in tf.paragraphs:
            t = (p.text or "").strip()
            if t:
                tokens.append(t)
    return " ".join(tokens)


def _is_ui_slide(slide) -> bool:
    return "Bluesky UI:" in _slide_text(slide)


def _is_receipt_slide(slide) -> bool:
    return "Data receipt:" in _slide_text(slide)


def _iter_all_shapes(slide) -> Iterable:
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


def _normalize_dashes(text: str) -> str:
    # Avoid glyph fallback boxes in PDF export by normalizing uncommon dash codepoints.
    return (
        text.replace("\u2011", "-")  # non-breaking hyphen
        .replace("\u2010", "-")  # hyphen
        .replace("\u2212", "-")  # minus sign
    )


def _normalize_text_frames(slide) -> int:
    changed = 0
    for shape in _iter_all_shapes(slide):
        if not getattr(shape, "has_text_frame", False):
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text:
                    continue
                out = _normalize_dashes(run.text)
                if out != run.text:
                    run.text = out
                    changed += 1
    return changed


def _maybe_swap_shape_fill(shape, *, old: RGBColor, new: RGBColor) -> bool:
    try:
        fill = shape.fill
    except Exception:
        return False
    if fill.type != MSO_FILL.SOLID:
        return False
    fore = fill.fore_color
    if fore.type != MSO_COLOR_TYPE.RGB:
        return False
    if fore.rgb != old:
        return False
    fore.rgb = new
    return True


def _maybe_swap_shape_line(shape, *, old: RGBColor, new: RGBColor) -> bool:
    try:
        line = shape.line
    except Exception:
        return False
    if line is None:
        return False
    color = line.color
    if color is None or color.type != MSO_COLOR_TYPE.RGB:
        return False
    if color.rgb != old:
        return False
    color.rgb = new
    return True


def _maybe_swap_text_colors(shape, *, old: RGBColor, new: RGBColor) -> int:
    if not getattr(shape, "has_text_frame", False):
        return 0
    swapped = 0
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            color = run.font.color
            if color is None or color.type != MSO_COLOR_TYPE.RGB:
                continue
            if color.rgb != old:
                continue
            color.rgb = new
            swapped += 1
    return swapped


def _set_line_rgb(shape, rgb: RGBColor) -> None:
    try:
        line = shape.line
    except Exception:
        return
    if line is None:
        return
    line.color.rgb = rgb


def _set_line_width_pt(shape, width_pt: float) -> None:
    try:
        line = shape.line
    except Exception:
        return
    if line is None:
        return
    line.width = Pt(width_pt)


def _set_solid_fill_rgb(shape, rgb: RGBColor) -> None:
    try:
        fill = shape.fill
    except Exception:
        return
    fill.solid()
    fill.fore_color.rgb = rgb


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found")


def _picture_shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id and sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            return sh
    raise KeyError(f"picture shape_id={shape_id} not found")


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001
    blip.set(qn("r:embed"), rid)


def _open_image_from_blob(blob: bytes) -> Image.Image:
    return Image.open(io.BytesIO(blob))


def _compute_card_bboxes_px_from_picture(pic_shape) -> list[BboxPx]:
    img = _open_image_from_blob(pic_shape.image.blob).convert(PIL_RGBA)
    alpha = img.getchannel("A")
    w, h = alpha.size

    px = alpha.load()
    row_has = [False] * h
    for y in range(h):
        for x in range(w):
            if px[x, y] != 0:
                row_has[y] = True
                break

    segments: list[tuple[int, int]] = []
    y = 0
    while y < h:
        if not row_has[y]:
            y += 1
            continue
        y0 = y
        while y < h and row_has[y]:
            y += 1
        y1 = y
        segments.append((y0, y1))

    bboxes: list[BboxPx] = []
    for seg_y0, seg_y1 in segments:
        bbox = alpha.crop((0, seg_y0, w, seg_y1)).getbbox()
        if bbox is None:
            continue
        x0, y0, x1, y1 = bbox
        bboxes.append(BboxPx(x0=x0, y0=seg_y0 + y0, x1=x1, y1=seg_y0 + y1))

    return bboxes


def _map_bbox_to_slide_in(*, bbox: BboxPx, pic_shape) -> BboxIn:
    img_w, img_h = pic_shape.image.size
    pic_left = _emu_to_in(int(pic_shape.left))
    pic_top = _emu_to_in(int(pic_shape.top))
    pic_w = _emu_to_in(int(pic_shape.width))
    pic_h = _emu_to_in(int(pic_shape.height))

    sx = pic_w / img_w
    sy = pic_h / img_h
    x = pic_left + bbox.x0 * sx
    y = pic_top + bbox.y0 * sy
    w = bbox.w * sx
    h = bbox.h * sy
    return BboxIn(x=x, y=y, w=w, h=h)


def align_ovals_to_rectangles(
    *,
    slide,
    rect_shape_ids: Sequence[int],
    oval_shape_ids: Sequence[int],
    oval_overlap_in: float,
    oval_top_offset_in: float,
) -> None:
    if len(rect_shape_ids) != len(oval_shape_ids):
        raise ValueError("rect_shape_ids and oval_shape_ids must have same length")

    for rect_id, oval_id in zip(rect_shape_ids, oval_shape_ids, strict=True):
        rect = _shape(slide, rect_id)
        oval = _shape(slide, oval_id)

        rect_left = _emu_to_in(int(rect.left))
        rect_top = _emu_to_in(int(rect.top))
        oval_w = _emu_to_in(int(oval.width))
        oval_h = _emu_to_in(int(oval.height))

        oval.left = _emu(rect_left - oval_w + oval_overlap_in)
        oval.top = _emu(rect_top + oval_top_offset_in)
        oval.width = _emu(oval_w)
        oval.height = _emu(oval_h)


def align_callouts_to_cards(
    *,
    slide,
    picture_shape_id: int,
    rect_shape_ids: Sequence[int],
    oval_shape_ids: Sequence[int],
    inset_x_in: float,
    inset_y_in: float,
    oval_overlap_in: float,
    oval_top_offset_in: float,
) -> None:
    if len(rect_shape_ids) != len(oval_shape_ids):
        raise ValueError("rect_shape_ids and oval_shape_ids must have same length")

    pic = _picture_shape(slide, picture_shape_id)
    bboxes_px = _compute_card_bboxes_px_from_picture(pic)
    if len(bboxes_px) < len(rect_shape_ids):
        raise RuntimeError(f"expected >= {len(rect_shape_ids)} card bboxes, got {len(bboxes_px)}")

    for rect_id, oval_id, bbox in zip(rect_shape_ids, oval_shape_ids, bboxes_px, strict=True):
        bbox_in = _map_bbox_to_slide_in(bbox=bbox, pic_shape=pic)

        rect = _shape(slide, rect_id)
        rect.left = _emu(bbox_in.x + inset_x_in)
        rect.top = _emu(bbox_in.y + inset_y_in)
        rect.width = _emu(max(0.0, bbox_in.w - 2 * inset_x_in))
        rect.height = _emu(max(0.0, bbox_in.h - 2 * inset_y_in))

        oval = _shape(slide, oval_id)
        oval_w = _emu_to_in(int(oval.width))
        oval_h = _emu_to_in(int(oval.height))
        oval.left = _emu(_emu_to_in(int(rect.left)) - oval_w + oval_overlap_in)
        oval.top = _emu(_emu_to_in(int(rect.top)) + oval_top_offset_in)
        oval.width = _emu(oval_w)
        oval.height = _emu(oval_h)


def _shape_bbox(shape) -> tuple[int, int, int, int]:
    left = int(getattr(shape, "left", 0))
    top = int(getattr(shape, "top", 0))
    width = int(getattr(shape, "width", 0))
    height = int(getattr(shape, "height", 0))
    return left, top, left + width, top + height


def _center_y(top: int, bottom: int) -> int:
    return int(round((top + bottom) / 2))


def _set_connectors_to_highlights(slide) -> None:
    # Connector 23: A chips -> rect 18
    chip_a0 = _shape(slide, 10)
    chip_a1 = _shape(slide, 11)
    rect_a = _shape(slide, 18)
    conn_a = _shape(slide, 23)

    a_l0, a_t0, a_r0, a_b0 = _shape_bbox(chip_a0)
    a_l1, a_t1, a_r1, a_b1 = _shape_bbox(chip_a1)
    a_right = max(a_r0, a_r1)
    a_top = min(a_t0, a_t1)
    a_bottom = max(a_b0, a_b1)

    ra_l, ra_t, _ra_r, ra_b = _shape_bbox(rect_a)
    conn_a.begin_x = a_right
    conn_a.begin_y = _center_y(a_top, a_bottom)
    conn_a.end_x = ra_l
    conn_a.end_y = _center_y(ra_t, ra_b)

    # Connector 24: B chips -> rect 20
    chip_b0 = _shape(slide, 12)
    chip_b1 = _shape(slide, 13)
    rect_b = _shape(slide, 20)
    conn_b = _shape(slide, 24)

    b_l0, b_t0, b_r0, b_b0 = _shape_bbox(chip_b0)
    b_l1, b_t1, b_r1, b_b1 = _shape_bbox(chip_b1)
    b_right = max(b_r0, b_r1)
    b_top = min(b_t0, b_t1)
    b_bottom = max(b_b0, b_b1)

    rb_l, rb_t, _rb_r, rb_b = _shape_bbox(rect_b)
    conn_b.begin_x = b_right
    conn_b.begin_y = _center_y(b_top, b_bottom)
    conn_b.end_x = rb_l
    conn_b.end_y = _center_y(rb_t, rb_b)


def _shift_top(shape, *, delta_in: float) -> None:
    if not hasattr(shape, "top"):
        return
    shape.top = int(shape.top) + _emu(delta_in)


def _reposition_main_picture_stack(slide, *, target_top_in: float) -> float:
    pic = _picture_shape(slide, 8)
    current_top_in = _emu_to_in(int(pic.top))
    delta_in = target_top_in - current_top_in
    if abs(delta_in) < 1e-6:
        return 0.0

    _shift_top(pic, delta_in=delta_in)

    if _is_ui_slide(slide):
        for sid in (17, 18, 19, 20, 21, 22):
            _shift_top(_shape(slide, sid), delta_in=delta_in)
        for sid in (23, 24):
            _shift_top(_shape(slide, sid), delta_in=delta_in)
    elif _is_receipt_slide(slide):
        for sid in (17, 18, 19, 20, 21, 22):
            _shift_top(_shape(slide, sid), delta_in=delta_in)

    return delta_in


def _target_main_pic_top_in(
    *,
    slide,
    slide_h_in: float,
    default_top_in: float,
    min_takeaway_gap_in: float,
    bottom_margin_in: float,
) -> float:
    # Keep a minimum gap under the takeaway line, without pushing the picture off-slide.
    pic = _picture_shape(slide, 8)
    pic_h_in = _emu_to_in(int(pic.height))

    desired = float(default_top_in)
    try:
        takeaway = _shape(slide, 7)
        takeaway_bottom_in = _emu_to_in(int(takeaway.top + takeaway.height))
        desired = max(desired, takeaway_bottom_in + float(min_takeaway_gap_in))
    except Exception:
        pass

    max_top = float(slide_h_in) - float(bottom_margin_in) - pic_h_in
    if max_top > 0:
        desired = min(desired, max_top)
    return desired


def _render_receipt_hd(*, pic_shape, out_png: Path, target_w: int, target_h: int, pad_px: int) -> Path:
    out_png.parent.mkdir(parents=True, exist_ok=True)

    img = _open_image_from_blob(pic_shape.image.blob).convert(PIL_RGBA)
    alpha = img.getchannel("A")
    bbox = alpha.getbbox()
    if bbox is None:
        raise ValueError("receipt image alpha bbox is empty")

    x0, y0, x1, y1 = bbox
    x0 = max(0, x0 - pad_px)
    y0 = max(0, y0 - pad_px)
    x1 = min(img.width, x1 + pad_px)
    y1 = min(img.height, y1 + pad_px)
    crop = img.crop((x0, y0, x1, y1))

    resampling = getattr(Image, "Resampling", Image)
    scaled = ImageOps.contain(crop, (target_w, target_h), method=resampling.LANCZOS)

    # Only sharpen if we upscaled.
    if scaled.width > crop.width or scaled.height > crop.height:
        scaled = scaled.filter(ImageFilter.UnsharpMask(radius=1.1, percent=120, threshold=3))

    canvas = Image.new(PIL_RGBA, (target_w, target_h), (0, 0, 0, 0))
    left = (target_w - scaled.width) // 2
    top = (target_h - scaled.height) // 2
    canvas.paste(scaled, (left, top), scaled)
    canvas.save(out_png, format="PNG")
    return out_png


def _normalize_global_colors(prs: Presentation, *, theme: ThemeColors) -> None:
    for slide in prs.slides:
        for shape in _iter_all_shapes(slide):
            _maybe_swap_shape_fill(shape, old=theme.card_fill_old, new=theme.card_fill_new)
            _maybe_swap_shape_line(shape, old=theme.card_line_old, new=theme.card_line_new)
            _maybe_swap_text_colors(shape, old=theme.muted_text_old, new=theme.muted_text_new)


def _normalize_abc_palette(slide, *, theme: ThemeColors) -> None:
    # Only touch slides with A/B/C callout grammar.
    if not (_is_ui_slide(slide) or _is_receipt_slide(slide)):
        return

    for shape in _iter_all_shapes(slide):
        _maybe_swap_shape_line(shape, old=theme.abc_a_old, new=theme.abc_a_new)
        _maybe_swap_shape_line(shape, old=theme.abc_b_old, new=theme.abc_b_new)
        _maybe_swap_shape_line(shape, old=theme.abc_c_old, new=theme.abc_c_new)

        _maybe_swap_shape_fill(shape, old=theme.abc_a_old, new=theme.abc_a_new)
        _maybe_swap_shape_fill(shape, old=theme.abc_b_old, new=theme.abc_b_new)
        _maybe_swap_shape_fill(shape, old=theme.abc_c_old, new=theme.abc_c_new)

        _maybe_swap_shape_line(shape, old=theme.connector_old, new=theme.connector_new)


def _fix_receipt_c_palette(slide, *, theme: ThemeColors) -> None:
    # Receipts in the baseline deck use A + (B/C) amber for some elements; force distinct C accents.
    if not _is_receipt_slide(slide):
        return

    # C group chips are 14-16; C callouts are oval 21 and rect 22.
    for sid in (14, 15, 16):
        _set_line_rgb(_shape(slide, sid), theme.abc_c_new)

    _set_solid_fill_rgb(_shape(slide, 21), theme.abc_c_new)
    _set_line_rgb(_shape(slide, 22), theme.abc_c_new)


def _tune_callout_line_weights(
    slide,
    *,
    chip_line_pt: float,
    highlight_line_pt: float,
    connector_line_pt: float,
) -> None:
    if not (_is_ui_slide(slide) or _is_receipt_slide(slide)):
        return

    for sid in range(10, 17):
        _set_line_width_pt(_shape(slide, sid), chip_line_pt)

    for sid in (18, 20, 22):
        _set_line_width_pt(_shape(slide, sid), highlight_line_pt)

    if _is_ui_slide(slide):
        for sid in (23, 24):
            _set_line_width_pt(_shape(slide, sid), connector_line_pt)


def _patch_texture_alpha(
    *,
    pptx_in: Path,
    pptx_out: Path,
    texture_descr: str,
    alpha_amt: int,
) -> None:
    if not (0 <= alpha_amt <= 100000):
        raise ValueError("alpha_amt must be in [0, 100000]")

    slide_re = etree.XPath(".//p:cNvPr[@descr=$descr]", namespaces=NSMAP)
    pic_ancestor = etree.XPath("ancestor::p:pic[1]", namespaces=NSMAP)

    with zipfile.ZipFile(pptx_in, "r") as zin:
        pptx_out.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if not info.filename.startswith("ppt/slides/slide") or not info.filename.endswith(".xml"):
                    zout.writestr(info, data)
                    continue

                root = etree.fromstring(data)
                changed = False
                for cNvPr in slide_re(root, descr=texture_descr):
                    pics = pic_ancestor(cNvPr)
                    if not pics:
                        continue
                    blip = pics[0].find(f".//{_qn(A_NS, 'blip')}")
                    if blip is None:
                        continue
                    alpha = blip.find(f"{_qn(A_NS, 'alphaModFix')}")
                    if alpha is None:
                        alpha = etree.SubElement(blip, _qn(A_NS, "alphaModFix"))
                    alpha.set("amt", str(alpha_amt))
                    changed = True

                if changed:
                    data = etree.tostring(root, encoding="UTF-8", xml_declaration=True, standalone=True)
                zout.writestr(info, data)


def polish_deck(
    *,
    pptx_in: Path,
    pptx_out: Path,
    cache_dir: Path,
    main_pic_top_in: float,
    min_takeaway_gap_in: float,
    bottom_margin_in: float,
    receipt_hd_size: tuple[int, int],
    receipt_crop_pad_px: int,
    texture_alpha_amt: int,
    chip_line_pt: float,
    highlight_line_pt: float,
    connector_line_pt: float,
) -> None:
    pptx_in = pptx_in.resolve()
    pptx_out = pptx_out.resolve()
    cache_dir = cache_dir.resolve()

    if not pptx_in.exists():
        raise FileNotFoundError(f"missing --in: {pptx_in}")

    prs = Presentation(str(pptx_in))
    _normalize_global_colors(prs, theme=DV_POLISH)
    slide_h_in = _emu_to_in(int(prs.slide_height))

    for slide_num, slide in enumerate(prs.slides, start=1):
        _normalize_text_frames(slide)
        _normalize_abc_palette(slide, theme=DV_POLISH)

        if _is_ui_slide(slide) or _is_receipt_slide(slide):
            target = _target_main_pic_top_in(
                slide=slide,
                slide_h_in=slide_h_in,
                default_top_in=main_pic_top_in,
                min_takeaway_gap_in=min_takeaway_gap_in,
                bottom_margin_in=bottom_margin_in,
            )
            _reposition_main_picture_stack(slide, target_top_in=target)

        if _is_ui_slide(slide):
            align_ovals_to_rectangles(
                slide=slide,
                rect_shape_ids=(18, 20, 22),
                oval_shape_ids=(17, 19, 21),
                oval_overlap_in=0.05,
                oval_top_offset_in=0.04,
            )
            _set_connectors_to_highlights(slide)

        if _is_receipt_slide(slide):
            pic = _picture_shape(slide, 8)
            target_w, target_h = receipt_hd_size
            out_png = cache_dir / f"receipt_hd_slide{slide_num:02d}.png"
            _render_receipt_hd(
                pic_shape=pic,
                out_png=out_png,
                target_w=target_w,
                target_h=target_h,
                pad_px=receipt_crop_pad_px,
            )
            _replace_picture_image(pic, slide, out_png)
            align_callouts_to_cards(
                slide=slide,
                picture_shape_id=8,
                rect_shape_ids=(18, 20, 22),
                oval_shape_ids=(17, 19, 21),
                inset_x_in=0.08,
                inset_y_in=0.06,
                oval_overlap_in=0.05,
                oval_top_offset_in=0.04,
            )
            _fix_receipt_c_palette(slide, theme=DV_POLISH)

        _tune_callout_line_weights(
            slide,
            chip_line_pt=chip_line_pt,
            highlight_line_pt=highlight_line_pt,
            connector_line_pt=connector_line_pt,
        )

    with tempfile.TemporaryDirectory(prefix="pptx_polish_") as tmp_dir:
        tmp_path = Path(tmp_dir) / "polished.tmp.pptx"
        prs.save(str(tmp_path))
        _patch_texture_alpha(
            pptx_in=tmp_path,
            pptx_out=pptx_out,
            texture_descr="A2_card_texture.png",
            alpha_amt=texture_alpha_amt,
        )


def main() -> None:
    here = Path(__file__).resolve().parent
    parser = argparse.ArgumentParser(description="Polish the professor deck (layout + palette + receipt HD).")
    parser.add_argument(
        "--in",
        dest="pptx_in",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v13.preanim.pptx",
        help="Input PPTX (typically pre-animations).",
    )
    parser.add_argument(
        "--out",
        dest="pptx_out",
        type=Path,
        default=here / "deck_versions/Prof_Meeting_Bluesky_RQs_Data_FINAL_more_examples_v15.preanim.pptx",
        help="Output PPTX (pre-animations).",
    )
    parser.add_argument(
        "--cache-dir",
        type=Path,
        default=here / "prof_build/assets_cache/polish_v15",
        help="Cache directory for derived high-res receipt images.",
    )
    parser.add_argument(
        "--main-pic-top-in",
        type=float,
        default=2.06,
        help="Target top (inches) for the main screenshot/receipt stack on UI + receipt slides.",
    )
    parser.add_argument(
        "--min-takeaway-gap-in",
        type=float,
        default=0.16,
        help="Minimum gap (inches) between takeaway bottom and the main screenshot/receipt stack.",
    )
    parser.add_argument(
        "--bottom-margin-in",
        type=float,
        default=0.5,
        help="Minimum bottom margin (inches) below the main screenshot/receipt stack.",
    )
    parser.add_argument(
        "--receipt-hd-size",
        type=str,
        default="1920x1080",
        help="Target WxH for re-rendered receipt stacks (e.g., 1920x1080).",
    )
    parser.add_argument(
        "--receipt-crop-pad-px",
        type=int,
        default=18,
        help="Padding around alpha bbox before scaling receipt images.",
    )
    parser.add_argument(
        "--texture-alpha-amt",
        type=int,
        default=9000,
        help="OOXML alphaModFix amt for A2_card_texture.png pictures (0-100000).",
    )
    parser.add_argument(
        "--chip-line-pt",
        type=float,
        default=1.9,
        help="Line width (pt) for A/B/C chips on UI + receipt slides.",
    )
    parser.add_argument(
        "--highlight-line-pt",
        type=float,
        default=2.6,
        help="Line width (pt) for highlight rectangles on UI + receipt slides.",
    )
    parser.add_argument(
        "--connector-line-pt",
        type=float,
        default=1.2,
        help="Line width (pt) for chip-to-highlight connectors on UI slides.",
    )
    args = parser.parse_args()

    if "x" not in args.receipt_hd_size:
        raise SystemExit("--receipt-hd-size must be like 1920x1080")
    w_s, h_s = args.receipt_hd_size.split("x", maxsplit=1)
    receipt_hd_size = (int(w_s), int(h_s))

    polish_deck(
        pptx_in=args.pptx_in,
        pptx_out=args.pptx_out,
        cache_dir=args.cache_dir,
        main_pic_top_in=float(args.main_pic_top_in),
        min_takeaway_gap_in=float(args.min_takeaway_gap_in),
        bottom_margin_in=float(args.bottom_margin_in),
        receipt_hd_size=receipt_hd_size,
        receipt_crop_pad_px=int(args.receipt_crop_pad_px),
        texture_alpha_amt=int(args.texture_alpha_amt),
        chip_line_pt=float(args.chip_line_pt),
        highlight_line_pt=float(args.highlight_line_pt),
        connector_line_pt=float(args.connector_line_pt),
    )
    print(f"OK: wrote {args.pptx_out.resolve()}")


if __name__ == "__main__":
    main()
