#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400


@dataclass(frozen=True)
class DeckQcResult:
    slide_count: int
    total_click_effects: int
    slides_missing_timing: list[str]
    out_of_bounds: list[str]
    asset_issues: list[str]


def _iter_slide_xml_parts(z: zipfile.ZipFile) -> list[str]:
    parts = [
        name
        for name in z.namelist()
        if name.startswith("ppt/slides/slide") and name.endswith(".xml") and "/_rels/" not in name
    ]
    parts.sort(key=lambda s: int(Path(s).stem.replace("slide", "")))
    return parts


def _count_click_effects(xml: str) -> int:
    # Most reliable indicator across template slides is p:cTn nodeType="clickEffect".
    return len(re.findall(r'nodeType="clickEffect"', xml))


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):  # group shapes
            for child in shape.shapes:
                yield child


def _check_bounds(*, pptx_path: Path, margin_in: float) -> list[str]:
    prs = Presentation(str(pptx_path))
    margin = int(margin_in * EMU_PER_INCH)
    w = int(prs.slide_width)
    h = int(prs.slide_height)

    issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide):
            left = int(getattr(shape, "left", 0))
            top = int(getattr(shape, "top", 0))
            width = int(getattr(shape, "width", 0))
            height = int(getattr(shape, "height", 0))
            right = left + width
            bottom = top + height

            if width <= 0 or height <= 0:
                continue

            if left < -margin or top < -margin or right > w + margin or bottom > h + margin:
                issues.append(
                    f"slide {slide_num:02d}: shape_id={shape.shape_id} "
                    f"bbox=({left},{top})-({right},{bottom}) slide=({w},{h})"
                )
    return issues


def _is_under(child: Path, parent: Path) -> bool:
    child_s = str(child.resolve())
    parent_s = str(parent.resolve())
    if child_s == parent_s:
        return True
    return child_s.startswith(parent_s + "/")


def _check_asset_manifest_with_provenance(
    *,
    asset_manifest_path: Path,
    expected_slides: int,
    min_assets: int,
    assert_dir: Path,
    assets_cache_dir: Path,
) -> list[str]:
    if not asset_manifest_path.exists():
        return [f"Missing asset manifest: {asset_manifest_path}"]

    try:
        payload = json.loads(asset_manifest_path.read_text(encoding="utf-8"))
    except json.JSONDecodeError as err:
        return [f"Invalid JSON in asset manifest: {err}"]

    if not isinstance(payload, list):
        return ["Asset manifest root is not a list"]

    slide_to_assets: dict[int, int] = {}
    assert_dir_resolved = assert_dir.resolve()
    assets_cache_dir_resolved = assets_cache_dir.resolve()
    for entry in payload:
        if not isinstance(entry, dict):
            continue
        slide = entry.get("slide")
        assets = entry.get("assets")
        if not isinstance(slide, int) or not isinstance(assets, list):
            continue
        slide_to_assets[slide] = len(assets)
        for idx, asset in enumerate(assets, start=1):
            if not isinstance(asset, dict):
                issues = [f"slide {slide:02d} asset {idx}: malformed asset entry"]
                return issues
            source = asset.get("source_pptx")
            derived = asset.get("derived_path")
            role = asset.get("role")
            if not isinstance(source, str) or not isinstance(derived, str):
                issues = [f"slide {slide:02d} asset {idx}: missing source/derived path"]
                return issues
            if role not in {"structural", "semantic"}:
                issues = [f"slide {slide:02d} asset {idx}: invalid role={role!r}"]
                return issues
            source_path = Path(source)
            derived_path = Path(derived)
            if not _is_under(source_path, assert_dir_resolved):
                issues = [
                    f"slide {slide:02d} asset {idx}: source not under assert dir: {source_path}",
                ]
                return issues
            if not _is_under(derived_path, assets_cache_dir_resolved):
                issues = [
                    f"slide {slide:02d} asset {idx}: derived asset not under cache dir: {derived_path}",
                ]
                return issues

    issues: list[str] = []
    for s in range(1, expected_slides + 1):
        n = slide_to_assets.get(s, 0)
        if n < min_assets:
            issues.append(f"slide {s:02d}: expected >= {min_assets} assets, found {n}")
    return issues


def qc_deck(
    *,
    pptx_path: Path,
    expected_slides: int,
    min_click_effects: int,
    margin_in: float,
    asset_manifest_path: Path | None,
    min_assets_per_slide: int,
    assert_dir: Path,
    assets_cache_dir: Path,
) -> DeckQcResult:
    with zipfile.ZipFile(pptx_path, "r") as z:
        slide_parts = _iter_slide_xml_parts(z)

        total_clicks = 0
        missing_timing: list[str] = []
        for part in slide_parts:
            xml = z.read(part).decode("utf-8", "ignore")
            total_clicks += _count_click_effects(xml)
            if "<p:timing" not in xml:
                missing_timing.append(part)

    if len(slide_parts) != expected_slides:
        missing_timing.append(f"(slide xml count mismatch: {len(slide_parts)} != {expected_slides})")

    out_of_bounds = _check_bounds(pptx_path=pptx_path, margin_in=margin_in)

    asset_issues: list[str] = []
    if asset_manifest_path is not None:
        asset_issues = _check_asset_manifest_with_provenance(
            asset_manifest_path=asset_manifest_path,
            expected_slides=expected_slides,
            min_assets=min_assets_per_slide,
            assert_dir=assert_dir,
            assets_cache_dir=assets_cache_dir,
        )

    if total_clicks < min_click_effects:
        asset_issues.append(f"total_click_effects {total_clicks} < min_click_effects {min_click_effects}")

    return DeckQcResult(
        slide_count=len(slide_parts),
        total_click_effects=total_clicks,
        slides_missing_timing=missing_timing,
        out_of_bounds=out_of_bounds,
        asset_issues=asset_issues,
    )


def _build_parser() -> argparse.ArgumentParser:
    p = argparse.ArgumentParser(description="QC Slide2_RQs_and_Data.pptx (timing, click reveals, bounds, assets).")
    p.add_argument("pptx", nargs="?", default="Slide2/Slide2_RQs_and_Data.pptx")
    p.add_argument("--expected-slides", type=int, default=54)
    p.add_argument("--min-click-effects", type=int, default=600)
    p.add_argument("--margin-in", type=float, default=0.05)
    p.add_argument("--asset-manifest", type=Path, default=Path("Slide2/Slide2_asset_manifest.json"))
    p.add_argument("--min-assets-per-slide", type=int, default=3)
    p.add_argument("--assert-dir", type=Path, default=Path("Slides/assert"))
    p.add_argument("--assets-cache-dir", type=Path, default=Path("Slide2/_assets"))
    return p


def main(argv: list[str] | None = None) -> int:
    args = _build_parser().parse_args(argv if argv is not None else sys.argv[1:])
    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Missing pptx: {pptx_path}", file=sys.stderr)
        return 2

    res = qc_deck(
        pptx_path=pptx_path,
        expected_slides=int(args.expected_slides),
        min_click_effects=int(args.min_click_effects),
        margin_in=float(args.margin_in),
        asset_manifest_path=Path(args.asset_manifest) if args.asset_manifest else None,
        min_assets_per_slide=int(args.min_assets_per_slide),
        assert_dir=Path(args.assert_dir),
        assets_cache_dir=Path(args.assets_cache_dir),
    )

    ok = True
    if res.slide_count != int(args.expected_slides):
        ok = False
        print(f"FAIL: slide_count={res.slide_count} expected={args.expected_slides}", file=sys.stderr)
    if res.slides_missing_timing:
        ok = False
        print("FAIL: slides missing <p:timing>:", file=sys.stderr)
        for s in res.slides_missing_timing:
            print("-", s, file=sys.stderr)
    if res.total_click_effects < int(args.min_click_effects):
        ok = False
        print(
            f"FAIL: total_click_effects={res.total_click_effects} min={args.min_click_effects}",
            file=sys.stderr,
        )
    if res.out_of_bounds:
        ok = False
        print("FAIL: out-of-bounds shapes:", file=sys.stderr)
        for line in res.out_of_bounds:
            print("-", line, file=sys.stderr)
    if res.asset_issues:
        ok = False
        print("FAIL: asset manifest issues:", file=sys.stderr)
        for line in res.asset_issues:
            print("-", line, file=sys.stderr)

    if not ok:
        return 1

    print(
        "OK:",
        f"slides={res.slide_count}",
        f"clickEffects={res.total_click_effects}",
        "timing=present",
        f"bounds_ok (margin_in={args.margin_in})",
        f"assets_ok (min_per_slide={args.min_assets_per_slide})",
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
