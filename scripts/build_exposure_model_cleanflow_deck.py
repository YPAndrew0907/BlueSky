#!/usr/bin/env python3
from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


@dataclass(frozen=True)
class Theme:
    bg: RGBColor = RGBColor(255, 255, 255)
    text: RGBColor = RGBColor(15, 15, 15)
    muted: RGBColor = RGBColor(96, 96, 96)
    rule: RGBColor = RGBColor(206, 206, 206)
    title_font: str = "Times New Roman"
    body_font: str = "Arial"
    math_font: str = "Times New Roman"


THEME = Theme()

SLIDE_NOTES: list[str] = [
    (
        "[Sources]\n"
        "- Open with the research questions first, then earn the notation.\n"
        "- Ashudeep Singh and Thorsten Joachims. Fairness of Exposure in Rankings. KDD 2018. DOI: 10.1145/3219819.3220088\n"
        "- Yunqi Li et al. Fairness in Recommendation: Foundations, Methods, and Applications. ACM TIST 2023. DOI: 10.1145/3610302\n"
        "- Meike Zehlike, Ke Yang, Julia Stoyanovich. Fairness in Ranking, Part I and Part II. ACM CSUR 2022.\n"
    ),
    (
        "[Sources]\n"
        "- Early concept slide: exposure is the scarce resource and rank position allocates it.\n"
        "- Ashudeep Singh and Thorsten Joachims. Fairness of Exposure in Rankings. KDD 2018. DOI: 10.1145/3219819.3220088\n"
    ),
    (
        "[Sources]\n"
        "- This slide defines monitored exposure, not total platform impressions.\n"
        "- Ashudeep Singh and Thorsten Joachims. Fairness of Exposure in Rankings. KDD 2018. DOI: 10.1145/3219819.3220088\n"
    ),
    (
        "[Sources]\n"
        "- We use a conservative comparison unit: local content-similarity within a competition window, not claimed identity of real-world events.\n"
        "- Meike Zehlike, Ke Yang, Julia Stoyanovich. Fairness in Ranking, Part II. ACM CSUR 2022.\n"
    ),
    (
        "[Sources]\n"
        "- Formal setup for RQ1: local context, risk set, and observed share.\n"
        "- Ashudeep Singh and Thorsten Joachims. Fairness of Exposure in Rankings. KDD 2018. DOI: 10.1145/3219819.3220088\n"
    ),
    (
        "[Sources]\n"
        "- The gap is a distance between observed and baseline-implied shares, not a direct point estimate of bias.\n"
        "- Dynamic/contestable factors are handled as a separate pathway rather than folded into pure merit.\n"
        "- Yunqi Li et al. Fairness in Recommendation: Foundations, Methods, and Applications. ACM TIST 2023. DOI: 10.1145/3610302\n"
    ),
    (
        "[Sources]\n"
        "- Bridge slide: the exposure accounting stays fixed while the fairness unit changes from copies to frames.\n"
        "- Repo context: /Volumes/T9/BlueSky/_build/2026-03-13/content_bias_pilot_readout.md\n"
    ),
    (
        "[Sources]\n"
        "- RQ2 begins with exposure share versus observed frame supply share inside one event/topic cluster.\n"
        "- Repo context: /Volumes/T9/BlueSky/_build/2026-03-13/content_bias_pilot_readout.md\n"
    ),
    (
        "[Sources]\n"
        "- Raw frame disparity is the supply-relative first-pass answer for RQ2.\n"
        "- Ashudeep Singh and Thorsten Joachims. Fairness of Exposure in Rankings. KDD 2018. DOI: 10.1145/3219819.3220088\n"
    ),
    (
        "[Sources]\n"
        "- Adjusted frame disparity reuses the same timing / aligned / dynamic ladder as RQ1 and then stratifies by surface.\n"
        "- Yunqi Li et al. Fairness in Recommendation: Foundations, Methods, and Applications. ACM TIST 2023. DOI: 10.1145/3610302\n"
    ),
    (
        "[Sources]\n"
        "- WYSIWYG versus WAE belongs at the interpretation stage, not the setup stage.\n"
        "- Meike Zehlike, Ke Yang, Julia Stoyanovich. Fairness in Ranking, Part I and Part II. ACM CSUR 2022.\n"
        "- Dynamic popularity and early trajectory cautions follow Li et al. 2023.\n"
    ),
]


def add_line(
    slide,
    *,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    width_pt: float = 1.0,
    color: RGBColor = THEME.rule,
) -> None:
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.width = Pt(width_pt)
    connector.line.color.rgb = color


def add_text(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    text: str,
    font_name: str,
    font_size: float,
    color: RGBColor = THEME.text,
    bold: bool = False,
    italic: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.1,
    margin_left_pt: float = 0,
    margin_right_pt: float = 0,
    margin_top_pt: float = 0,
    margin_bottom_pt: float = 0,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    frame.margin_left = Pt(margin_left_pt)
    frame.margin_right = Pt(margin_right_pt)
    frame.margin_top = Pt(margin_top_pt)
    frame.margin_bottom = Pt(margin_bottom_pt)
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.line_spacing = line_spacing
    run = paragraph.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(font_size)
    font.bold = bold
    font.italic = italic
    font.color.rgb = color


def add_paragraph_block(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    paragraphs: list[str],
    font_name: str,
    font_size: float,
    color: RGBColor = THEME.text,
    bold: bool = False,
    italic: bool = False,
    line_spacing: float = 1.12,
    space_after_pt: float = 9,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    frame.margin_left = Pt(0)
    frame.margin_right = Pt(0)
    frame.margin_top = Pt(0)
    frame.margin_bottom = Pt(0)
    first = True
    for text in paragraphs:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.alignment = PP_ALIGN.LEFT
        paragraph.line_spacing = line_spacing
        paragraph.space_after = Pt(space_after_pt)
        run = paragraph.add_run()
        run.text = text
        font = run.font
        font.name = font_name
        font.size = Pt(font_size)
        font.bold = bold
        font.italic = italic
        font.color.rgb = color


def add_fraction(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    prefix: str,
    numerator: str,
    denominator: str,
    prefix_w: float = 1.15,
    prefix_size: float = 24,
    frac_size: float = 28,
    caption: str | None = None,
    caption_size: float = 14,
) -> None:
    add_text(
        slide,
        x=x,
        y=y + 0.33,
        w=prefix_w,
        h=0.55,
        text=prefix,
        font_name=THEME.math_font,
        font_size=prefix_size,
        color=THEME.text,
        italic=True,
    )
    frac_x = x + prefix_w + 0.18
    frac_w = max(w - prefix_w - 0.18, 0.8)
    add_text(
        slide,
        x=frac_x,
        y=y,
        w=frac_w,
        h=0.48,
        text=numerator,
        font_name=THEME.math_font,
        font_size=frac_size,
        color=THEME.text,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_line(
        slide,
        x1=frac_x + 0.05,
        y1=y + 0.54,
        x2=frac_x + frac_w - 0.05,
        y2=y + 0.54,
        width_pt=1.4,
        color=THEME.text,
    )
    add_text(
        slide,
        x=frac_x,
        y=y + 0.57,
        w=frac_w,
        h=0.48,
        text=denominator,
        font_name=THEME.math_font,
        font_size=frac_size - 1,
        color=THEME.text,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    if caption:
        add_text(
            slide,
            x=x,
            y=y + 1.08,
            w=w,
            h=0.28,
            text=caption,
            font_name=THEME.body_font,
            font_size=caption_size,
            color=THEME.muted,
            align=PP_ALIGN.CENTER,
        )


def add_scaffold(slide, *, kicker: str, title: str, page_number: int) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = THEME.bg
    add_text(
        slide,
        x=0.82,
        y=0.18,
        w=3.2,
        h=0.26,
        text=kicker.upper(),
        font_name=THEME.body_font,
        font_size=13,
        color=THEME.muted,
        bold=True,
    )
    add_text(
        slide,
        x=0.82,
        y=0.44,
        w=11.2,
        h=0.52,
        text=title,
        font_name=THEME.title_font,
        font_size=30,
        color=THEME.text,
    )
    add_line(slide, x1=0.82, y1=0.97, x2=12.58, y2=0.97, width_pt=1.0, color=THEME.rule)
    add_text(
        slide,
        x=12.28,
        y=0.18,
        w=0.36,
        h=0.22,
        text=str(page_number),
        font_name=THEME.body_font,
        font_size=12,
        color=THEME.muted,
        align=PP_ALIGN.RIGHT,
    )


def add_section_label(slide, *, x: float, y: float, text: str) -> None:
    add_text(
        slide,
        x=x,
        y=y,
        w=2.9,
        h=0.26,
        text=text.upper(),
        font_name=THEME.body_font,
        font_size=14,
        color=THEME.muted,
        bold=True,
    )


def slide_1(slide) -> None:
    add_scaffold(
        slide,
        kicker="Bluesky exposure audit",
        title="Research questions before the machinery",
        page_number=1,
    )
    add_text(
        slide,
        x=0.92,
        y=1.30,
        w=10.9,
        h=0.54,
        text="We ask how monitored top-K visibility is split across near-identical copies and across frames inside the same event.",
        font_name=THEME.body_font,
        font_size=20,
        color=THEME.text,
    )
    add_section_label(slide, x=0.92, y=2.12, text="RQ1")
    add_paragraph_block(
        slide,
        x=0.92,
        y=2.45,
        w=11.0,
        h=1.05,
        paragraphs=[
            "After timing and local panel context are accounted for, do near-duplicate post versions still receive systematically unequal monitored exposure opportunities?"
        ],
        font_name=THEME.title_font,
        font_size=20,
        color=THEME.text,
        line_spacing=1.08,
        space_after_pt=0,
    )
    add_line(slide, x1=0.92, y1=3.62, x2=12.05, y2=3.62, width_pt=0.9, color=THEME.rule)
    add_section_label(slide, x=0.92, y=3.90, text="RQ2")
    add_paragraph_block(
        slide,
        x=0.92,
        y=4.23,
        w=11.0,
        h=1.15,
        paragraphs=[
            "Within the same event or topic cluster, does monitored exposure across narrative frames track frame supply, or does it drift in a patterned way after adjustment?"
        ],
        font_name=THEME.title_font,
        font_size=20,
        color=THEME.text,
        line_spacing=1.08,
        space_after_pt=0,
    )
    add_text(
        slide,
        x=0.92,
        y=6.42,
        w=11.15,
        h=0.32,
        text="Flow: why exposure is the unit, then RQ1, then the bridge to RQ2, then how to read the remainder.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_2(slide) -> None:
    add_scaffold(
        slide,
        kicker="Step 1",
        title="Why exposure is the right object",
        page_number=2,
    )
    add_paragraph_block(
        slide,
        x=0.92,
        y=1.46,
        w=5.25,
        h=2.55,
        paragraphs=[
            "A monitored top-K appearance is not a full impression log. It is an observed opportunity to be seen.",
            "The front of the panel matters more than the back because viewers attend to higher ranks first.",
            "So our accounting object is not just “did the post appear?” It is “how much monitored visibility did that appearance carry?”",
        ],
        font_name=THEME.body_font,
        font_size=20,
        color=THEME.text,
        line_spacing=1.14,
        space_after_pt=11,
    )
    add_line(slide, x1=6.45, y1=1.38, x2=6.45, y2=5.88, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=6.82, y=1.46, text="One simple intuition")
    add_text(
        slide,
        x=6.82,
        y=1.95,
        w=5.0,
        h=0.38,
        text="Front ranks are worth more than lower ranks.",
        font_name=THEME.body_font,
        font_size=20,
        color=THEME.text,
    )
    add_text(
        slide,
        x=7.00,
        y=2.72,
        w=4.6,
        h=0.35,
        text="rank 1  ->  highest monitored attention",
        font_name=THEME.math_font,
        font_size=26,
        color=THEME.text,
    )
    add_text(
        slide,
        x=7.00,
        y=3.48,
        w=4.6,
        h=0.35,
        text="rank 5  ->  less",
        font_name=THEME.math_font,
        font_size=26,
        color=THEME.text,
    )
    add_text(
        slide,
        x=7.00,
        y=4.24,
        w=4.6,
        h=0.35,
        text="rank 10 ->  less still",
        font_name=THEME.math_font,
        font_size=26,
        color=THEME.text,
    )
    add_text(
        slide,
        x=0.92,
        y=6.16,
        w=11.15,
        h=0.36,
        text="This is the Singh-Joachims idea in plain English: rankings allocate visibility, and position is the mechanism that allocates it.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_3(slide) -> None:
    add_scaffold(
        slide,
        kicker="Step 1",
        title="Define monitored exposure",
        page_number=3,
    )
    add_line(slide, x1=6.45, y1=1.38, x2=6.45, y2=5.88, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=1.44, text="Position weight")
    add_fraction(
        slide,
        x=1.08,
        y=2.12,
        w=4.60,
        prefix="w(r) =",
        numerator="1",
        denominator="log₂(1 + r)",
        caption="Higher ranks carry more monitored attention.",
    )
    add_section_label(slide, x=6.82, y=1.44, text="Post total")
    add_text(
        slide,
        x=6.98,
        y=2.22,
        w=5.0,
        h=0.42,
        text="Eᵢ = Σₘ∈Mᵢ w(rₘ)",
        font_name=THEME.math_font,
        font_size=28,
        color=THEME.text,
        italic=True,
    )
    add_text(
        slide,
        x=6.98,
        y=2.90,
        w=4.92,
        h=0.68,
        text="Add every monitored appearance of post i, but give more credit to appearances that occur higher in the panel.",
        font_name=THEME.body_font,
        font_size=17,
        color=THEME.text,
        line_spacing=1.12,
    )
    add_line(slide, x1=0.92, y1=4.12, x2=12.00, y2=4.12, width_pt=0.9, color=THEME.rule)
    add_section_label(slide, x=0.92, y=4.34, text="Within one context")
    add_text(
        slide,
        x=1.08,
        y=4.82,
        w=5.35,
        h=0.42,
        text="yᵢg = Σₘ∈Mᵢg w(rₘ)",
        font_name=THEME.math_font,
        font_size=27,
        color=THEME.text,
        italic=True,
    )
    add_text(
        slide,
        x=7.04,
        y=4.74,
        w=4.86,
        h=0.78,
        text="g = (feed, hour, viewer, vantage)\nKeep the accounting local before you compare two versions.",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        line_spacing=1.16,
    )


def slide_4(slide) -> None:
    add_scaffold(
        slide,
        kicker="Step 2A",
        title="RQ1 starts with a fair local comparison",
        page_number=4,
    )
    add_paragraph_block(
        slide,
        x=0.92,
        y=1.46,
        w=11.0,
        h=1.06,
        paragraphs=[
            "We do not claim the posts describe the same real-world event. We use a weaker and more defensible unit: versions that were similar enough, and close enough in time, to compete for the same monitored opportunity."
        ],
        font_name=THEME.body_font,
        font_size=19,
        color=THEME.text,
        line_spacing=1.12,
        space_after_pt=0,
    )
    add_line(slide, x1=4.40, y1=2.92, x2=4.40, y2=5.30, width_pt=0.9, color=THEME.rule)
    add_line(slide, x1=8.06, y1=2.92, x2=8.06, y2=5.30, width_pt=0.9, color=THEME.rule)
    add_text(
        slide,
        x=1.04,
        y=3.10,
        w=2.9,
        h=0.72,
        text="same\ncontent family",
        font_name=THEME.title_font,
        font_size=23,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=4.72,
        y=3.10,
        w=2.9,
        h=0.72,
        text="same local\npanel context",
        font_name=THEME.title_font,
        font_size=23,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=8.38,
        y=3.10,
        w=2.9,
        h=0.72,
        text="same competition\nwindow",
        font_name=THEME.title_font,
        font_size=23,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=1.02,
        y=4.26,
        w=10.96,
        h=0.74,
        text="Main analysis: exact duplicates. Robustness tiers: same-link duplicates and near-text duplicates.",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=0.92,
        y=6.18,
        w=11.20,
        h=0.34,
        text="First fix the comparison unit. Only then introduce the symbols that describe it.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_5(slide) -> None:
    add_scaffold(
        slide,
        kicker="Step 2A",
        title="RQ1 formal setup",
        page_number=5,
    )
    add_line(slide, x1=6.45, y1=1.38, x2=6.45, y2=5.94, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=1.46, text="Context and risk set")
    add_text(
        slide,
        x=1.08,
        y=2.06,
        w=5.02,
        h=0.42,
        text="g = (feed, hour, viewer, vantage)",
        font_name=THEME.math_font,
        font_size=24,
        color=THEME.text,
        italic=True,
    )
    add_text(
        slide,
        x=1.08,
        y=2.90,
        w=5.02,
        h=0.58,
        text="R_cg = { i : c(i)=c, tᵢ ≤ τg, τg − tᵢ ≤ H }",
        font_name=THEME.math_font,
        font_size=24,
        color=THEME.text,
        italic=True,
    )
    add_paragraph_block(
        slide,
        x=1.08,
        y=3.80,
        w=5.05,
        h=1.10,
        paragraphs=[
            "At snapshot time τg, compare only the duplicate versions that were already available and still inside the local comparison horizon H.",
        ],
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        line_spacing=1.12,
        space_after_pt=0,
    )
    add_section_label(slide, x=6.82, y=1.46, text="Observed share")
    add_fraction(
        slide,
        x=7.00,
        y=2.12,
        w=4.70,
        prefix="qᵢg =",
        numerator="Eᵢg",
        denominator="Σᵤ∈Rcg Eᵤg",
        caption="The share that version i actually receives inside that duplicate-local comparison.",
    )
    add_text(
        slide,
        x=0.92,
        y=6.18,
        w=11.15,
        h=0.34,
        text="RQ1 is about how that local share gets split across copies of essentially the same content.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_6(slide) -> None:
    add_scaffold(
        slide,
        kicker="RQ1 target",
        title="RQ1 asks how far the split moves away from a baseline",
        page_number=6,
    )
    add_text(
        slide,
        x=1.00,
        y=1.58,
        w=10.96,
        h=0.40,
        text="D_M(c,g) = ½ Σᵢ∈Rcg | qᵢg − πᵢg(M) |",
        font_name=THEME.math_font,
        font_size=29,
        color=THEME.text,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=1.40,
        y=2.28,
        w=10.16,
        h=0.52,
        text="qᵢg is the observed share. πᵢg(M) is the share implied by a baseline M. The gap asks how much the actual split departs from that baseline.",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x1=0.92, y1=3.18, x2=12.02, y2=3.18, width_pt=0.9, color=THEME.rule)
    add_section_label(slide, x=0.92, y=3.34, text="Baseline ladder")
    add_paragraph_block(
        slide,
        x=1.08,
        y=3.84,
        w=5.86,
        h=1.50,
        paragraphs=[
            "M₀  : timing only",
            "M_A : timing + declared-objective factors",
            "M_D : timing + dynamic / contestable factors",
        ],
        font_name=THEME.math_font,
        font_size=24,
        color=THEME.text,
        italic=True,
        line_spacing=1.18,
        space_after_pt=9,
    )
    add_section_label(slide, x=7.10, y=3.34, text="What we read from it")
    add_paragraph_block(
        slide,
        x=7.24,
        y=3.84,
        w=4.55,
        h=1.72,
        paragraphs=[
            "timing share",
            "aligned share",
            "dynamic / contestable share",
            "unexplained share",
        ],
        font_name=THEME.body_font,
        font_size=19,
        color=THEME.text,
        line_spacing=1.22,
        space_after_pt=9,
    )
    add_text(
        slide,
        x=0.92,
        y=6.16,
        w=11.15,
        h=0.34,
        text="Important reading rule: early trajectory is evidence of a dynamic pathway, not a free pass that proves merit.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_7(slide) -> None:
    add_scaffold(
        slide,
        kicker="Bridge",
        title="Same exposure engine, new comparison object",
        page_number=7,
    )
    add_line(slide, x1=6.45, y1=1.46, x2=6.45, y2=5.88, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=1.52, text="RQ1")
    add_text(
        slide,
        x=1.06,
        y=2.06,
        w=4.94,
        h=0.56,
        text="Compare copies inside one duplicate-local risk set.",
        font_name=THEME.title_font,
        font_size=24,
        color=THEME.text,
    )
    add_text(
        slide,
        x=1.10,
        y=3.08,
        w=4.90,
        h=0.42,
        text="question: which copy gets the monitored visibility?",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
    )
    add_text(
        slide,
        x=1.10,
        y=4.12,
        w=4.86,
        h=0.42,
        text="local object: qᵢg",
        font_name=THEME.math_font,
        font_size=28,
        color=THEME.text,
        italic=True,
    )
    add_section_label(slide, x=6.82, y=1.52, text="RQ2")
    add_text(
        slide,
        x=6.96,
        y=2.06,
        w=4.94,
        h=0.56,
        text="Compare frames inside one event or topic cluster.",
        font_name=THEME.title_font,
        font_size=24,
        color=THEME.text,
    )
    add_text(
        slide,
        x=7.00,
        y=3.08,
        w=4.90,
        h=0.42,
        text="question: which narrative frame gets the monitored visibility?",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
    )
    add_text(
        slide,
        x=7.00,
        y=4.12,
        w=4.86,
        h=0.42,
        text="local objects: x_fhg and s_fhg",
        font_name=THEME.math_font,
        font_size=28,
        color=THEME.text,
        italic=True,
    )
    add_text(
        slide,
        x=0.92,
        y=6.18,
        w=11.15,
        h=0.34,
        text="Nothing about the exposure accounting changes here. Only the fairness unit changes.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_8(slide) -> None:
    add_scaffold(
        slide,
        kicker="Step 2B",
        title="RQ2 starts with two shares inside one event",
        page_number=8,
    )
    add_line(slide, x1=6.45, y1=1.38, x2=6.45, y2=5.88, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=1.46, text="Frame exposure share")
    add_fraction(
        slide,
        x=1.04,
        y=2.12,
        w=5.10,
        prefix="x_fhg =",
        numerator="Y_fhg",
        denominator="Σᵤ Y_uhg",
        caption="Share of monitored exposure that goes to frame f inside event h and context g.",
    )
    add_section_label(slide, x=6.82, y=1.46, text="Frame supply share")
    add_fraction(
        slide,
        x=6.98,
        y=2.12,
        w=5.06,
        prefix="s_fhg =",
        numerator="N_fhg",
        denominator="Σᵤ N_uhg",
        caption="Share of observed candidate posts that belong to frame f inside the same event and context.",
    )
    add_text(
        slide,
        x=1.02,
        y=4.78,
        w=11.0,
        h=0.62,
        text="Y_fhg measures where the monitored exposure went. N_fhg measures what was available to receive it.",
        font_name=THEME.body_font,
        font_size=19,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=0.92,
        y=6.16,
        w=11.15,
        h=0.34,
        text="RQ2 begins with a supply-relative comparison: exposure share versus frame supply share.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_9(slide) -> None:
    add_scaffold(
        slide,
        kicker="RQ2 raw disparity",
        title="RQ2 asks whether frame exposure matches frame supply",
        page_number=9,
    )
    add_text(
        slide,
        x=1.00,
        y=1.62,
        w=10.96,
        h=0.40,
        text="D_supply(h,g) = ½ Σ_f∈Fₕ | x_fhg − s_fhg |",
        font_name=THEME.math_font,
        font_size=29,
        color=THEME.text,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=1.36,
        y=2.28,
        w=10.22,
        h=0.50,
        text="D_supply(h,g) = 0 means the monitored exposure split matches the observed frame supply split inside that event and context.",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x1=0.92, y1=3.18, x2=12.02, y2=3.18, width_pt=0.9, color=THEME.rule)
    add_section_label(slide, x=0.92, y=3.38, text="Per-frame signed reading")
    add_text(
        slide,
        x=1.06,
        y=4.02,
        w=5.12,
        h=0.40,
        text="L_fhg = log [ (x_fhg + ε) / (s_fhg + ε) ]",
        font_name=THEME.math_font,
        font_size=25,
        color=THEME.text,
        italic=True,
    )
    add_paragraph_block(
        slide,
        x=6.90,
        y=3.92,
        w=4.80,
        h=1.42,
        paragraphs=[
            "L_fhg > 0  -> over-exposed",
            "L_fhg < 0  -> under-exposed",
            "L_fhg = 0  -> supply-matched",
        ],
        font_name=THEME.body_font,
        font_size=19,
        color=THEME.text,
        line_spacing=1.2,
        space_after_pt=8,
    )
    add_text(
        slide,
        x=0.92,
        y=6.16,
        w=11.15,
        h=0.34,
        text="This is the clean first-pass answer to RQ2 before any adjusted baseline explains part of the disparity away.",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_10(slide) -> None:
    add_scaffold(
        slide,
        kicker="RQ2 adjusted gap",
        title="RQ2 then asks how much frame disparity remains after adjustment",
        page_number=10,
    )
    add_text(
        slide,
        x=0.98,
        y=1.56,
        w=11.02,
        h=0.42,
        text="D_M^frame(h,g) = ½ Σ_f | x_fhg − π_fhg(M) |",
        font_name=THEME.math_font,
        font_size=29,
        color=THEME.text,
        italic=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        x=1.24,
        y=2.22,
        w=10.52,
        h=0.52,
        text="π_fhg(M) is the model-implied frame share produced by the same post-level ladder we already used for RQ1.",
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, x1=0.92, y1=3.14, x2=12.02, y2=3.14, width_pt=0.9, color=THEME.rule)
    add_line(slide, x1=6.45, y1=3.34, x2=6.45, y2=5.90, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=3.34, text="Same ladder")
    add_paragraph_block(
        slide,
        x=1.08,
        y=3.88,
        w=5.16,
        h=1.48,
        paragraphs=[
            "M₀  : timing only",
            "M_A : timing + declared-objective factors",
            "M_D : timing + dynamic / contestable factors",
        ],
        font_name=THEME.math_font,
        font_size=22,
        color=THEME.text,
        italic=True,
        line_spacing=1.18,
        space_after_pt=9,
    )
    add_section_label(slide, x=6.82, y=3.34, text="What we stratify")
    add_paragraph_block(
        slide,
        x=6.98,
        y=3.88,
        w=4.92,
        h=1.48,
        paragraphs=[
            "by feed bucket",
            "by viewer mode",
            "by surface or moderation environment",
        ],
        font_name=THEME.body_font,
        font_size=19,
        color=THEME.text,
        line_spacing=1.2,
        space_after_pt=8,
    )
    add_text(
        slide,
        x=0.92,
        y=6.22,
        w=11.15,
        h=0.34,
        text="RQ2 is not just “does supply differ?” It is “after baseline M, which frame gaps are still left?”",
        font_name=THEME.body_font,
        font_size=15,
        color=THEME.muted,
    )


def slide_11(slide) -> None:
    add_scaffold(
        slide,
        kicker="Interpretation",
        title="How to read the remainder",
        page_number=11,
    )
    add_line(slide, x1=6.45, y1=1.40, x2=6.45, y2=5.92, width_pt=1.0, color=THEME.rule)
    add_section_label(slide, x=0.92, y=1.48, text="WYSIWYG")
    add_paragraph_block(
        slide,
        x=1.06,
        y=2.00,
        w=4.98,
        h=2.48,
        paragraphs=[
            "observable signals ≈ legitimate merit",
            "S_U = unexplained after observed-factor adjustment",
            "S_U is not automatically bias",
            "The remainder may still include omitted legitimate signals, hidden platform objectives, stochasticity, or measurement error.",
        ],
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        line_spacing=1.17,
        space_after_pt=9,
    )
    add_section_label(slide, x=6.82, y=1.48, text="WAE")
    add_paragraph_block(
        slide,
        x=6.96,
        y=2.00,
        w=4.96,
        h=2.52,
        paragraphs=[
            "observed popularity and early trajectory may already contain prior unequal exposure",
            "what looks “explained” can still be part of dynamic amplification",
            "bias-consistent remainder is larger than S_U alone",
            "under this reading, dynamic / contestable share and unexplained share both matter",
        ],
        font_name=THEME.body_font,
        font_size=18,
        color=THEME.text,
        line_spacing=1.17,
        space_after_pt=9,
    )
    add_line(slide, x1=0.92, y1=5.72, x2=12.02, y2=5.72, width_pt=0.9, color=THEME.rule)
    add_text(
        slide,
        x=0.92,
        y=5.96,
        w=11.18,
        h=0.44,
        text="Final reading rule: report the last remainder as conditional on the observed factor set, not as a point estimate of “bias.”",
        font_name=THEME.body_font,
        font_size=16,
        color=THEME.muted,
        align=PP_ALIGN.CENTER,
    )


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide_builders = [
        slide_1,
        slide_2,
        slide_3,
        slide_4,
        slide_5,
        slide_6,
        slide_7,
        slide_8,
        slide_9,
        slide_10,
        slide_11,
    ]
    for index, builder in enumerate(slide_builders):
        slide = prs.slides.add_slide(blank)
        builder(slide)
        notes_frame = slide.notes_slide.notes_text_frame
        notes_frame.text = SLIDE_NOTES[index]
    return prs


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Build a cleaner white-background editable deck for the Bluesky exposure model slides."
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=Path(
            "/Volumes/T9/BlueSky/Bluesky_RQ1_RQ2_exposure_model_restructured_minimal_math_v5_cleanflow_editable.pptx"
        ),
        help="Output PPTX path.",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    prs = build_presentation()
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))
    print(f"OK: wrote {args.out}")


if __name__ == "__main__":
    main()
