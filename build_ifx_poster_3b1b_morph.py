#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
import shutil
import tempfile
import zipfile
from dataclasses import dataclass, replace
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

import build_ifx_poster_svg as base_svg


PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
P159_NS = "http://schemas.microsoft.com/office/powerpoint/2015/09/main"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
BG_HEX = "0B1320"


@dataclass(frozen=True)
class SvgAsset:
    name: str
    slide_num: int
    order: int
    x_in: float
    y_in: float
    w_in: float
    h_in: float
    svg_text: str


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _emu_to_in(emu: int) -> float:
    return float(emu) / 914400.0


def _extract_svg_assets(pptx_path: Path) -> dict[str, SvgAsset]:
    assets: dict[str, SvgAsset] = {}

    with zipfile.ZipFile(pptx_path) as zf:
        slide_parts = sorted(
            [name for name in zf.namelist() if re.fullmatch(r"ppt/slides/slide\d+\.xml", name)],
            key=lambda s: int(re.search(r"slide(\d+)\.xml", s).group(1)),
        )

        for slide_part in slide_parts:
            slide_num = int(re.search(r"slide(\d+)\.xml", slide_part).group(1))
            rels_part = f"ppt/slides/_rels/slide{slide_num}.xml.rels"

            slide_root = etree.fromstring(zf.read(slide_part))
            rels_root = etree.fromstring(zf.read(rels_part))
            rel_by_id = {
                rel.get("Id"): rel.get("Target")
                for rel in rels_root.findall(f"{{{PKG_REL_NS}}}Relationship")
            }

            pics = slide_root.xpath("./p:cSld/p:spTree/p:pic", namespaces={"p": PML_NS})
            for order, pic in enumerate(pics):
                c_nv_pr = pic.find("./p:nvPicPr/p:cNvPr", namespaces={"p": PML_NS})
                blip = pic.find(".//a:blip", namespaces={"a": A_NS})
                off = pic.find(".//a:xfrm/a:off", namespaces={"a": A_NS})
                ext = pic.find(".//a:xfrm/a:ext", namespaces={"a": A_NS})
                if c_nv_pr is None or blip is None or off is None or ext is None:
                    continue

                name = c_nv_pr.get("name") or ""
                rid = blip.get(f"{{{R_NS}}}embed")
                target = rel_by_id.get(rid or "")
                if not name or not target:
                    continue
                if not target.lower().endswith(".svg"):
                    raise ValueError(f"Asset {name} is not SVG-backed: {target}")

                media_part = str(Path("ppt/slides").joinpath(target).resolve()).replace("\\", "/")
                # Resolve /ppt/slides/../media/... style paths back into package-relative form.
                media_part = re.sub(r"^[A-Za-z]:", "", media_part)
                media_part = media_part.replace("/ppt/slides/../", "/ppt/")
                media_part = media_part.lstrip("/")
                if media_part not in zf.namelist():
                    alt = f"ppt/media/{Path(target).name}"
                    if alt in zf.namelist():
                        media_part = alt
                if media_part not in zf.namelist():
                    raise ValueError(f"Missing SVG media for {name}: {target}")

                assets[name] = SvgAsset(
                    name=name,
                    slide_num=slide_num,
                    order=order,
                    x_in=_emu_to_in(int(off.get("x"))),
                    y_in=_emu_to_in(int(off.get("y"))),
                    w_in=_emu_to_in(int(ext.get("cx"))),
                    h_in=_emu_to_in(int(ext.get("cy"))),
                    svg_text=zf.read(media_part).decode("utf-8"),
                )

    return assets


def _prefixed(name: str) -> str:
    return f"!!{name}"


def _clone(asset: SvgAsset, **kwargs: float | int | str) -> SvgAsset:
    return replace(asset, **kwargs)


def _bbox(scene: dict[str, SvgAsset], names: list[str]) -> tuple[float, float, float, float]:
    xs = [scene[name].x_in for name in names]
    ys = [scene[name].y_in for name in names]
    rights = [scene[name].x_in + scene[name].w_in for name in names]
    bottoms = [scene[name].y_in + scene[name].h_in for name in names]
    x0 = min(xs)
    y0 = min(ys)
    return x0, y0, max(rights) - x0, max(bottoms) - y0


def _scale_group(scene: dict[str, SvgAsset], names: list[str], factor: float) -> None:
    x0, y0, w, h = _bbox(scene, names)
    cx = x0 + w / 2.0
    cy = y0 + h / 2.0
    for name in names:
        asset = scene[name]
        new_w = asset.w_in * factor
        new_h = asset.h_in * factor
        old_cx = asset.x_in + asset.w_in / 2.0
        old_cy = asset.y_in + asset.h_in / 2.0
        new_cx = cx + (old_cx - cx) * factor
        new_cy = cy + (old_cy - cy) * factor
        scene[name] = _clone(
            asset,
            x_in=new_cx - new_w / 2.0,
            y_in=new_cy - new_h / 2.0,
            w_in=new_w,
            h_in=new_h,
        )


def _shift_group(scene: dict[str, SvgAsset], names: list[str], dx: float = 0.0, dy: float = 0.0) -> None:
    for name in names:
        asset = scene[name]
        scene[name] = _clone(asset, x_in=asset.x_in + dx, y_in=asset.y_in + dy)


def _move_group_offscreen(scene: dict[str, SvgAsset], names: list[str], direction: str, margin: float = 0.4) -> None:
    x0, y0, w, h = _bbox(scene, names)
    if direction == "left":
        dx = -x0 - w - margin
        dy = 0.0
    elif direction == "right":
        dx = SLIDE_W_IN + margin - x0
        dy = 0.0
    elif direction == "top":
        dx = 0.0
        dy = -y0 - h - margin
    elif direction == "bottom":
        dx = 0.0
        dy = SLIDE_H_IN + margin - y0
    else:
        raise ValueError(f"Unknown direction: {direction}")
    _shift_group(scene, names, dx=dx, dy=dy)


def _names_by_prefix(all_names: list[str], *prefixes: str) -> list[str]:
    return [name for name in all_names if any(name.startswith(prefix) for prefix in prefixes)]


def _subset(base_assets: dict[str, SvgAsset], names: list[str]) -> dict[str, SvgAsset]:
    return {name: _clone(base_assets[name]) for name in names}


def _build_scenes(base_assets: dict[str, SvgAsset]) -> list[tuple[str, str]]:
    names = sorted(base_assets)

    g = {
        "S1_FRAME": [
            "_BG_S1_GRID",
            "STATIC_S1_TITLE",
            "_BG_S1_PANEL_THEORY",
            "_BG_S1_PANEL_DATA",
            "STATIC_S1_LBL_THEORY",
            "STATIC_S1_LBL_DATA",
        ],
        "S1_EQ": ["_BG_S1_EQ_LINKS"] + _names_by_prefix(names, "S1_EQSTRIP_"),
        "S1_MATRIX": ["STATIC_S1_PILL_P", "_BG_S1_P_ROW_HI"] + _names_by_prefix(names, "S1_P_"),
        "S1_V": ["STATIC_S1_PILL_V"] + _names_by_prefix(names, "S1_V_") + ["S1_ARROW_PV", "S1_OUT_E"],
        "S1_DS": ["_BG_S1_DS_CARD", "STATIC_S1_PILL_DS"] + _names_by_prefix(names, "S1_DS_"),
        "S1_DATA": (
            _names_by_prefix(names, "S1_Q_")
            + ["STATIC_S1_PILL_TAU"]
            + _names_by_prefix(names, "S1_RANK_")
            + [
                "_BG_S1_CARD_PIPE",
                "STATIC_S1_PILL_PIPE_SNAP",
                "S1_PIPE_EQ_E",
                "STATIC_S1_PIPE_NOTE_E0",
                "STATIC_S1_PILL_PIPE_PHAT",
                "S1_PIPE_EQ_PH1",
                "S1_PIPE_EQ_PH2",
                "S1_JOIN_GRAPH",
                "_BG_S1_CARD_PROMO",
                "STATIC_S1_PILL_PROMO",
                "S1_EQ_TOPK_Y",
                "S1_PROMO_BAR",
            ]
        ),
        "S2_FRAME": [
            "_BG_S2_GRID",
            "STATIC_S2_TITLE",
            "_BG_S2_PANEL_LEFT",
            "_BG_S2_PANEL_RIGHT",
            "STATIC_S2_LBL_SIM",
            "STATIC_S2_LBL_METRICS",
        ],
        "S2_LEFT": ["_BG_S2_LINKS_CLUSTER", "S2_POST_ALL", "S2_CLUSTER", "S2_EQ_SIM", "S2_AUDITS"],
        "S2_CARDS": [
            "_BG_S2_CARD_D",
            "STATIC_S2_CARD_01_PILL",
            "STATIC_S2_CARD_01_SUB",
            "_BG_S2_CARD_LIP",
            "STATIC_S2_CARD_02_PILL",
            "STATIC_S2_CARD_02_SUB",
            "_BG_S2_CARD_DELTA",
            "STATIC_S2_CARD_03_PILL",
            "STATIC_S2_CARD_03_SUB",
            "S2_EQ_D",
            "S2_EQ_LIP",
            "S2_EQ_DELTA",
        ],
        "S2_HEAT": ["STATIC_S2_PILL_HEAT", "S2_HEATMAP", "_BG_S2_HEAT_OVERLAY"],
        "S2_TIMELINE": ["S2_TIMELINE"],
        "S2_OPS": ["S2_OPS"],
    }

    scene_specs: list[tuple[str, dict[str, SvgAsset]]] = []

    # Scene 1: open frame, equation staged above the slide.
    s1 = _subset(base_assets, g["S1_FRAME"] + g["S1_EQ"])
    _move_group_offscreen(s1, g["S1_EQ"], "top")
    scene_specs.append(("none", s1))

    # Scene 2: equation enters; matrix staged below for the next beat.
    s2 = _subset(base_assets, g["S1_FRAME"] + g["S1_EQ"] + g["S1_MATRIX"])
    _move_group_offscreen(s2, g["S1_MATRIX"], "bottom")
    scene_specs.append(("morph", s2))

    # Scene 3: theory focus with the P matrix enlarged; v and exposure staged at right.
    s3 = _subset(base_assets, g["S1_FRAME"] + g["S1_EQ"] + g["S1_MATRIX"] + g["S1_V"] + g["S1_DS"])
    _scale_group(s3, g["S1_EQ"] + g["S1_MATRIX"], 1.08)
    _shift_group(s3, g["S1_EQ"] + g["S1_MATRIX"], dx=0.10, dy=0.12)
    _move_group_offscreen(s3, g["S1_V"] + g["S1_DS"], "right")
    scene_specs.append(("morph", s3))

    # Scene 4: position bias and exposure mechanics arrive; data mapping staged at right.
    s4 = _subset(base_assets, g["S1_FRAME"] + g["S1_EQ"] + g["S1_MATRIX"] + g["S1_V"] + g["S1_DS"] + g["S1_DATA"])
    _move_group_offscreen(s4, g["S1_DATA"], "right")
    scene_specs.append(("morph", s4))

    # Scene 5: full application view.
    s5 = _subset(base_assets, g["S1_FRAME"] + g["S1_EQ"] + g["S1_MATRIX"] + g["S1_V"] + g["S1_DS"] + g["S1_DATA"])
    scene_specs.append(("morph", s5))

    # Scene 6: topic cut to the second paper, left-side content staged from the left.
    s6 = _subset(base_assets, g["S2_FRAME"] + g["S2_LEFT"])
    _move_group_offscreen(s6, g["S2_LEFT"], "left")
    scene_specs.append(("fade", s6))

    # Scene 7: similarity framing; gap cards staged from the right.
    s7 = _subset(base_assets, g["S2_FRAME"] + g["S2_LEFT"] + g["S2_CARDS"])
    _scale_group(s7, g["S2_LEFT"], 1.05)
    _shift_group(s7, g["S2_LEFT"], dx=0.06, dy=0.03)
    _move_group_offscreen(s7, g["S2_CARDS"], "right")
    scene_specs.append(("morph", s7))

    # Scene 8: gap metrics enter; heatmap and timeline staged from below.
    s8 = _subset(base_assets, g["S2_FRAME"] + g["S2_LEFT"] + g["S2_CARDS"] + g["S2_HEAT"] + g["S2_TIMELINE"])
    _move_group_offscreen(s8, g["S2_HEAT"] + g["S2_TIMELINE"], "bottom")
    scene_specs.append(("morph", s8))

    # Scene 9: right panel focus with heatmap and trend line enlarged; ops staged below.
    s9 = _subset(base_assets, g["S2_FRAME"] + g["S2_LEFT"] + g["S2_CARDS"] + g["S2_HEAT"] + g["S2_TIMELINE"] + g["S2_OPS"])
    _scale_group(s9, g["S2_CARDS"] + g["S2_HEAT"] + g["S2_TIMELINE"], 1.04)
    _shift_group(s9, g["S2_CARDS"] + g["S2_HEAT"] + g["S2_TIMELINE"], dx=0.05, dy=0.02)
    _move_group_offscreen(s9, g["S2_OPS"], "bottom")
    scene_specs.append(("morph", s9))

    # Scene 10: final integrated slide.
    s10 = _subset(base_assets, g["S2_FRAME"] + g["S2_LEFT"] + g["S2_CARDS"] + g["S2_HEAT"] + g["S2_TIMELINE"] + g["S2_OPS"])
    scene_specs.append(("morph", s10))

    return [(kind, name) for kind, name in scene_specs]


def _scene_assets_list(scene: dict[str, SvgAsset]) -> list[SvgAsset]:
    return sorted(scene.values(), key=lambda asset: (asset.slide_num, asset.order, asset.name))


def _set_transition(root: etree._Element, kind: str, dur_ms: int = 1100) -> None:
    for child in list(root):
        if etree.QName(child).localname == "transition":
            root.remove(child)
        elif child.tag == f"{{{MC_NS}}}AlternateContent":
            root.remove(child)

    if kind == "none":
        return

    insert_at = len(root)
    for idx, child in enumerate(root):
        if etree.QName(child).localname == "timing":
            insert_at = idx
            break

    if kind == "fade":
        trans = etree.Element(f"{{{PML_NS}}}transition", spd="slow")
        etree.SubElement(trans, f"{{{PML_NS}}}fade")
        root.insert(insert_at, trans)
        return

    if kind != "morph":
        raise ValueError(f"Unknown transition kind: {kind}")

    alt = etree.Element(
        f"{{{MC_NS}}}AlternateContent",
        nsmap={"mc": MC_NS, "p14": P14_NS, "p159": P159_NS},
    )
    choice = etree.SubElement(alt, f"{{{MC_NS}}}Choice", Requires="p159")
    trans = etree.SubElement(choice, f"{{{PML_NS}}}transition")
    trans.set(f"{{{P14_NS}}}dur", str(dur_ms))
    etree.SubElement(trans, f"{{{P159_NS}}}morph", option="byObject")

    fallback = etree.SubElement(alt, f"{{{MC_NS}}}Fallback")
    fallback_trans = etree.SubElement(fallback, f"{{{PML_NS}}}transition", spd="slow")
    etree.SubElement(fallback_trans, f"{{{PML_NS}}}fade")
    root.insert(insert_at, alt)


def _patch_transitions(pptx_in: Path, pptx_out: Path, transition_kinds: list[str]) -> None:
    with zipfile.ZipFile(pptx_in) as zin, zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
        for info in zin.infolist():
            data = zin.read(info.filename)
            m = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", info.filename)
            if m:
                slide_num = int(m.group(1))
                root = etree.fromstring(data)
                kind = transition_kinds[slide_num - 1] if slide_num - 1 < len(transition_kinds) else "none"
                _set_transition(root, kind=kind)
                data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
            zout.writestr(info, data)


def build_deck(pptx_in: Path, out_path: Path) -> None:
    base_assets = _extract_svg_assets(pptx_in)
    scene_specs = _build_scenes(base_assets)

    work_dir = Path(tempfile.mkdtemp(prefix="ifx_morph_"))
    try:
        svg_dir = work_dir / "svgs"
        png_dir = work_dir / "pngs"
        svg_dir.mkdir(parents=True, exist_ok=True)
        png_dir.mkdir(parents=True, exist_ok=True)

        shape_svg_map: dict[str, Path] = {}
        placeholder_png_map: dict[str, Path] = {}

        for asset in base_assets.values():
            svg_path = svg_dir / f"{asset.name}.svg"
            svg_path.write_text(asset.svg_text, encoding="utf-8")
            shape_name = _prefixed(asset.name)
            shape_svg_map[shape_name] = svg_path
            png_path = png_dir / f"{asset.name}.png"
            base_svg._write_unique_png(out_path=png_path, seed=asset.name)
            placeholder_png_map[shape_name] = png_path

        prs = Presentation()
        prs.slide_width = Inches(SLIDE_W_IN)
        prs.slide_height = Inches(SLIDE_H_IN)
        blank = prs.slide_layouts[6]

        for _, scene in scene_specs:
            slide = prs.slides.add_slide(blank)
            fill = slide.background.fill
            fill.solid()
            fill.fore_color.rgb = _rgb(BG_HEX)

            for asset in _scene_assets_list(scene):
                shape_name = _prefixed(asset.name)
                pic = slide.shapes.add_picture(
                    str(placeholder_png_map[shape_name]),
                    Inches(asset.x_in),
                    Inches(asset.y_in),
                    width=Inches(asset.w_in),
                    height=Inches(asset.h_in),
                )
                pic.name = shape_name

        pre_svg = work_dir / "pre_svg.pptx"
        prs.save(str(pre_svg))

        svg_pptx = work_dir / "svg.pptx"
        base_svg._patch_pptx_replace_images_with_svg(
            pptx_in=pre_svg,
            pptx_out=svg_pptx,
            shape_svg_map=shape_svg_map,
        )
        _patch_transitions(svg_pptx, out_path, [kind for kind, _ in scene_specs])
    finally:
        shutil.rmtree(work_dir, ignore_errors=True)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a morph-driven 3Blue1Brown-style IFX poster deck.")
    parser.add_argument(
        "pptx_in",
        type=Path,
        nargs="?",
        default=Path("_build/ifx_poster_animated_02_workcopy_20260305.pptx"),
        help="Input PPTX used as the visual asset source.",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("_build/ifx_poster_animated_02_3b1b_morph.pptx"),
        help="Output PPTX path.",
    )
    args = parser.parse_args()

    pptx_in = args.pptx_in.resolve()
    out_path = args.out.resolve()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    build_deck(pptx_in=pptx_in, out_path=out_path)
    print(f"OK: wrote {out_path}")


if __name__ == "__main__":
    main()
