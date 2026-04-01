#!/usr/bin/env python3
from __future__ import annotations

"""
Poster-style 2-slide deck for:
  - Singh & Joachims (Fairness of Exposure in Rankings)
  - Zehlike et al. (Fairness in Ranking, Part II)

Key constraints from user:
  - Use real SVG assets (embedded as SVG picture parts in PPTX).
  - Zero overlap collisions among non-background shapes.
  - Lots of click-to-reveal animations.
"""

import argparse
import hashlib
import io
import re
import sys
import zipfile
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches


# Canvas (widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Margins / safe area
SAFE_L_IN = 0.65
SAFE_R_IN = 0.65
SAFE_T_IN = 0.45
SAFE_B_IN = 0.45

SAFE_W_IN = SLIDE_W_IN - SAFE_L_IN - SAFE_R_IN
SAFE_H_IN = SLIDE_H_IN - SAFE_T_IN - SAFE_B_IN


COLORS_HEX: dict[str, str] = {
    "bg": "0B1320",
    "ink": "F7F9FC",
    "muted": "C9D2E3",
    "card": "162033",
    "card2": "0F1B2E",
    "card_line": "3A4D66",
    "cyan": "58B7E6",
    "amber": "F4C55D",
    "coral": "FF7D77",
    "green": "54C687",
    "violet": "8B7CFF",
}

FONT_SANS = "Calibri"


PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
PKG_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _escape_xml(text: str) -> str:
    return (
        text.replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
        .replace("'", "&apos;")
    )


def _in_to_px(inches: float) -> int:
    return int(round(inches * 96))


def _scale_to_fit(*, w_in: float, h_in: float, max_w_in: float, max_h_in: float | None = None) -> tuple[float, float]:
    if w_in <= 0 or h_in <= 0:
        return w_in, h_in
    scale = max_w_in / w_in
    if max_h_in is not None:
        scale = min(scale, max_h_in / h_in)
    scale = min(scale, 1.0)
    return w_in * scale, h_in * scale


def _svg_root(*, w_px: int, h_px: int, body: str) -> str:
    return (
        f'<svg xmlns="http://www.w3.org/2000/svg" width="{w_px}" height="{h_px}" '
        f'viewBox="0 0 {w_px} {h_px}" version="1.1" overflow="hidden">{body}</svg>'
    )


def _svg_rounded_rect(
    *, x: int, y: int, w: int, h: int, r: int, fill: str, stroke: str, stroke_w: int, opacity: float = 1.0
) -> str:
    op = f' opacity="{opacity:.4f}"' if opacity < 1.0 else ""
    return (
        f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="{r}" ry="{r}" '
        f'fill="{fill}" stroke="{stroke}" stroke-width="{stroke_w}"{op}/>'
    )


def _svg_text(
    *,
    x: int,
    y: int,
    text: str,
    size_px: int,
    color: str,
    weight: int = 400,
    anchor: str = "start",
    family: str = FONT_SANS,
    opacity: float = 1.0,
    baseline: str = "middle",
) -> str:
    op = f' opacity="{opacity:.4f}"' if opacity < 1.0 else ""
    return (
        f'<text x="{x}" y="{y}" fill="{color}" font-family="{_escape_xml(family)}" '
        f'font-size="{size_px}" font-weight="{weight}" text-anchor="{anchor}" '
        f'dominant-baseline="{baseline}"{op}>{_escape_xml(text)}</text>'
    )


def _svg_multiline(
    *,
    x: int,
    y: int,
    lines: list[str],
    size_px: int,
    color: str,
    weight: int = 400,
    anchor: str = "start",
    family: str = FONT_SANS,
    line_gap_px: int | None = None,
    baseline: str = "hanging",
) -> str:
    if line_gap_px is None:
        line_gap_px = int(round(size_px * 1.25))
    tspans = []
    for i, line in enumerate(lines):
        dy = 0 if i == 0 else line_gap_px
        tspans.append(f'<tspan x="{x}" dy="{dy}">{_escape_xml(line)}</tspan>')
    return (
        f'<text x="{x}" y="{y}" fill="{color}" font-family="{_escape_xml(family)}" '
        f'font-size="{size_px}" font-weight="{weight}" text-anchor="{anchor}" '
        f'dominant-baseline="{baseline}">{"".join(tspans)}</text>'
    )


def _svg_grid_bg(*, w_px: int, h_px: int) -> str:
    step = 48
    lines = []
    for x in range(0, w_px + 1, step):
        lines.append(
            f'<line x1="{x}" y1="0" x2="{x}" y2="{h_px}" stroke="#FFFFFF" stroke-opacity="0.035" stroke-width="1"/>'
        )
    for y in range(0, h_px + 1, step):
        lines.append(
            f'<line x1="0" y1="{y}" x2="{w_px}" y2="{y}" stroke="#FFFFFF" stroke-opacity="0.035" stroke-width="1"/>'
        )
    return _svg_root(w_px=w_px, h_px=h_px, body="".join(lines))


def _svg_title_block(*, w_px: int, h_px: int, title: str, subtitle: str) -> str:
    pad = 22
    body = _svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=18,
        fill=f"#{COLORS_HEX['card']}",
        stroke=f"#{COLORS_HEX['card_line']}",
        stroke_w=2,
        opacity=1.0,
    )
    # Subtle accent spine + highlight (poster polish).
    spine_w = 12
    body += f'<rect x="0" y="0" width="{spine_w}" height="{h_px}" rx="18" ry="18" fill="#{COLORS_HEX["cyan"]}" opacity="0.85"/>'
    body += f'<rect x="0" y="{int(round(h_px*0.52))}" width="{spine_w}" height="{h_px - int(round(h_px*0.52))}" rx="18" ry="18" fill="#{COLORS_HEX["coral"]}" opacity="0.55"/>'
    body += f'<line x1="{spine_w}" y1="1" x2="{w_px-1}" y2="1" stroke="#FFFFFF" stroke-opacity="0.06" stroke-width="2"/>'

    body += _svg_text(x=pad, y=int(h_px * 0.40), text=title, size_px=42, color=f"#{COLORS_HEX['ink']}", weight=800)
    body += _svg_text(
        x=pad,
        y=int(h_px * 0.77),
        text=subtitle,
        size_px=18,
        color=f"#{COLORS_HEX['muted']}",
        weight=500,
    )
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_panel_label(*, w_px: int, h_px: int, text: str, accent_hex: str) -> str:
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=14, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{accent_hex}", stroke_w=2)
    body += f'<rect x="0" y="0" width="10" height="{h_px}" rx="14" ry="14" fill="#{accent_hex}" opacity="0.18"/>'
    body += f'<circle cx="18" cy="{h_px//2}" r="{max(4, int(round(h_px*0.18)))}" fill="#{accent_hex}" opacity="0.95"/>'
    body += _svg_text(x=32, y=h_px // 2, text=text, size_px=16, color=f"#{COLORS_HEX['ink']}", weight=800)
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_pill(*, w_px: int, h_px: int, text: str, accent_hex: str) -> str:
    r = max(8, int(round(h_px / 2)))
    pad = int(round(h_px * 0.36))
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=r, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{accent_hex}", stroke_w=2)
    # icon dot
    body += f'<circle cx="{pad}" cy="{h_px//2}" r="{max(4, int(round(h_px*0.18)))}" fill="#{accent_hex}" opacity="0.9"/>'
    body += _svg_text(x=pad + int(round(h_px * 0.42)), y=h_px // 2, text=text, size_px=max(12, int(round(h_px * 0.55))), color=f"#{COLORS_HEX['ink']}", weight=800)
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_chip(*, w_px: int, h_px: int, head: str, val: str, accent_hex: str) -> str:
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=16, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    body += f'<rect x="0" y="0" width="10" height="{h_px}" rx="16" ry="16" fill="#{accent_hex}" opacity="0.16"/>'
    body += _svg_text(x=20, y=int(h_px * 0.38), text=head, size_px=14, color=f"#{COLORS_HEX['muted']}", weight=750)
    body += _svg_text(x=20, y=int(h_px * 0.72), text=val, size_px=16, color=f"#{accent_hex}", weight=800)
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_rank_row(*, w_px: int, h_px: int, rank: int, label: str, accent_hex: str) -> str:
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=max(8, int(round(h_px / 2))), fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    badge_w = max(int(round(w_px * 0.15)), int(round(h_px * 1.55)))
    body += _svg_rounded_rect(x=0, y=0, w=badge_w, h=h_px, r=max(8, int(round(h_px / 2))), fill=f"#{accent_hex}", stroke=f"#{accent_hex}", stroke_w=0)
    body += _svg_text(x=badge_w // 2, y=h_px // 2, text=str(rank), size_px=max(11, int(round(h_px * 0.66))), color=f"#{COLORS_HEX['bg']}", weight=900, anchor="middle")

    # avatar stub
    av_r = max(4, int(round(h_px * 0.28)))
    av_cx = badge_w + int(round(av_r * 1.6))
    body += f'<circle cx="{av_cx}" cy="{h_px//2}" r="{av_r}" fill="#{COLORS_HEX["card_line"]}" opacity="0.95"/>'
    body += f'<circle cx="{av_cx-1}" cy="{h_px//2-1}" r="{max(2, av_r-2)}" fill="#{COLORS_HEX["card"]}" opacity="0.85"/>'

    tx = av_cx + av_r + max(8, int(round(h_px * 0.22)))
    label_size = max(10, int(round(h_px * 0.56)))
    body += _svg_text(x=tx, y=h_px // 2, text=label, size_px=label_size, color=f"#{COLORS_HEX['ink']}", weight=650)

    # Tiny exposure bar hint (ties rank rows to v(rank) visually).
    pad = 10
    bar_area_w = max(36, int(round(w_px * 0.22)))
    bar_h = max(4, int(round(h_px * 0.42)))
    bar_x = w_px - pad - bar_area_w
    bar_y = int(round((h_px - bar_h) / 2))
    frac = 1.0 / (1.0 + 0.18 * (max(1, rank) - 1))
    body += f'<rect x="{bar_x}" y="{bar_y}" width="{bar_area_w}" height="{bar_h}" rx="8" ry="8" fill="#{COLORS_HEX["card"]}" opacity="0.92"/>'
    body += f'<rect x="{bar_x}" y="{bar_y}" width="{int(round(bar_area_w * frac))}" height="{bar_h}" rx="8" ry="8" fill="#{accent_hex}" opacity="0.85"/>'
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_post_card(
    *,
    w_px: int,
    h_px: int,
    author: str,
    lines: list[str],
    rank: int,
    accent_hex: str,
) -> str:
    pad = max(10, int(round(h_px * 0.12)))
    r = max(14, int(round(min(w_px, h_px) * 0.18)))
    stripe_w = max(8, int(round(w_px * 0.03)))

    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=r, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    body += _svg_rounded_rect(x=0, y=0, w=stripe_w, h=h_px, r=r, fill=f"#{accent_hex}", stroke=f"#{accent_hex}", stroke_w=0, opacity=0.95)

    # Type scale (relative) for consistent readability across card sizes.
    author_size = max(13, min(17, int(round(h_px * 0.17))))
    body_size = max(12, min(15, int(round(h_px * 0.15))))
    footer_size = max(11, min(14, int(round(h_px * 0.13))))
    line_gap = max(14, int(round(body_size * 1.25)))

    # Header: avatar + author
    av_r = max(9, min(13, int(round(h_px * 0.13))))
    av_cx = pad + av_r + 4
    av_cy = pad + av_r + 1
    body += f'<circle cx="{av_cx}" cy="{av_cy}" r="{av_r}" fill="#{accent_hex}" opacity="0.95"/>'
    body += f'<circle cx="{av_cx}" cy="{av_cy}" r="{max(2, av_r-3)}" fill="#{COLORS_HEX["card"]}" opacity="0.95"/>'
    body += _svg_text(x=av_cx + av_r + 8, y=av_cy, text=author, size_px=author_size, color=f"#{COLORS_HEX['ink']}", weight=800)

    # Footer region (reserve height so body never collides)
    footer_h = max(18, int(round(h_px * 0.24)))
    sep_y = h_px - footer_h
    body += f'<line x1="{pad}" y1="{sep_y}" x2="{w_px-pad}" y2="{sep_y}" stroke="#{COLORS_HEX["card_line"]}" stroke-opacity="0.55" stroke-width="1"/>'

    # Quote body (shrink if needed to fit above footer)
    body_lines = lines[:3]
    ty = av_cy + av_r + max(6, int(round(h_px * 0.06)))
    while body_lines and (ty + (len(body_lines) - 1) * line_gap + body_size) > (sep_y - 6) and body_size > 11:
        body_size -= 1
        line_gap = max(13, int(round(body_size * 1.18)))
    body += _svg_multiline(
        x=pad + 2,
        y=ty,
        lines=body_lines,
        size_px=body_size,
        color=f"#{COLORS_HEX['muted']}",
        weight=650,
        line_gap_px=line_gap,
        baseline="hanging",
    )

    # Footer: rank badge + v(rank) bar (exposure proxy)
    foot_cy = sep_y + footer_h // 2
    badge_h = max(16, int(round(footer_h * 0.56)))
    badge_w = max(66, int(round(badge_h * 3.4)))
    bx = pad + 2
    by = int(round(foot_cy - badge_h / 2))
    body += f'<rect x="{bx}" y="{by}" width="{badge_w}" height="{badge_h}" rx="{badge_h//2}" ry="{badge_h//2}" fill="#{COLORS_HEX["card"]}" stroke="#{accent_hex}" stroke-width="2" opacity="0.98"/>'
    body += _svg_text(x=bx + badge_w // 2, y=foot_cy, text=f"rank {int(rank)}", size_px=footer_size, color=f"#{COLORS_HEX['ink']}", weight=850, anchor="middle", baseline="middle")

    bar_x = bx + badge_w + 10
    bar_w = max(44, (w_px - pad) - bar_x)
    bar_h = max(6, int(round(footer_h * 0.28)))
    bar_y = int(round(foot_cy - bar_h / 2))
    frac = 1.0 / (1.0 + 0.18 * (max(1, int(rank)) - 1))
    frac = max(0.0, min(1.0, frac))

    body += _svg_text(x=bar_x, y=sep_y + 4, text="v(rank)", size_px=max(11, footer_size - 1), color=f"#{COLORS_HEX['muted']}", weight=750, baseline="hanging")
    body += f'<rect x="{bar_x}" y="{bar_y}" width="{bar_w}" height="{bar_h}" rx="8" ry="8" fill="#{COLORS_HEX["card"]}" opacity="0.92"/>'
    body += f'<rect x="{bar_x}" y="{bar_y}" width="{int(round(bar_w * frac))}" height="{bar_h}" rx="8" ry="8" fill="#{accent_hex}" opacity="0.88"/>'

    # Eye marker at the right end of the bar.
    eye_rx = max(9, int(round(bar_h * 1.20)))
    eye_ry = max(6, int(round(bar_h * 0.75)))
    ex = bar_x + bar_w - eye_rx - 8
    ey = foot_cy
    body += f'<ellipse cx="{int(ex)}" cy="{int(ey)}" rx="{eye_rx}" ry="{eye_ry}" fill="none" stroke="#{accent_hex}" stroke-width="2" opacity="0.85"/>'
    body += f'<circle cx="{int(ex)}" cy="{int(ey)}" r="{max(3, int(round(eye_ry*0.55)))}" fill="#{accent_hex}" opacity="0.85"/>'

    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_cluster_bubble(*, w_px: int, h_px: int, label: str, accent_hex: str) -> str:
    cx = w_px // 2
    cy = h_px // 2
    r = min(cx, cy) - 10
    body = f'<circle cx="{cx}" cy="{cy}" r="{r}" fill="#{COLORS_HEX["card2"]}" stroke="#{accent_hex}" stroke-width="3" opacity="0.98"/>'
    body += f'<circle cx="{cx}" cy="{cy}" r="{max(6, r-8)}" fill="none" stroke="#{accent_hex}" stroke-opacity="0.35" stroke-width="2" stroke-dasharray="6 6"/>'

    # inner nodes + light connections
    nodes = [(-0.26, -0.10), (0.22, -0.18), (0.10, 0.24), (-0.18, 0.22)]
    pts: list[tuple[int, int]] = []
    for dx, dy in nodes:
        x = int(round(cx + dx * r))
        y = int(round(cy + dy * r))
        pts.append((x, y))
        body += f'<circle cx="{x}" cy="{y}" r="{max(8, int(round(r*0.12)))}" fill="#{accent_hex}" opacity="0.32"/>'
        body += f'<circle cx="{x}" cy="{y}" r="{max(5, int(round(r*0.08)))}" fill="#{accent_hex}" opacity="0.92"/>'
    for (x1, y1), (x2, y2) in zip(pts, pts[1:] + pts[:1]):
        body += f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="#{accent_hex}" stroke-opacity="0.22" stroke-width="2"/>'

    # Small "hash" badge (visual: clustering via normalization/fingerprints).
    badge_w, badge_h = 44, 18
    bx, by = cx - badge_w // 2, cy - r + 16
    body += f'<rect x="{bx}" y="{by}" width="{badge_w}" height="{badge_h}" rx="9" ry="9" fill="#{COLORS_HEX["card"]}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.96"/>'
    body += _svg_text(x=cx, y=by + badge_h // 2, text="# hash", size_px=11, color=f"#{COLORS_HEX['muted']}", weight=800, anchor="middle")

    body += _svg_text(x=cx, y=cy - 6, text=label, size_px=18, color=f"#{COLORS_HEX['ink']}", weight=900, anchor="middle")
    body += _svg_text(x=cx, y=cy + 14, text="same / near-same text", size_px=13, color=f"#{COLORS_HEX['muted']}", weight=700, anchor="middle")
    body += _svg_text(x=cx, y=cy + 34, text="+ similar time (optional)", size_px=12, color=f"#{COLORS_HEX['muted']}", weight=650, anchor="middle", opacity=0.92)
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_timeline_card(*, w_px: int, h_px: int, accent_hex: str) -> str:
    pad = 16
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=18, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    body += f'<rect x="0" y="0" width="10" height="{h_px}" rx="18" ry="18" fill="#{accent_hex}" opacity="0.10"/>'
    body += f'<line x1="0" y1="1" x2="{w_px}" y2="1" stroke="#FFFFFF" stroke-opacity="0.05" stroke-width="2"/>'
    body += _svg_text(x=pad, y=pad + 4, text="Amortized attention over snapshots", size_px=14, color=f"#{COLORS_HEX['ink']}", weight=900, baseline="hanging")
    body += _svg_text(
        x=pad,
        y=pad + 24,
        text="A(p)=Σ_t v(rank_t(p))     R(p)=Σ_t u(p|q_t)",
        size_px=12,
        color=f"#{COLORS_HEX['muted']}",
        weight=700,
        baseline="hanging",
    )

    x0 = pad
    y0 = pad + 48
    x1 = w_px - pad
    y1 = h_px - pad - 8
    body += f'<line x1="{x0}" y1="{y1}" x2="{x1}" y2="{y1}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.8"/>'
    body += f'<line x1="{x0}" y1="{y0}" x2="{x0}" y2="{y1}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.8"/>'

    # Snapshot ticks on the x-axis.
    ticks = 6
    for i in range(ticks):
        tx = int(round(x0 + i * (x1 - x0) / (ticks - 1)))
        body += f'<line x1="{tx}" y1="{y1}" x2="{tx}" y2="{y1+6}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.65"/>'
        if i in (0, ticks - 1):
            body += _svg_text(x=tx, y=y1 + 10, text=f"t{i+1}", size_px=10, color=f"#{COLORS_HEX['muted']}", weight=700, anchor="middle", baseline="hanging")

    # two smooth-ish polylines (A and R)
    pts_a = [(x0, y1), (x0 + 0.22 * (x1 - x0), y1 - 0.22 * (y1 - y0)), (x0 + 0.55 * (x1 - x0), y1 - 0.48 * (y1 - y0)), (x1, y0 + 0.18 * (y1 - y0))]
    pts_r = [(x0, y1), (x0 + 0.25 * (x1 - x0), y1 - 0.12 * (y1 - y0)), (x0 + 0.60 * (x1 - x0), y1 - 0.34 * (y1 - y0)), (x1, y0 + 0.28 * (y1 - y0))]
    poly_a = " ".join(f"{int(x)},{int(y)}" for x, y in pts_a)
    poly_r = " ".join(f"{int(x)},{int(y)}" for x, y in pts_r)
    body += f'<polygon points="{poly_a} {int(x1)},{int(y1)} {int(x0)},{int(y1)}" fill="#{accent_hex}" opacity="0.10"/>'
    body += f'<polygon points="{poly_r} {int(x1)},{int(y1)} {int(x0)},{int(y1)}" fill="#{COLORS_HEX["muted"]}" opacity="0.06"/>'
    body += f'<polyline fill="none" stroke="#{accent_hex}" stroke-width="3" opacity="0.95" points="{poly_a}"/>'
    body += f'<polyline fill="none" stroke="#{COLORS_HEX["muted"]}" stroke-width="3" opacity="0.85" points="{poly_r}"/>'
    for x, y in pts_a[1:]:
        body += f'<circle cx="{int(x)}" cy="{int(y)}" r="3" fill="#{accent_hex}" opacity="0.95"/>'
    for x, y in pts_r[1:]:
        body += f'<circle cx="{int(x)}" cy="{int(y)}" r="3" fill="#{COLORS_HEX["muted"]}" opacity="0.85"/>'
    body += _svg_text(x=x0 + 8, y=y0 + 10, text="A(p)", size_px=13, color=f"#{accent_hex}", weight=900, baseline="hanging")
    body += _svg_text(x=x0 + 54, y=y0 + 10, text="R(p)", size_px=13, color=f"#{COLORS_HEX['muted']}", weight=900, baseline="hanging")

    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_join_graph(*, w_px: int, h_px: int, accent_hex: str) -> str:
    pad = 10
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=18, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    body += f'<rect x="0" y="0" width="10" height="{h_px}" rx="18" ry="18" fill="#{accent_hex}" opacity="0.10"/>'
    body += f'<line x1="0" y1="1" x2="{w_px}" y2="1" stroke="#FFFFFF" stroke-opacity="0.05" stroke-width="2"/>'
    body += _svg_text(x=pad, y=pad, text="Audit table (one row per impression)", size_px=14, color=f"#{COLORS_HEX['ink']}", weight=900, baseline="hanging")

    nodes = ["feed_items", "posts_first_seen", "post_metrics", "author_profiles"]
    nx0 = pad
    ny = pad + 18
    gap = 8
    box_h = 18
    # pack nodes across
    box_w = int((w_px - 2 * pad - gap * (len(nodes) - 1)) / len(nodes))
    for i, name in enumerate(nodes):
        x = nx0 + i * (box_w + gap)
        body += _svg_rounded_rect(x=x, y=ny, w=box_w, h=box_h, r=10, fill=f"#{COLORS_HEX['card']}", stroke=f"#{accent_hex if i==0 else COLORS_HEX['card_line']}", stroke_w=2, opacity=0.98)
        body += _svg_text(x=x + box_w // 2, y=ny + box_h // 2, text=name, size_px=12, color=f"#{COLORS_HEX['muted']}", weight=800, anchor="middle")
        if i < len(nodes) - 1:
            x2 = x + box_w
            y2 = ny + box_h // 2
            body += f'<line x1="{x2}" y1="{y2}" x2="{x2 + gap}" y2="{y2}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.85"/>'
            body += f'<circle cx="{x2 + gap}" cy="{y2}" r="2" fill="#{accent_hex}" opacity="0.9"/>'

    body += _svg_text(
        x=pad,
        y=ny + box_h + 6,
        text="Keys: post_uri · author_did   (+ q: f,t,m,ν)",
        size_px=11,
        color=f"#{COLORS_HEX['muted']}",
        weight=700,
        baseline="hanging",
    )
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_metric_card_bg(*, w_px: int, h_px: int, stroke_hex: str, kind: str) -> str:
    """Background for metric cards (no <text>), with a small icon hint."""
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=18, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{stroke_hex}", stroke_w=2)
    body += f'<line x1="0" y1="1" x2="{w_px}" y2="1" stroke="#FFFFFF" stroke-opacity="0.05" stroke-width="2"/>'

    # Icon area (bottom-right; light opacity so it doesn't fight the equation).
    pad = 14
    iw = max(64, int(round(w_px * 0.34)))
    ih = max(36, int(round(h_px * 0.34)))
    ix = w_px - pad - iw
    iy = h_px - pad - ih
    icon_op = 0.20

    if kind.upper() == "D":
        # Two bars + bracket (|E - E'|).
        b1w = int(round(iw * 0.86))
        b2w = int(round(iw * 0.46))
        bh = 8
        y1 = iy + 10
        y2 = iy + 24
        body += f'<rect x="{ix}" y="{y1}" width="{b1w}" height="{bh}" rx="6" ry="6" fill="#{stroke_hex}" opacity="{icon_op:.3f}"/>'
        body += f'<rect x="{ix}" y="{y2}" width="{b2w}" height="{bh}" rx="6" ry="6" fill="#{stroke_hex}" opacity="{icon_op:.3f}"/>'
        bx = ix + int(round(iw * 0.90))
        body += f'<line x1="{bx}" y1="{y1}" x2="{bx}" y2="{y2+bh}" stroke="#{stroke_hex}" stroke-width="3" stroke-opacity="{icon_op:.3f}"/>'
        body += f'<line x1="{bx-8}" y1="{y1}" x2="{bx}" y2="{y1}" stroke="#{stroke_hex}" stroke-width="3" stroke-opacity="{icon_op:.3f}"/>'
        body += f'<line x1="{bx-8}" y1="{y2+bh}" x2="{bx}" y2="{y2+bh}" stroke="#{stroke_hex}" stroke-width="3" stroke-opacity="{icon_op:.3f}"/>'
    elif kind.upper() == "IF":
        # Two nodes + similarity edge.
        cy = iy + ih // 2
        c1x = ix + int(round(iw * 0.28))
        c2x = ix + int(round(iw * 0.76))
        body += f'<line x1="{c1x}" y1="{cy}" x2="{c2x}" y2="{cy}" stroke="#{stroke_hex}" stroke-width="3" stroke-opacity="{icon_op:.3f}" stroke-dasharray="6 6"/>'
        body += f'<circle cx="{c1x}" cy="{cy}" r="10" fill="#{stroke_hex}" opacity="{icon_op:.3f}"/>'
        body += f'<circle cx="{c2x}" cy="{cy}" r="10" fill="#{stroke_hex}" opacity="{icon_op:.3f}"/>'
        body += f'<circle cx="{c1x}" cy="{cy}" r="5" fill="#{stroke_hex}" opacity="{min(0.55, icon_op*2.0):.3f}"/>'
        body += f'<circle cx="{c2x}" cy="{cy}" r="5" fill="#{stroke_hex}" opacity="{min(0.55, icon_op*2.0):.3f}"/>'
    else:
        # Δ: two context badges + minus.
        bw = int(round(iw * 0.62))
        bh = 12
        rx = ix + int(round(iw * 0.10))
        r1y = iy + 10
        r2y = iy + 26
        body += f'<rect x="{rx}" y="{r1y}" width="{bw}" height="{bh}" rx="8" ry="8" fill="#{stroke_hex}" opacity="{icon_op:.3f}"/>'
        body += f'<rect x="{rx}" y="{r2y}" width="{int(round(bw*0.78))}" height="{bh}" rx="8" ry="8" fill="#{stroke_hex}" opacity="{icon_op*0.70:.3f}"/>'
        mx = ix + int(round(iw * 0.86))
        my = iy + ih // 2
        body += f'<line x1="{mx-10}" y1="{my}" x2="{mx+10}" y2="{my}" stroke="#{stroke_hex}" stroke-width="4" stroke-opacity="{icon_op:.3f}"/>'

    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_bar_row(*, w_px: int, h_px: int, j: int, frac: float, accent_hex: str) -> str:
    r = max(6, int(round(h_px / 2)))
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=r, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    pad = max(8, int(round(h_px * 0.70)))
    bar_w = int(round((w_px - 2 * pad) * max(0.0, min(1.0, frac))))
    bar_h = max(4, int(round(h_px * 0.50)))
    bar_y = max(0, int(round((h_px - bar_h) / 2)))
    body += (
        f'<rect x="{pad}" y="{bar_y}" width="{bar_w}" height="{bar_h}" '
        f'rx="8" ry="8" fill="#{accent_hex}" fill-opacity="0.85"/>'
    )
    body += _svg_text(
        x=w_px - pad,
        y=h_px // 2,
        text=f"v{j}",
        size_px=max(11, int(round(h_px * 0.72))),
        color=f"#{COLORS_HEX['muted']}",
        weight=800,
        anchor="end",
    )
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_scalar_bar(*, w_px: int, h_px: int, label: str, frac: float, accent_hex: str, value: str | None = None) -> str:
    frac = max(0.0, min(1.0, frac))
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=16, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    pad = 14
    body += _svg_text(x=pad, y=int(h_px * 0.30), text=label, size_px=max(12, int(round(h_px * 0.34))), color=f"#{COLORS_HEX['muted']}", weight=800, baseline="middle")

    bx = pad
    by = int(h_px * 0.52)
    bw = w_px - 2 * pad
    bh = max(10, int(h_px * 0.28))
    body += f'<rect x="{bx}" y="{by}" width="{bw}" height="{bh}" rx="10" ry="10" fill="#{COLORS_HEX["card"]}" opacity="0.95"/>'
    body += f'<rect x="{bx}" y="{by}" width="{int(round(bw * frac))}" height="{bh}" rx="10" ry="10" fill="#{accent_hex}" opacity="0.9"/>'
    if value is not None:
        body += _svg_text(x=w_px - pad, y=int(h_px * 0.30), text=value, size_px=max(12, int(round(h_px * 0.34))), color=f"#{COLORS_HEX['ink']}", weight=900, anchor="end", baseline="middle")
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_cell(*, w_px: int, h_px: int, val: str | None, heat: float, accent_hex: str) -> str:
    heat = max(0.0, min(1.0, heat))

    def _mix(a: int, b: int, t: float) -> int:
        return int(round(a * (1 - t) + b * t))

    base = COLORS_HEX["card2"]
    br, bg, bb = int(base[0:2], 16), int(base[2:4], 16), int(base[4:6], 16)
    ar, ag, ab = int(accent_hex[0:2], 16), int(accent_hex[2:4], 16), int(accent_hex[4:6], 16)
    t = 0.18 + 0.72 * heat
    fill = f"#{_mix(br, ar, t):02X}{_mix(bg, ag, t):02X}{_mix(bb, ab, t):02X}"

    rr = max(3, int(round(min(w_px, h_px) * 0.25)))
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=rr, fill=fill, stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    if val:
        # Keep numbers readable across different cell sizes (matrix vs tiny heatmaps).
        cell = max(1, min(w_px, h_px))
        size = int(round(cell * 0.60))
        size = max(11, min(15, size))
        body += _svg_text(x=w_px // 2, y=h_px // 2, text=val, size_px=size, color=f"#{COLORS_HEX['ink']}", weight=800, anchor="middle")
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_arrow(*, w_px: int, h_px: int, accent_hex: str, text: str) -> str:
    body = f'<path d="M0 {h_px*0.20} H{w_px*0.72} V0 L{w_px} {h_px/2} L{w_px*0.72} {h_px} V{h_px*0.80} H0 Z" fill="#{accent_hex}" fill-opacity="0.92"/>'
    body += _svg_text(x=int(w_px * 0.34), y=h_px // 2, text=text, size_px=16, color=f"#{COLORS_HEX['bg']}", weight=900, anchor="middle")
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_caption(*, w_px: int, h_px: int, lines: list[str], accent_hex: str | None = None) -> str:
    body = _svg_rounded_rect(x=0, y=0, w=w_px, h=h_px, r=16, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
    pad = 18
    if accent_hex:
        body += f'<rect x="0" y="0" width="10" height="{h_px}" rx="16" ry="16" fill="#{accent_hex}" opacity="0.14"/>'
    body += f'<line x1="0" y1="1" x2="{w_px}" y2="1" stroke="#FFFFFF" stroke-opacity="0.05" stroke-width="2"/>'

    if not lines:
        return _svg_root(w_px=w_px, h_px=h_px, body=body)

    header = lines[0]
    rest = lines[1:]
    header_size = 16
    body_size = 14
    line_gap = 18

    while rest and (pad + 26 + len(rest) * line_gap + 10) > h_px and body_size > 12:
        body_size -= 1
        line_gap = max(15, int(round(body_size * 1.25)))

    body += _svg_text(x=pad, y=pad + 2, text=header, size_px=header_size, color=f"#{COLORS_HEX['ink']}", weight=900, baseline="hanging")

    y0 = pad + 26
    bullet_r = 3
    bcol = f"#{accent_hex}" if accent_hex else f"#{COLORS_HEX['muted']}"
    for i, line in enumerate(rest):
        y = y0 + i * line_gap
        body += f'<circle cx="{pad}" cy="{y + 7}" r="{bullet_r}" fill="{bcol}" opacity="0.95"/>'
        body += _svg_text(
            x=pad + 10,
            y=y,
            text=line.lstrip("\u2022 ").strip(),
            size_px=body_size,
            color=f"#{COLORS_HEX['muted']}",
            weight=650,
            baseline="hanging",
        )
    return _svg_root(w_px=w_px, h_px=h_px, body=body)


def _mpl_math_svg(expr: str, *, color_hex: str, fontsize_pt: float) -> tuple[str, float, float]:
    # True vector: glyphs as paths, not <text>.
    import matplotlib

    matplotlib.use("Agg")  # noqa: S104
    import matplotlib as mpl
    from matplotlib import pyplot as plt

    mpl.rcParams["svg.fonttype"] = "path"
    mpl.rcParams["figure.dpi"] = 96

    fig = plt.figure(figsize=(0.1, 0.1))
    fig.patch.set_alpha(0.0)
    ax = fig.add_axes([0, 0, 1, 1])
    ax.axis("off")

    txt = ax.text(0.0, 0.0, expr, fontsize=fontsize_pt, color=f"#{color_hex}", ha="left", va="baseline")
    fig.canvas.draw()
    bbox = txt.get_window_extent(renderer=fig.canvas.get_renderer())
    w_in = (bbox.width + 6) / fig.dpi
    h_in = (bbox.height + 6) / fig.dpi

    fig.set_size_inches(w_in, h_in)
    ax.set_position([0, 0, 1, 1])
    txt.set_position((0.0, 0.0))

    buf = io.StringIO()
    fig.savefig(buf, format="svg", transparent=True, bbox_inches="tight", pad_inches=0.0)
    plt.close(fig)
    return buf.getvalue(), w_in, h_in


def _math_svg_baseline_ratio(svg_text: str) -> float | None:
    """
    Return baseline position as a fraction of the SVG height (distance from top / height).

    Matplotlib writes mathtext as a <g ... transform="translate(x y) scale(...)"> where the translate y
    corresponds to the text baseline in the SVG coordinate system (y grows downward). We parse that
    to align token baselines in multi-token equation strips.
    """

    m_h = re.search(r'height="([0-9.]+)pt"', svg_text)
    if not m_h:
        return None
    h_pt = float(m_h.group(1))
    if h_pt <= 0:
        return None

    # Prefer the transform inside the first text group (more robust than grabbing any random transform).
    m_t = re.search(r'id="text_1".*?transform="translate\([^\s]+\s+([-0-9.]+)\)', svg_text, flags=re.DOTALL)
    if not m_t:
        m_t = re.search(r'transform="translate\([^\s]+\s+([-0-9.]+)\)', svg_text)
    if not m_t:
        return None

    baseline_pt = float(m_t.group(1))
    r = baseline_pt / h_pt
    # Guard against small rounding/formatting drift.
    return max(0.0, min(1.0, r))


def _write_unique_png(*, out_path: Path, seed: str) -> None:
    # Ensure each picture gets a distinct media blob so python-pptx doesn't de-dupe image parts.
    h = hashlib.sha1(seed.encode("utf-8")).digest()
    rgb = (h[0], h[1], h[2], 255)
    img = Image.new("RGBA", (1, 1), rgb)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)


@dataclass(frozen=True)
class SvgTile:
    slide_num: int
    name: str
    x_in: float
    y_in: float
    w_in: float
    h_in: float
    svg_text: str
    static: bool = False


def _assert_tiles_within_slide_bounds(tiles: list[SvgTile]) -> None:
    eps = 1e-4
    for t in tiles:
        if t.w_in <= 0 or t.h_in <= 0:
            raise ValueError(f"Invalid tile size: slide={t.slide_num} name={t.name} w_in={t.w_in} h_in={t.h_in}")
        if t.x_in < -eps or t.y_in < -eps:
            raise ValueError(f"Tile out of bounds (neg): slide={t.slide_num} name={t.name} x={t.x_in} y={t.y_in}")
        if t.x_in + t.w_in > SLIDE_W_IN + eps or t.y_in + t.h_in > SLIDE_H_IN + eps:
            raise ValueError(
                f"Tile out of bounds: slide={t.slide_num} name={t.name} "
                f"right={t.x_in + t.w_in:.3f} bottom={t.y_in + t.h_in:.3f} "
                f"slide_w={SLIDE_W_IN:.3f} slide_h={SLIDE_H_IN:.3f}"
            )


def _add_math_strip(
    *,
    tiles: list[SvgTile],
    slide_num: int,
    name_prefix: str,
    x_in: float,
    y_in: float,
    tokens: list[str],
    color_hex: str,
    fontsize_pt: float,
    gap_in: float = 0.06,
    token_color_hexes: list[str] | None = None,
) -> tuple[float, float]:
    if token_color_hexes is not None and len(token_color_hexes) != len(tokens):
        raise ValueError(f"token_color_hexes length {len(token_color_hexes)} != tokens length {len(tokens)}")

    rendered: list[tuple[str, float, float, float]] = []
    baseline_top_max = 0.0
    for idx, tok in enumerate(tokens):
        tok_color = token_color_hexes[idx] if token_color_hexes is not None else color_hex
        svg, w_in, h_in = _mpl_math_svg(tok, color_hex=tok_color, fontsize_pt=fontsize_pt)
        r = _math_svg_baseline_ratio(svg)
        baseline_top = h_in * (r if r is not None else 0.82)
        rendered.append((svg, w_in, h_in, baseline_top))
        baseline_top_max = max(baseline_top_max, baseline_top)

    x = x_in
    strip_h = 0.0
    for svg, _, h_in, baseline_top in rendered:
        y_off = baseline_top_max - baseline_top
        strip_h = max(strip_h, y_off + h_in)

    for idx, (svg, w_in, h_in, baseline_top) in enumerate(rendered, start=1):
        y_off = baseline_top_max - baseline_top
        tiles.append(
            SvgTile(
                slide_num=slide_num,
                name=f"{name_prefix}_{idx:02d}",
                x_in=x,
                y_in=y_in + y_off,
                w_in=w_in,
                h_in=h_in,
                svg_text=svg,
            )
        )
        x += w_in + gap_in

    total_w = max(0.0, x - x_in - gap_in)
    return total_w, strip_h


def _build_slide1_legacy() -> list[SvgTile]:
    tiles: list[SvgTile] = []

    # Background grid (ignored by overlap QC by name prefix).
    tiles.append(
        SvgTile(
            slide_num=1,
            name="_BG_GRID_S1",
            x_in=0.0,
            y_in=0.0,
            w_in=SLIDE_W_IN,
            h_in=SLIDE_H_IN,
            svg_text=_svg_grid_bg(w_px=_in_to_px(SLIDE_W_IN), h_px=_in_to_px(SLIDE_H_IN)),
            static=True,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_TITLE",
            x_in=SAFE_L_IN,
            y_in=0.28,
            w_in=SAFE_W_IN,
            h_in=0.95,
            svg_text=_svg_title_block(
                w_px=_in_to_px(SAFE_W_IN),
                h_px=_in_to_px(0.95),
                title="Exposure = expected attention",
                subtitle="Singh & Joachims (KDD’18): P (rank marginals) · v (position bias)  →  Exposure",
            ),
            static=True,
        )
    )

    col_gap = 0.28
    left_w = 7.25
    right_w = SAFE_W_IN - left_w - col_gap
    left_x = SAFE_L_IN
    right_x = SAFE_L_IN + left_w + col_gap
    body_y = 1.38
    panel_h = SAFE_H_IN - (body_y - SAFE_T_IN)

    # Background panels.
    for name, x, w in [("_BG_PANEL_THEORY", left_x, left_w), ("_BG_PANEL_DATA", right_x, right_w)]:
        tiles.append(
            SvgTile(
                slide_num=1,
                name=name,
                x_in=x,
                y_in=body_y,
                w_in=w,
                h_in=panel_h,
                svg_text=_svg_root(
                    w_px=_in_to_px(w),
                    h_px=_in_to_px(panel_h),
                    body=_svg_rounded_rect(
                        x=0,
                        y=0,
                        w=_in_to_px(w),
                        h=_in_to_px(panel_h),
                        r=22,
                        fill=f"#{COLORS_HEX['card']}",
                        stroke=f"#{COLORS_HEX['card_line']}",
                        stroke_w=2,
                        opacity=0.98,
                    ),
                ),
                static=True,
            )
        )

    # Left panel: theory
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_LBL_THEORY",
            x_in=left_x + 0.22,
            y_in=body_y + 0.12,
            w_in=3.80,
            h_in=0.35,
            svg_text=_svg_panel_label(w_px=_in_to_px(3.80), h_px=_in_to_px(0.35), text="Theory: probabilistic ranking", accent_hex=COLORS_HEX["cyan"]),
            static=True,
        )
    )

    # P grid
    grid_n = 13
    cell_in = 0.26
    gap_in = 0.02
    grid_w_in = grid_n * cell_in + (grid_n - 1) * gap_in
    grid_x = left_x + 0.30
    grid_y = body_y + 0.70

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_CAP_P",
            x_in=grid_x,
            y_in=grid_y - 0.20,
            w_in=grid_w_in,
            h_in=0.20,
            svg_text=_svg_root(
                w_px=_in_to_px(grid_w_in),
                h_px=_in_to_px(0.20),
                body=_svg_text(x=0, y=_in_to_px(0.15), text="P (rank marginals) — pick a row i", size_px=16, color=f"#{COLORS_HEX['muted']}", weight=700),
            ),
            static=True,
        )
    )

    hi_r = 5
    probs = [0.40, 0.22, 0.14, 0.09, 0.06, 0.04, 0.03, 0.02, 0.00, 0.00, 0.00, 0.00, 0.00]
    for r in range(grid_n):
        for c in range(grid_n):
            x = grid_x + c * (cell_in + gap_in)
            y = grid_y + r * (cell_in + gap_in)
            heat = probs[c] / max(probs) if r == hi_r else 0.05
            val = f"{probs[c]:.2f}" if (r == hi_r and probs[c] > 0) else None
            tiles.append(
                SvgTile(
                    slide_num=1,
                    name=f"S1_P_{r:02d}_{c:02d}",
                    x_in=x,
                    y_in=y,
                    w_in=cell_in,
                    h_in=cell_in,
                    svg_text=_svg_cell(w_px=_in_to_px(cell_in), h_px=_in_to_px(cell_in), val=val, heat=heat, accent_hex=COLORS_HEX["cyan"]),
                )
            )

    # v bars (12)
    v_x = grid_x + grid_w_in + 0.38
    v_y = grid_y
    bar_w_in = 2.65
    bar_h_in = 0.26
    bar_gap_in = 0.06

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_CAP_V",
            x_in=v_x,
            y_in=v_y - 0.38,
            w_in=bar_w_in,
            h_in=0.30,
            svg_text=_svg_root(
                w_px=_in_to_px(bar_w_in),
                h_px=_in_to_px(0.30),
                body=_svg_text(x=0, y=_in_to_px(0.15), text="v (position bias)", size_px=16, color=f"#{COLORS_HEX['muted']}", weight=700),
            ),
            static=True,
        )
    )

    for j in range(1, grid_n + 1):
        frac = 1.0 / (1.0 + 0.18 * (j - 1))
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_V_{j:02d}",
                x_in=v_x,
                y_in=v_y + (j - 1) * (bar_h_in + bar_gap_in),
                w_in=bar_w_in,
                h_in=bar_h_in,
                svg_text=_svg_bar_row(w_px=_in_to_px(bar_w_in), h_px=_in_to_px(bar_h_in), j=j, frac=frac, accent_hex=COLORS_HEX["green"]),
            )
        )

    # Arrow and exposure caption (place below the full v-stack to avoid collisions)
    v_total_h = grid_n * bar_h_in + (grid_n - 1) * bar_gap_in
    arrow_y = v_y + v_total_h + 0.18
    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_ARROW_PV",
            x_in=v_x,
            y_in=arrow_y,
            w_in=bar_w_in,
            h_in=0.48,
            svg_text=_svg_arrow(w_px=_in_to_px(bar_w_in), h_px=_in_to_px(0.48), accent_hex=COLORS_HEX["amber"], text="P·v"),
        )
    )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_NOTE_EXPOSURE",
            x_in=left_x + 0.30,
            y_in=arrow_y + 0.48 + 0.18,
            w_in=left_w - 0.60,
            h_in=0.92,
            svg_text=_svg_caption(
                w_px=_in_to_px(left_w - 0.60),
                h_px=_in_to_px(0.92),
                lines=[
                    "Exposure(d_i | P) = Σ_j P_{i,j} · v_j",
                    "Utility: U(P|q)=uᵀPv  (anchors the model).",
                ],
            ),
        )
    )

    # Right panel: Bluesky mapping
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_LBL_DATA",
            x_in=right_x + 0.22,
            y_in=body_y + 0.20,
            w_in=right_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(
                w_px=_in_to_px(right_w - 0.44),
                h_px=_in_to_px(0.35),
                text="Your measurement: hourly feed snapshots",
                accent_hex=COLORS_HEX["coral"],
            ),
            static=True,
        )
    )

    chip_w_in = (right_w - 0.22 * 2 - 0.12) / 2
    chip_h_in = 0.58
    chip_x0 = right_x + 0.22
    chip_y0 = body_y + 0.70
    chips = [
        ("feed_uri", "f", COLORS_HEX["cyan"]),
        ("snapshot_hour_utc", "t", COLORS_HEX["amber"]),
        ("viewer_mode", "m", COLORS_HEX["green"]),
        ("vantage_id", "ν", COLORS_HEX["violet"]),
    ]
    for idx, (head, val, accent) in enumerate(chips):
        cx = chip_x0 + (idx % 2) * (chip_w_in + 0.12)
        cy = chip_y0 + (idx // 2) * (chip_h_in + 0.12)
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_Q_{idx+1:02d}",
                x_in=cx,
                y_in=cy,
                w_in=chip_w_in,
                h_in=chip_h_in,
                svg_text=_svg_chip(w_px=_in_to_px(chip_w_in), h_px=_in_to_px(chip_h_in), head=head, val=val, accent_hex=accent),
            )
        )

    list_x = chip_x0
    list_y = chip_y0 + 2 * (chip_h_in + 0.12) + 0.32
    list_w_in = right_w - 0.44
    row_h_in = 0.26
    row_gap_in = 0.06
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_CAP_RANK",
            x_in=list_x,
            y_in=list_y - 0.36,
            w_in=list_w_in,
            h_in=0.30,
            svg_text=_svg_root(
                w_px=_in_to_px(list_w_in),
                h_px=_in_to_px(0.30),
                body=_svg_text(x=0, y=_in_to_px(0.15), text="Observed top‑K list (one snapshot)", size_px=16, color=f"#{COLORS_HEX['muted']}", weight=700),
            ),
            static=True,
        )
    )

    for r in range(1, 26):
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_RANK_{r:02d}",
                x_in=list_x,
                y_in=list_y + (r - 1) * (row_h_in + row_gap_in),
                w_in=list_w_in,
                h_in=row_h_in,
                svg_text=_svg_rank_row(
                    w_px=_in_to_px(list_w_in),
                    h_px=_in_to_px(row_h_in),
                    rank=r,
                    label="post_uri / author_did …",
                    accent_hex=COLORS_HEX["coral"] if r <= 5 else COLORS_HEX["card_line"],
                ),
            )
        )

    map_note_y = list_y + 25 * (row_h_in + row_gap_in) + 0.18
    eq_svg, eq_w_in, eq_h_in = _mpl_math_svg(
        r"$E(\mathrm{post}\mid q)=v(\mathrm{rank})\quad(\mathrm{absent}\Rightarrow E=0)$",
        color_hex=COLORS_HEX["ink"],
        fontsize_pt=22.0,
    )
    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_EQ_MAP",
            x_in=list_x,
            y_in=map_note_y,
            w_in=min(list_w_in, eq_w_in),
            h_in=eq_h_in,
            svg_text=eq_svg,
        )
    )

    phat_svg, phat_w_in, phat_h_in = _mpl_math_svg(
        r"$\hat P_{p,j}=\mathrm{freq}_t[\mathrm{rank}_t(p)=j]\;\Rightarrow\;\hat E(p)=\sum\,_{j} \hat P_{p,j}v_j$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=20.0,
    )
    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_EQ_PHAT",
            x_in=list_x,
            y_in=map_note_y + eq_h_in + 0.12,
            w_in=min(list_w_in, phat_w_in),
            h_in=phat_h_in,
            svg_text=phat_svg,
        )
    )

    return tiles


def _build_slide2_legacy() -> list[SvgTile]:
    tiles: list[SvgTile] = []

    tiles.append(
        SvgTile(
            slide_num=2,
            name="_BG_GRID_S2",
            x_in=0.0,
            y_in=0.0,
            w_in=SLIDE_W_IN,
            h_in=SLIDE_H_IN,
            svg_text=_svg_grid_bg(w_px=_in_to_px(SLIDE_W_IN), h_px=_in_to_px(SLIDE_H_IN)),
            static=True,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_TITLE",
            x_in=SAFE_L_IN,
            y_in=0.28,
            w_in=SAFE_W_IN,
            h_in=0.95,
            svg_text=_svg_title_block(
                w_px=_in_to_px(SAFE_W_IN),
                h_px=_in_to_px(0.95),
                title="Individual fairness audit: same content → different exposure",
                subtitle="Zehlike et al. (CSUR’22 Part II): define similarity, then measure exposure discrepancies.",
            ),
            static=True,
        )
    )

    col_gap = 0.28
    left_w = 6.10
    right_w = SAFE_W_IN - left_w - col_gap
    left_x = SAFE_L_IN
    right_x = SAFE_L_IN + left_w + col_gap
    body_y = 1.38
    panel_h = SAFE_H_IN - (body_y - SAFE_T_IN)

    for name, x, w in [("_BG_PANEL_LEFT", left_x, left_w), ("_BG_PANEL_RIGHT", right_x, right_w)]:
        tiles.append(
            SvgTile(
                slide_num=2,
                name=name,
                x_in=x,
                y_in=body_y,
                w_in=w,
                h_in=panel_h,
                svg_text=_svg_root(
                    w_px=_in_to_px(w),
                    h_px=_in_to_px(panel_h),
                    body=_svg_rounded_rect(
                        x=0,
                        y=0,
                        w=_in_to_px(w),
                        h=_in_to_px(panel_h),
                        r=22,
                        fill=f"#{COLORS_HEX['card']}",
                        stroke=f"#{COLORS_HEX['card_line']}",
                        stroke_w=2,
                        opacity=0.98,
                    ),
                ),
                static=True,
            )
        )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_LBL_SIM",
            x_in=left_x + 0.22,
            y_in=body_y + 0.20,
            w_in=left_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(w_px=_in_to_px(left_w - 0.44), h_px=_in_to_px(0.35), text="Define similarity d(p,p')", accent_hex=COLORS_HEX["cyan"]),
            static=True,
        )
    )
    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_LBL_METRICS",
            x_in=right_x + 0.22,
            y_in=body_y + 0.20,
            w_in=right_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(w_px=_in_to_px(right_w - 0.44), h_px=_in_to_px(0.35), text="Measure exposure gaps", accent_hex=COLORS_HEX["coral"]),
            static=True,
        )
    )

    # Left: 3 near-duplicate posts
    post_w_in = left_w - 0.44
    post_h_in = 1.05
    post_x = left_x + 0.22
    post_y0 = body_y + 0.70
    authors = [("A", COLORS_HEX["cyan"]), ("B", COLORS_HEX["amber"]), ("C", COLORS_HEX["green"])]
    for i, (a, accent) in enumerate(authors, start=1):
        svg = _svg_caption(
            w_px=_in_to_px(post_w_in),
            h_px=_in_to_px(post_h_in),
            lines=[f"author_did = {a}", 'text ≈ “same / near‑same content”', "Outcome: different ranks → different exposure"],
        ).replace(COLORS_HEX["card_line"], accent)
        tiles.append(
            SvgTile(
                slide_num=2,
                name=f"S2_POST_{i:02d}",
                x_in=post_x,
                y_in=post_y0 + (i - 1) * (post_h_in + 0.14),
                w_in=post_w_in,
                h_in=post_h_in,
                svg_text=svg,
            )
        )

    sim_svg, sim_w_in, sim_h_in = _mpl_math_svg(
        r"$d(p,p')=\alpha\,d_{\text{text}}(p,p')+(1-\alpha)\,d_{\text{time}}(p,p')$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=20.0,
    )
    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_EQ_SIM",
            x_in=post_x,
            y_in=post_y0 + 3 * (post_h_in + 0.14) + 0.10,
            w_in=min(post_w_in, sim_w_in),
            h_in=sim_h_in,
            svg_text=sim_svg,
        )
    )

    # Right: equations + heatmap
    fx = right_x + 0.22
    fy = body_y + 0.70
    fw = right_w - 0.44

    eq1, w1, h1 = _mpl_math_svg(
        r"$D(p,p')=\left|\,\mathrm{Exposure}(p)-\mathrm{Exposure}(p')\,\right|$",
        color_hex=COLORS_HEX["amber"],
        fontsize_pt=24.0,
    )
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_D", x_in=fx, y_in=fy, w_in=min(fw, w1), h_in=h1, svg_text=eq1))

    eq2, w2, h2 = _mpl_math_svg(
        r"$d(p,p')\approx 0\;\Rightarrow\;|E(p)-E(p')|\leq \varepsilon$",
        color_hex=COLORS_HEX["ink"],
        fontsize_pt=22.0,
    )
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_LIP", x_in=fx, y_in=fy + h1 + 0.16, w_in=min(fw, w2), h_in=h2, svg_text=eq2))

    eq3, w3, h3 = _mpl_math_svg(
        r"$\Delta E(p;q_1,q_2)=v(\mathrm{rank}_{q_1}(p)) - v(\mathrm{rank}_{q_2}(p))$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=20.0,
    )
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_DELTA", x_in=fx, y_in=fy + h1 + h2 + 0.30, w_in=min(fw, w3), h_in=h3, svg_text=eq3))

    eq4, w4, h4 = _mpl_math_svg(
        r"$A(p)=\sum\,_{t} v(\mathrm{rank}_t(p))\quad\mathrm{vs}\quad R(p)=\sum\,_{t} u(p\mid q_t)$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=20.0,
    )
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_AMORT", x_in=fx, y_in=fy + h1 + h2 + h3 + 0.44, w_in=min(fw, w4), h_in=h4, svg_text=eq4))

    heat_nx = 14
    heat_ny = 14
    cell_in = 0.17
    gap_in = 0.02
    heat_w_in = heat_nx * cell_in + (heat_nx - 1) * gap_in
    heat_h_in = heat_ny * cell_in + (heat_ny - 1) * gap_in
    heat_x = fx
    heat_y = fy + h1 + h2 + h3 + h4 + 0.90

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_CAP_HEAT",
            x_in=heat_x,
            y_in=heat_y - 0.36,
            w_in=heat_w_in,
            h_in=0.30,
            svg_text=_svg_root(
                w_px=_in_to_px(heat_w_in),
                h_px=_in_to_px(0.30),
                body=_svg_text(x=0, y=_in_to_px(0.15), text="Within‑cluster exposure gaps (illustration)", size_px=16, color=f"#{COLORS_HEX['muted']}", weight=700),
            ),
            static=True,
        )
    )

    for ry in range(heat_ny):
        for cx in range(heat_nx):
            heat = ((cx + 1) * (ry + 2)) % 13 / 13.0
            tiles.append(
                SvgTile(
                    slide_num=2,
                    name=f"S2_HEAT_{ry:02d}_{cx:02d}",
                    x_in=heat_x + cx * (cell_in + gap_in),
                    y_in=heat_y + ry * (cell_in + gap_in),
                    w_in=cell_in,
                    h_in=cell_in,
                    svg_text=_svg_cell(w_px=_in_to_px(cell_in), h_px=_in_to_px(cell_in), val=None, heat=heat, accent_hex=COLORS_HEX["coral"]),
                )
            )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_OPS",
            x_in=fx,
            y_in=heat_y + heat_h_in + 0.18,
            w_in=fw,
            h_in=0.92,
            svg_text=_svg_caption(
                w_px=_in_to_px(fw),
                h_px=_in_to_px(0.92),
                lines=[
                    "Operationalize in Bluesky tables:",
                    "• cluster_id from posts_first_seen.text (exact/near‑dup)",
                    "• E(post|q)=v(rank) from feed_items",
                    "• test ΔE within cluster + across viewer_mode/vantage_id",
                ],
            ),
        )
    )

    return tiles


def _build_slide1() -> list[SvgTile]:
    tiles: list[SvgTile] = []

    tiles.append(
        SvgTile(
            slide_num=1,
            name="_BG_S1_GRID",
            x_in=0.0,
            y_in=0.0,
            w_in=SLIDE_W_IN,
            h_in=SLIDE_H_IN,
            svg_text=_svg_grid_bg(w_px=_in_to_px(SLIDE_W_IN), h_px=_in_to_px(SLIDE_H_IN)),
            static=True,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_TITLE",
            x_in=SAFE_L_IN,
            y_in=0.28,
            w_in=SAFE_W_IN,
            h_in=0.92,
            svg_text=_svg_title_block(
                w_px=_in_to_px(SAFE_W_IN),
                h_px=_in_to_px(0.92),
                title="Exposure = expected attention",
                subtitle="Singh & Joachims (KDD '18): rank marginals P and position bias v define Exposure.",
            ),
            static=True,
        )
    )

    col_gap = 0.28
    left_w = 6.85
    right_w = SAFE_W_IN - left_w - col_gap
    left_x = SAFE_L_IN
    right_x = SAFE_L_IN + left_w + col_gap
    body_y = 1.34
    panel_h = SAFE_H_IN - (body_y - SAFE_T_IN)

    for name, x, w in [
        ("_BG_S1_PANEL_THEORY", left_x, left_w),
        ("_BG_S1_PANEL_DATA", right_x, right_w),
    ]:
        tiles.append(
            SvgTile(
                slide_num=1,
                name=name,
                x_in=x,
                y_in=body_y,
                w_in=w,
                h_in=panel_h,
                svg_text=_svg_root(
                    w_px=_in_to_px(w),
                    h_px=_in_to_px(panel_h),
                    body=_svg_rounded_rect(
                        x=0,
                        y=0,
                        w=_in_to_px(w),
                        h=_in_to_px(panel_h),
                        r=22,
                        fill=f"#{COLORS_HEX['card']}",
                        stroke=f"#{COLORS_HEX['card_line']}",
                        stroke_w=2,
                        opacity=0.98,
                    ),
                ),
                static=True,
            )
        )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_LBL_THEORY",
            x_in=left_x + 0.22,
            y_in=body_y + 0.12,
            w_in=4.15,
            h_in=0.35,
            svg_text=_svg_panel_label(
                w_px=_in_to_px(4.15),
                h_px=_in_to_px(0.35),
                text="Theory: probabilistic ranking",
                accent_hex=COLORS_HEX["cyan"],
            ),
            static=True,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_LBL_DATA",
            x_in=right_x + 0.22,
            y_in=body_y + 0.12,
            w_in=right_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(
                w_px=_in_to_px(right_w - 0.44),
                h_px=_in_to_px(0.35),
                text="Bluesky mapping: snapshots → exposure",
                accent_hex=COLORS_HEX["coral"],
            ),
            static=True,
        )
    )

    # ---------- Left panel content ----------
    lx = left_x + 0.30
    left_inner_w = left_w - 0.60

    eq_y = body_y + 0.50
    _, strip_h = _add_math_strip(
        tiles=tiles,
        slide_num=1,
        name_prefix="S1_EQSTRIP",
        x_in=lx,
        y_in=eq_y,
        tokens=[
            r"$\mathrm{Exposure}(d_i\mid P)$",
            r"$=$",
            r"$\sum\,_{j}$",
            r"$P_{i,j}$",
            r"$v_j$",
        ],
        color_hex=COLORS_HEX["ink"],
        token_color_hexes=[
            COLORS_HEX["ink"],
            COLORS_HEX["muted"],
            COLORS_HEX["ink"],
            COLORS_HEX["cyan"],
            COLORS_HEX["green"],
        ],
        fontsize_pt=20.0,
    )

    # Small, readable example (dense, doubly-stochastic by construction).
    grid_n = 13
    cell_in = 0.185
    gap_in = 0.02
    grid_w_in = grid_n * cell_in + (grid_n - 1) * gap_in
    grid_h_in = grid_w_in
    grid_x = lx
    grid_y = eq_y + strip_h + 0.30

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_P",
            x_in=grid_x,
            y_in=grid_y - 0.30,
            w_in=2.05,
            h_in=0.22,
            svg_text=_svg_pill(w_px=_in_to_px(2.05), h_px=_in_to_px(0.22), text="P (rank marginals)", accent_hex=COLORS_HEX["cyan"]),
            static=True,
        )
    )

    # Circulant weights: each row is a distribution over ranks; each column also sums to 1.
    #
    # Important: the *labels* in the grid are rounded, and we want the "rows/cols sum to 1"
    # property to remain visually true even at 2 decimals. So we round using a largest-remainder
    # method to get 2-decimal probabilities that still sum to exactly 1.00.
    raw = list(range(grid_n, 0, -1))  # 13..1
    s = float(sum(raw))
    weights_f = [w / s for w in raw]
    cents_f = [w * 100.0 for w in weights_f]
    cents = [int(w) for w in cents_f]  # floor since all positive
    remainder = 100 - sum(cents)
    for idx in sorted(range(grid_n), key=lambda i: cents_f[i] - cents[i], reverse=True)[: max(0, remainder)]:
        cents[idx] += 1
    weights = [c / 100.0 for c in cents]
    hi_r = grid_n // 2  # visually emphasize "pick a row i"
    tiles.append(
        SvgTile(
            slide_num=1,
            name="_BG_S1_P_ROW_HI",
            x_in=grid_x - 0.01,
            y_in=grid_y + hi_r * (cell_in + gap_in) - 0.01,
            w_in=grid_w_in + 0.02,
            h_in=cell_in + 0.02,
            svg_text=_svg_root(
                w_px=_in_to_px(grid_w_in + 0.02),
                h_px=_in_to_px(cell_in + 0.02),
                body=_svg_rounded_rect(
                    x=0,
                    y=0,
                    w=_in_to_px(grid_w_in + 0.02),
                    h=_in_to_px(cell_in + 0.02),
                    r=10,
                    fill=f"#{COLORS_HEX['cyan']}",
                    stroke=f"#{COLORS_HEX['cyan']}",
                    stroke_w=0,
                    opacity=0.14,
                ),
            ),
            static=True,
        )
    )
    max_w = max(weights)
    for r in range(grid_n):
        for c in range(grid_n):
            x = grid_x + c * (cell_in + gap_in)
            y = grid_y + r * (cell_in + gap_in)
            p = weights[(c - r) % grid_n]
            heat = p / max_w
            val = f"{p:.2f}"
            tiles.append(
                SvgTile(
                    slide_num=1,
                    name=f"S1_P_{r:02d}_{c:02d}",
                    x_in=x,
                    y_in=y,
                    w_in=cell_in,
                    h_in=cell_in,
                    svg_text=_svg_cell(w_px=_in_to_px(cell_in), h_px=_in_to_px(cell_in), val=val, heat=heat, accent_hex=COLORS_HEX["cyan"]),
                )
            )

    v_x = grid_x + grid_w_in + 0.30
    v_y = grid_y
    bar_w_in = 2.55
    bar_h_in = 0.15
    bar_gap_in = 0.04
    v_total_h = grid_n * bar_h_in + (grid_n - 1) * bar_gap_in

    # Visual integration: connect equation tokens to P and v (no text, drawn behind).
    conn_h_in = max(0.0, grid_y - eq_y)
    p_tok = next((t for t in tiles if t.slide_num == 1 and t.name == "S1_EQSTRIP_04"), None)
    v_tok = next((t for t in tiles if t.slide_num == 1 and t.name == "S1_EQSTRIP_05"), None)
    if p_tok is not None and v_tok is not None and conn_h_in > 0.12:
        conn_x = left_x
        conn_y = eq_y
        conn_w = left_w
        conn_h = conn_h_in

        def _rel_px(x_in: float, y_in: float) -> tuple[int, int]:
            return _in_to_px(x_in - conn_x), _in_to_px(y_in - conn_y)

        spx, spy = _rel_px(p_tok.x_in + p_tok.w_in / 2, p_tok.y_in + p_tok.h_in / 2)
        svx, svy = _rel_px(v_tok.x_in + v_tok.w_in / 2, v_tok.y_in + v_tok.h_in / 2)

        epy_x, epy_y = _rel_px(grid_x + grid_w_in * 0.55, grid_y - 0.01)
        evy_x, evy_y = _rel_px(v_x + bar_w_in * 0.55, v_y - 0.01)

        def _curve(x1: int, y1: int, x2: int, y2: int) -> str:
            mx = int(round((x1 + x2) / 2))
            return f"M{x1} {y1} C {mx} {y1} {mx} {y2} {x2} {y2}"

        conn_body = ""
        conn_body += f'<path d="{_curve(spx, spy, epy_x, epy_y)}" fill="none" stroke="#{COLORS_HEX["cyan"]}" stroke-width="2" stroke-opacity="0.30"/>'
        conn_body += f'<circle cx="{epy_x}" cy="{epy_y}" r="3" fill="#{COLORS_HEX["cyan"]}" opacity="0.30"/>'
        conn_body += f'<path d="{_curve(svx, svy, evy_x, evy_y)}" fill="none" stroke="#{COLORS_HEX["green"]}" stroke-width="2" stroke-opacity="0.30"/>'
        conn_body += f'<circle cx="{evy_x}" cy="{evy_y}" r="3" fill="#{COLORS_HEX["green"]}" opacity="0.30"/>'

        conn_tile = SvgTile(
            slide_num=1,
            name="_BG_S1_EQ_LINKS",
            x_in=conn_x,
            y_in=conn_y,
            w_in=conn_w,
            h_in=conn_h,
            svg_text=_svg_root(w_px=_in_to_px(conn_w), h_px=_in_to_px(conn_h), body=conn_body),
            static=True,
        )
        for ins_idx, tt in enumerate(tiles):
            if tt.slide_num == 1 and tt.name == "S1_EQSTRIP_01":
                tiles.insert(ins_idx, conn_tile)
                break

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_V",
            x_in=v_x,
            y_in=v_y - 0.30,
            w_in=1.75,
            h_in=0.22,
            svg_text=_svg_pill(w_px=_in_to_px(1.75), h_px=_in_to_px(0.22), text="v (position bias)", accent_hex=COLORS_HEX["green"]),
            static=True,
        )
    )

    for j in range(1, grid_n + 1):
        frac = 1.0 / (1.0 + 0.18 * (j - 1))
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_V_{j:02d}",
                x_in=v_x,
                y_in=v_y + (j - 1) * (bar_h_in + bar_gap_in),
                w_in=bar_w_in,
                h_in=bar_h_in,
                svg_text=_svg_bar_row(w_px=_in_to_px(bar_w_in), h_px=_in_to_px(bar_h_in), j=j, frac=frac, accent_hex=COLORS_HEX["green"]),
            )
        )

    arrow_y = max(grid_y + grid_h_in, v_y + v_total_h) + 0.06
    arrow_h = 0.34
    out_gap = 0.06
    out_h = 0.58
    out_y = arrow_y + arrow_h + out_gap

    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_ARROW_PV",
            x_in=v_x,
            y_in=arrow_y,
            w_in=bar_w_in,
            h_in=arrow_h,
            svg_text=_svg_arrow(w_px=_in_to_px(bar_w_in), h_px=_in_to_px(arrow_h), accent_hex=COLORS_HEX["amber"], text="P·v"),
        )
    )

    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_OUT_E",
            x_in=v_x,
            y_in=out_y,
            w_in=bar_w_in,
            h_in=out_h,
            svg_text=_svg_scalar_bar(
                w_px=_in_to_px(bar_w_in),
                h_px=_in_to_px(out_h),
                label="Exposure(d_i)",
                frac=0.62,
                accent_hex=COLORS_HEX["amber"],
                value="E_i",
            ),
        )
    )

    ds_h = (out_y + out_h) - arrow_y
    ds_bg = _svg_root(
        w_px=_in_to_px(grid_w_in),
        h_px=_in_to_px(ds_h),
        body=_svg_rounded_rect(
            x=0,
            y=0,
            w=_in_to_px(grid_w_in),
            h=_in_to_px(ds_h),
            r=18,
            fill=f"#{COLORS_HEX['card2']}",
            stroke=f"#{COLORS_HEX['card_line']}",
            stroke_w=2,
        ),
    )
    tiles.append(
        SvgTile(
            slide_num=1,
            name="_BG_S1_DS_CARD",
            x_in=grid_x,
            y_in=arrow_y,
            w_in=grid_w_in,
            h_in=ds_h,
            svg_text=ds_bg,
            static=True,
        )
    )
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_DS",
            x_in=grid_x + 0.10,
            y_in=arrow_y + 0.04,
            w_in=grid_w_in - 0.20,
            h_in=0.16,
            svg_text=_svg_pill(
                w_px=_in_to_px(grid_w_in - 0.20),
                h_px=_in_to_px(0.16),
                text="Doubly‑stochastic P",
                accent_hex=COLORS_HEX["cyan"],
            ),
            static=True,
        )
    )
    util_svg, util_w, util_h = _mpl_math_svg(r"$U(P\mid q)=u^\top P\,v$", color_hex=COLORS_HEX["muted"], fontsize_pt=16.0)
    util_w2, util_h2 = _scale_to_fit(w_in=util_w, h_in=util_h, max_w_in=grid_w_in - 0.20, max_h_in=0.22)
    tiles.append(SvgTile(slide_num=1, name="S1_DS_UTIL", x_in=grid_x + 0.10, y_in=arrow_y + 0.22, w_in=util_w2, h_in=util_h2, svg_text=util_svg))
    row_svg, row_w, row_h = _mpl_math_svg(r"$\sum\,_{j} P_{i,j}=1$", color_hex=COLORS_HEX["cyan"], fontsize_pt=16.0)
    row_w2, row_h2 = _scale_to_fit(w_in=row_w, h_in=row_h, max_w_in=grid_w_in - 0.20, max_h_in=0.18)
    tiles.append(SvgTile(slide_num=1, name="S1_DS_ROW", x_in=grid_x + 0.10, y_in=arrow_y + 0.50, w_in=row_w2, h_in=row_h2, svg_text=row_svg))
    col_svg, col_w, col_h = _mpl_math_svg(r"$\sum\,_{i} P_{i,j}=1$", color_hex=COLORS_HEX["cyan"], fontsize_pt=16.0)
    col_w2, col_h2 = _scale_to_fit(w_in=col_w, h_in=col_h, max_w_in=grid_w_in - 0.20, max_h_in=0.18)
    tiles.append(SvgTile(slide_num=1, name="S1_DS_COL", x_in=grid_x + 0.10, y_in=arrow_y + 0.72, w_in=col_w2, h_in=col_h2, svg_text=col_svg))

    # Where P comes from: distribution over rankings (π) → rank marginals (P)
    panel_bottom = body_y + panel_h
    pi_x = grid_x
    pi_y = out_y + out_h + 0.18
    pi_h = panel_bottom - pi_y - 0.10
    if pi_h > 0.60:
        pi_w = left_inner_w
        pi_bg = _svg_root(
            w_px=_in_to_px(pi_w),
            h_px=_in_to_px(pi_h),
            body=_svg_rounded_rect(
                x=0,
                y=0,
                w=_in_to_px(pi_w),
                h=_in_to_px(pi_h),
                r=18,
                fill=f"#{COLORS_HEX['card2']}",
                stroke=f"#{COLORS_HEX['card_line']}",
                stroke_w=2,
            ),
        )
        tiles.append(SvgTile(slide_num=1, name="_BG_S1_CARD_PI", x_in=pi_x, y_in=pi_y, w_in=pi_w, h_in=pi_h, svg_text=pi_bg, static=True))
        tiles.append(
            SvgTile(
                slide_num=1,
                name="STATIC_S1_PILL_PI",
                x_in=pi_x + 0.14,
                y_in=pi_y + 0.06,
                w_in=2.55,
                h_in=0.18,
                svg_text=_svg_pill(w_px=_in_to_px(2.55), h_px=_in_to_px(0.18), text="From π → P", accent_hex=COLORS_HEX["violet"]),
                static=True,
            )
        )
        eq_pi1, w_pi1, h_pi1 = _mpl_math_svg(r"$P_{i,j}=\mathrm{Pr}_{\tau\sim\pi}[\tau(i)=j]$", color_hex=COLORS_HEX["ink"], fontsize_pt=18.0)
        w_pi1b, h_pi1b = _scale_to_fit(w_in=w_pi1, h_in=h_pi1, max_w_in=pi_w - 0.28, max_h_in=0.34)
        eq1_y = pi_y + 0.28
        tiles.append(SvgTile(slide_num=1, name="S1_EQ_PI_1", x_in=pi_x + 0.14, y_in=eq1_y, w_in=w_pi1b, h_in=h_pi1b, svg_text=eq_pi1))

        eq_pi2, w_pi2, h_pi2 = _mpl_math_svg(r"$P=\mathrm{E}_{\tau\sim\pi}[\Pi_{\tau}]$", color_hex=COLORS_HEX["muted"], fontsize_pt=18.0)
        w_pi2b, h_pi2b = _scale_to_fit(w_in=w_pi2, h_in=h_pi2, max_w_in=pi_w - 0.28, max_h_in=0.34)
        eq2_y = eq1_y + h_pi1b + 0.06
        if eq2_y + h_pi2b <= pi_y + pi_h - 0.10:
            tiles.append(SvgTile(slide_num=1, name="S1_EQ_PI_2", x_in=pi_x + 0.14, y_in=eq2_y, w_in=w_pi2b, h_in=h_pi2b, svg_text=eq_pi2))

    # ---------- Right panel content ----------
    rx = right_x + 0.22
    right_inner_w = right_w - 0.44

    chip_w_in = (right_inner_w - 0.12) / 2
    chip_h_in = 0.52
    chip_y0 = body_y + 0.60
    chips = [
        ("feed_uri", "f", COLORS_HEX["cyan"]),
        ("hour_utc", "t", COLORS_HEX["amber"]),
        ("viewer_mode", "m", COLORS_HEX["green"]),
        ("vantage_id", "ν", COLORS_HEX["violet"]),
    ]
    for idx, (head, val, accent) in enumerate(chips):
        cx = rx + (idx % 2) * (chip_w_in + 0.12)
        cy = chip_y0 + (idx // 2) * (chip_h_in + 0.12)
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_Q_{idx+1:02d}",
                x_in=cx,
                y_in=cy,
                w_in=chip_w_in,
                h_in=chip_h_in,
                svg_text=_svg_chip(w_px=_in_to_px(chip_w_in), h_px=_in_to_px(chip_h_in), head=head, val=val, accent_hex=accent),
            )
        )

    list_y = chip_y0 + 2 * (chip_h_in + 0.12) + 0.26
    split_gap = 0.18
    list_col_w = min(2.35, right_inner_w * 0.58)
    side_w = right_inner_w - list_col_w - split_gap

    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_TAU",
            x_in=rx,
            y_in=list_y - 0.30,
            w_in=list_col_w,
            h_in=0.22,
            svg_text=_svg_pill(w_px=_in_to_px(list_col_w), h_px=_in_to_px(0.22), text="Observed ranking τ_t (top‑K)", accent_hex=COLORS_HEX["coral"]),
            static=True,
        )
    )

    row_h_in = 0.22
    row_gap_in = 0.04
    k = 6
    for r in range(1, k + 1):
        tiles.append(
            SvgTile(
                slide_num=1,
                name=f"S1_RANK_{r:02d}",
                x_in=rx,
                y_in=list_y + (r - 1) * (row_h_in + row_gap_in),
                w_in=list_col_w,
                h_in=row_h_in,
                svg_text=_svg_rank_row(
                    w_px=_in_to_px(list_col_w),
                    h_px=_in_to_px(row_h_in),
                    rank=r,
                    label="post_uri · author",
                    accent_hex=COLORS_HEX["coral"] if r <= 5 else COLORS_HEX["card_line"],
                ),
            )
        )

    list_h = k * row_h_in + (k - 1) * row_gap_in
    sx = rx + list_col_w + split_gap
    pipe_h = list_h
    pipe_bg = _svg_root(
        w_px=_in_to_px(side_w),
        h_px=_in_to_px(pipe_h),
        body=_svg_rounded_rect(x=0, y=0, w=_in_to_px(side_w), h=_in_to_px(pipe_h), r=18, fill=f"#{COLORS_HEX['card2']}", stroke=f"#{COLORS_HEX['card_line']}", stroke_w=2)
        + f'<rect x="0" y="0" width="10" height="{_in_to_px(pipe_h)}" fill="#{COLORS_HEX["coral"]}" opacity="0.16"/>'
        + f'<rect x="0" y="{int(round(_in_to_px(pipe_h)*0.50))}" width="10" height="{_in_to_px(pipe_h) - int(round(_in_to_px(pipe_h)*0.50))}" fill="#{COLORS_HEX["violet"]}" opacity="0.20"/>'
        + f'<line x1="10" y1="{int(round(_in_to_px(pipe_h)*0.50))}" x2="{_in_to_px(side_w)-10}" y2="{int(round(_in_to_px(pipe_h)*0.50))}" stroke="#{COLORS_HEX["card_line"]}" stroke-opacity="0.70" stroke-width="1"/>',
    )
    tiles.append(SvgTile(slide_num=1, name="_BG_S1_CARD_PIPE", x_in=sx, y_in=list_y, w_in=side_w, h_in=pipe_h, svg_text=pipe_bg, static=True))

    # Top half: per snapshot
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_PIPE_SNAP",
            x_in=sx + 0.09,
            y_in=list_y + 0.04,
            w_in=side_w - 0.18,
            h_in=0.16,
            svg_text=_svg_pill(w_px=_in_to_px(side_w - 0.18), h_px=_in_to_px(0.16), text="One snapshot", accent_hex=COLORS_HEX["coral"]),
            static=True,
        )
    )
    eq_svg, eq_w_in, eq_h_in = _mpl_math_svg(r"$E(p\mid q)=v(\mathrm{rank})$", color_hex=COLORS_HEX["ink"], fontsize_pt=18.5)
    eq_w2, eq_h2 = _scale_to_fit(w_in=eq_w_in, h_in=eq_h_in, max_w_in=side_w - 0.18, max_h_in=0.26)
    tiles.append(SvgTile(slide_num=1, name="S1_PIPE_EQ_E", x_in=sx + 0.09, y_in=list_y + 0.24, w_in=eq_w2, h_in=eq_h2, svg_text=eq_svg))
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PIPE_NOTE_E0",
            x_in=sx + 0.09,
            y_in=list_y + 0.52,
            w_in=side_w - 0.18,
            h_in=0.14,
            svg_text=_svg_root(
                w_px=_in_to_px(side_w - 0.18),
                h_px=_in_to_px(0.14),
                body=_svg_text(x=0, y=_in_to_px(0.07), text="absent in top‑K ⇒ E=0 (lower bound)", size_px=12, color=f"#{COLORS_HEX['muted']}", weight=650, baseline="middle"),
            ),
            static=True,
        )
    )

    # Bottom half: empirical P̂ across snapshots
    mid_y = list_y + pipe_h * 0.50
    tiles.append(
        SvgTile(
            slide_num=1,
            name="STATIC_S1_PILL_PIPE_PHAT",
            x_in=sx + 0.09,
            y_in=mid_y + 0.02,
            w_in=side_w - 0.18,
            h_in=0.16,
            svg_text=_svg_pill(w_px=_in_to_px(side_w - 0.18), h_px=_in_to_px(0.16), text="Across snapshots", accent_hex=COLORS_HEX["violet"]),
            static=True,
        )
    )
    ph1_svg, ph1_w_in, ph1_h_in = _mpl_math_svg(r"$\hat P_{p,j}=\mathrm{freq}_t[\mathrm{rank}_t(p)=j]$", color_hex=COLORS_HEX["muted"], fontsize_pt=15.8)
    ph1_w2, ph1_h2 = _scale_to_fit(w_in=ph1_w_in, h_in=ph1_h_in, max_w_in=side_w - 0.18, max_h_in=0.22)
    ph1_y = mid_y + 0.22
    tiles.append(SvgTile(slide_num=1, name="S1_PIPE_EQ_PH1", x_in=sx + 0.09, y_in=ph1_y, w_in=ph1_w2, h_in=ph1_h2, svg_text=ph1_svg))
    ph2_svg, ph2_w_in, ph2_h_in = _mpl_math_svg(r"$\hat E(p)=\sum\,_{j} \hat P_{p,j}v_j$", color_hex=COLORS_HEX["muted"], fontsize_pt=16.8)
    ph2_w2, ph2_h2 = _scale_to_fit(w_in=ph2_w_in, h_in=ph2_h_in, max_w_in=side_w - 0.18, max_h_in=0.22)
    tiles.append(SvgTile(slide_num=1, name="S1_PIPE_EQ_PH2", x_in=sx + 0.09, y_in=ph1_y + ph1_h2 + 0.04, w_in=ph2_w2, h_in=ph2_h2, svg_text=ph2_svg))

    join_y = list_y + list_h + 0.18
    join_h = 0.68
    tiles.append(
        SvgTile(
            slide_num=1,
            name="S1_JOIN_GRAPH",
            x_in=rx,
            y_in=join_y,
            w_in=right_inner_w,
            h_in=join_h,
            svg_text=_svg_join_graph(w_px=_in_to_px(right_inner_w), h_px=_in_to_px(join_h), accent_hex=COLORS_HEX["coral"]),
        )
    )

    # Saved space: promotion probability (top‑k) as an exposure-friendly outcome.
    panel_bottom = body_y + panel_h
    promo_x = rx
    promo_y = join_y + join_h + 0.12
    promo_h = panel_bottom - promo_y - 0.10
    if promo_h > 0.62:
        promo_w = right_inner_w
        promo_bg = _svg_root(
            w_px=_in_to_px(promo_w),
            h_px=_in_to_px(promo_h),
            body=_svg_rounded_rect(
                x=0,
                y=0,
                w=_in_to_px(promo_w),
                h=_in_to_px(promo_h),
                r=18,
                fill=f"#{COLORS_HEX['card2']}",
                stroke=f"#{COLORS_HEX['card_line']}",
                stroke_w=2,
            ),
        )
        tiles.append(SvgTile(slide_num=1, name="_BG_S1_CARD_PROMO", x_in=promo_x, y_in=promo_y, w_in=promo_w, h_in=promo_h, svg_text=promo_bg, static=True))
        tiles.append(
            SvgTile(
                slide_num=1,
                name="STATIC_S1_PILL_PROMO",
                x_in=promo_x + 0.14,
                y_in=promo_y + 0.06,
                w_in=promo_w - 0.28,
                h_in=0.18,
                svg_text=_svg_pill(w_px=_in_to_px(promo_w - 0.28), h_px=_in_to_px(0.18), text="Promotion probability (top‑k)", accent_hex=COLORS_HEX["amber"]),
                static=True,
            )
        )

        eq_y1, w_y1, h_y1 = _mpl_math_svg(
            r"$y_t=\mathbf{1}[\mathrm{rank}_t(p)\leq k]$",
            color_hex=COLORS_HEX["ink"],
            fontsize_pt=16.5,
        )
        w_y1b, h_y1b = _scale_to_fit(w_in=w_y1, h_in=h_y1, max_w_in=promo_w - 0.28, max_h_in=0.24)
        y1_x = promo_x + 0.14
        y1_y = promo_y + 0.26
        tiles.append(SvgTile(slide_num=1, name="S1_EQ_TOPK_Y", x_in=y1_x, y_in=y1_y, w_in=w_y1b, h_in=h_y1b, svg_text=eq_y1))

        bar_h = 0.26
        bar_y = y1_y + h_y1b + 0.06
        if bar_y + bar_h <= promo_y + promo_h - 0.10:
            tiles.append(
                SvgTile(
                    slide_num=1,
                    name="S1_PROMO_BAR",
                    x_in=promo_x + 0.14,
                    y_in=bar_y,
                    w_in=promo_w - 0.28,
                    h_in=bar_h,
                    svg_text=_svg_scalar_bar(
                        w_px=_in_to_px(promo_w - 0.28),
                        h_px=_in_to_px(bar_h),
                        label="p̂(p) = freq_t[rank≤k]",
                        frac=0.37,
                        accent_hex=COLORS_HEX["amber"],
                        value="p̂",
                    ),
                )
            )

    return tiles


def _build_slide2() -> list[SvgTile]:
    tiles: list[SvgTile] = []

    tiles.append(
        SvgTile(
            slide_num=2,
            name="_BG_S2_GRID",
            x_in=0.0,
            y_in=0.0,
            w_in=SLIDE_W_IN,
            h_in=SLIDE_H_IN,
            svg_text=_svg_grid_bg(w_px=_in_to_px(SLIDE_W_IN), h_px=_in_to_px(SLIDE_H_IN)),
            static=True,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_S2_TITLE",
            x_in=SAFE_L_IN,
            y_in=0.28,
            w_in=SAFE_W_IN,
            h_in=0.92,
            svg_text=_svg_title_block(
                w_px=_in_to_px(SAFE_W_IN),
                h_px=_in_to_px(0.92),
                title="Individual fairness: same content → different exposure",
                subtitle="Zehlike et al. (CSUR '22, Part II): choose similarity d(p,p'), then audit exposure gaps.",
            ),
            static=True,
        )
    )

    col_gap = 0.28
    left_w = 6.10
    right_w = SAFE_W_IN - left_w - col_gap
    left_x = SAFE_L_IN
    right_x = SAFE_L_IN + left_w + col_gap
    body_y = 1.34
    panel_h = SAFE_H_IN - (body_y - SAFE_T_IN)

    for name, x, w in [
        ("_BG_S2_PANEL_LEFT", left_x, left_w),
        ("_BG_S2_PANEL_RIGHT", right_x, right_w),
    ]:
        tiles.append(
            SvgTile(
                slide_num=2,
                name=name,
                x_in=x,
                y_in=body_y,
                w_in=w,
                h_in=panel_h,
                svg_text=_svg_root(
                    w_px=_in_to_px(w),
                    h_px=_in_to_px(panel_h),
                    body=_svg_rounded_rect(
                        x=0,
                        y=0,
                        w=_in_to_px(w),
                        h=_in_to_px(panel_h),
                        r=22,
                        fill=f"#{COLORS_HEX['card']}",
                        stroke=f"#{COLORS_HEX['card_line']}",
                        stroke_w=2,
                        opacity=0.98,
                    ),
                ),
                static=True,
            )
        )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_S2_LBL_SIM",
            x_in=left_x + 0.22,
            y_in=body_y + 0.12,
            w_in=left_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(w_px=_in_to_px(left_w - 0.44), h_px=_in_to_px(0.35), text="Define similarity over posts", accent_hex=COLORS_HEX["cyan"]),
            static=True,
        )
    )
    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_S2_LBL_METRICS",
            x_in=right_x + 0.22,
            y_in=body_y + 0.12,
            w_in=right_w - 0.44,
            h_in=0.35,
            svg_text=_svg_panel_label(w_px=_in_to_px(right_w - 0.44), h_px=_in_to_px(0.35), text="Measure exposure gaps", accent_hex=COLORS_HEX["coral"]),
            static=True,
        )
    )

    # ---------- Left panel content: post cards + cluster ----------
    lx = left_x + 0.22
    lw = left_w - 0.44
    col_gap2 = 0.18
    col1_w = 3.70
    col2_w = lw - col1_w - col_gap2

    y0 = body_y + 0.60
    post_h = 0.92
    post_gap = 0.12
    posts = [
        ("author_did = A", ["“same / near‑same text …”", "Δt small (optional)"], 3, COLORS_HEX["cyan"]),
        ("author_did = B", ["“same / near‑same text …”", "Δt small (optional)"], 12, COLORS_HEX["amber"]),
        ("author_did = C", ["“same / near‑same text …”", "Δt small (optional)"], 25, COLORS_HEX["green"]),
    ]

    # Visual links: each near-duplicate post → cluster g (no text; sits behind cards).
    link_h = max(1.45, 3 * post_h + 2 * post_gap)
    link_body = ""
    x_start = _in_to_px(col1_w) - 2
    x_end = _in_to_px(col1_w + col_gap2) + 2
    for i, (_, _, rank_i, accent) in enumerate(posts, start=1):
        y_start = _in_to_px((i - 1) * (post_h + post_gap) + post_h / 2)
        y_end = _in_to_px(1.45 / 2 + (i - 2) * 0.18)
        mx = int(round((x_start + x_end) / 2))
        path = f"M{x_start} {int(round(y_start))} C {mx} {int(round(y_start))} {mx} {int(round(y_end))} {x_end} {int(round(y_end))}"
        op = 0.22 + 0.55 * (1.0 / (1.0 + 0.18 * (max(1, int(rank_i)) - 1)))
        link_body += f'<path d="{path}" fill="none" stroke="#{accent}" stroke-width="2" stroke-opacity="{op:.3f}"/>'
        link_body += f'<circle cx="{x_end}" cy="{int(round(y_end))}" r="3" fill="#{accent}" opacity="{op:.3f}"/>'

    tiles.append(
        SvgTile(
            slide_num=2,
            name="_BG_S2_LINKS_CLUSTER",
            x_in=lx,
            y_in=y0,
            w_in=lw,
            h_in=link_h,
            svg_text=_svg_root(w_px=_in_to_px(lw), h_px=_in_to_px(link_h), body=link_body),
            static=True,
        )
    )

    for i, (author, lines, rank_i, accent) in enumerate(posts, start=1):
        tiles.append(
            SvgTile(
                slide_num=2,
                name=f"S2_POST_{i:02d}",
                x_in=lx,
                y_in=y0 + (i - 1) * (post_h + post_gap),
                w_in=col1_w,
                h_in=post_h,
                svg_text=_svg_post_card(w_px=_in_to_px(col1_w), h_px=_in_to_px(post_h), author=author, lines=lines, rank=int(rank_i), accent_hex=accent),
            )
        )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_CLUSTER",
            x_in=lx + col1_w + col_gap2,
            y_in=y0,
            w_in=col2_w,
            h_in=1.45,
            svg_text=_svg_cluster_bubble(w_px=_in_to_px(col2_w), h_px=_in_to_px(1.45), label="cluster g", accent_hex=COLORS_HEX["cyan"]),
        )
    )

    sim_svg, sim_w, sim_h = _mpl_math_svg(
        r"$d(p,p')=\alpha\,d_{\text{text}}(p,p')+(1-\alpha)\,d_{\text{time}}(p,p')$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=18.0,
    )
    sim_w2, sim_h2 = _scale_to_fit(w_in=sim_w, h_in=sim_h, max_w_in=col2_w, max_h_in=0.60)
    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_EQ_SIM",
            x_in=lx + col1_w + col_gap2,
            y_in=y0 + 1.55,
            w_in=sim_w2,
            h_in=sim_h2,
            svg_text=sim_svg,
        )
    )

    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_AUDITS",
            x_in=lx,
            y_in=y0 + 3 * (post_h + post_gap) + 0.10,
            w_in=lw,
            h_in=1.40,
            svg_text=_svg_caption(
                w_px=_in_to_px(lw),
                h_px=_in_to_px(1.40),
                lines=[
                    "Two audits you can run from the same exposure math:",
                    "A) same post_uri across contexts → ΔE(p; q1,q2)",
                    "B) same/near‑same text within cluster g → |E(p)-E(p')|",
                    "Control similarity: include posting time so claims ≠ “it’s just time”.",
                ],
                accent_hex=COLORS_HEX["cyan"],
            ),
        )
    )

    # ---------- Right panel content: metric cards + heatmap + timeline ----------
    fx = right_x + 0.22
    fw = right_w - 0.44
    fy0 = body_y + 0.60

    gapx = 0.14
    card_w = (fw - 2 * gapx) / 3
    card_h = 1.08

    for i, (bg_name, x_off, title, subtitle, stroke_hex, kind) in enumerate(
        [
            ("_BG_S2_CARD_D", 0.0, "Duplicate gap", "D(p,p')", COLORS_HEX["amber"], "D"),
            ("_BG_S2_CARD_LIP", card_w + gapx, "Individual fairness", "d≈0 ⇒ |E−E'|≤ε", COLORS_HEX["coral"], "IF"),
            ("_BG_S2_CARD_DELTA", 2 * (card_w + gapx), "Context ΔE", "same post, different q", COLORS_HEX["violet"], "DELTA"),
        ],
        start=1,
    ):
        bg = _svg_metric_card_bg(w_px=_in_to_px(card_w), h_px=_in_to_px(card_h), stroke_hex=stroke_hex, kind=kind)
        card_x = fx + x_off
        tiles.append(SvgTile(slide_num=2, name=bg_name, x_in=card_x, y_in=fy0, w_in=card_w, h_in=card_h, svg_text=bg, static=True))
        tiles.append(
            SvgTile(
                slide_num=2,
                name=f"STATIC_S2_CARD_{i:02d}_PILL",
                x_in=card_x + 0.08,
                y_in=fy0 + 0.06,
                w_in=card_w - 0.16,
                h_in=0.18,
                svg_text=_svg_pill(w_px=_in_to_px(card_w - 0.16), h_px=_in_to_px(0.18), text=title, accent_hex=stroke_hex),
                static=True,
            )
        )
        tiles.append(
            SvgTile(
                slide_num=2,
                name=f"STATIC_S2_CARD_{i:02d}_SUB",
                x_in=card_x + 0.10,
                y_in=fy0 + 0.30,
                w_in=card_w - 0.20,
                h_in=0.14,
                svg_text=_svg_root(
                    w_px=_in_to_px(card_w - 0.20),
                    h_px=_in_to_px(0.14),
                    body=_svg_text(x=0, y=_in_to_px(0.07), text=subtitle, size_px=12, color=f"#{COLORS_HEX['muted']}", weight=700, baseline="middle"),
                ),
                static=True,
            )
        )

    eq_d, w_d, h_d = _mpl_math_svg(r"$D(p,p')=\left|E(p)-E(p')\right|$", color_hex=COLORS_HEX["amber"], fontsize_pt=20.0)
    w_d2, h_d2 = _scale_to_fit(w_in=w_d, h_in=h_d, max_w_in=card_w - 0.22, max_h_in=0.42)
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_D", x_in=fx + 0.11, y_in=fy0 + 0.56, w_in=w_d2, h_in=h_d2, svg_text=eq_d))

    eq_l, w_l, h_l = _mpl_math_svg(r"$d\approx 0\Rightarrow |E-E'|\leq \varepsilon$", color_hex=COLORS_HEX["ink"], fontsize_pt=18.0)
    w_l2, h_l2 = _scale_to_fit(w_in=w_l, h_in=h_l, max_w_in=card_w - 0.22, max_h_in=0.42)
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_LIP", x_in=fx + (card_w + gapx) + 0.11, y_in=fy0 + 0.56, w_in=w_l2, h_in=h_l2, svg_text=eq_l))

    eq_de, w_de, h_de = _mpl_math_svg(
        r"$\Delta E(p)=v(\mathrm{rank}_{q_1})-v(\mathrm{rank}_{q_2})$",
        color_hex=COLORS_HEX["muted"],
        fontsize_pt=16.5,
    )
    w_de2, h_de2 = _scale_to_fit(w_in=w_de, h_in=h_de, max_w_in=card_w - 0.22, max_h_in=0.42)
    tiles.append(SvgTile(slide_num=2, name="S2_EQ_DELTA", x_in=fx + 2 * (card_w + gapx) + 0.11, y_in=fy0 + 0.56, w_in=w_de2, h_in=h_de2, svg_text=eq_de))

    heat_y = fy0 + card_h + 0.22
    heat_n = 14
    cell_in = 0.15
    gap_in = 0.015
    heat_w = heat_n * cell_in + (heat_n - 1) * gap_in
    heat_h = heat_w

    tiles.append(
        SvgTile(
            slide_num=2,
            name="STATIC_S2_PILL_HEAT",
            x_in=fx,
            y_in=heat_y - 0.30,
            w_in=heat_w,
            h_in=0.22,
            svg_text=_svg_pill(w_px=_in_to_px(heat_w), h_px=_in_to_px(0.22), text="Within‑cluster gap heatmap (illustration)", accent_hex=COLORS_HEX["coral"]),
            static=True,
        )
    )

    for ry in range(heat_n):
        for cx in range(heat_n):
            heat = ((cx + 1) * (ry + 2)) % 13 / 13.0
            tiles.append(
                SvgTile(
                    slide_num=2,
                    name=f"S2_HEAT_{ry:02d}_{cx:02d}",
                    x_in=fx + cx * (cell_in + gap_in),
                    y_in=heat_y + ry * (cell_in + gap_in),
                    w_in=cell_in,
                    h_in=cell_in,
                    svg_text=_svg_cell(w_px=_in_to_px(cell_in), h_px=_in_to_px(cell_in), val=None, heat=heat, accent_hex=COLORS_HEX["coral"]),
                )
            )

    # Heatmap overlay frame (drawn on top; no text).
    hw_px = _in_to_px(heat_w)
    hh_px = _in_to_px(heat_h)
    heat_overlay = _svg_root(
        w_px=hw_px,
        h_px=hh_px,
        body=(
            f'<rect x="0" y="0" width="{hw_px}" height="{hh_px}" rx="10" ry="10" fill="none" '
            f'stroke="#{COLORS_HEX["coral"]}" stroke-width="2" opacity="0.85"/>'
            + f'<line x1="0" y1="0" x2="{hw_px}" y2="{hh_px}" stroke="#{COLORS_HEX["card_line"]}" stroke-width="2" opacity="0.30"/>'
            + f'<circle cx="0" cy="0" r="3" fill="#{COLORS_HEX["coral"]}" opacity="0.45"/>'
            + f'<circle cx="{hw_px}" cy="{hh_px}" r="3" fill="#{COLORS_HEX["coral"]}" opacity="0.45"/>'
        ),
    )
    tiles.append(SvgTile(slide_num=2, name="_BG_S2_HEAT_OVERLAY", x_in=fx, y_in=heat_y, w_in=heat_w, h_in=heat_h, svg_text=heat_overlay, static=True))

    timeline_x = fx + heat_w + 0.18
    timeline_w = fw - heat_w - 0.18
    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_TIMELINE",
            x_in=timeline_x,
            y_in=heat_y,
            w_in=timeline_w,
            h_in=heat_h,
            svg_text=_svg_timeline_card(w_px=_in_to_px(timeline_w), h_px=_in_to_px(heat_h), accent_hex=COLORS_HEX["violet"]),
        )
    )

    ops_y = heat_y + heat_h + 0.18
    ops_h = 0.94
    tiles.append(
        SvgTile(
            slide_num=2,
            name="S2_OPS",
            x_in=fx,
            y_in=ops_y,
            w_in=fw,
            h_in=ops_h,
            svg_text=_svg_caption(
                w_px=_in_to_px(fw),
                h_px=_in_to_px(ops_h),
                lines=[
                    "Operationalize in your tables: cluster_id from posts_first_seen.text (exact/near‑dup),",
                    "E(post|q)=v(rank) from feed_items, then summarize gaps by author + context.",
                ],
                accent_hex=COLORS_HEX["coral"],
            ),
        )
    )

    return tiles


def _iter_slide_xml_paths(pptx: zipfile.ZipFile) -> Iterable[tuple[int, str]]:
    for name in pptx.namelist():
        m = re.match(r"^ppt/slides/slide(\d+)\.xml$", name)
        if not m:
            continue
        yield int(m.group(1)), name


def _sanitize_media_name(raw: str) -> str:
    s = re.sub(r"[^A-Za-z0-9_\\-]+", "_", raw.strip())
    s = re.sub(r"_+", "_", s).strip("_")
    return (s or "asset")[:64]


def _patch_pptx_replace_images_with_svg(*, pptx_in: Path, pptx_out: Path, shape_svg_map: dict[str, Path]) -> None:
    """Rewrite slide image rel targets to SVG media parts and add SVG files to the PPTX."""

    ns_slide = {"p": PML_NS, "a": A_NS, "r": R_NS}
    rel_tag = f"{{{PKG_REL_NS}}}Relationship"

    svg_media: dict[str, bytes] = {}
    svg_default_needed = True

    with zipfile.ZipFile(pptx_in, "r") as zin:
        ct_root = etree.fromstring(zin.read("[Content_Types].xml"))
        for el in ct_root.findall(f"{{{CT_NS}}}Default"):
            if el.get("Extension") == "svg":
                svg_default_needed = False
                break

        slide_xml_by_num: dict[int, etree._Element] = {}
        slide_rels_by_num: dict[int, etree._Element] = {}

        for slide_num, slide_path in _iter_slide_xml_paths(zin):
            slide_xml_by_num[slide_num] = etree.fromstring(zin.read(slide_path))
            rels_path = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            if rels_path in zin.namelist():
                slide_rels_by_num[slide_num] = etree.fromstring(zin.read(rels_path))

        for slide_num, slide_root in slide_xml_by_num.items():
            rels_root = slide_rels_by_num.get(slide_num)
            if rels_root is None:
                continue
            rid_to_rel = {rel.get("Id"): rel for rel in rels_root.findall(rel_tag)}

            for pic in slide_root.xpath(".//p:pic", namespaces=ns_slide):
                c_nv = pic.find(".//p:cNvPr", namespaces=ns_slide)
                if c_nv is None:
                    continue
                shape_name = (c_nv.get("name") or "").strip()
                if not shape_name or shape_name not in shape_svg_map:
                    continue

                blip = pic.find(".//a:blip", namespaces=ns_slide)
                if blip is None:
                    continue
                rid = blip.get(f"{{{R_NS}}}embed")
                if not rid:
                    continue
                rel = rid_to_rel.get(rid)
                if rel is None:
                    continue

                svg_src = shape_svg_map[shape_name]
                base = _sanitize_media_name(shape_name)
                media_name = f"{base}.svg"
                if media_name in svg_media:
                    media_name = f"{base}_{hashlib.sha1(shape_name.encode('utf-8')).hexdigest()[:8]}.svg"
                rel.set("Target", f"../media/{media_name}")
                svg_media[media_name] = svg_src.read_bytes()

        with zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)

                if info.filename == "[Content_Types].xml":
                    if svg_default_needed:
                        ct_root = etree.fromstring(data)
                        etree.SubElement(ct_root, f"{{{CT_NS}}}Default", Extension="svg", ContentType="image/svg+xml")
                        data = etree.tostring(ct_root, xml_declaration=True, encoding="UTF-8", standalone=True)
                    zout.writestr(info, data)
                    continue

                m = re.match(r"^ppt/slides/slide(\d+)\.xml$", info.filename)
                if m:
                    slide_num = int(m.group(1))
                    root = slide_xml_by_num.get(slide_num)
                    if root is not None:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                    zout.writestr(info, data)
                    continue

                m = re.match(r"^ppt/slides/_rels/slide(\d+)\.xml\.rels$", info.filename)
                if m:
                    slide_num = int(m.group(1))
                    root = slide_rels_by_num.get(slide_num)
                    if root is not None:
                        data = etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True)
                    zout.writestr(info, data)
                    continue

                zout.writestr(info, data)

            for media_name, blob in svg_media.items():
                zout.writestr(f"ppt/media/{media_name}", blob)


def _compute_exclude_spids_for_static_by_slide(pptx_path: Path) -> dict[int, set[int]]:
    prs = Presentation(str(pptx_path))
    exclude_by_slide: dict[int, set[int]] = {}
    for slide_num, slide in enumerate(prs.slides, start=1):
        exclude: set[int] = set()
        for sh in slide.shapes:
            name = (getattr(sh, "name", "") or "").strip()
            if not name:
                continue
            if name.startswith("_BG") or name.startswith("STATIC_"):
                exclude.add(int(sh.shape_id))
        exclude_by_slide[slide_num] = exclude
    return exclude_by_slide


def build_deck(*, out_pre: Path, out_svg: Path, out_animated: Path, effect_dur_ms: int, keep_intermediates: bool) -> None:
    out_pre = out_pre.resolve()
    out_svg = out_svg.resolve()
    out_animated = out_animated.resolve()
    out_pre.parent.mkdir(parents=True, exist_ok=True)
    out_svg.parent.mkdir(parents=True, exist_ok=True)
    out_animated.parent.mkdir(parents=True, exist_ok=True)

    tiles = _build_slide1() + _build_slide2()

    # Hard requirements: no off-slide placements + no duplicate names (name is the SVG lookup key).
    _assert_tiles_within_slide_bounds(tiles)
    seen_names: set[str] = set()
    for t in tiles:
        if t.name in seen_names:
            raise ValueError(f"Duplicate tile name (must be unique across deck): {t.name}")
        seen_names.add(t.name)

    tmp_dir = out_pre.parent / "_ifx_svg_tmp"
    if tmp_dir.exists():
        import shutil

        shutil.rmtree(tmp_dir, ignore_errors=True)
    svg_dir = tmp_dir / "svgs"
    png_dir = tmp_dir / "pngs"
    svg_dir.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)

    shape_to_svg: dict[str, Path] = {}
    shape_to_png: dict[str, Path] = {}
    for t in tiles:
        svg_path = svg_dir / f"{t.name}_s{t.slide_num:02d}.svg"
        svg_path.write_text(t.svg_text, encoding="utf-8")
        shape_to_svg[t.name] = svg_path

        png_path = png_dir / f"{t.name}_s{t.slide_num:02d}.png"
        _write_unique_png(out_path=png_path, seed=f"{t.slide_num}:{t.name}")
        shape_to_png[t.name] = png_path

    # Build PPTX using placeholder PNG picture parts (one unique blob per shape), then patch to SVG.
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    s2 = prs.slides.add_slide(blank)

    for slide in [s1, s2]:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(COLORS_HEX["bg"])

    for t in tiles:
        slide = s1 if t.slide_num == 1 else s2
        pic = slide.shapes.add_picture(
            str(shape_to_png[t.name]),
            Inches(t.x_in),
            Inches(t.y_in),
            width=Inches(t.w_in),
            height=Inches(t.h_in),
        )
        pic.name = t.name

    prs.save(str(out_pre))
    _patch_pptx_replace_images_with_svg(pptx_in=out_pre, pptx_out=out_svg, shape_svg_map=shape_to_svg)

    # Inject animations.
    repo_root = Path(__file__).resolve().parent
    slide2_dir = repo_root / "Slide2"
    sys.path.insert(0, str(slide2_dir))
    from pptx_click_animations import inject_click_reveals  # noqa: PLC0415

    exclude_by_slide = _compute_exclude_spids_for_static_by_slide(out_svg)
    # On Windows, PowerPoint may hold an exclusive lock on the last output file. Fall back to a new name.
    animated_target = out_animated
    for attempt in range(0, 25):
        try:
            tmp_target = animated_target.with_name(f"{animated_target.stem}__tmp{animated_target.suffix}")
            res1 = inject_click_reveals(
                pptx_in=out_svg,
                pptx_out=tmp_target,
                slide_nums={1},
                exclude_spids=exclude_by_slide.get(1, set()),
                effect_dur_ms=effect_dur_ms,
            )
            res2 = inject_click_reveals(
                pptx_in=tmp_target,
                pptx_out=animated_target,
                slide_nums={2},
                exclude_spids=exclude_by_slide.get(2, set()),
                effect_dur_ms=effect_dur_ms,
            )
            click_effects: dict[int, int] = {}
            click_effects.update(res1.slide_click_effects)
            click_effects.update(res2.slide_click_effects)
            try:
                tmp_target.unlink()
            except OSError:
                pass
            break
        except PermissionError:
            try:
                tmp_target.unlink()
            except OSError:
                pass
            animated_target = out_animated.with_name(f"{out_animated.stem}_{attempt+1:02d}{out_animated.suffix}")
    else:
        raise

    print(f"OK: wrote {res2.pptx_out}")
    for s, n in sorted(click_effects.items()):
        print(f"OK: slide {s} clickEffects={n}")
    print(f"OK: total clickEffects={sum(click_effects.values())}")

    if not keep_intermediates:
        for p in [out_pre, out_svg]:
            try:
                p.unlink()
            except OSError:
                pass


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--out-pre", type=Path, default=Path("_build/ifx_poster_preanim.pptx"))
    parser.add_argument("--out-svg", type=Path, default=Path("_build/ifx_poster_svg.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/ifx_poster_animated.pptx"))
    parser.add_argument("--dur-ms", type=int, default=260)
    parser.add_argument("--keep-intermediates", action="store_true", help="Keep preanim/svg PPTX artifacts in _build/.")
    args = parser.parse_args()

    build_deck(out_pre=args.out_pre, out_svg=args.out_svg, out_animated=args.out, effect_dur_ms=int(args.dur_ms), keep_intermediates=bool(args.keep_intermediates))


if __name__ == "__main__":
    main()
