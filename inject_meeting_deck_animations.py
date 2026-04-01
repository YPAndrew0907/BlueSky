#!/usr/bin/env python3
from __future__ import annotations

import argparse
import sys
import tempfile
from pathlib import Path

from pptx import Presentation


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


def _compute_exclude_spids_for_slide6(pptx_path: Path) -> set[int]:
    prs = Presentation(str(pptx_path))
    if len(prs.slides) < 6:
        raise ValueError("Deck must have at least 6 slides")
    slide6 = prs.slides[5]

    exclude: set[int] = set()
    for sh in _iter_shapes(slide6):
        name = (getattr(sh, "name", "") or "").strip()
        if not name:
            continue
        if name == "TITLE" or name == "TAKEAWAY":
            exclude.add(int(sh.shape_id))
            continue
        if name.startswith("_BG"):
            exclude.add(int(sh.shape_id))
            continue
        # Optional equation mode (not used unless assets are provided): keep EQ_STEP_* static by default.
        if name.startswith("EQ_") or name.startswith("EQ_STEP_"):
            exclude.add(int(sh.shape_id))
            continue

    return exclude


def main() -> None:
    parser = argparse.ArgumentParser(description="Inject meeting-deck animations (structured slides 1–5; legacy slide 6).")
    parser.add_argument("pptx_in", type=Path, nargs="?", default=Path("_build/bluesky_meeting_rqs_papers_data_6slides_strict_preanim.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/bluesky_meeting_rqs_papers_data_6slides_strict_animated.pptx"))
    args = parser.parse_args()

    repo_root = Path(__file__).resolve().parent
    slide2_dir = repo_root / "Slide2"
    sys.path.insert(0, str(slide2_dir))
    from pptx_click_animations import inject_click_reveals, inject_structured_reveals  # noqa: PLC0415

    pptx_in = args.pptx_in.resolve()
    if not pptx_in.exists():
        raise SystemExit(f"Missing input PPTX: {pptx_in}")

    out = args.out.resolve()
    out.parent.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory(prefix="meeting_anim_") as tmp_dir:
        tmp1 = Path(tmp_dir) / "step1_structured.pptx"

        res_struct = inject_structured_reveals(
            pptx_in=pptx_in,
            pptx_out=tmp1,
            slide_nums={1, 2, 3, 4, 5},
            exclude_spids=set(),
            effect_dur_ms=220,
        )

        exclude6 = _compute_exclude_spids_for_slide6(tmp1)
        res_legacy6 = inject_click_reveals(
            pptx_in=tmp1,
            pptx_out=out,
            slide_nums={6},
            exclude_spids=exclude6,
            effect_dur_ms=180,
        )

    print(f"OK: wrote {out}")
    for s in range(1, 6):
        if s in res_struct.slide_click_effects:
            print(f"OK: slide {s} clickEffects={res_struct.slide_click_effects[s]}")
    if 6 in res_legacy6.slide_click_effects:
        print(f"OK: slide 6 clickEffects={res_legacy6.slide_click_effects[6]}")


if __name__ == "__main__":
    main()

