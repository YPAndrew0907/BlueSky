#!/usr/bin/env python3
from __future__ import annotations

import argparse
import os
import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# Canvas (widescreen)
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Margins (keep everything inside)
SAFE_L_IN = 0.65
SAFE_R_IN = 0.65
SAFE_T_IN = 0.45
SAFE_B_IN = 0.45

SAFE_W_IN = SLIDE_W_IN - SAFE_L_IN - SAFE_R_IN
SAFE_H_IN = SLIDE_H_IN - SAFE_T_IN - SAFE_B_IN


COLORS_HEX: dict[str, str] = {
    "bg": "0B1320",
    "ink": "F7F9FC",
    "muted": "C9D2E3",
    "card": "162033",
    "card_line": "3A4D66",
    "cyan": "58B7E6",
    "amber": "F4C55D",
    "coral": "FF7D77",
    "green": "54C687",
}

FONT_UI = "Calibri"
FONT_MONO = "Consolas"


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(COLORS_HEX["bg"])


def _style_text_frame(tf, *, ml: float = 0.0, mr: float = 0.0, mt: float = 0.0, mb: float = 0.0) -> None:
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)


def _set_text(
    shape,
    *,
    text: str,
    size_pt: float,
    bold: bool = False,
    italic: bool = False,
    color_hex: str = "F7F9FC",
    font_name: str = FONT_UI,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    _style_text_frame(tf, ml=0.0, mr=0.0, mt=0.0, mb=0.0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size_pt)
    font.bold = bold
    font.italic = italic
    font.color.rgb = _rgb(color_hex)


@dataclass(frozen=True)
class MonoMetrics:
    char_w_in: float
    line_h_in: float


def _mono_metrics(*, font_pt: float, font_name: str = FONT_MONO) -> MonoMetrics:
    # Use a real TTF on Windows so widths are stable.
    # Pillow uses pixel sizes; approximate points at 96dpi.
    windir = Path(os.environ.get("WINDIR", "C:/Windows"))
    if font_name.lower() == "consolas":
        ttf = windir / "Fonts" / "consola.ttf"
    else:
        # Fall back: hope font exists; if not, Pillow raises.
        ttf = windir / "Fonts" / f"{font_name}.ttf"

    px = int(round(font_pt * 96 / 72))
    font = ImageFont.truetype(str(ttf), px)
    # Consolas is monospaced; measure a single glyph width.
    char_px = float(font.getlength("M"))
    char_w_in = char_px / 96.0

    # Line height: a bit of breathing room for "chalk" look.
    line_h_in = (font_pt * 1.30) / 72.0
    return MonoMetrics(char_w_in=char_w_in, line_h_in=line_h_in)


def _add_mono_reveal_line(
    slide,
    *,
    x: float,
    y: float,
    text: str,
    font_pt: float,
    color_hex: str,
    name_prefix: str,
    bold: bool = False,
) -> tuple[int, float]:
    """Adds one textbox per non-space character (spaces advance cursor). Returns (#shapes, y_next)."""
    mm = _mono_metrics(font_pt=font_pt)
    cx = x
    count = 0
    for idx, ch in enumerate(text):
        if ch == "\n":
            cx = x
            y += mm.line_h_in
            continue
        if ch == " ":
            cx += mm.char_w_in
            continue

        box = slide.shapes.add_textbox(Inches(cx), Inches(y), Inches(mm.char_w_in), Inches(mm.line_h_in))
        box.name = f"{name_prefix}_{idx:04d}"
        _set_text(
            box,
            text=ch,
            size_pt=font_pt,
            bold=bold,
            color_hex=color_hex,
            font_name=FONT_MONO,
            align=PP_ALIGN.LEFT,
        )
        cx += mm.char_w_in
        count += 1
    return count, (y + mm.line_h_in)


def _add_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    title: str,
    name: str,
    fill_hex: str = COLORS_HEX["card"],
    line_hex: str = COLORS_HEX["card_line"],
    title_hex: str = COLORS_HEX["ink"],
) -> object:
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.name = name
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb(fill_hex)
    card.line.color.rgb = _rgb(line_hex)
    card.line.width = Pt(1.5)

    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.14), Inches(w - 0.44), Inches(0.45))
    title_box.name = f"{name}_TITLE"
    _set_text(title_box, text=title, size_pt=18, bold=True, color_hex=title_hex, font_name=FONT_UI)
    return card


def _add_matrix_p_diagram(
    slide,
    *,
    x: float,
    y: float,
    size: float,
    n: int,
    name_prefix: str,
) -> None:
    cell = size / n

    # Outer frame
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(size), Inches(size))
    frame.name = f"{name_prefix}_FRAME"
    frame.fill.background()
    frame.line.color.rgb = _rgb(COLORS_HEX["muted"])
    frame.line.width = Pt(1.25)

    # Grid cells (each cell is its own shape for click granularity)
    for r in range(n):
        for c in range(n):
            sx = x + c * cell
            sy = y + r * cell
            sq = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(sx), Inches(sy), Inches(cell), Inches(cell))
            sq.name = f"{name_prefix}_CELL_{r:02d}_{c:02d}"
            sq.fill.solid()
            sq.fill.fore_color.rgb = _rgb(COLORS_HEX["card"])
            sq.line.color.rgb = _rgb(COLORS_HEX["card_line"])
            sq.line.width = Pt(0.8)

    # Highlight a single row i (row 2)
    i = 2
    hi = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y + i * cell), Inches(size), Inches(cell))
    hi.name = f"{name_prefix}_ROW_I"
    hi.fill.solid()
    hi.fill.fore_color.rgb = _rgb("0F2A3D")
    hi.fill.transparency = 0.25
    hi.line.color.rgb = _rgb(COLORS_HEX["cyan"])
    hi.line.width = Pt(1.2)

    # Put a simple probability mass on that row (illustrative)
    probs = [0.60, 0.25, 0.10, 0.05, 0.00, 0.00][:n]
    for c, p in enumerate(probs):
        if p <= 0:
            continue
        tx = x + c * cell
        ty = y + i * cell
        t = slide.shapes.add_textbox(Inches(tx), Inches(ty + cell * 0.16), Inches(cell), Inches(cell * 0.70))
        t.name = f"{name_prefix}_PVAL_{c:02d}"
        _set_text(t, text=f"{p:.2f}", size_pt=12, bold=True, color_hex=COLORS_HEX["amber"], font_name=FONT_UI, align=PP_ALIGN.CENTER)

    # Labels
    lab = slide.shapes.add_textbox(Inches(x), Inches(y - 0.42), Inches(size), Inches(0.40))
    lab.name = f"{name_prefix}_LABEL"
    _set_text(lab, text="P  (Pr[item i at rank j])", size_pt=16, bold=True, color_hex=COLORS_HEX["muted"], font_name=FONT_UI)


def _add_v_vector_diagram(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    n: int,
    name_prefix: str,
) -> None:
    # Decreasing bars for position bias v_j.
    bar_h = h / n
    for j in range(n):
        frac = max(0.12, 1.0 - 0.13 * j)
        bw = w * frac
        by = y + j * bar_h
        bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(by), Inches(bw), Inches(bar_h * 0.75))
        bar.name = f"{name_prefix}_BAR_{j+1:02d}"
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(COLORS_HEX["green"])
        bar.fill.transparency = 0.15
        bar.line.color.rgb = _rgb(COLORS_HEX["green"])
        bar.line.width = Pt(1.0)

        t = slide.shapes.add_textbox(Inches(x + w + 0.05), Inches(by - 0.02), Inches(0.60), Inches(bar_h))
        t.name = f"{name_prefix}_LBL_{j+1:02d}"
        _set_text(t, text=f"v{j+1}", size_pt=12, bold=True, color_hex=COLORS_HEX["muted"], font_name=FONT_UI)

    lab = slide.shapes.add_textbox(Inches(x), Inches(y - 0.42), Inches(w + 0.9), Inches(0.40))
    lab.name = f"{name_prefix}_LABEL"
    _set_text(lab, text="v  (position bias / exposure weight)", size_pt=16, bold=True, color_hex=COLORS_HEX["muted"], font_name=FONT_UI)


def _add_arrow(slide, *, x: float, y: float, w: float, h: float, color_hex: str, name: str) -> None:
    arr = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    arr.name = name
    arr.fill.solid()
    arr.fill.fore_color.rgb = _rgb(color_hex)
    arr.line.color.rgb = _rgb(color_hex)
    arr.line.width = Pt(1.0)


def _add_connector(slide, *, x1: float, y1: float, x2: float, y2: float, color_hex: str, name: str) -> None:
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.name = name
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(2.0)


def _slide1(slide) -> int:
    _set_slide_bg(slide)
    shapes = 0

    # Title
    title = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(0.25), Inches(SAFE_W_IN), Inches(0.75))
    title.name = "TITLE"
    _set_text(
        title,
        text="Exposure in rankings = expected attention (Singh & Joachims, KDD’18)",
        size_pt=36,
        bold=True,
        color_hex=COLORS_HEX["ink"],
        font_name=FONT_UI,
    )
    shapes += 1

    subtitle = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(1.02), Inches(SAFE_W_IN), Inches(0.45))
    subtitle.name = "TAKEAWAY"
    _set_text(
        subtitle,
        text="Core object: Exposure(item i) = Σ_j P_{i,j} · v_j  (then map P to Bluesky snapshot rankings).",
        size_pt=18,
        bold=False,
        color_hex=COLORS_HEX["muted"],
        font_name=FONT_UI,
    )
    shapes += 1

    # Left visuals: P matrix and v vector
    left_x = SAFE_L_IN
    top_y = 1.65
    _add_matrix_p_diagram(slide, x=left_x, y=top_y, size=3.65, n=6, name_prefix="P")
    shapes += 1 + (6 * 6) + 1 + 4 + 1  # frame + cells + row highlight + up to 4 pvals + label (rough)

    _add_v_vector_diagram(slide, x=left_x + 4.10, y=top_y + 0.25, w=1.85, h=3.05, n=6, name_prefix="V")
    shapes += (6 * 2) + 1  # bars + labels + label

    # Right side: blackboard-style derivation (monospace, per-char reveals)
    # Keep the chalk region narrow enough to sit beside the diagram, so lines must be short/wrapped.
    chalk_x = SAFE_L_IN + 7.05
    chalk_y = 1.70
    label = slide.shapes.add_textbox(Inches(chalk_x), Inches(chalk_y - 0.42), Inches(4.60), Inches(0.35))
    label.name = "CHALK_LABEL"
    _set_text(label, text="Write it like a blackboard:", size_pt=16, bold=True, color_hex=COLORS_HEX["coral"], font_name=FONT_UI)
    shapes += 1

    line_pt = 18
    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y,
        text="q = (feed_uri, snapshot_hour_utc,\nviewer_mode, vantage_id)",
        font_pt=line_pt,
        color_hex=COLORS_HEX["ink"],
        name_prefix="EQ1_Q",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.10,
        text="P_{i,j} = Pr[item i at rank j]\nP is doubly stochastic",
        font_pt=line_pt,
        color_hex=COLORS_HEX["ink"],
        name_prefix="EQ2_P",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.10,
        text="Exposure(i|P) = Σ_j P_{i,j} v_j\n          = (P·v)_i",
        font_pt=20,
        color_hex=COLORS_HEX["amber"],
        name_prefix="EQ3_E",
        bold=True,
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.10,
        text="U(P|q) = u^T P v",
        font_pt=18,
        color_hex=COLORS_HEX["muted"],
        name_prefix="EQ4_U",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.12,
        text="Snapshot: E(post|q)=v(rank)\nAbsent ⇒ E=0  (top-K truncation)",
        font_pt=18,
        color_hex=COLORS_HEX["ink"],
        name_prefix="EQ5_MAP",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.06,
        text="Across t: P̂_{p,j}=freq(rank=j)\nÊ(p)=Σ_j P̂_{p,j} v_j",
        font_pt=18,
        color_hex=COLORS_HEX["ink"],
        name_prefix="EQ6_PHAT",
    )
    shapes += c

    foot = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(SLIDE_H_IN - SAFE_B_IN - 0.30), Inches(SAFE_W_IN), Inches(0.30))
    foot.name = "FOOT"
    _set_text(
        foot,
        text="Notation: P from Singh & Joachims (2018). v_j = position bias / examination probability at rank j.",
        size_pt=12,
        bold=False,
        color_hex=COLORS_HEX["muted"],
        font_name=FONT_UI,
        align=PP_ALIGN.LEFT,
    )
    shapes += 1

    return shapes


def _slide2(slide) -> int:
    _set_slide_bg(slide)
    shapes = 0

    # Title
    title = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(0.25), Inches(SAFE_W_IN), Inches(0.75))
    title.name = "TITLE"
    _set_text(
        title,
        text="Individual fairness: similar posts should get similar exposure (Zehlike et al., CSUR’22 Part II)",
        size_pt=32,
        bold=True,
        color_hex=COLORS_HEX["ink"],
        font_name=FONT_UI,
    )
    shapes += 1

    subtitle = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(1.02), Inches(SAFE_W_IN), Inches(0.45))
    subtitle.name = "TAKEAWAY"
    _set_text(
        subtitle,
        text="Make “same content → different exposure” audit-ready by defining a similarity metric over posts.",
        size_pt=18,
        bold=False,
        color_hex=COLORS_HEX["muted"],
        font_name=FONT_UI,
    )
    shapes += 1

    # Left: content cluster diagram (vector shapes)
    cluster_x = SAFE_L_IN
    cluster_y = 1.70
    _add_card(slide, x=cluster_x, y=cluster_y, w=5.90, h=4.90, title="Content cluster g (near-duplicate text)", name="CLUSTER_CARD")
    shapes += 2  # card + title

    # Three post cards inside
    posts = [
        ("author_did = A", "text ≈ “hello world …”", COLORS_HEX["cyan"]),
        ("author_did = B", "text ≈ “hello world …”", COLORS_HEX["amber"]),
        ("author_did = C", "text ≈ “hello world …”", COLORS_HEX["green"]),
    ]
    py = cluster_y + 0.80
    for i, (a, t, accent) in enumerate(posts, start=1):
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(cluster_x + 0.35), Inches(py), Inches(5.20), Inches(1.10))
        box.name = f"POST_{i:02d}"
        box.fill.solid()
        box.fill.fore_color.rgb = _rgb("0F1B2E")
        box.line.color.rgb = _rgb(accent)
        box.line.width = Pt(1.5)
        shapes += 1

        at = slide.shapes.add_textbox(Inches(cluster_x + 0.55), Inches(py + 0.15), Inches(5.00), Inches(0.35))
        at.name = f"POST_{i:02d}_AUTH"
        _set_text(at, text=a, size_pt=16, bold=True, color_hex=COLORS_HEX["ink"], font_name=FONT_UI)
        shapes += 1

        tt = slide.shapes.add_textbox(Inches(cluster_x + 0.55), Inches(py + 0.55), Inches(5.00), Inches(0.40))
        tt.name = f"POST_{i:02d}_TEXT"
        _set_text(tt, text=t, size_pt=15, bold=False, color_hex=COLORS_HEX["muted"], font_name=FONT_UI)
        shapes += 1

        py += 1.32

    # Arrow to ranking/exposure box
    diff = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(cluster_x + 4.05), Inches(cluster_y + 4.05), Inches(1.55), Inches(0.85))
    diff.name = "DIFF_BOX"
    diff.fill.solid()
    diff.fill.fore_color.rgb = _rgb("112033")
    diff.line.color.rgb = _rgb(COLORS_HEX["coral"])
    diff.line.width = Pt(1.5)
    shapes += 1

    diff_t = slide.shapes.add_textbox(Inches(cluster_x + 4.12), Inches(cluster_y + 4.17), Inches(1.40), Inches(0.65))
    diff_t.name = "DIFF_TEXT"
    _set_text(diff_t, text="≠ exposure", size_pt=16, bold=True, color_hex=COLORS_HEX["ink"], font_name=FONT_UI, align=PP_ALIGN.CENTER)
    shapes += 1

    # Right: equations / audit metrics (monospace, per-char reveals)
    chalk_x = cluster_x + 6.25
    chalk_y = 1.70
    label = slide.shapes.add_textbox(Inches(chalk_x), Inches(chalk_y - 0.42), Inches(4.60), Inches(0.35))
    label.name = "CHALK_LABEL"
    _set_text(label, text="Individual-fairness math:", size_pt=16, bold=True, color_hex=COLORS_HEX["coral"], font_name=FONT_UI)
    shapes += 1

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y,
        text="Define similarity d(p,p'):\ntext + time ⇒ d = α·d_text\n          + (1−α)·d_time",
        font_pt=18,
        color_hex=COLORS_HEX["ink"],
        name_prefix="IF1_D",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.10,
        text="Zehlike (exposure-based):\nD(p,p') = |Exposure(p)\n        − Exposure(p')|",
        font_pt=20,
        color_hex=COLORS_HEX["amber"],
        name_prefix="IF2_DISC",
        bold=True,
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.10,
        text="Individual fairness:  d(p,p')≈0\n⇒  |E(p) − E(p')| ≤ ε",
        font_pt=18,
        color_hex=COLORS_HEX["muted"],
        name_prefix="IF3_LIP",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.12,
        text="Amortized attention (over snapshots):\nA(p)=Σ_t v(rank_t(p))\nR(p)=Σ_t u(p|q_t)",
        font_pt=18,
        color_hex=COLORS_HEX["ink"],
        name_prefix="IF4_AMORT",
    )
    shapes += c

    c, chalk_y = _add_mono_reveal_line(
        slide,
        x=chalk_x,
        y=chalk_y + 0.06,
        text="Context gap (same post_uri):\nΔE(p; q1,q2)=v(rank_q1)−v(rank_q2)",
        font_pt=18,
        color_hex=COLORS_HEX["ink"],
        name_prefix="IF5_DELTA",
    )
    shapes += c

    # Quick operationalization bullets (as UI text; still revealable as shapes)
    ops = slide.shapes.add_textbox(Inches(chalk_x), Inches(5.62), Inches(5.05), Inches(1.00))
    ops.name = "OPS_BOX"
    _set_text(
        ops,
        text="Operationalize with your tables:\n• build cluster_id from posts_first_seen.text\n• E(post|q)=v(rank) using feed_items(rank, post_uri, author_did)\n• compare within-cluster across authors + across viewer_mode/vantage_id",
        size_pt=15,
        bold=False,
        color_hex=COLORS_HEX["ink"],
        font_name=FONT_UI,
    )
    # This one textbox counts as 1 shape; but it will still click-reveal as a unit.
    shapes += 1

    foot = slide.shapes.add_textbox(Inches(SAFE_L_IN), Inches(SLIDE_H_IN - SAFE_B_IN - 0.30), Inches(SAFE_W_IN), Inches(0.30))
    foot.name = "FOOT"
    _set_text(
        foot,
        text="Exposure definitions: Singh & Joachims (2018). Individual-fairness framing + amortized attention: Zehlike et al. (2022) + Biega et al.",
        size_pt=12,
        bold=False,
        color_hex=COLORS_HEX["muted"],
        font_name=FONT_UI,
    )
    shapes += 1

    return shapes


def build_deck(*, out_pptx: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank = prs.slide_layouts[6]
    s1 = prs.slides.add_slide(blank)
    _slide1(s1)
    s2 = prs.slides.add_slide(blank)
    _slide2(s2)

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def _inject_click_reveals(*, pptx_in: Path, pptx_out: Path) -> None:
    repo_root = Path(__file__).resolve().parent
    slide2_dir = repo_root / "Slide2"
    sys.path.insert(0, str(slide2_dir))
    from pptx_click_animations import inject_click_reveals  # noqa: PLC0415

    res = inject_click_reveals(
        pptx_in=pptx_in,
        pptx_out=pptx_out,
        slide_nums={1, 2},
        exclude_spids=set(),
        effect_dur_ms=140,
    )
    print(f"OK: wrote {pptx_out}")
    for s, n in sorted(res.slide_click_effects.items()):
        print(f"OK: slide {s} clickEffects={n}")
    print(f"OK: total clickEffects={res.total_click_effects}")


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a 2-slide individual-fairness exposure deck (with optional click animations).")
    parser.add_argument("--out-pre", type=Path, default=Path("_build/individual_fairness_exposure_2slides_preanim.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/individual_fairness_exposure_2slides_animated.pptx"))
    parser.add_argument("--animate", action="store_true", help="Inject per-shape click-to-reveal animations.")
    args = parser.parse_args()

    out_pre = args.out_pre.resolve()
    build_deck(out_pptx=out_pre)
    print(f"OK: built {out_pre}")

    if args.animate:
        _inject_click_reveals(pptx_in=out_pre, pptx_out=args.out.resolve())


if __name__ == "__main__":
    main()
