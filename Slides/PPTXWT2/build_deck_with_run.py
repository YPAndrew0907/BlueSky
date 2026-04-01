#!/usr/bin/env python3
from __future__ import annotations

import csv
import gzip
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Iterator, Sequence

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


RUN_DIR = Path("out/bsky_fair_run_20260201T090454Z")
RUN_CSV_DIR = RUN_DIR / "csv"

BASE_DECK = Path("Slides/PPTXWT2/bluesky-data-collection-pipeline-final.pptx")
OUT_DECK = Path("Slides/PPTXWT2/bluesky-data-collection-pipeline-run-20260201.pptx")

ASSET_DIR = Path("/tmp/bsky_pptx_receipts_20260201")


COLORS = {
    "bg": "0B0B0F",
    "card": "14141C",
    "ink": "F5F5FA",
    "muted": "B4B4C8",
    "muted2": "46465A",
    "cyan": "00D4FF",
    "purple": "A855F7",
    "pink": "FF5D9E",
    "amber": "F59E0B",
    "green": "22C55E",
}


FONT_SANS = Path("/System/Library/Fonts/SFNS.ttf")
FONT_MONO = Path("/System/Library/Fonts/SFNSMono.ttf")


RTL_MARKS_RE = re.compile(r"[\u202A-\u202E\u2066-\u2069]")


def _clean_text(s: str) -> str:
    return RTL_MARKS_RE.sub("", s).strip()


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _truncate_middle(s: str, max_len: int) -> str:
    if len(s) <= max_len:
        return s
    if max_len < 6:
        return s[:max_len]
    keep_head = (max_len - 1) // 2
    keep_tail = max_len - 1 - keep_head
    return f"{s[:keep_head]}…{s[-keep_tail:]}"


def _compact_did(did: str) -> str:
    did = _clean_text(did)
    if not did.startswith("did:"):
        return _truncate_middle(did, 22)
    if len(did) <= 20:
        return did
    return f"{did[:12]}…{did[-6:]}"


def _compact_cid(cid: str) -> str:
    cid = _clean_text(cid)
    if len(cid) <= 22:
        return cid
    return f"{cid[:10]}…{cid[-10:]}"


def _compact_at_uri(uri: str) -> str:
    uri = _clean_text(uri)
    # Keep the collection + rkey tail when possible.
    idx = uri.rfind("/app.bsky.")
    if idx != -1:
        return f"…{uri[idx:]}"
    return _truncate_middle(uri, 34)


def _compact_sha256(h: str) -> str:
    h = _clean_text(h)
    if len(h) <= 18:
        return h
    return f"{h[:10]}…{h[-10:]}"


@dataclass(frozen=True)
class RunAnchors:
    run_id: str
    started_at_utc: str
    finished_at_utc: str
    auth_mode: str
    num_feeds_panel: int
    num_feed_items: int


def _read_single_row_csv(path: Path) -> dict[str, str]:
    with path.open("r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        return next(reader)


def load_run_anchors() -> RunAnchors:
    meta = _read_single_row_csv(RUN_CSV_DIR / "run_metadata.csv")
    summ = _read_single_row_csv(RUN_CSV_DIR / "run_summary.csv")
    return RunAnchors(
        run_id=_clean_text(meta["run_id"]),
        started_at_utc=_clean_text(meta["started_at_utc"]),
        finished_at_utc=_clean_text(meta["finished_at_utc"]),
        auth_mode=_clean_text(meta["auth_mode"]),
        num_feeds_panel=int(summ["num_feeds_panel"]),
        num_feed_items=int(summ["num_feed_items"]),
    )


def _iter_csv_rows(path: Path) -> Iterator[dict[str, str]]:
    with path.open("r", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        yield from reader


def _iter_csv_gz_rows(path: Path) -> Iterator[dict[str, str]]:
    with gzip.open(path, "rt", encoding="utf-8") as f:
        reader = csv.DictReader(f)
        yield from reader


def _find_row(path: Path, *, key: str, value: str) -> dict[str, str]:
    for row in _iter_csv_rows(path):
        if row.get(key) == value:
            return row
    raise RuntimeError(f"Row not found in {path.name}: {key}={value!r}")


def _find_row_gz(path: Path, *, predicate) -> dict[str, str]:
    for row in _iter_csv_gz_rows(path):
        if predicate(row):
            return row
    raise RuntimeError(f"Row not found in {path.name} (predicate)")


def pick_named_starterpack_with_inclusion() -> tuple[dict[str, str], dict[str, str]]:
    # Find a safe ASCII-ish named pack, then take its first inclusion.
    name_map: dict[str, dict[str, str]] = {}
    for row in _iter_csv_rows(RUN_CSV_DIR / "starterpacks.csv"):
        name = _clean_text(row.get("name", ""))
        if not name:
            continue
        name_map[row["starterpack_uri"]] = row

    ascii_re = re.compile(r"^[\x20-\x7E]+$")
    for inc in _iter_csv_rows(RUN_CSV_DIR / "starterpack_feeds.csv"):
        sp_uri = inc["starterpack_uri"]
        sp = name_map.get(sp_uri)
        if not sp:
            continue
        name = _clean_text(sp.get("name", ""))
        if not ascii_re.match(name):
            continue
        return sp, inc

    raise RuntimeError("No suitable named starter pack with feed inclusion found.")


def render_kv_card(
    *,
    out_path: Path,
    title: str,
    items: Sequence[tuple[str, str]],
    accent_hex: str,
    width_px: int = 1200,
    height_px: int = 260,
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    pad = 28
    radius = 28
    border = 4

    # Card container
    draw.rounded_rectangle(
        (0, 0, width_px - 1, height_px - 1),
        radius=radius,
        fill=f"#{COLORS['card']}",
        outline=f"#{accent_hex}",
        width=border,
    )

    title_font = ImageFont.truetype(str(FONT_SANS), 28)
    label_font = ImageFont.truetype(str(FONT_SANS), 16)
    value_font = ImageFont.truetype(str(FONT_SANS), 19)
    value_mono = ImageFont.truetype(str(FONT_MONO), 18)

    draw.text((pad, 18), _clean_text(title), fill=f"#{COLORS['ink']}", font=title_font)

    # Grid layout (2–4 columns)
    cols = max(2, min(4, len(items)))
    rows = (len(items) + cols - 1) // cols
    col_w = (width_px - 2 * pad) / cols
    base_y = 78
    row_h = 74

    for idx, (k, v) in enumerate(items):
        r = idx // cols
        c = idx % cols
        x = int(pad + c * col_w)
        y = int(base_y + r * row_h)
        draw.text((x, y), _clean_text(k), fill=f"#{COLORS['muted']}", font=label_font)

        vv = _clean_text(v)
        is_mono = any(t in vv for t in ("did:", "at://", "…", "bafy", "/"))
        vf = value_mono if is_mono else value_font
        draw.text((x, y + 22), vv, fill=f"#{COLORS['ink']}", font=vf)

    # Small caption
    cap = "Example row (truncated)"
    cap_font = ImageFont.truetype(str(FONT_SANS), 13)
    cap_w = draw.textlength(cap, font=cap_font)
    draw.text((width_px - pad - cap_w, height_px - 26), cap, fill=f"#{COLORS['muted2']}", font=cap_font)

    img.save(out_path)


def render_list_card(
    *,
    out_path: Path,
    title: str,
    lines: Sequence[str],
    accent_hex: str,
    width_px: int = 1200,
    height_px: int = 320,
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    pad = 28
    radius = 28
    border = 4

    draw.rounded_rectangle(
        (0, 0, width_px - 1, height_px - 1),
        radius=radius,
        fill=f"#{COLORS['card']}",
        outline=f"#{accent_hex}",
        width=border,
    )

    title_font = ImageFont.truetype(str(FONT_SANS), 28)
    mono = ImageFont.truetype(str(FONT_MONO), 18)
    draw.text((pad, 18), _clean_text(title), fill=f"#{COLORS['ink']}", font=title_font)

    y = 78
    for line in lines[:8]:
        draw.text((pad, y), _clean_text(line), fill=f"#{COLORS['ink']}", font=mono)
        y += 28

    cap = "Excerpt"
    cap_font = ImageFont.truetype(str(FONT_SANS), 13)
    cap_w = draw.textlength(cap, font=cap_font)
    draw.text((width_px - pad - cap_w, height_px - 26), cap, fill=f"#{COLORS['muted2']}", font=cap_font)

    img.save(out_path)


def build_receipts(anchors: RunAnchors) -> dict[str, Path]:
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    receipts: dict[str, Path] = {}

    # Coherent exemplar anchored on the top popular feed "whats-hot".
    popular_row = next(_iter_csv_rows(RUN_CSV_DIR / "popular_feeds.csv"))
    exemplar_feed_uri = popular_row["feed_uri"]

    # Feed generator index (canonical metadata)
    feed_idx = _find_row(RUN_CSV_DIR / "feed_generators_index.csv", key="feed_uri", value=exemplar_feed_uri)
    p_feed_idx = ASSET_DIR / "feed_index_row.png"
    render_kv_card(
        out_path=p_feed_idx,
        title="Feed generator index",
        accent_hex=COLORS["cyan"],
        items=(
            ("Feed", _clean_text(feed_idx.get("display_name", "")) or "(no name)"),
            ("Provider", _clean_text(feed_idx.get("provider_bucket", "")) or "(unknown)"),
            ("URI", _compact_at_uri(exemplar_feed_uri)),
        ),
    )
    receipts["feed_index"] = p_feed_idx

    # Starter pack + inclusion (pick a safe named pack)
    sp, inc = pick_named_starterpack_with_inclusion()
    p_sp = ASSET_DIR / "starterpack_row.png"
    render_kv_card(
        out_path=p_sp,
        title="Starter pack (metadata)",
        accent_hex=COLORS["pink"],
        items=(
            ("Name", _truncate_middle(_clean_text(sp["name"]), 34)),
            ("Creator DID", _compact_did(sp["creator_did"])),
            ("Pack URI", _compact_at_uri(sp["starterpack_uri"])),
        ),
        height_px=280,
    )
    receipts["starterpack"] = p_sp

    p_inc = ASSET_DIR / "starterpack_inclusion_row.png"
    render_kv_card(
        out_path=p_inc,
        title="Starter pack → feed inclusion",
        accent_hex=COLORS["pink"],
        items=(
            ("Pack", _compact_at_uri(inc["starterpack_uri"])),
            ("Feed", _compact_at_uri(inc["feed_uri"])),
        ),
        height_px=240,
    )
    receipts["starterpack_inclusion"] = p_inc

    # Popular row (rank surface)
    p_pop = ASSET_DIR / "popular_row.png"
    render_kv_card(
        out_path=p_pop,
        title="Popular feeds (example row)",
        accent_hex=COLORS["amber"],
        items=(
            ("Rank", "top"),
            ("Feed", _clean_text(feed_idx.get("display_name", "")) or "(no name)"),
            ("URI", _compact_at_uri(exemplar_feed_uri)),
        ),
    )
    receipts["popular"] = p_pop

    # Panel entry
    panel_row = _find_row(RUN_CSV_DIR / "feed_panel.csv", key="feed_uri", value=exemplar_feed_uri)

    def _human_reason(reason: str) -> str:
        reason = _clean_text(reason)
        if reason == "top_popular_feed_generators":
            return "Top of popular list"
        if reason == "starterpack_top_by_slot_count":
            return "Most included in starter packs"
        if reason == "provider_balanced_long_tail":
            return "Provider-balanced long tail"
        return reason.replace("_", " ")

    p_panel = ASSET_DIR / "panel_row.png"
    render_kv_card(
        out_path=p_panel,
        title="Feed panel (chosen feed)",
        accent_hex=COLORS["cyan"],
        items=(
            ("Feed", _clean_text(panel_row.get("display_name", "")) or "(no name)"),
            ("Group", _clean_text(panel_row.get("feed_group", "")).replace("_", "-").title()),
            ("Chosen by", _human_reason(panel_row.get("selection_reason", ""))),
            ("URI", _compact_at_uri(exemplar_feed_uri)),
        ),
        height_px=300,
    )
    receipts["panel"] = p_panel

    # Impression (ranked exposure)
    feed_items_path = RUN_CSV_DIR / "feed_items.csv.gz"
    impression_row = _find_row_gz(
        feed_items_path,
        predicate=lambda r: r.get("feed_uri") == exemplar_feed_uri and r.get("viewer_mode") == "unauth" and r.get("rank") == "1",
    )
    p_imp = ASSET_DIR / "impression_row.png"
    render_kv_card(
        out_path=p_imp,
        title="Ranked impression (example)",
        accent_hex=COLORS["purple"],
        items=(
            ("Viewer", "unauth (public)"),
            ("Rank", "top"),
            ("Post", _compact_at_uri(impression_row["post_uri"])),
            ("Author DID", _compact_did(impression_row["author_did"])),
        ),
        height_px=300,
    )
    receipts["impression"] = p_imp

    # Post metadata (no text)
    post_uri = impression_row["post_uri"]
    post_cid = impression_row["post_cid"]
    post_row = _find_row_gz(
        RUN_CSV_DIR / "posts.csv.gz",
        predicate=lambda r: r.get("post_uri") == post_uri and r.get("post_cid") == post_cid,
    )

    embed_type = _clean_text(post_row.get("embed_type", ""))
    embed_short = embed_type.replace("app.bsky.embed.", "") if embed_type else "(none)"
    has_images = "yes" if _clean_text(post_row.get("image_count", "0")) not in ("", "0") else "no"
    ext_domain = _clean_text(post_row.get("external_domain", "")) or "(none)"
    langs = _clean_text(post_row.get("langs_json", "")) or "(unknown)"

    p_post = ASSET_DIR / "post_row.png"
    render_kv_card(
        out_path=p_post,
        title="Post metadata (no text)",
        accent_hex=COLORS["pink"],
        items=(
            ("Post", _compact_at_uri(post_uri)),
            ("Embed", embed_short),
            ("Images", has_images),
            ("External", ext_domain),
            ("Langs", _truncate_middle(langs, 22)),
        ),
        height_px=340,
        width_px=1400,
    )
    receipts["post"] = p_post

    # Author profile (no counts)
    author_did = post_row["author_did"]
    author_row = _find_row_gz(
        RUN_CSV_DIR / "authors.csv.gz",
        predicate=lambda r: r.get("author_did") == author_did,
    )
    p_author = ASSET_DIR / "author_row.png"
    render_kv_card(
        out_path=p_author,
        title="Author profile (no counts)",
        accent_hex=COLORS["green"],
        items=(
            ("Handle", _truncate_middle(_clean_text(author_row.get("handle", "")), 30)),
            ("Name", _truncate_middle(_clean_text(author_row.get("display_name", "")) or "(none)", 30)),
            ("DID", _compact_did(author_did)),
        ),
        height_px=280,
    )
    receipts["author"] = p_author

    # Appendix: manifest excerpt + data dictionary excerpt
    want_files = {"feed_panel.csv", "feed_items.csv.gz", "posts.csv.gz", "authors.csv.gz"}
    manifest_rows = [r for r in _iter_csv_rows(RUN_CSV_DIR / "manifest.csv") if r.get("file_name") in want_files]
    manifest_rows.sort(key=lambda r: r["file_name"])
    manifest_lines = [f"{r['file_name']}: sha256 {_compact_sha256(r['sha256'])}" for r in manifest_rows[:4]]
    p_manifest = ASSET_DIR / "manifest_excerpt.png"
    render_list_card(
        out_path=p_manifest,
        title="Manifest (files + hashes)",
        accent_hex=COLORS["purple"],
        lines=manifest_lines,
        height_px=260,
        width_px=1400,
    )
    receipts["manifest"] = p_manifest

    dd_rows = [
        r
        for r in _iter_csv_rows(RUN_CSV_DIR / "data_dictionary.csv")
        if r.get("file_name") == "feed_items.csv.gz"
        and r.get("column_name") in {"feed_uri", "viewer_mode", "rank", "post_uri"}
    ]
    dd_rows.sort(key=lambda r: ["feed_uri", "viewer_mode", "rank", "post_uri"].index(r["column_name"]))
    dd_lines = [f"{r['column_name']}: {r['description']}" for r in dd_rows]
    p_dd = ASSET_DIR / "data_dictionary_excerpt.png"
    render_list_card(
        out_path=p_dd,
        title="Data dictionary (selected columns)",
        accent_hex=COLORS["green"],
        lines=dd_lines,
        height_px=300,
        width_px=1400,
    )
    receipts["data_dictionary"] = p_dd

    return receipts


def _add_bg(slide, prs: Presentation) -> None:
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _rgb(COLORS["bg"])
    bg.line.fill.background()


def _add_title(slide, *, title: str, subtitle: str | None) -> None:
    # Match deck style: ~34pt title, 16pt subtitle
    t = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(8.6), Inches(0.6))
    tf = t.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Calibri"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = _rgb(COLORS["ink"])

    if subtitle:
        s = slide.shapes.add_textbox(Inches(0.72), Inches(1.15), Inches(8.6), Inches(0.4))
        sf = s.text_frame
        sf.clear()
        sp = sf.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Calibri"
        sp.font.size = Pt(16)
        sp.font.color.rgb = _rgb(COLORS["muted"])


def _add_card(slide, *, x: float, y: float, w: float, h: float, border_hex: str) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(COLORS["card"])
    shape.line.color.rgb = _rgb(border_hex)
    shape.line.width = Pt(2)


def _add_text_in_box(slide, *, x: float, y: float, w: float, h: float, text: str, size: int, bold: bool = False, color_hex: str = "F5F5FA") -> None:
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Calibri"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(color_hex)


def _move_last_slide_to(prs: Presentation, pos: int) -> None:
    # pos is 0-based index in the slide list (insert position).
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001 (python-pptx internal)
    last = sld_id_lst[-1]
    sld_id_lst.remove(last)
    sld_id_lst.insert(pos, last)


def add_appendix_repro(prs: Presentation, receipts: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_title(slide, title="Reproducible outputs", subtitle="Every export is hashed; every column is documented.")

    # Two large cards
    _add_card(slide, x=0.9, y=1.75, w=8.2, h=1.7, border_hex=COLORS["purple"])
    slide.shapes.add_picture(str(receipts["manifest"]), Inches(1.05), Inches(1.92), height=Inches(1.35))

    _add_card(slide, x=0.9, y=3.65, w=8.2, h=1.75, border_hex=COLORS["green"])
    slide.shapes.add_picture(str(receipts["data_dictionary"]), Inches(1.05), Inches(3.82), height=Inches(1.4))

    _move_last_slide_to(prs, 56)  # after original slide 56


def add_stitched_example(prs: Presentation, receipts: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_title(
        slide,
        title="One stitched example (keys line up)",
        subtitle="Panel entry → ranked impression → post metadata → author profile (no post text).",
    )

    y = 1.75
    h = 3.6
    w = 2.05
    g = 0.2
    x0 = 0.7
    xs = [x0 + i * (w + g) for i in range(4)]
    keys = ["panel", "impression", "post", "author"]
    borders = [COLORS["cyan"], COLORS["purple"], COLORS["pink"], COLORS["green"]]

    for x, key, border in zip(xs, keys, borders, strict=True):
        _add_card(slide, x=x, y=y, w=w, h=h, border_hex=border)
        # Slight inset so card border isn't clipped by image alpha.
        slide.shapes.add_picture(str(receipts[key]), Inches(x + 0.12), Inches(y + 0.15), width=Inches(w - 0.24))

    # Arrows between cards
    for i in range(3):
        ax = xs[i] + w + 0.02
        ay = y + 1.65
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(ax), Inches(ay), Inches(g - 0.04), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = _rgb(COLORS["muted2"])
        arrow.line.fill.background()

    _move_last_slide_to(prs, 54)  # after original slide 54


def add_discovery_receipts(prs: Presentation, receipts: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_title(
        slide,
        title="Discovery leaves artifacts",
        subtitle="Starter packs and popular lists become joinable tables (IDs we reuse later).",
    )

    # Two big sections
    _add_card(slide, x=0.9, y=1.75, w=4.05, h=3.75, border_hex=COLORS["pink"])
    _add_text_in_box(slide, x=1.1, y=1.9, w=3.7, h=0.3, text="Starter packs", size=16, bold=True)
    slide.shapes.add_picture(str(receipts["starterpack"]), Inches(1.1), Inches(2.2), width=Inches(3.75))
    slide.shapes.add_picture(str(receipts["starterpack_inclusion"]), Inches(1.1), Inches(3.65), width=Inches(3.75))
    _add_text_in_box(
        slide,
        x=1.1,
        y=4.95,
        w=3.7,
        h=0.45,
        text="We also aggregate pack inclusions into a feed-level ranking.",
        size=12,
        color_hex=COLORS["muted"],
    )

    _add_card(slide, x=5.05, y=1.75, w=4.05, h=3.75, border_hex=COLORS["amber"])
    _add_text_in_box(slide, x=5.25, y=1.9, w=3.7, h=0.3, text="Popular feeds", size=16, bold=True)
    slide.shapes.add_picture(str(receipts["popular"]), Inches(5.25), Inches(2.2), width=Inches(3.75))
    _add_text_in_box(
        slide,
        x=5.25,
        y=3.65,
        w=3.7,
        h=0.45,
        text="Stores a ranked list we can sample from for the panel.",
        size=12,
        color_hex=COLORS["muted"],
    )

    _move_last_slide_to(prs, 31)  # after original slide 31


def add_run_slide(prs: Presentation, anchors: RunAnchors, receipts: dict[str, Path]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, prs)
    _add_title(slide, title="The run behind this story", subtitle="A single unauth run turns public APIs into analysis-ready datasets.")

    # Top-left: run context
    _add_card(slide, x=0.9, y=1.75, w=4.05, h=1.45, border_hex=COLORS["green"])
    _add_text_in_box(slide, x=1.1, y=1.92, w=3.7, h=0.3, text="Viewer context", size=16, bold=True)
    _add_text_in_box(
        slide,
        x=1.1,
        y=2.25,
        w=3.7,
        h=0.8,
        text=f"auth_mode: {anchors.auth_mode}\nrun_id: {_truncate_middle(anchors.run_id, 18)}",
        size=13,
        color_hex=COLORS["muted"],
    )

    # Top-right: anchors (only 2 numbers)
    _add_card(slide, x=5.05, y=1.75, w=4.05, h=1.45, border_hex=COLORS["cyan"])
    _add_text_in_box(slide, x=5.25, y=1.92, w=3.7, h=0.3, text="Anchors (minimal)", size=16, bold=True)
    _add_text_in_box(
        slide,
        x=5.25,
        y=2.25,
        w=3.7,
        h=0.8,
        text=f"Panel: {anchors.num_feeds_panel:,} feeds\nCaptured: ~{anchors.num_feed_items//1000:,}k ranked impressions",
        size=18,
        bold=True,
        color_hex=COLORS["ink"],
    )

    # Bottom: what gets exported + example index row
    _add_card(slide, x=0.9, y=3.4, w=8.2, h=2.05, border_hex=COLORS["purple"])
    _add_text_in_box(slide, x=1.1, y=3.56, w=4.5, h=0.3, text="Exports (conceptual)", size=16, bold=True)
    _add_text_in_box(
        slide,
        x=1.1,
        y=3.9,
        w=4.6,
        h=1.4,
        text="• Feed generator index\n• Starter packs + inclusions\n• Popular feed list\n• Feed panel\n• Ranked impressions\n• Posts (metadata)\n• Authors (profiles)",
        size=13,
        color_hex=COLORS["muted"],
    )
    slide.shapes.add_picture(str(receipts["feed_index"]), Inches(5.05), Inches(3.78), width=Inches(3.85))
    _add_text_in_box(
        slide,
        x=5.05,
        y=5.2,
        w=3.9,
        h=0.25,
        text="(Example: one row from the index)",
        size=11,
        color_hex=COLORS["muted2"],
    )

    _move_last_slide_to(prs, 10)  # after original slide 10


def main() -> None:
    if not BASE_DECK.exists():
        raise SystemExit(f"Base deck not found: {BASE_DECK}")

    anchors = load_run_anchors()
    receipts = build_receipts(anchors)

    prs = Presentation(str(BASE_DECK))

    # Insert in descending order so earlier positions remain stable.
    add_appendix_repro(prs, receipts)
    add_stitched_example(prs, receipts)
    add_discovery_receipts(prs, receipts)
    add_run_slide(prs, anchors, receipts)

    prs.save(str(OUT_DECK))
    print(f"Wrote: {OUT_DECK}")


if __name__ == "__main__":
    main()
