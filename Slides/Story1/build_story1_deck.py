#!/usr/bin/env python3
from __future__ import annotations

import io
import json
import math
import shutil
import zipfile
from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Literal

from lxml import etree
from PIL import Image, ImageDraw, ImageEnhance, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


EMU_PER_INCH = 914400

PPTX_MIME_SLIDE = "application/vnd.openxmlformats-officedocument.presentationml.slide+xml"
REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
CT_NS = "http://schemas.openxmlformats.org/package/2006/content-types"

NSMAP = {"p": PML_NS, "r": R_NS}
CT_NSMAP = {"ct": CT_NS}
REL_NSMAP = {"rel": REL_NS}

COLORS = {
    "bg": "0B0B0F",
    "card": "14141C",
    "ink": "F5F5FA",
    "muted": "B4B4C8",
    "muted2": "46465A",
    "cyan": "00D4FF",
    "purple": "A855F7",
    "pink": "FF5D9E",
    "amber": "F59E0B",
    "green": "22C55E",
}

# System fonts used by the existing decks/scripts.
FONT_SANS = Path("/System/Library/Fonts/SFNS.ttf")
FONT_MONO = Path("/System/Library/Fonts/SFNSMono.ttf")


# The blackboard font used in the template deck is missing several Unicode dash variants.
# Normalize text to ASCII so it renders consistently in LibreOffice/PowerPoint.
_DASH_TRANSLATION = str.maketrans(
    {
        "\u2010": "-",  # hyphen
        "\u2011": "-",  # non-breaking hyphen
        "\u2012": "-",  # figure dash
        "\u2013": "-",  # en dash
        "\u2014": "-",  # em dash
        "\u2212": "-",  # minus sign
        "\u00A0": " ",  # non-breaking space
    }
)


def _norm_text(text: str) -> str:
    return text.translate(_DASH_TRANSLATION)


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


AssertAssetRole = Literal["structural", "semantic"]


@dataclass(frozen=True)
class AssertAssetRef:
    source_pptx: Path
    media_path: str
    role: AssertAssetRole
    derived_path: Path


@dataclass(frozen=True)
class SlideManifestEntry:
    slide_num: int
    title: str
    assets: list[AssertAssetRef]


@dataclass(frozen=True)
class PreparedAssets:
    bg_blueprint_overlay: AssertAssetRef
    bg_blueprint_card: AssertAssetRef

    icon_pawn: AssertAssetRef
    icon_pawn_wm: AssertAssetRef
    icon_knight: AssertAssetRef
    icon_knight_wm: AssertAssetRef
    icon_king: AssertAssetRef
    icon_king_wm: AssertAssetRef
    icon_rook_or_bishop: AssertAssetRef
    icon_rook_or_bishop_wm: AssertAssetRef
    icon_queen_or_rook: AssertAssetRef
    icon_queen_or_rook_wm: AssertAssetRef
    icon_chessboard: AssertAssetRef
    icon_chessboard_wm: AssertAssetRef

    icon_satellite: AssertAssetRef
    icon_satellite_wm: AssertAssetRef
    icon_brain: AssertAssetRef
    icon_brain_wm: AssertAssetRef

    icon_fragile_rot: AssertAssetRef


@dataclass(frozen=True)
class BuildPaths:
    template_pptx: Path
    working_pptx: Path
    plots_dir: Path
    assets_dir: Path


def _read_zip_xml(z: zipfile.ZipFile, name: str) -> etree._Element:
    return etree.fromstring(z.read(name))


def _write_zip_xml(z: zipfile.ZipFile, name: str, root: etree._Element) -> None:
    z.writestr(name, etree.tostring(root, xml_declaration=True, encoding="UTF-8", standalone=True))


def _max_slide_number(z: zipfile.ZipFile) -> int:
    max_n = 0
    for name in z.namelist():
        if not name.startswith("ppt/slides/slide") or not name.endswith(".xml"):
            continue
        stem = Path(name).stem
        try:
            n = int(stem.replace("slide", ""))
        except ValueError:
            continue
        max_n = max(max_n, n)
    return max_n


def _next_rid(pres_rels_root: etree._Element) -> str:
    rids: list[int] = []
    for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
        rid = rel.get("Id", "")
        if rid.startswith("rId"):
            try:
                rids.append(int(rid[3:]))
            except ValueError:
                continue
    nxt = (max(rids) + 1) if rids else 1
    return f"rId{nxt}"


def _next_slide_id(pres_root: etree._Element) -> int:
    sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
    if sld_id_lst is None:
        raise RuntimeError("presentation.xml missing p:sldIdLst")
    ids: list[int] = []
    for sld_id in sld_id_lst.findall("p:sldId", namespaces=NSMAP):
        try:
            ids.append(int(sld_id.get("id", "0")))
        except ValueError:
            continue
    return (max(ids) + 1) if ids else 256


def _remove_notes_relationship(slide_rels_root: etree._Element) -> None:
    for rel in list(slide_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP)):
        if rel.get("Type") == "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide":
            slide_rels_root.remove(rel)


def build_working_structure(paths: BuildPaths) -> None:
    """Create the working deck by duplicating the bullet slide and reordering into the Story1 arc."""
    paths.working_pptx.parent.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(paths.template_pptx, "r") as zin:
        max_slide = _max_slide_number(zin)

        # Duplicate slide 2 ("What you should believe") 8x for narrative bullet slides.
        duplicate_slide_num = 2
        duplicate_count = 8
        slide_src = f"ppt/slides/slide{duplicate_slide_num}.xml"
        slide_rels_src = f"ppt/slides/_rels/slide{duplicate_slide_num}.xml.rels"
        slide_src_bytes = zin.read(slide_src)
        slide_rels_root = _read_zip_xml(zin, slide_rels_src)
        _remove_notes_relationship(slide_rels_root)
        slide_rels_bytes = etree.tostring(slide_rels_root, xml_declaration=True, encoding="UTF-8", standalone=True)

        pres_root = _read_zip_xml(zin, "ppt/presentation.xml")
        pres_rels_root = _read_zip_xml(zin, "ppt/_rels/presentation.xml.rels")
        ct_root = _read_zip_xml(zin, "[Content_Types].xml")

        new_slide_nums = list(range(max_slide + 1, max_slide + 1 + duplicate_count))
        sld_id_lst = pres_root.find("p:sldIdLst", namespaces=NSMAP)
        if sld_id_lst is None:
            raise RuntimeError("presentation.xml missing p:sldIdLst")

        next_slide_id = _next_slide_id(pres_root)
        for n in new_slide_nums:
            rid = _next_rid(pres_rels_root)

            rel_el = etree.Element(f"{{{REL_NS}}}Relationship")
            rel_el.set("Id", rid)
            rel_el.set("Type", "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide")
            rel_el.set("Target", f"slides/slide{n}.xml")
            pres_rels_root.append(rel_el)

            sld_id_el = etree.Element(f"{{{PML_NS}}}sldId")
            sld_id_el.set("id", str(next_slide_id))
            sld_id_el.set(f"{{{R_NS}}}id", rid)
            sld_id_lst.append(sld_id_el)
            next_slide_id += 1

            override = etree.Element(f"{{{CT_NS}}}Override")
            override.set("PartName", f"/ppt/slides/slide{n}.xml")
            override.set("ContentType", PPTX_MIME_SLIDE)
            ct_root.append(override)

        # Desired slide part order (original 1–18 + new 19–26).
        desired_nums = [
            1,
            2,
            4,
            14,
            18,
            19,
            20,
            21,
            3,
            13,
            15,
            6,
            7,
            8,
            9,
            10,
            11,
            12,
            5,
            22,
            23,
            16,
            24,
            25,
            26,
            17,
        ]

        rid_to_target: dict[str, str] = {}
        for rel in pres_rels_root.findall("rel:Relationship", namespaces=REL_NSMAP):
            if rel.get("Type") != "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide":
                continue
            rid = rel.get("Id")
            tgt = rel.get("Target")
            if rid and tgt:
                rid_to_target[rid] = tgt

        target_to_sld_id: dict[str, etree._Element] = {}
        for sld_id in list(sld_id_lst.findall("p:sldId", namespaces=NSMAP)):
            rid = sld_id.get(f"{{{R_NS}}}id")
            tgt = rid_to_target.get(rid or "")
            if tgt:
                target_to_sld_id[tgt] = sld_id

        for child in list(sld_id_lst):
            sld_id_lst.remove(child)
        for n in desired_nums:
            target = f"slides/slide{n}.xml"
            el = target_to_sld_id.get(target)
            if el is None:
                raise RuntimeError(f"Missing slide target in presentation: {target}")
            sld_id_lst.append(deepcopy(el))

        tmp_path = paths.working_pptx.with_suffix(".tmp.pptx")
        if tmp_path.exists():
            tmp_path.unlink()

        with zipfile.ZipFile(tmp_path, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                if item.filename in {"ppt/presentation.xml", "ppt/_rels/presentation.xml.rels", "[Content_Types].xml"}:
                    continue
                zout.writestr(item, zin.read(item.filename))

            for n in new_slide_nums:
                zout.writestr(f"ppt/slides/slide{n}.xml", slide_src_bytes)
                zout.writestr(f"ppt/slides/_rels/slide{n}.xml.rels", slide_rels_bytes)

            _write_zip_xml(zout, "ppt/presentation.xml", pres_root)
            _write_zip_xml(zout, "ppt/_rels/presentation.xml.rels", pres_rels_root)
            _write_zip_xml(zout, "[Content_Types].xml", ct_root)

        shutil.move(tmp_path, paths.working_pptx)


def _shape(slide, shape_id: int):
    for sh in slide.shapes:
        if sh.shape_id == shape_id:
            return sh
    raise KeyError(f"shape_id={shape_id} not found on slide")


def _set_run_text(shape, text: str) -> None:
    if not shape.has_text_frame:
        return
    text = _norm_text(text)
    tf = shape.text_frame
    if not tf.paragraphs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = text
        for run in p.runs[1:]:
            run.text = ""
    else:
        r = p.add_run()
        r.text = text

    # Clear any leftover text from other paragraphs in the template shape.
    for para in tf.paragraphs[1:]:
        for run in para.runs:
            run.text = ""


def _set_border_color(shape, hex6: str) -> None:
    if not hasattr(shape, "line"):
        return
    shape.line.color.rgb = _rgb(hex6)


def _fill_shape_with_picture(shape, slide, image_path: Path) -> None:
    # Keep the shape (and its animations) and only change its fill to a:blipFill.
    _image_part, rid = slide.part.get_or_add_image_part(str(image_path))

    sp = shape._element  # noqa: SLF001 (pptx oxml)
    sp_pr = sp.xpath("./p:spPr")[0]

    fill_tags = {"noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}
    for child in list(sp_pr):
        # child.tag is like '{ns}solidFill'
        local = etree.QName(child).localname
        if local in fill_tags:
            sp_pr.remove(child)

    blip_fill = OxmlElement("a:blipFill")
    blip = OxmlElement("a:blip")
    blip.set(qn("r:embed"), rid)
    blip_fill.append(blip)
    stretch = OxmlElement("a:stretch")
    stretch.append(OxmlElement("a:fillRect"))
    blip_fill.append(stretch)

    ln_els = sp_pr.xpath("./a:ln")
    insert_at = sp_pr.index(ln_els[0]) if ln_els else len(sp_pr)
    sp_pr.insert(insert_at, blip_fill)


def _clear_text(shape) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        for run in para.runs:
            run.text = ""


def _set_run_font_size(shape, size_pt: float) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    if not tf.paragraphs:
        return
    p = tf.paragraphs[0]
    if not p.runs:
        return
    p.runs[0].font.size = Pt(size_pt)


def _set_shape_line_nofill(shape) -> None:
    sp = shape._element  # noqa: SLF001 (pptx oxml)
    sp_pr = sp.xpath("./p:spPr")[0]
    ln_els = sp_pr.xpath("./a:ln")
    if not ln_els:
        return
    ln = ln_els[0]
    # Remove existing fill children and set noFill.
    for child in list(ln):
        local = etree.QName(child).localname
        if local in {"solidFill", "gradFill", "pattFill", "noFill"}:
            ln.remove(child)
    ln.append(OxmlElement("a:noFill"))


def _set_shape_fill_nofill(shape) -> None:
    sp = shape._element  # noqa: SLF001 (pptx oxml)
    sp_pr = sp.xpath("./p:spPr")[0]
    fill_tags = {"noFill", "solidFill", "gradFill", "blipFill", "pattFill", "grpFill"}
    for child in list(sp_pr):
        local = etree.QName(child).localname
        if local in fill_tags:
            sp_pr.remove(child)
    sp_pr.insert(0, OxmlElement("a:noFill"))


def _hide_shape(shape) -> None:
    _clear_text(shape)
    _set_shape_fill_nofill(shape)
    _set_shape_line_nofill(shape)


def _emu(inches: float) -> int:
    return int(inches * EMU_PER_INCH)


def _slide_wh_in(prs: Presentation) -> tuple[float, float]:
    return (prs.slide_width / EMU_PER_INCH, prs.slide_height / EMU_PER_INCH)


def _move_shape_to_back(slide, shape, *, after_background: bool) -> None:
    sp_tree = slide.shapes._spTree  # noqa: SLF001 (pptx oxml)
    insert_at = 3 if after_background else 2
    sp_tree.insert(insert_at, shape._element)  # noqa: SLF001 (pptx oxml)


def _add_fullslide_background(slide, prs: Presentation, image_path: Path) -> None:
    pic = slide.shapes.add_picture(str(image_path), 0, 0, width=prs.slide_width, height=prs.slide_height)
    _move_shape_to_back(slide, pic, after_background=False)


def _add_icon_watermark(
    *,
    slide,
    prs: Presentation,
    image_path: Path,
    size_in: float,
    right_in: float = 0.7,
    top_in: float = 4.2,
) -> None:
    slide_w_in, _slide_h_in = _slide_wh_in(prs)
    left_in = max(0.0, slide_w_in - right_in - size_in)
    pic = slide.shapes.add_picture(str(image_path), _emu(left_in), _emu(top_in), width=_emu(size_in), height=_emu(size_in))
    _move_shape_to_back(slide, pic, after_background=True)


def _add_picture_icon(
    *,
    slide,
    image_path: Path,
    left_in: float,
    top_in: float,
    size_in: float,
    send_to_back: bool,
) -> None:
    pic = slide.shapes.add_picture(
        str(image_path),
        _emu(left_in),
        _emu(top_in),
        width=_emu(size_in),
        height=_emu(size_in),
    )
    if send_to_back:
        _move_shape_to_back(slide, pic, after_background=True)


def _extract_media(*, source_pptx: Path, media_path: str, out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    with zipfile.ZipFile(source_pptx, "r") as z:
        out_path.write_bytes(z.read(media_path))
    return out_path


def _apply_alpha_multiplier(img: Image.Image, alpha_mult: float) -> Image.Image:
    if alpha_mult >= 0.999:
        return img
    if img.mode != "RGBA":
        img = img.convert("RGBA")
    r, g, b, a = img.split()
    a = a.point(lambda v: int(v * alpha_mult))
    return Image.merge("RGBA", (r, g, b, a))


def _prep_blueprint_overlay(
    *, src_path: Path, out_path: Path, tint_hex: str, alpha_max: float, size_px: tuple[int, int]
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    img = Image.open(src_path).convert("RGB")
    g = ImageOps.grayscale(img)
    g = ImageOps.invert(g)
    g = ImageOps.autocontrast(g, cutoff=2)
    g = ImageEnhance.Contrast(g).enhance(1.7)
    g = ImageEnhance.Brightness(g).enhance(0.9)
    g = g.resize(size_px, resample=Image.LANCZOS)

    tint_rgb = tuple(int(tint_hex[i : i + 2], 16) for i in (0, 2, 4))
    alpha_scale = max(0.0, min(1.0, alpha_max))
    a = g.point(lambda v: int(v * alpha_scale))
    overlay = Image.new("RGBA", size_px, tint_rgb + (0,))
    overlay.putalpha(a)
    overlay.save(out_path)
    return out_path


def _prep_blueprint_card_texture(
    *,
    src_path: Path,
    out_path: Path,
    card_hex: str,
    tint_hex: str,
    alpha_max: float,
    size_px: tuple[int, int],
) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    overlay_path = out_path.with_name(out_path.stem + "_overlay_tmp.png")
    _prep_blueprint_overlay(src_path=src_path, out_path=overlay_path, tint_hex=tint_hex, alpha_max=alpha_max, size_px=size_px)

    bg_rgb = tuple(int(card_hex[i : i + 2], 16) for i in (0, 2, 4))
    bg = Image.new("RGB", size_px, bg_rgb)
    overlay = Image.open(overlay_path).convert("RGBA")
    bg_rgba = bg.convert("RGBA")
    comp = Image.alpha_composite(bg_rgba, overlay)
    comp.convert("RGB").save(out_path)
    try:
        overlay_path.unlink()
    except OSError:
        pass
    return out_path


def _prep_icon(*, src_path: Path, out_path: Path, size_px: int, alpha_mult: float) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    img = Image.open(src_path).convert("RGBA")
    bbox = img.split()[-1].getbbox()
    if bbox:
        img = img.crop(bbox)
    img = _apply_alpha_multiplier(img, alpha_mult=alpha_mult)

    w, h = img.size
    if w == 0 or h == 0:
        raise ValueError(f"empty icon image: {src_path}")
    scale = min(size_px / w, size_px / h)
    new_w = max(1, int(w * scale))
    new_h = max(1, int(h * scale))
    img = img.resize((new_w, new_h), resample=Image.LANCZOS)

    canvas = Image.new("RGBA", (size_px, size_px), (0, 0, 0, 0))
    canvas.paste(img, ((size_px - new_w) // 2, (size_px - new_h) // 2), img)
    canvas.save(out_path)
    return out_path


def _prep_rotated_icon(*, src_path: Path, out_path: Path, size_px: int, alpha_mult: float, degrees: float) -> Path:
    tmp = out_path.with_name(out_path.stem + "_tmp.png")
    if out_path.exists() and out_path.stat().st_size > 0:
        return out_path
    _prep_icon(src_path=src_path, out_path=tmp, size_px=size_px, alpha_mult=alpha_mult)
    img = Image.open(tmp).convert("RGBA").rotate(degrees, resample=Image.BICUBIC, expand=True)
    # Center-crop back to square.
    w, h = img.size
    if w > size_px or h > size_px:
        left = max(0, (w - size_px) // 2)
        top = max(0, (h - size_px) // 2)
        img = img.crop((left, top, left + size_px, top + size_px))
    else:
        canvas = Image.new("RGBA", (size_px, size_px), (0, 0, 0, 0))
        canvas.paste(img, ((size_px - w) // 2, (size_px - h) // 2), img)
        img = canvas
    img.save(out_path)
    try:
        tmp.unlink()
    except OSError:
        pass
    return out_path


def prepare_assert_assets(*, assert_dir: Path, cache_dir: Path) -> PreparedAssets:
    # Structural blueprint overlay source.
    arch_pptx = assert_dir / "Architect Infographics by Slidesgo.pptx"
    arch_media = "ppt/media/image12.png"
    src_blueprint = _extract_media(
        source_pptx=arch_pptx, media_path=arch_media, out_path=cache_dir / "raw_architect_image12.png"
    )

    overlay_path = cache_dir / "bg_blueprint_overlay.png"
    card_path = cache_dir / "bg_blueprint_card.png"
    _prep_blueprint_overlay(src_path=src_blueprint, out_path=overlay_path, tint_hex=COLORS["muted2"], alpha_max=0.11, size_px=(1920, 1080))
    _prep_blueprint_card_texture(
        src_path=src_blueprint,
        out_path=card_path,
        card_hex=COLORS["card"],
        tint_hex=COLORS["muted2"],
        alpha_max=0.22,
        size_px=(900, 1120),
    )

    strategy_pptx = assert_dir / "Copy of Strategy Infographics by Slidesgo.pptx"
    chess = {
        "pawn": "ppt/media/image1.png",
        "knight": "ppt/media/image2.png",
        "king": "ppt/media/image3.png",
        "rook_or_bishop": "ppt/media/image4.png",
        "queen_or_rook": "ppt/media/image5.png",
        "chessboard": "ppt/media/image6.png",
    }
    raw_chess: dict[str, Path] = {}
    for k, media in chess.items():
        raw_chess[k] = _extract_media(source_pptx=strategy_pptx, media_path=media, out_path=cache_dir / f"raw_{k}.png")

    icons: dict[str, Path] = {}
    icons_wm: dict[str, Path] = {}
    for k, src in raw_chess.items():
        icons[k] = cache_dir / f"icon_{k}.png"
        icons_wm[k] = cache_dir / f"icon_{k}_wm.png"
        _prep_icon(src_path=src, out_path=icons[k], size_px=1024, alpha_mult=1.0)
        _prep_icon(src_path=src, out_path=icons_wm[k], size_px=1024, alpha_mult=0.14)

    satellite_pptx = assert_dir / "EN What Is a Satellite_ by Slidesgo.pptx"
    sat_media = "ppt/media/image2.png"
    raw_sat = _extract_media(source_pptx=satellite_pptx, media_path=sat_media, out_path=cache_dir / "raw_satellite.png")
    sat_icon = cache_dir / "icon_satellite.png"
    sat_icon_wm = cache_dir / "icon_satellite_wm.png"
    _prep_icon(src_path=raw_sat, out_path=sat_icon, size_px=1024, alpha_mult=1.0)
    _prep_icon(src_path=raw_sat, out_path=sat_icon_wm, size_px=1024, alpha_mult=0.14)

    brain_pptx = assert_dir / "Copy of Brain Infographics by Slidesgo.pptx"
    brain_media = "ppt/media/image1.png"
    raw_brain = _extract_media(source_pptx=brain_pptx, media_path=brain_media, out_path=cache_dir / "raw_brain.png")
    brain_icon = cache_dir / "icon_brain.png"
    brain_icon_wm = cache_dir / "icon_brain_wm.png"
    _prep_icon(src_path=raw_brain, out_path=brain_icon, size_px=1024, alpha_mult=1.0)
    _prep_icon(src_path=raw_brain, out_path=brain_icon_wm, size_px=1024, alpha_mult=0.12)

    fragile_rot = cache_dir / "icon_fragile_rot.png"
    _prep_rotated_icon(src_path=raw_chess["queen_or_rook"], out_path=fragile_rot, size_px=1024, alpha_mult=0.16, degrees=-18)

    return PreparedAssets(
        bg_blueprint_overlay=AssertAssetRef(arch_pptx, arch_media, "structural", overlay_path),
        bg_blueprint_card=AssertAssetRef(arch_pptx, arch_media, "structural", card_path),
        icon_pawn=AssertAssetRef(strategy_pptx, chess["pawn"], "semantic", icons["pawn"]),
        icon_pawn_wm=AssertAssetRef(strategy_pptx, chess["pawn"], "semantic", icons_wm["pawn"]),
        icon_knight=AssertAssetRef(strategy_pptx, chess["knight"], "semantic", icons["knight"]),
        icon_knight_wm=AssertAssetRef(strategy_pptx, chess["knight"], "semantic", icons_wm["knight"]),
        icon_king=AssertAssetRef(strategy_pptx, chess["king"], "semantic", icons["king"]),
        icon_king_wm=AssertAssetRef(strategy_pptx, chess["king"], "semantic", icons_wm["king"]),
        icon_rook_or_bishop=AssertAssetRef(strategy_pptx, chess["rook_or_bishop"], "semantic", icons["rook_or_bishop"]),
        icon_rook_or_bishop_wm=AssertAssetRef(strategy_pptx, chess["rook_or_bishop"], "semantic", icons_wm["rook_or_bishop"]),
        icon_queen_or_rook=AssertAssetRef(strategy_pptx, chess["queen_or_rook"], "semantic", icons["queen_or_rook"]),
        icon_queen_or_rook_wm=AssertAssetRef(strategy_pptx, chess["queen_or_rook"], "semantic", icons_wm["queen_or_rook"]),
        icon_chessboard=AssertAssetRef(strategy_pptx, chess["chessboard"], "semantic", icons["chessboard"]),
        icon_chessboard_wm=AssertAssetRef(strategy_pptx, chess["chessboard"], "semantic", icons_wm["chessboard"]),
        icon_satellite=AssertAssetRef(satellite_pptx, sat_media, "semantic", sat_icon),
        icon_satellite_wm=AssertAssetRef(satellite_pptx, sat_media, "semantic", sat_icon_wm),
        icon_brain=AssertAssetRef(brain_pptx, brain_media, "semantic", brain_icon),
        icon_brain_wm=AssertAssetRef(brain_pptx, brain_media, "semantic", brain_icon_wm),
        icon_fragile_rot=AssertAssetRef(strategy_pptx, chess["queen_or_rook"], "semantic", fragile_rot),
    )


def build_assert_inventory(*, assert_dir: Path, out_json: Path, out_md: Path, curated: list[AssertAssetRef]) -> None:
    items: list[dict[str, object]] = []
    for pptx in sorted(assert_dir.glob("*.pptx")):
        with zipfile.ZipFile(pptx, "r") as z:
            media = sorted([n for n in z.namelist() if n.startswith("ppt/media/")])
            for name in media:
                data = z.read(name)
                ext = Path(name).suffix.lower()
                width = height = 0
                mode = ""
                has_alpha = False
                if ext in {".png", ".jpg", ".jpeg", ".gif"}:
                    try:
                        img = Image.open(io.BytesIO(data))
                    except Exception:
                        img = None
                    if img is not None:
                        width, height = img.size
                        mode = img.mode
                        has_alpha = "A" in img.mode
                items.append(
                    {
                        "pptx": str(pptx),
                        "media_path": name,
                        "filename": Path(name).name,
                        "ext": ext,
                        "bytes": len(data),
                        "width": width,
                        "height": height,
                        "mode": mode,
                        "has_alpha": has_alpha,
                    }
                )

    curated_rows = [
        {
            "role": a.role,
            "source_pptx": str(a.source_pptx),
            "media_path": a.media_path,
            "derived_path": str(a.derived_path),
        }
        for a in curated
    ]
    out_json.parent.mkdir(parents=True, exist_ok=True)
    out_json.write_text(
        json.dumps(
            {"assert_dir": str(assert_dir), "items": items, "curated": curated_rows},
            indent=2,
            sort_keys=True,
        )
        + "\n",
        encoding="utf-8",
    )

    md_lines = [
        "# /Slides/assert inventory",
        "",
        "## Curated asset set (Blueprint + Chess)",
        "",
        "| role | source pptx | media | derived |",
        "|---|---|---|---|",
    ]
    for row in curated_rows:
        md_lines.append(
            f"| {row['role']} | `{Path(str(row['source_pptx'])).name}` | `{row['media_path']}` | `{Path(str(row['derived_path'])).name}` |"
        )
    md_lines += [
        "",
        "## Full media listing (flattened)",
        "",
        "This is a machine-generated index of `ppt/media/*` across all `/Slides/assert/*.pptx`.",
        "",
        "| pptx | media | type | px | alpha | bytes |",
        "|---|---|---:|---:|---:|---:|",
    ]
    for it in items:
        px = f"{it['width']}×{it['height']}" if it["width"] and it["height"] else ""
        md_lines.append(
            f"| `{Path(str(it['pptx'])).name}` | `{it['media_path']}` | `{it['ext']}` | {px} | {str(it['has_alpha']).lower()} | {it['bytes']} |"
        )
    out_md.write_text("\n".join(md_lines) + "\n", encoding="utf-8")


def _record_asset(
    manifest: dict[int, list[AssertAssetRef]], slide_num: int, asset: AssertAssetRef
) -> None:
    assets = manifest.setdefault(slide_num, [])
    key = (str(asset.source_pptx), asset.media_path, asset.role, str(asset.derived_path))
    for existing in assets:
        if (str(existing.source_pptx), existing.media_path, existing.role, str(existing.derived_path)) == key:
            return
    assets.append(asset)


def _slide_title(slide) -> str:
    try:
        title_shape = _shape(slide, 2)
    except KeyError:
        return ""
    if not title_shape.has_text_frame:
        return ""
    return title_shape.text_frame.text.strip()


def _write_slide_index(
    *,
    prs: Presentation,
    out_md: Path,
    out_json: Path | None,
    goals: dict[int, str],
) -> None:
    lines: list[str] = ["# Story1 slide index", ""]
    items: list[dict[str, object]] = []
    for i, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide) or f"Slide {i}"
        goal = goals.get(i, "")
        lines.append(f"{i:02d}. {title} — {goal}".rstrip())
        items.append({"slide": i, "title": title, "goal": goal})
    out_md.write_text("\n".join(lines) + "\n", encoding="utf-8")
    if out_json is not None:
        out_json.write_text(json.dumps(items, indent=2) + "\n", encoding="utf-8")


def _write_asset_manifest(*, entries: list[SlideManifestEntry], out_json: Path) -> None:
    payload: list[dict[str, object]] = []
    for e in entries:
        payload.append(
            {
                "slide": e.slide_num,
                "title": e.title,
                "assets": [
                    {
                        "role": a.role,
                        "source_pptx": str(a.source_pptx),
                        "media_path": a.media_path,
                        "derived_path": str(a.derived_path),
                    }
                    for a in e.assets
                ],
            }
        )
    out_json.write_text(json.dumps(payload, indent=2) + "\n", encoding="utf-8")


def _render_reference_block(
    *,
    out_path: Path,
    lines: Iterable[str],
    width_px: int = 1600,
    height_px: int = 1100,
    font_size: int = 28,
) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGBA", (width_px, height_px), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    font = ImageFont.truetype(str(FONT_MONO), font_size)

    x = 20
    y = 10
    line_h = int(font_size * 1.35)
    for line in lines:
        draw.text((x, y), line, fill=f"#{COLORS['ink']}", font=font)
        y += line_h
        if y > height_px - line_h:
            break
    img.save(out_path)


def _replace_picture_image(pic_shape, slide, image_path: Path) -> None:
    # Replace the picture's r:embed to point at a new image relationship.
    image_part, rid = slide.part.get_or_add_image_part(str(image_path))
    blip = pic_shape._element.xpath(".//a:blip")[0]  # noqa: SLF001 (pptx oxml)
    blip.set(f"{{{R_NS}}}embed", rid)


def _make_attack_surface_diagram(out_path: Path) -> None:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    w, h = 1200, 1500
    img = Image.new("RGBA", (w, h), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)

    font_title = ImageFont.truetype(str(FONT_SANS), 36)
    font = ImageFont.truetype(str(FONT_SANS), 26)
    font_small = ImageFont.truetype(str(FONT_SANS), 22)

    def box(x0: int, y0: int, x1: int, y1: int, label: str, border: str) -> None:
        r = 28
        draw.rounded_rectangle(
            (x0, y0, x1, y1),
            radius=r,
            outline=f"#{border}",
            width=5,
            fill=f"#{COLORS['card']}",
        )
        tw = draw.textlength(label, font=font)
        draw.text(((x0 + x1 - tw) / 2, y0 + 18), label, fill=f"#{COLORS['ink']}", font=font)

    def arrow(x0: int, y0: int, x1: int, y1: int) -> None:
        draw.line((x0, y0, x1, y1), fill=f"#{COLORS['muted2']}", width=6)
        # arrowhead
        ang = math.atan2(y1 - y0, x1 - x0)
        size = 20
        p1 = (x1, y1)
        p2 = (x1 - size * math.cos(ang - 0.5), y1 - size * math.sin(ang - 0.5))
        p3 = (x1 - size * math.cos(ang + 0.5), y1 - size * math.sin(ang + 0.5))
        draw.polygon([p1, p2, p3], fill=f"#{COLORS['muted2']}")

    draw.text((40, 20), "Attack surface (marketplace)", fill=f"#{COLORS['ink']}", font=font_title)

    box(90, 140, 1110, 320, "Discovery surfaces", COLORS["pink"])
    draw.text((140, 230), "starter packs, popular lists", fill=f"#{COLORS['muted']}", font=font_small)

    box(90, 420, 1110, 600, "Feed generators (rankers)", COLORS["cyan"])
    draw.text((140, 510), "3rd-party ranking policies", fill=f"#{COLORS['muted']}", font=font_small)

    box(90, 700, 1110, 880, "Exposure outcomes", COLORS["purple"])
    draw.text((140, 790), "who gets seen (and how often)", fill=f"#{COLORS['muted']}", font=font_small)

    box(90, 980, 1110, 1160, "Labelers (safety signals)", COLORS["green"])
    draw.text((140, 1070), "labels/actions vary by provider", fill=f"#{COLORS['muted']}", font=font_small)

    arrow(600, 320, 600, 420)
    arrow(600, 600, 600, 700)
    arrow(600, 880, 600, 980)

    # Attacker callouts
    draw.rounded_rectangle(
        (60, 1230, 1140, 1450),
        radius=28,
        outline=f"#{COLORS['amber']}",
        width=5,
        fill=(0, 0, 0, 0),
    )
    draw.text((90, 1255), "Adversaries:", fill=f"#{COLORS['ink']}", font=font)
    lines = [
        "• discovery capture / poisoning",
        "• growth-gaming feeds + sybils",
        "• label evasion + coverage gaps",
    ]
    yy = 1300
    for line in lines:
        draw.text((120, yy), line, fill=f"#{COLORS['ink']}", font=font_small)
        yy += 42

    img.save(out_path)


def apply_story_content(
    *,
    paths: BuildPaths,
    prepared: PreparedAssets,
    out_manifest: Path,
    out_slide_index_md: Path,
    out_slide_index_json: Path | None,
) -> None:
    prs = Presentation(str(paths.working_pptx))
    manifest: dict[int, list[AssertAssetRef]] = {}

    # Global structural asset: blueprint overlay background on every slide.
    for slide_num, slide in enumerate(prs.slides, start=1):
        _add_fullslide_background(slide, prs, prepared.bg_blueprint_overlay.derived_path)
        _record_asset(manifest, slide_num, prepared.bg_blueprint_overlay)

    plots = {
        "h1": paths.plots_dir / "h1_discovery_lorenz.png",
        "h2": paths.plots_dir / "h2_provider_leverage.png",
        "h3": paths.plots_dir / "h3_exposure_lorenz.png",
        "h4": paths.plots_dir / "h4_overlap_box.png",
        "h5": paths.plots_dir / "h5_deciles.png",
        "h6": paths.plots_dir / "h6_label_variability.png",
    }

    metrics_path = paths.plots_dir / "metrics.json"
    metrics = json.loads(metrics_path.read_text(encoding="utf-8")) if metrics_path.exists() else {}

    def pct(key: str) -> str:
        v = float(metrics[key])
        return f"{round(v * 100):d}%"

    def times_int(key: str) -> str:
        v = float(metrics[key])
        return f"×{round(v):d}"

    # Slide 1 — Title
    s1 = prs.slides[0]
    _set_run_text(_shape(s1, 2), "Algorithmic Choice in Bluesky")
    _set_run_text(_shape(s1, 3), "Responsible ranking + security in an open feed ecosystem")
    _set_run_text(_shape(s1, 4), "Measured: discovery → panel → 600,651 ranked impressions (Feb 1–2, 2026)")
    _add_icon_watermark(slide=s1, prs=prs, image_path=prepared.icon_king_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.0)
    _record_asset(manifest, 1, prepared.icon_king_wm)

    # Slide 2 — Hook
    s2 = prs.slides[1]
    _set_run_text(_shape(s2, 2), "Algorithmic choice ≠ decentralized power")
    _set_run_text(_shape(s2, 3), "• Feeds + labelers turn ranking/moderation into a marketplace.")
    _set_run_text(_shape(s2, 4), "• Discovery + hosting create chokepoints + new attack surfaces.")
    _hide_shape(_shape(s2, 5))
    _set_run_text(_shape(s2, 6), "Thesis: pluralism shifts power to discovery + infrastructure.")
    _add_icon_watermark(slide=s2, prs=prs, image_path=prepared.icon_queen_or_rook_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.0)
    _record_asset(manifest, 2, prepared.icon_queen_or_rook_wm)

    # Slide 3 — System background (repurpose “Where requests go”)
    s3 = prs.slides[2]
    _set_run_text(_shape(s3, 2), "Bluesky’s open feed ecosystem")
    _set_run_text(_shape(s3, 3), "Ranking + moderation become third‑party services (a marketplace).")
    _set_run_text(_shape(s3, 5), "Client")
    _set_run_text(_shape(s3, 7), "Discovery")
    _set_run_text(_shape(s3, 8), "Starter packs + popular lists")
    _set_run_text(_shape(s3, 10), "Feed generators")
    _set_run_text(_shape(s3, 11), "Third‑party ranking")
    _set_run_text(_shape(s3, 13), "Labelers")
    _set_run_text(_shape(s3, 14), "Safety labels + actions")
    _set_run_text(_shape(s3, 18), "Chokepoints emerge where users find feeds + who hosts them.  [ICWSM’24 10.1609/icwsm.v18i1.31293]")
    # Faint chessboard watermark behind nodes.
    _add_picture_icon(
        slide=s3,
        image_path=prepared.icon_chessboard_wm.derived_path,
        left_in=4.7,
        top_in=2.1,
        size_in=4.7,
        send_to_back=True,
    )
    _record_asset(manifest, 3, prepared.icon_chessboard_wm)

    # Slide 4 — Discovery surfaces
    s4 = prs.slides[3]
    _set_run_text(_shape(s4, 2), "Discovery is measurable (not anecdotes)")
    _set_run_text(_shape(s4, 3), "Starter packs + popular feeds are explicit lists that seed attention.")
    _set_run_text(_shape(s4, 5), "Starter packs (onboarding bundles)")
    _set_run_text(_shape(s4, 9), "Popular feeds (in‑app list)")
    _set_run_text(_shape(s4, 11), "These surfaces seed attention → panel sampling.  [ICWSM’25 10.1609/icwsm.v19i1.35810]")
    # Pawn icon on starter packs card.
    left_card = _shape(s4, 4)
    _add_picture_icon(
        slide=s4,
        image_path=prepared.icon_pawn.derived_path,
        left_in=(left_card.left + left_card.width - _emu(0.75)) / EMU_PER_INCH,
        top_in=(left_card.top + _emu(0.25)) / EMU_PER_INCH,
        size_in=0.55,
        send_to_back=False,
    )
    _record_asset(manifest, 4, prepared.icon_pawn)

    # Slide 5 — What we measured (H1–H6)
    s5 = prs.slides[4]
    _set_run_text(_shape(s5, 2), "What we measured (H1–H6)")
    _set_run_text(_shape(s5, 3), "Discovery → provider leverage → exposure → overlap → rich-get-richer → safety labels")
    _set_run_text(_shape(s5, 9), "Which providers get outsized discovery?")
    _set_run_text(_shape(s5, 12), "How concentrated is exposure?")
    _set_run_text(_shape(s5, 15), "Do feeds show the same posts?")
    _set_run_text(_shape(s5, 19), "H6 Risk exposure variability (labels)")
    # Add small semantic chess icons to each H card.
    card_specs = [
        (4, prepared.icon_pawn),
        (7, prepared.icon_rook_or_bishop),
        (10, prepared.icon_knight),
        (13, prepared.icon_chessboard),
        (16, prepared.icon_queen_or_rook),
    ]
    for box_id, asset in card_specs:
        box = _shape(s5, box_id)
        left_in = (box.left + box.width - _emu(0.75)) / EMU_PER_INCH
        top_in = (box.top + _emu(0.20)) / EMU_PER_INCH
        _add_picture_icon(
            slide=s5,
            image_path=asset.derived_path,
            left_in=left_in,
            top_in=top_in,
            size_in=0.5,
            send_to_back=False,
        )
        _record_asset(manifest, 5, asset)
    # H6 icon by the footer line.
    footer = _shape(s5, 19)
    _add_picture_icon(
        slide=s5,
        image_path=prepared.icon_king.derived_path,
        left_in=(footer.left + footer.width - _emu(0.75)) / EMU_PER_INCH,
        top_in=(footer.top + _emu(0.05)) / EMU_PER_INCH,
        size_in=0.45,
        send_to_back=False,
    )
    _record_asset(manifest, 5, prepared.icon_king)

    # Slide 6 — Responsible ranking lens (dup bullet)
    s6 = prs.slides[5]
    _set_run_text(_shape(s6, 2), "Problem: responsible ranking meets open markets")
    _set_run_text(_shape(s6, 3), "• Users choose algorithms — but discovery chooses defaults.")
    _set_run_text(_shape(s6, 4), "• Responsible ranking needs: security, fairness, safety, accountability.")
    _hide_shape(_shape(s6, 5))
    _set_run_text(_shape(s6, 6), "")
    _add_icon_watermark(slide=s6, prs=prs, image_path=prepared.icon_pawn_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.1)
    _record_asset(manifest, 6, prepared.icon_pawn_wm)

    # Slide 7 — Threat model (dup bullet)
    s7 = prs.slides[6]
    _set_run_text(_shape(s7, 2), "Threat model: marketplace adversaries")
    _set_run_text(_shape(s7, 3), "• Malicious feeds/rankers can steer exposure or surface explicit content.")
    _set_run_text(_shape(s7, 4), "• Growth gaming + label gaps enable evasion, bias, and targeting.")
    _hide_shape(_shape(s7, 5))
    _set_run_text(
        _shape(s7, 6),
        "[NDSS’21 10.14722/ndss.2021.24525; USENIX Sec’23 10.5555/3620237.3620333]",
    )
    _add_icon_watermark(slide=s7, prs=prs, image_path=prepared.icon_knight_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.1)
    _record_asset(manifest, 7, prepared.icon_knight_wm)

    # Slide 8 — Contributions (dup bullet)
    s8 = prs.slides[7]
    _set_run_text(_shape(s8, 2), "Our contributions (S&P tone)")
    _set_run_text(_shape(s8, 3), "• Threat model for ranking+moderation marketplaces.")
    _set_run_text(_shape(s8, 4), "• Reproducible measurement + evidence + mitigations (H1–H6).")
    _hide_shape(_shape(s8, 5))
    _set_run_text(_shape(s8, 6), "")
    _add_icon_watermark(slide=s8, prs=prs, image_path=prepared.icon_chessboard_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.1)
    _record_asset(manifest, 8, prepared.icon_chessboard_wm)

    # Slide 9 — Methods overview
    s9 = prs.slides[8]
    _set_run_text(_shape(s9, 2), "Methods credibility (read‑only snapshot)")
    _set_run_text(_shape(s9, 3), "Discovery → fixed panel → ranked impressions → joinable tables")
    _set_run_text(
        _shape(s9, 38),
        "Read-only XRPC to Relay/AppView (only POST is session auth/refresh).",
    )
    # Hide steps 5–7 (bottom row) to keep the credibility segment to 2 slides.
    for sid in list(range(20, 32)) + [35, 36, 37]:
        _hide_shape(_shape(s9, sid))
    _add_icon_watermark(slide=s9, prs=prs, image_path=prepared.icon_rook_or_bishop_wm.derived_path, size_in=2.8, right_in=0.8, top_in=4.2)
    _record_asset(manifest, 9, prepared.icon_rook_or_bishop_wm)

    # Slide 10 — Run receipts
    s10 = prs.slides[9]
    _set_run_text(_shape(s10, 11), "Captured: 600,651 ranked impressions")
    _set_run_text(_shape(s10, 13), "Every stage leaves a trace (exports)")
    _set_run_text(_shape(s10, 21), "Keys line up: feed_uri / post_cid / author_did")
    _add_icon_watermark(slide=s10, prs=prs, image_path=prepared.icon_chessboard_wm.derived_path, size_in=2.8, right_in=0.8, top_in=4.2)
    _record_asset(manifest, 10, prepared.icon_chessboard_wm)

    # Slide 11 — Adversary playbook (repurpose joinability slide)
    s11 = prs.slides[10]
    _set_run_text(_shape(s11, 2), "Adversary playbook: 4 ways markets get gamed")
    _set_run_text(_shape(s11, 3), "Market incentives create predictable attack patterns.")
    quad = [
        (4, 8, 12, prepared.icon_pawn, "Discovery poisoning", "Capture/poison starter packs + popular lists."),
        (5, 9, 13, prepared.icon_knight, "Growth gaming", "Sybil + engagement pumps to climb ranks."),
        (6, 10, 14, prepared.icon_queen_or_rook, "Malicious rankers", "3rd-party policies steer exposure by design."),
        (7, 11, 15, prepared.icon_king, "Label gaps/evasion", "Coverage gaps + inconsistent semantics."),
    ]
    for box_id, head_id, pic_id, asset, head, body in quad:
        box = _shape(s11, box_id)
        if box.has_text_frame:
            box.text_frame.margin_top = _emu(0.55)
            box.text_frame.margin_left = _emu(0.35)
        _set_run_text(box, body)
        _set_run_text(_shape(s11, head_id), head)
        _replace_picture_image(_shape(s11, pic_id), s11, asset.derived_path)
        _record_asset(manifest, 11, asset)
    # Remove arrows that imply a pipeline.
    for sid in (16, 17):
        _hide_shape(_shape(s11, sid))

    # Slides 12–17 — H1–H6 results (stage-slide template)
    result_specs = [
        (
            11,
            "H1 — Discovery slots are concentrated",
            "Top 1% share of starter-pack slots",
            pct("h1_top_1pct_share"),
            "Defaults become gatekeepers → cheaper discovery capture.",
            plots["h1"],
            COLORS["pink"],
            prepared.icon_pawn,
            "[ICWSM’25 10.1609/icwsm.v19i1.35810]",
        ),
        (
            12,
            "H2 — Discovery creates provider leverage",
            "Max leverage: discovery share vs hosting footprint",
            times_int("h2_max_leverage_ratio"),
            "Hosting chokepoints become targets → ecosystem-wide risk.",
            plots["h2"],
            COLORS["cyan"],
            prepared.icon_rook_or_bishop,
            "",
        ),
        (
            13,
            "H3 — Exposure is concentrated among authors",
            "Top 10% of authors get this share of impressions",
            pct("h3_top_10pct_share"),
            "Concentrated exposure lowers manipulation/targeting cost.",
            plots["h3"],
            COLORS["purple"],
            prepared.icon_knight,
            "",
        ),
        (
            14,
            "H4 — Many feeds show the same posts",
            "Share of impressions for posts in ≥2 feeds",
            pct("h4_share_impressions_post_in_ge2_feeds"),
            "Overlap amplifies winners → pluralism without diversity.",
            plots["h4"],
            COLORS["amber"],
            prepared.icon_chessboard,
            "",
        ),
        (
            15,
            "H5 — Bigger accounts get outsized exposure",
            "Top follower decile gets this share of impressions",
            pct("h5_top_decile_share"),
            "Popularity bias reappears → unfair + gameable outcomes.",
            plots["h5"],
            COLORS["green"],
            prepared.icon_queen_or_rook,
            "[KDD’21 10.1145/3447548.3467376]",
        ),
        (
            16,
            "H6 — Safety label exposure varies by feed",
            "Feeds (≥100 impressions) with any adult labels",
            pct("h6_share_feeds_with_any_adult_labels"),
            "Safety boundaries depend on feed choice → uneven risk.",
            plots["h6"],
            COLORS["pink"],
            prepared.icon_king,
            "[NAACL’22 10.18653/v1/2022.naacl-main.431]",
        ),
    ]

    for slide_idx, title, claim, stat, why, plot_path, border, icon_asset, foot in result_specs:
        sl = prs.slides[slide_idx]
        slide_num = slide_idx + 1
        _set_run_text(_shape(sl, 2), title)

        frame = _shape(sl, 3)
        plot = _shape(sl, 6)
        icon = _shape(sl, 5)
        _clear_text(plot)
        _clear_text(icon)

        _set_border_color(frame, border)
        _fill_shape_with_picture(frame, sl, prepared.bg_blueprint_card.derived_path)
        _record_asset(manifest, slide_num, prepared.bg_blueprint_card)

        # Plot appears later in the existing animation order by repurposing shape 6.
        margin = _emu(0.15)
        plot.left = frame.left + margin
        plot.top = frame.top + margin
        plot.width = frame.width - 2 * margin
        plot.height = frame.height - 2 * margin
        _fill_shape_with_picture(plot, sl, plot_path)
        _set_shape_line_nofill(plot)

        # Semantic icon on the left column.
        icon.left = _emu(1.0)
        icon.top = _emu(2.25)
        icon.width = _emu(0.8)
        icon.height = _emu(0.8)
        _fill_shape_with_picture(icon, sl, icon_asset.derived_path)
        _set_shape_line_nofill(icon)
        _record_asset(manifest, slide_num, icon_asset)

        # Left column text.
        claim_shape = _shape(sl, 4)
        stat_shape = _shape(sl, 7)
        why_shape = _shape(sl, 10)

        _set_run_text(claim_shape, claim)
        _set_run_font_size(claim_shape, 28)

        stat_shape.left = _emu(2.0)
        stat_shape.top = _emu(2.20)
        stat_shape.width = _emu(5.0)
        stat_shape.height = _emu(0.9)
        _set_run_text(stat_shape, stat)
        _set_run_font_size(stat_shape, 72)
        if stat_shape.has_text_frame and stat_shape.text_frame.paragraphs and stat_shape.text_frame.paragraphs[0].runs:
            stat_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = _rgb(border)

        why_shape.top = _emu(3.35)
        _set_run_text(why_shape, why)
        _set_run_font_size(why_shape, 26)

        # Remove unused scaffolding shapes so nothing overlaps the plot.
        for sid in (8, 9, 11, 12):
            _hide_shape(_shape(sl, sid))

        _set_run_text(_shape(sl, 13), foot)

    # Slide 18 — Security implications (stage slide)
    s18 = prs.slides[17]
    diagram = paths.assets_dir / "attack_surface.png"
    _make_attack_surface_diagram(diagram)
    _set_run_text(_shape(s18, 2), "Security implications: a new attack surface")
    _set_border_color(_shape(s18, 3), COLORS["amber"])
    _fill_shape_with_picture(_shape(s18, 3), s18, diagram)
    _set_run_text(_shape(s18, 4), "• Concentrated discovery reduces attack cost.")
    _set_run_text(_shape(s18, 7), "• 3rd-party rankers can be adversarial → marketplace-wide audit + guardrails.")
    _hide_shape(_shape(s18, 10))
    for sid in (5, 6, 8, 9, 11, 12):
        _clear_text(_shape(s18, sid))
    _set_run_text(
        _shape(s18, 13),
        "[NDSS’21 10.14722/ndss.2021.24525; CCS’21 10.1145/3460120.3484770]",
    )
    _add_icon_watermark(slide=s18, prs=prs, image_path=prepared.icon_rook_or_bishop_wm.derived_path, size_in=2.8, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 18, prepared.icon_rook_or_bishop_wm)

    # Slide 19 — Privacy (repurpose XRPC slide)
    s19 = prs.slides[18]
    _set_run_text(_shape(s19, 2), "Privacy: third-party services can profile")
    _set_run_text(_shape(s19, 3), "Request patterns + outputs can reveal preferences and identity.")
    _set_run_text(_shape(s19, 6), "https://…/xrpc/app.bsky.feed.getFeed?feed=at://…&cursor=…")
    _set_run_text(_shape(s19, 8), "What a service sees")
    _set_run_text(
        _shape(s19, 9),
        "• IP + timing + feed URI (and viewer DID if auth)\n• repeated queries + outputs reveal preferences",
    )
    _set_run_text(_shape(s19, 11), "Why it matters")
    _set_run_text(
        _shape(s19, 12),
        "• profiling / tracking + output leakage\n• cross-service linkage risk",
    )
    _set_run_text(_shape(s19, 13), "[CCS’21 10.1145/3460120.3484770]")
    _add_icon_watermark(slide=s19, prs=prs, image_path=prepared.icon_satellite_wm.derived_path, size_in=2.9, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 19, prepared.icon_satellite_wm)

    # Slide 20 — Moderation fragility (dup bullet)
    s20 = prs.slides[19]
    _set_run_text(_shape(s20, 2), "Implication: decentralized moderation is fragile")
    _set_run_text(_shape(s20, 3), "• Label semantics are subjective → disagreement across labelers.")
    _set_run_text(_shape(s20, 4), "• Feed choice can bypass/amplify safety boundaries (coverage gaps + bias).")
    _hide_shape(_shape(s20, 5))
    _set_run_text(
        _shape(s20, 6),
        "[NAACL’22 10.18653/v1/2022.naacl-main.431; Findings ACL’22 10.18653/v1/2022.findings-acl.176]",
    )
    _add_icon_watermark(slide=s20, prs=prs, image_path=prepared.icon_fragile_rot.derived_path, size_in=3.2, right_in=0.9, top_in=3.9)
    _record_asset(manifest, 20, prepared.icon_fragile_rot)

    # Slide 21 — Mitigation 1 (dup bullet)
    s21 = prs.slides[20]
    _set_run_text(_shape(s21, 2), "Mitigation: discovery transparency + diversification")
    _set_run_text(_shape(s21, 3), "• Publish logs for starter packs + popular lists.")
    _set_run_text(_shape(s21, 4), "• Diversify defaults + detect manipulation (growth + overlap spikes).")
    _hide_shape(_shape(s21, 5))
    _set_run_text(_shape(s21, 6), "")
    _add_icon_watermark(slide=s21, prs=prs, image_path=prepared.icon_pawn_wm.derived_path, size_in=2.8, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 21, prepared.icon_pawn_wm)

    # Slide 22 — Mitigation 2 (repurpose resumability slide)
    s22 = prs.slides[21]
    _set_run_text(_shape(s22, 2), "Mitigation: auditable feed marketplace")
    _set_run_text(_shape(s22, 3), "Logs + reproducible snapshots enable accountability")
    _set_run_text(_shape(s22, 5), "Add audit hooks around discovery, ranking, and labelers.")
    _set_run_text(_shape(s22, 6), "Expose stable identities + policies so clients can warn and researchers can audit.")
    _set_run_text(_shape(s22, 8), "Transparency logs")
    _set_run_text(_shape(s22, 10), "Audit APIs")
    _set_run_text(_shape(s22, 12), "Reproducible snapshots")
    _set_run_text(_shape(s22, 14), "Why this matters")
    _set_run_text(_shape(s22, 15), "Choice only helps if users can compare risk, diversity, and trust.")
    _set_run_text(_shape(s22, 17), "What to disclose (minimal)")
    _set_run_text(_shape(s22, 18), "• provider identity + hosting footprint")
    _set_run_text(_shape(s22, 19), "• discovery surfaces (who promoted what)")
    _set_run_text(_shape(s22, 20), "• ranking output summaries (exposure / overlap)")
    _set_run_text(_shape(s22, 21), "• labeler policies + coverage")
    _set_run_text(_shape(s22, 22), "• incident response + takedown hooks")
    _set_run_text(_shape(s22, 23), "")
    _add_icon_watermark(slide=s22, prs=prs, image_path=prepared.icon_chessboard_wm.derived_path, size_in=2.8, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 22, prepared.icon_chessboard_wm)

    # Slide 23 — Mitigation 3 (dup bullet)
    s23 = prs.slides[22]
    _set_run_text(_shape(s23, 2), "Mitigation: “feed nutrition labels” (user-facing)")
    _set_run_text(_shape(s23, 3), "• Show provider, overlap, and exposure concentration at a glance.")
    _set_run_text(_shape(s23, 4), "• Show label exposure + applied labelers (warnings/filters/defaults).")
    _hide_shape(_shape(s23, 5))
    _set_run_text(_shape(s23, 6), "")
    _add_icon_watermark(slide=s23, prs=prs, image_path=prepared.icon_brain_wm.derived_path, size_in=2.7, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 23, prepared.icon_brain_wm)

    # Slide 24 — Mitigation 4 (dup bullet)
    s24 = prs.slides[23]
    _set_run_text(_shape(s24, 2), "Mitigation: labeler interoperability + secure sharing")
    _set_run_text(_shape(s24, 3), "• Standardize vocabularies + transparency for labeler policies.")
    _set_run_text(_shape(s24, 4), "• Secure/private sharing + multi-labeler consensus to resist poisoning.")
    _hide_shape(_shape(s24, 5))
    _set_run_text(_shape(s24, 6), "[IEEE S&P’22 10.1109/SP46214.2022.9833647]")
    _add_icon_watermark(slide=s24, prs=prs, image_path=prepared.icon_satellite_wm.derived_path, size_in=2.9, right_in=0.8, top_in=4.1)
    _record_asset(manifest, 24, prepared.icon_satellite_wm)

    # Slide 25 — Closing (dup bullet)
    s25 = prs.slides[24]
    _set_run_text(_shape(s25, 2), "Takeaways")
    _set_run_text(_shape(s25, 3), "• Algorithmic choice shifts power to discovery + infrastructure.")
    _set_run_text(_shape(s25, 4), "• Open markets create new adversarial surfaces + inequalities.")
    _hide_shape(_shape(s25, 5))
    _set_run_text(_shape(s25, 6), "Questions?")
    _add_icon_watermark(slide=s25, prs=prs, image_path=prepared.icon_chessboard_wm.derived_path, size_in=3.0, right_in=0.7, top_in=4.0)
    _record_asset(manifest, 25, prepared.icon_chessboard_wm)

    # Slide 26 — References (repurpose reproducible outputs slide)
    s26 = prs.slides[25]
    _set_run_text(_shape(s26, 2), "References (venue-filtered, 2021–2025)")
    _set_run_text(_shape(s26, 3), "CCS / IEEE S&P / NDSS / USENIX Sec / WWW / CSCW / ICWSM / ACL / KDD")
    _set_run_text(_shape(s26, 5), "Decentralization + discovery")
    _set_run_text(_shape(s26, 8), "Ranking + security + labels")

    # Resize picture placeholders to fill the cards.
    pic_left = _shape(s26, 6)
    pic_right = _shape(s26, 9)
    for pic in (pic_left, pic_right):
        pic.height = int(3.9 * EMU_PER_INCH)

    left_refs = [
        "ICWSM’25 Starter Packs",
        "https://doi.org/10.1609/icwsm.v19i1.35810",
        "",
        "ICWSM’24 Decentralised Moderation",
        "https://doi.org/10.1609/icwsm.v18i1.31293",
        "",
        "ICWSM’24 Group Sanctions",
        "https://doi.org/10.1609/icwsm.v18i1.31316",
        "",
        "WWW’23 Will Admins Cope?",
        "https://doi.org/10.1145/3543507.3583487",
        "",
        "CSCW’24 Trouble in Paradise",
        "https://doi.org/10.1145/3687059",
    ]
    right_refs = [
        "NDSS’21 Poisoning Recommenders",
        "https://doi.org/10.14722/ndss.2021.24525",
        "",
        "USENIX’23 PORE",
        "https://doi.org/10.5555/3620237.3620333",
        "",
        "CCS’21 Membership Inference",
        "https://doi.org/10.1145/3460120.3484770",
        "",
        "IEEE S&P’22 FL Poisoning Eval",
        "https://doi.org/10.1109/SP46214.2022.9833647",
        "",
        "KDD’21 Popularity Bias (Dynamic)",
        "https://doi.org/10.1145/3447548.3467376",
        "",
        "NAACL’22 Annotators w/ Attitudes",
        "https://doi.org/10.18653/v1/2022.naacl-main.431",
        "",
        "Findings ACL’22 Fairness Varies",
        "https://doi.org/10.18653/v1/2022.findings-acl.176",
        "",
        "Findings EMNLP’23 Hate Bias (Causal)",
        "https://doi.org/10.18653/v1/2023.findings-emnlp.440",
    ]

    left_img = paths.assets_dir / "refs_left.png"
    right_img = paths.assets_dir / "refs_right.png"
    _render_reference_block(out_path=left_img, lines=left_refs)
    _render_reference_block(out_path=right_img, lines=right_refs)
    _replace_picture_image(pic_left, s26, left_img)
    _replace_picture_image(pic_right, s26, right_img)

    prs.save(str(paths.working_pptx))

    # Write manifest + slide index.
    entries: list[SlideManifestEntry] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide) or f"Slide {slide_num}"
        assets = manifest.get(slide_num, [])
        entries.append(SlideManifestEntry(slide_num=slide_num, title=title, assets=assets))

    # Validate asset usage requirements.
    for e in entries:
        if not e.assets:
            raise RuntimeError(f"slide {e.slide_num} missing assert assets")
    for n in range(12, 18):
        assets = manifest.get(n, [])
        roles = {a.role for a in assets}
        if len(assets) < 2 or not {"structural", "semantic"}.issubset(roles):
            raise RuntimeError(f"results slide {n} missing structural+semantic assert assets")

    _write_asset_manifest(entries=entries, out_json=out_manifest)

    goals = {
        1: "Hook + scale of measurement.",
        2: "Thesis: choice ≠ decentralized power.",
        3: "Map the feed marketplace + chokepoints.",
        4: "Show discovery is explicit + measurable.",
        5: "Preview the H1–H6 arc.",
        6: "Define the responsible ranking stakes.",
        7: "Define marketplace adversaries.",
        8: "State contributions (threat model + evidence).",
        9: "Methods credibility in 4 steps.",
        10: "Anchors + exports + joinability.",
        11: "Show 4 predictable gaming paths.",
        12: "H1: quantify discovery concentration.",
        13: "H2: quantify provider leverage.",
        14: "H3: quantify exposure concentration.",
        15: "H4: quantify feed overlap redundancy.",
        16: "H5: show popularity bias by author size.",
        17: "H6: show safety-label variability.",
        18: "Connect concentration to attack surface.",
        19: "Explain privacy leakage from services.",
        20: "Argue moderation fragility in markets.",
        21: "Mitigation: transparency + diversification.",
        22: "Mitigation: auditable marketplace hooks.",
        23: "Mitigation: feed nutrition labels.",
        24: "Mitigation: labeler interoperability.",
        25: "Crisp closing + Q&A.",
        26: "References for grounding.",
    }
    _write_slide_index(prs=prs, out_md=out_slide_index_md, out_json=out_slide_index_json, goals=goals)


def main() -> None:
    story_dir = Path("Slides/Story1")
    assert_dir = Path("Slides/assert")
    cache_dir = story_dir / "assets/assert_cache"
    prepared = prepare_assert_assets(assert_dir=assert_dir, cache_dir=cache_dir)
    curated = [
        prepared.bg_blueprint_overlay,
        prepared.bg_blueprint_card,
        prepared.icon_pawn,
        prepared.icon_pawn_wm,
        prepared.icon_knight,
        prepared.icon_knight_wm,
        prepared.icon_king,
        prepared.icon_king_wm,
        prepared.icon_rook_or_bishop,
        prepared.icon_rook_or_bishop_wm,
        prepared.icon_queen_or_rook,
        prepared.icon_queen_or_rook_wm,
        prepared.icon_chessboard,
        prepared.icon_chessboard_wm,
        prepared.icon_satellite,
        prepared.icon_satellite_wm,
        prepared.icon_brain,
        prepared.icon_brain_wm,
        prepared.icon_fragile_rot,
    ]
    build_assert_inventory(
        assert_dir=assert_dir,
        out_json=story_dir / "assets/assert_inventory.json",
        out_md=story_dir / "assets/assert_inventory.md",
        curated=curated,
    )

    paths = BuildPaths(
        template_pptx=Path("Slides/PPTXWT2/bluesky-data-collection-pipeline-blackboard-animated-run-20260201.pptx"),
        working_pptx=story_dir / "Story1_Working.pptx",
        plots_dir=story_dir / "assets/plots",
        assets_dir=story_dir / "assets/story1_gen",
    )

    build_working_structure(paths)
    apply_story_content(
        paths=paths,
        prepared=prepared,
        out_manifest=story_dir / "Story1_asset_manifest.json",
        out_slide_index_md=story_dir / "Story1_slide_index.md",
        out_slide_index_json=story_dir / "Story1_slide_index.json",
    )
    final_path = story_dir / "Story1_Final.pptx"
    shutil.copy2(paths.working_pptx, final_path)
    print(f"Wrote {paths.working_pptx} and {final_path}")


if __name__ == "__main__":
    main()
