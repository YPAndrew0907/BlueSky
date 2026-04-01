#!/usr/bin/env python3
from __future__ import annotations

import argparse
import sys
from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):  # group shapes
            for child in shape.shapes:
                yield child


def check_bounds(*, pptx_path: Path, margin_in: float) -> list[str]:
    prs = Presentation(str(pptx_path))
    margin = int(margin_in * EMU_PER_INCH)
    w = int(prs.slide_width)
    h = int(prs.slide_height)

    issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape in _iter_shapes(slide):
            left = int(getattr(shape, "left", 0))
            top = int(getattr(shape, "top", 0))
            width = int(getattr(shape, "width", 0))
            height = int(getattr(shape, "height", 0))
            right = left + width
            bottom = top + height

            if width <= 0 or height <= 0:
                continue

            if left < -margin or top < -margin or right > w + margin or bottom > h + margin:
                issues.append(
                    f"slide {slide_num:02d}: shape_id={shape.shape_id} "
                    f"bbox=({left},{top})-({right},{bottom}) slide=({w},{h})"
                )
    return issues


def main(argv: list[str]) -> int:
    parser = argparse.ArgumentParser(description="Basic PPTX bounds QC (no out-of-slide elements).")
    parser.add_argument("pptx", nargs="?", default="Slides/Story1/Story1_Final.pptx")
    parser.add_argument("--margin-in", type=float, default=0.05)
    args = parser.parse_args(argv)

    pptx_path = Path(args.pptx)
    if not pptx_path.exists():
        print(f"Missing pptx: {pptx_path}", file=sys.stderr)
        return 2

    issues = check_bounds(pptx_path=pptx_path, margin_in=args.margin_in)
    if issues:
        print("OUT-OF-BOUNDS SHAPES:")
        for line in issues:
            print("-", line)
        return 1

    print("OK: no out-of-bounds shapes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv[1:]))

