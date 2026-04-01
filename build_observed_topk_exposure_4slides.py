#!/usr/bin/env python3
from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

BG = "07111C"
GRID = "0F2438"
INK = "F4F8FD"
MUTED = "C5D3E4"
CARD = "102033"
CARD_2 = "0B1828"
LINE = "335575"
CYAN = "64C9FF"
AMBER = "E4BC67"
GREEN = "6BD59C"
CORAL = "D68A57"
VIOLET = "9B8BFF"
RED = "E67C73"

FONT_UI = "Trebuchet MS"
FONT_MONO = "Consolas"

PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


def _qn(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _style_text_frame(tf, *, ml: float = 0.02, mr: float = 0.02, mt: float = 0.01, mb: float = 0.01) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(ml)
    tf.margin_right = Inches(mr)
    tf.margin_top = Inches(mt)
    tf.margin_bottom = Inches(mb)


def _textbox(slide, *, x: float, y: float, w: float, h: float, name: str):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    _style_text_frame(shape.text_frame)
    return shape


def _set_text(
    shape,
    *,
    text: str,
    size: float,
    color: str = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str = FONT_UI,
) -> None:
    tf = shape.text_frame
    tf.clear()
    _style_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = _rgb(color)


def _set_lines(
    shape,
    *,
    lines: list[str],
    size: float,
    color: str = MUTED,
    bold_first: bool = False,
    font_name: str = FONT_UI,
    line_spacing: float = 1.15,
) -> None:
    tf = shape.text_frame
    tf.clear()
    _style_text_frame(tf, ml=0.06, mr=0.04, mt=0.03, mb=0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        font = run.font
        font.name = font_name
        font.size = Pt(size)
        font.bold = bool(bold_first and i == 0)
        font.color.rgb = _rgb(color)


def _card(slide, *, x: float, y: float, w: float, h: float, name: str, fill_hex: str = CARD, line_hex: str = LINE):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(line_hex)
    shape.line.width = Pt(1.75)
    return shape


def _set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(BG)

    grid = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(SLIDE_W_IN), Inches(SLIDE_H_IN)
    )
    grid.name = "_BG_GRID"
    grid.fill.solid()
    grid.fill.fore_color.rgb = _rgb(BG)
    grid.line.fill.background()

    for i in range(1, 27):
        x = i * SLIDE_W_IN / 26.0
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x), Inches(0), Inches(x), Inches(SLIDE_H_IN))
        line.name = f"_BG_V_{i:02d}"
        line.line.color.rgb = _rgb(GRID)
        line.line.width = Pt(0.7)
        line.line.transparency = 0.55

    for i in range(1, 15):
        y = i * SLIDE_H_IN / 14.0
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(y), Inches(SLIDE_W_IN), Inches(y))
        line.name = f"_BG_H_{i:02d}"
        line.line.color.rgb = _rgb(GRID)
        line.line.width = Pt(0.7)
        line.line.transparency = 0.55

    for x, y, w, h, color, alpha, name in [
        (0.55, 0.35, 3.1, 2.1, CYAN, 0.90, "_BG_GLOW_1"),
        (9.55, 4.75, 2.9, 1.9, CORAL, 0.88, "_BG_GLOW_2"),
    ]:
        glow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        glow.name = name
        glow.fill.solid()
        glow.fill.fore_color.rgb = _rgb(color)
        glow.fill.transparency = alpha
        glow.line.fill.background()


def _title_bar(slide, *, slide_idx: int, title: str, subtitle: str) -> None:
    _card(slide, x=0.62, y=0.30, w=12.15, h=0.95, name=f"S{slide_idx}_TITLE_BAR", fill_hex="112540")
    spine = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.30), Inches(0.12), Inches(0.95))
    spine.name = f"S{slide_idx}_TITLE_SPINE"
    spine.fill.solid()
    spine.fill.fore_color.rgb = _rgb(CYAN)
    spine.line.fill.background()
    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(0.80), Inches(0.12), Inches(0.45)
    )
    accent.name = f"S{slide_idx}_TITLE_SPINE_2"
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(CORAL)
    accent.line.fill.background()
    title_box = _textbox(slide, x=0.86, y=0.38, w=11.5, h=0.38, name=f"S{slide_idx}_TITLE")
    _set_text(title_box, text=title, size=30, bold=True)
    sub_box = _textbox(slide, x=0.86, y=0.80, w=11.0, h=0.23, name=f"S{slide_idx}_SUB")
    _set_text(sub_box, text=subtitle, size=14.5, color=MUTED, bold=True)


def _panel_label(slide, *, x: float, y: float, w: float, text: str, accent: str, name: str) -> None:
    _card(slide, x=x, y=y, w=w, h=0.34, name=name, fill_hex=CARD_2, line_hex=accent)
    dot = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 0.08), Inches(y + 0.09), Inches(0.10), Inches(0.10))
    dot.name = f"{name}_DOT"
    dot.fill.solid()
    dot.fill.fore_color.rgb = _rgb(accent)
    dot.line.fill.background()
    text_box = _textbox(slide, x=x + 0.24, y=y + 0.05, w=w - 0.30, h=0.20, name=f"{name}_TEXT")
    _set_text(text_box, text=text, size=13.5, color=INK, bold=True)


def _mini_chip(slide, *, x: float, y: float, w: float, h: float, head: str, val: str, accent: str, name: str) -> None:
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD_2, line_hex=LINE)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.name = f"{name}_STRIPE"
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(accent)
    stripe.fill.transparency = 0.15
    stripe.line.fill.background()
    head_box = _textbox(slide, x=x + 0.16, y=y + 0.07, w=w - 0.22, h=0.16, name=f"{name}_HEAD")
    _set_text(head_box, text=head, size=10.2, color=MUTED, bold=True)
    val_box = _textbox(slide, x=x + 0.16, y=y + 0.25, w=w - 0.22, h=0.18, name=f"{name}_VAL")
    _set_text(val_box, text=val, size=14, color=accent, bold=True, font_name=FONT_MONO)


def _code_card(slide, *, x: float, y: float, w: float, h: float, lines: list[str], name: str, accent: str) -> None:
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD_2, line_hex=accent)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    stripe.name = f"{name}_STRIPE"
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(accent)
    stripe.fill.transparency = 0.12
    stripe.line.fill.background()
    body = _textbox(slide, x=x + 0.20, y=y + 0.16, w=w - 0.30, h=h - 0.24, name=f"{name}_TEXT")
    tf = body.text_frame
    tf.clear()
    _style_text_frame(tf, ml=0.02, mr=0.02, mt=0.02, mb=0.02)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.05
        run = p.add_run()
        run.text = line
        font = run.font
        font.name = FONT_MONO
        font.size = Pt(13.5)
        font.bold = i == 0 and "=" in line
        font.color.rgb = _rgb(INK if i == 0 else MUTED)


def _rank_row(slide, *, x: float, y: float, w: float, rank: int, author: str, accent: str, name: str) -> None:
    h = 0.34
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD_2, line_hex=LINE)
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(0.58), Inches(h))
    badge.name = f"{name}_BADGE"
    badge.fill.solid()
    badge.fill.fore_color.rgb = _rgb(accent)
    badge.line.fill.background()
    badge_text = _textbox(slide, x=x + 0.16, y=y + 0.05, w=0.26, h=0.16, name=f"{name}_RANK")
    _set_text(badge_text, text=str(rank), size=12.5, color=BG, bold=True, align=PP_ALIGN.CENTER)
    label = _textbox(slide, x=x + 0.72, y=y + 0.05, w=w - 1.75, h=0.16, name=f"{name}_AUTHOR")
    _set_text(label, text=author, size=12.2, color=INK, bold=True)
    bar_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + w - 0.95), Inches(y + 0.11), Inches(0.70), Inches(0.10))
    bar_bg.name = f"{name}_BAR_BG"
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = _rgb(CARD)
    bar_bg.line.fill.background()
    frac = max(0.16, 1.0 / (1.0 + 0.18 * (rank - 1)))
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x + w - 0.95),
        Inches(y + 0.11),
        Inches(0.70 * frac),
        Inches(0.10),
    )
    bar.name = f"{name}_BAR"
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(accent)
    bar.line.fill.background()


def _stat_box(slide, *, x: float, y: float, w: float, h: float, title: str, body: list[str], accent: str, name: str) -> None:
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD, line_hex=accent)
    _panel_label(slide, x=x + 0.10, y=y + 0.08, w=min(w - 0.20, 2.50), text=title, accent=accent, name=f"{name}_LBL")
    box = _textbox(slide, x=x + 0.18, y=y + 0.48, w=w - 0.30, h=h - 0.58, name=f"{name}_BODY")
    _set_lines(box, lines=body, size=12.2, color=MUTED)


def _claim_column(
    slide, *, x: float, y: float, w: float, title: str, accent: str, bullets: list[str], name: str, fill_hex: str
) -> None:
    _card(slide, x=x, y=y, w=w, h=3.65, name=name, fill_hex=fill_hex, line_hex=accent)
    _panel_label(slide, x=x + 0.12, y=y + 0.12, w=2.35, text=title, accent=accent, name=f"{name}_LBL")
    body = _textbox(slide, x=x + 0.18, y=y + 0.62, w=w - 0.30, h=2.72, name=f"{name}_BODY")
    _set_lines(body, lines=[f"- {b}" for b in bullets], size=15, color=INK if fill_hex == CARD else MUTED, line_spacing=1.18)


def _title_names(slide_idx: int) -> list[str]:
    return [
        f"S{slide_idx}_TITLE_BAR",
        f"S{slide_idx}_TITLE_SPINE",
        f"S{slide_idx}_TITLE_SPINE_2",
        f"S{slide_idx}_TITLE",
        f"S{slide_idx}_SUB",
    ]


def _mini_chip_names(name: str) -> list[str]:
    return [name, f"{name}_STRIPE", f"{name}_HEAD", f"{name}_VAL"]


def _code_card_names(name: str) -> list[str]:
    return [name, f"{name}_STRIPE", f"{name}_TEXT"]


def _rank_row_names(name: str) -> list[str]:
    return [name, f"{name}_BADGE", f"{name}_RANK", f"{name}_AUTHOR", f"{name}_BAR_BG", f"{name}_BAR"]


def _stat_box_names(name: str) -> list[str]:
    return [name, f"{name}_LBL", f"{name}_LBL_DOT", f"{name}_LBL_TEXT", f"{name}_BODY"]


def _claim_column_names(name: str) -> list[str]:
    return [name, f"{name}_LBL", f"{name}_LBL_DOT", f"{name}_LBL_TEXT", f"{name}_BODY"]


def _post_card(
    slide,
    *,
    x: float,
    y: float,
    w: float,
    h: float,
    author: str,
    stamp: str,
    text: str,
    accent: str,
    name: str,
) -> None:
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD_2, line_hex=LINE)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    stripe.name = f"{name}_STRIPE"
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(accent)
    stripe.fill.transparency = 0.16
    stripe.line.fill.background()

    author_box = _textbox(slide, x=x + 0.18, y=y + 0.08, w=w - 1.72, h=0.18, name=f"{name}_AUTHOR")
    _set_text(author_box, text=author, size=12.2, color=INK, bold=True)
    stamp_box = _textbox(slide, x=x + w - 1.38, y=y + 0.08, w=1.10, h=0.18, name=f"{name}_TIME")
    _set_text(stamp_box, text=stamp, size=10.6, color=accent, bold=True, align=PP_ALIGN.RIGHT, font_name=FONT_MONO)

    divider = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x + 0.18),
        Inches(y + 0.33),
        Inches(x + w - 0.18),
        Inches(y + 0.33),
    )
    divider.name = f"{name}_DIV"
    divider.line.color.rgb = _rgb(LINE)
    divider.line.width = Pt(1.0)

    body = _textbox(slide, x=x + 0.18, y=y + 0.40, w=w - 0.32, h=h - 0.50, name=f"{name}_TEXT")
    _set_lines(body, lines=[text], size=12.0, color=MUTED, line_spacing=1.10)


def _post_card_names(name: str) -> list[str]:
    return [name, f"{name}_STRIPE", f"{name}_AUTHOR", f"{name}_TIME", f"{name}_DIV", f"{name}_TEXT"]


def _spec_card(
    slide, *, x: float, y: float, w: float, h: float, title: str, body: list[str], accent: str, name: str
) -> None:
    _card(slide, x=x, y=y, w=w, h=h, name=name, fill_hex=CARD_2, line_hex=accent)
    stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.08), Inches(h))
    stripe.name = f"{name}_STRIPE"
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _rgb(accent)
    stripe.fill.transparency = 0.12
    stripe.line.fill.background()
    head = _textbox(slide, x=x + 0.18, y=y + 0.08, w=w - 0.26, h=0.18, name=f"{name}_HEAD")
    _set_text(head, text=title, size=12.8, color=accent, bold=True)
    body_box = _textbox(slide, x=x + 0.18, y=y + 0.32, w=w - 0.26, h=h - 0.40, name=f"{name}_BODY")
    _set_lines(body_box, lines=body, size=11.4, color=MUTED, line_spacing=1.12)


def _spec_card_names(name: str) -> list[str]:
    return [name, f"{name}_STRIPE", f"{name}_HEAD", f"{name}_BODY"]


def _slide1(slide) -> list[list[str]]:
    _set_bg(slide)
    _title_bar(
        slide,
        slide_idx=1,
        title="Observed top-K opportunity, not platform-wide fairness",
        subtitle="Fixed viewer, monitored feed panel, sampled top-50, and explicit scope boundaries",
    )
    steps: list[list[str]] = [_title_names(1)]

    _card(slide, x=0.72, y=1.48, w=4.18, h=4.48, name="S1_HERO", fill_hex="0C1B2D", line_hex=CYAN)
    hero_stripe = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.72), Inches(1.48), Inches(0.11), Inches(4.48))
    hero_stripe.name = "S1_HERO_STRIPE"
    hero_stripe.fill.solid()
    hero_stripe.fill.fore_color.rgb = _rgb(CYAN)
    hero_stripe.fill.transparency = 0.10
    hero_stripe.line.fill.background()
    hero_head = _textbox(slide, x=0.98, y=1.76, w=3.60, h=0.62, name="S1_HERO_HEAD")
    _set_text(hero_head, text="Same content.\nDifferent authors.\nScarce display opportunity.", size=24, color=INK, bold=True)
    hero_body = _textbox(slide, x=0.98, y=2.70, w=3.46, h=1.56, name="S1_HERO_BODY")
    _set_lines(
        hero_body,
        lines=[
            "The clean object is not total platform fairness.",
            "It is how monitored top-K opportunity gets allocated when content is effectively held fixed.",
            "That is the level our panel can actually identify.",
        ],
        size=14.2,
        color=MUTED,
        line_spacing=1.16,
    )
    _mini_chip(slide, x=0.96, y=4.72, w=1.08, h=0.58, head="unit", val="post x q", accent=CYAN, name="S1_HERO_CHIP_1")
    _mini_chip(slide, x=2.12, y=4.72, w=1.16, h=0.58, head="viewer", val="fixed C", accent=GREEN, name="S1_HERO_CHIP_2")
    _mini_chip(slide, x=3.38, y=4.72, w=1.18, h=0.58, head="panel", val="top-50", accent=AMBER, name="S1_HERO_CHIP_3")
    hero_note = _textbox(slide, x=0.98, y=5.45, w=3.56, h=0.30, name="S1_HERO_NOTE")
    _set_text(hero_note, text="Observed opportunity is the defensible proxy.", size=13.0, color=CYAN, bold=True)
    steps.append(
        [
            "S1_HERO",
            "S1_HERO_STRIPE",
            "S1_HERO_HEAD",
            "S1_HERO_BODY",
            "S1_HERO_NOTE",
            *_mini_chip_names("S1_HERO_CHIP_1"),
            *_mini_chip_names("S1_HERO_CHIP_2"),
            *_mini_chip_names("S1_HERO_CHIP_3"),
        ]
    )

    _code_card(
        slide,
        x=5.24,
        y=1.58,
        w=7.26,
        h=1.34,
        lines=[
            "q = (feed_uri, bucket, snapshot_hour, viewer_mode, vantage_id)",
            "I_iq = 1{post i appears in the observed",
            "sampled top-K panel for q}",
            "O_iq = I_iq * w(rank_iq)",
        ],
        name="S1_QUERY",
        accent=CYAN,
    )
    _mini_chip(slide, x=5.34, y=3.06, w=1.42, h=0.56, head="feed", val="generator", accent=CYAN, name="S1_QCHIP_1")
    _mini_chip(slide, x=6.86, y=3.06, w=1.38, h=0.56, head="bucket", val="sample", accent=AMBER, name="S1_QCHIP_2")
    _mini_chip(slide, x=8.34, y=3.06, w=1.72, h=0.56, head="mode", val="auth/unauth", accent=GREEN, name="S1_QCHIP_3")
    _mini_chip(slide, x=10.16, y=3.06, w=2.02, h=0.56, head="hour", val="captured q", accent=VIOLET, name="S1_QCHIP_4")
    steps.append(
        [
            *_code_card_names("S1_QUERY"),
            *_mini_chip_names("S1_QCHIP_1"),
            *_mini_chip_names("S1_QCHIP_2"),
            *_mini_chip_names("S1_QCHIP_3"),
            *_mini_chip_names("S1_QCHIP_4"),
        ]
    )

    _card(slide, x=5.24, y=3.82, w=7.26, h=1.62, name="S1_PANEL", fill_hex=CARD, line_hex=AMBER)
    _panel_label(slide, x=5.38, y=3.96, w=2.18, text="monitored panel", accent=AMBER, name="S1_PANEL_LBL")
    panel_note = _textbox(slide, x=7.74, y=4.00, w=4.38, h=0.18, name="S1_PANEL_NOTE")
    _set_text(panel_note, text="what competes is entry into an observed top-K slice", size=11.4, color=MUTED, bold=False)
    _rank_row(slide, x=5.50, y=4.32, w=6.70, rank=1, author="author A / duplicate version", accent=CYAN, name="S1_ROW_1")
    _rank_row(slide, x=5.50, y=4.68, w=6.70, rank=7, author="author B / duplicate version", accent=GREEN, name="S1_ROW_2")
    _rank_row(slide, x=5.50, y=5.04, w=6.70, rank=19, author="author C / duplicate version", accent=AMBER, name="S1_ROW_3")
    steps.append(
        [
            "S1_PANEL",
            "S1_PANEL_LBL",
            "S1_PANEL_LBL_DOT",
            "S1_PANEL_LBL_TEXT",
            "S1_PANEL_NOTE",
            *_rank_row_names("S1_ROW_1"),
            *_rank_row_names("S1_ROW_2"),
            *_rank_row_names("S1_ROW_3"),
        ]
    )

    _stat_box(
        slide,
        x=5.24,
        y=5.72,
        w=2.20,
        h=1.10,
        title="1. inclusion",
        body=["entry into sampled top-K", "not full internal eligibility"],
        accent=CYAN,
        name="S1_OUTCOME_1",
    )
    _stat_box(
        slide,
        x=7.65,
        y=5.72,
        w=2.20,
        h=1.10,
        title="2. rank",
        body=["conditional order in the same q", "small score gaps can blow up"],
        accent=GREEN,
        name="S1_OUTCOME_2",
    )
    _stat_box(
        slide,
        x=10.06,
        y=5.72,
        w=2.44,
        h=1.10,
        title="3. opportunity",
        body=["position-weighted proxy", "report multiple w(r) choices"],
        accent=AMBER,
        name="S1_OUTCOME_3",
    )
    _card(slide, x=0.72, y=6.18, w=4.18, h=0.64, name="S1_BOUNDARY", fill_hex="1C1214", line_hex=RED)
    boundary_text = _textbox(slide, x=0.92, y=6.34, w=3.82, h=0.22, name="S1_BOUNDARY_TEXT")
    _set_text(
        boundary_text,
        text="Not actual attention. Not full candidate-set eligibility. Not user-wide fairness.",
        size=11.4,
        color="FFD6D2",
        bold=True,
    )
    steps.append(
        [
            *_stat_box_names("S1_OUTCOME_1"),
            *_stat_box_names("S1_OUTCOME_2"),
            *_stat_box_names("S1_OUTCOME_3"),
            "S1_BOUNDARY",
            "S1_BOUNDARY_TEXT",
        ]
    )
    return steps


def _slide2(slide) -> list[list[str]]:
    _set_bg(slide)
    _title_bar(
        slide,
        slide_idx=2,
        title="Local risk set inside a duplicate cluster",
        subtitle="We do not observe the platform's global candidate set, so we compare indexed versions within the same content cluster",
    )
    steps: list[list[str]] = [_title_names(2)]

    _card(slide, x=0.72, y=1.50, w=4.28, h=4.86, name="S2_CLUSTER", fill_hex=CARD, line_hex=CYAN)
    _panel_label(slide, x=0.86, y=1.64, w=2.18, text="duplicate cluster c", accent=CYAN, name="S2_CLUSTER_LBL")
    cluster_note = _textbox(slide, x=0.92, y=2.02, w=3.78, h=0.26, name="S2_CLUSTER_NOTE")
    _set_text(cluster_note, text="same text or same link + high similarity + narrow time window", size=11.5, color=MUTED)
    _post_card(
        slide,
        x=0.92,
        y=2.40,
        w=3.84,
        h=0.94,
        author="author A",
        stamp="06:05",
        text="Breaking link + same caption + same news item.",
        accent=CYAN,
        name="S2_POST_1",
    )
    _post_card(
        slide,
        x=0.92,
        y=3.48,
        w=3.84,
        h=0.94,
        author="author B",
        stamp="06:11",
        text="Breaking link + same caption + same news item.",
        accent=GREEN,
        name="S2_POST_2",
    )
    _post_card(
        slide,
        x=0.92,
        y=4.56,
        w=3.84,
        h=0.94,
        author="author C",
        stamp="06:22",
        text="Breaking link + same caption + same news item.",
        accent=AMBER,
        name="S2_POST_3",
    )
    steps.append(
        [
            "S2_CLUSTER",
            "S2_CLUSTER_LBL",
            "S2_CLUSTER_LBL_DOT",
            "S2_CLUSTER_LBL_TEXT",
            "S2_CLUSTER_NOTE",
            *_post_card_names("S2_POST_1"),
            *_post_card_names("S2_POST_2"),
            *_post_card_names("S2_POST_3"),
        ]
    )

    chevron_1 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(5.18), Inches(2.90), Inches(0.58), Inches(0.62))
    chevron_1.name = "S2_FLOW_1"
    chevron_1.fill.solid()
    chevron_1.fill.fore_color.rgb = _rgb(CYAN)
    chevron_1.fill.transparency = 0.14
    chevron_1.line.fill.background()
    _code_card(
        slide,
        x=5.92,
        y=2.48,
        w=2.92,
        h=1.38,
        lines=[
            "R_cq = {i in c : indexed_at_i <= captured_at_q}",
            "compare only versions already indexable at q",
        ],
        name="S2_RISK",
        accent=VIOLET,
    )
    risk_note = _textbox(slide, x=6.00, y=3.98, w=2.74, h=0.30, name="S2_RISK_NOTE")
    _set_text(risk_note, text="This converts missing global eligibility into within-content observed competition.", size=11.2, color=VIOLET, bold=True)
    steps.append(["S2_FLOW_1", *_code_card_names("S2_RISK"), "S2_RISK_NOTE"])

    chevron_2 = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, Inches(8.98), Inches(2.90), Inches(0.58), Inches(0.62))
    chevron_2.name = "S2_FLOW_2"
    chevron_2.fill.solid()
    chevron_2.fill.fore_color.rgb = _rgb(AMBER)
    chevron_2.fill.transparency = 0.14
    chevron_2.line.fill.background()
    _card(slide, x=9.70, y=1.50, w=2.82, h=4.86, name="S2_COMP", fill_hex=CARD, line_hex=AMBER)
    _panel_label(slide, x=9.84, y=1.64, w=2.16, text="competition in q", accent=AMBER, name="S2_COMP_LBL")
    comp_note = _textbox(slide, x=9.90, y=2.02, w=2.36, h=0.34, name="S2_COMP_NOTE")
    _set_lines(
        comp_note,
        lines=["When at least one cluster member appears,", "which indexed version enters and where?"],
        size=10.8,
        color=MUTED,
        line_spacing=1.06,
    )
    _rank_row(slide, x=9.94, y=2.62, w=2.20, rank=3, author="author B", accent=GREEN, name="S2_ROW_1")
    _rank_row(slide, x=9.94, y=3.08, w=2.20, rank=11, author="author A", accent=CYAN, name="S2_ROW_2")
    _rank_row(slide, x=9.94, y=3.54, w=2.20, rank=27, author="author C", accent=AMBER, name="S2_ROW_3")
    _card(slide, x=9.94, y=4.18, w=2.20, h=1.00, name="S2_COMP_CALL", fill_hex=CARD_2, line_hex=LINE)
    call_text = _textbox(slide, x=10.10, y=4.36, w=1.84, h=0.56, name="S2_COMP_CALL_TEXT")
    _set_lines(
        call_text,
        lines=[
            "Outcome 1:",
            "entry into panel",
            "Outcome 2:",
            "relative rank",
        ],
        size=11.6,
        color=INK,
        line_spacing=1.06,
    )
    steps.append(
        [
            "S2_FLOW_2",
            "S2_COMP",
            "S2_COMP_LBL",
            "S2_COMP_LBL_DOT",
            "S2_COMP_LBL_TEXT",
            "S2_COMP_NOTE",
            *_rank_row_names("S2_ROW_1"),
            *_rank_row_names("S2_ROW_2"),
            *_rank_row_names("S2_ROW_3"),
            "S2_COMP_CALL",
            "S2_COMP_CALL_TEXT",
        ]
    )

    _card(slide, x=0.72, y=6.52, w=12.00, h=0.54, name="S2_FILTERS", fill_hex="101C2D", line_hex=LINE)
    _mini_chip(slide, x=0.92, y=6.48, w=3.12, h=0.58, head="main sample", val="exact duplicates first", accent=CYAN, name="S2_FILTER_1")
    _mini_chip(slide, x=4.26, y=6.48, w=3.64, h=0.58, head="drop", val="pure URL / hashtag / hard suppression", accent=RED, name="S2_FILTER_2")
    _mini_chip(slide, x=8.08, y=6.48, w=4.26, h=0.58, head="later", val="near-duplicate for external validity", accent=AMBER, name="S2_FILTER_3")
    steps.append(
        [
            "S2_FILTERS",
            *_mini_chip_names("S2_FILTER_1"),
            *_mini_chip_names("S2_FILTER_2"),
            *_mini_chip_names("S2_FILTER_3"),
        ]
    )
    return steps


def _slide3(slide) -> list[list[str]]:
    _set_bg(slide)
    _title_bar(
        slide,
        slide_idx=3,
        title="Three outcomes and nested specifications",
        subtitle="Selection, conditional rank, and reinforcement are different layers and should not be collapsed into one fairness number",
    )
    steps: list[list[str]] = [_title_names(3)]

    _stat_box(
        slide,
        x=0.72,
        y=1.56,
        w=4.10,
        h=1.34,
        title="1. observed inclusion",
        body=["Pr(I_iq = 1 | R_cq)", "who gets into the monitored top-K slice?"],
        accent=CYAN,
        name="S3_BOX_1",
    )
    steps.append([*_stat_box_names("S3_BOX_1")])

    _stat_box(
        slide,
        x=0.72,
        y=3.06,
        w=4.10,
        h=1.34,
        title="2. conditional rank",
        body=["within the same q and cluster c", "who is placed above whom once shown?"],
        accent=GREEN,
        name="S3_BOX_2",
    )
    opportunity = _code_card(
        slide,
        x=5.14,
        y=1.70,
        w=7.32,
        h=1.02,
        lines=[
            "O_it = sum over q in hour t of I_iq * w(rank_iq)",
            "report flat, log, and geometric w(r)",
        ],
        name="S3_OPP",
        accent=AMBER,
    )
    steps.append([*_stat_box_names("S3_BOX_2"), *_code_card_names("S3_OPP")])

    _stat_box(
        slide,
        x=0.72,
        y=4.56,
        w=4.10,
        h=1.52,
        title="3. within-content reinforcement",
        body=["does earlier opportunity or lagged engagement", "shift the next-period opportunity share?"],
        accent=AMBER,
        name="S3_BOX_3",
    )
    _code_card(
        slide,
        x=5.14,
        y=3.08,
        w=7.32,
        h=1.70,
        lines=[
            "s_ict+1 = rho * s_ict + gamma * lagged metrics + author priors + labels + timing + FE",
            "rho > 0 => early panel opportunity amplifies later share within the same cluster",
        ],
        name="S3_DYN",
        accent=AMBER,
    )
    steps.append([*_stat_box_names("S3_BOX_3"), *_code_card_names("S3_DYN")])

    _spec_card(slide, x=5.14, y=5.10, w=1.66, h=0.94, title="S0", body=["cluster FE + time"], accent=CYAN, name="S3_SPEC_0")
    _spec_card(slide, x=6.96, y=5.10, w=1.66, h=0.94, title="S1", body=["+ freshness + labels"], accent=GREEN, name="S3_SPEC_1")
    _spec_card(slide, x=8.78, y=5.10, w=1.66, h=0.94, title="S2", body=["+ author priors"], accent=VIOLET, name="S3_SPEC_2")
    _spec_card(slide, x=10.60, y=5.10, w=1.86, h=0.94, title="S3", body=["+ lagged engagement"], accent=AMBER, name="S3_SPEC_3")
    _card(slide, x=0.72, y=6.34, w=12.00, h=0.50, name="S3_CAUTION", fill_hex="1C1214", line_hex=RED)
    caution_text = _textbox(slide, x=0.94, y=6.47, w=11.54, h=0.22, name="S3_CAUTION_TEXT")
    _set_text(
        caution_text,
        text="Strict rule: use lagged metrics only. Contemporaneous counts are post-treatment. Dynamic results need prior observability.",
        size=11.5,
        color="FFD6D2",
        bold=True,
    )
    steps.append(
        [
            *_spec_card_names("S3_SPEC_0"),
            *_spec_card_names("S3_SPEC_1"),
            *_spec_card_names("S3_SPEC_2"),
            *_spec_card_names("S3_SPEC_3"),
            "S3_CAUTION",
            "S3_CAUTION_TEXT",
        ]
    )
    return steps


def _slide4(slide) -> list[list[str]]:
    _set_bg(slide)
    _title_bar(
        slide,
        slide_idx=4,
        title="What this deck can claim, and what it should not claim",
        subtitle="Tight scope is the strength: producer-side opportunity diagnostics under one viewer context",
    )
    ribbon = _card(slide, x=3.76, y=1.42, w=5.80, h=0.44, name="S4_RIBBON", fill_hex="102846", line_hex=CYAN)
    ribbon_text = _textbox(slide, x=4.04, y=1.50, w=5.26, h=0.18, name="S4_RIBBON_TEXT")
    _set_text(ribbon_text, text="This is a defensible methods paper, not a total platform verdict.", size=12.4, color=CYAN, bold=True, align=PP_ALIGN.CENTER)
    steps: list[list[str]] = [_title_names(4), ["S4_RIBBON", "S4_RIBBON_TEXT"]]

    _claim_column(
        slide,
        x=0.72,
        y=1.98,
        w=5.74,
        title="We can claim",
        accent=GREEN,
        bullets=[
            "viewer-conditional concentration in monitored feed panels",
            "producer-side opportunity allocation with content held fixed",
            "selection, rank, and reinforcement inside duplicate clusters",
        ],
        name="S4_CLAIM",
        fill_hex=CARD,
    )
    steps.append([*_claim_column_names("S4_CLAIM")])

    _claim_column(
        slide,
        x=6.86,
        y=1.98,
        w=5.86,
        title="We should not claim",
        accent=RED,
        bullets=[
            "platform-wide candidate-set eligibility fairness",
            "actual human attention or full feed consumption",
            "full ideological echo chamber without extra topic or graph layers",
        ],
        name="S4_NOCLAIM",
        fill_hex="15171F",
    )
    steps.append([*_claim_column_names("S4_NOCLAIM")])

    _card(slide, x=0.72, y=5.78, w=12.00, h=0.74, name="S4_PAPER", fill_hex="0C1B2D", line_hex=CYAN)
    paper_title = _textbox(slide, x=0.96, y=5.92, w=11.48, h=0.20, name="S4_PAPER_TITLE")
    _set_text(
        paper_title,
        text="Content-controlled observed top-K opportunity allocation in Bluesky feed panels",
        size=18.8,
        color=INK,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    paper_sub = _textbox(slide, x=1.16, y=6.18, w=11.08, h=0.16, name="S4_PAPER_SUB")
    _set_text(
        paper_sub,
        text="Selection, rank, and reinforcement among duplicate posts under a fixed viewer context",
        size=12.2,
        color=MUTED,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    steps.append(["S4_PAPER", "S4_PAPER_TITLE", "S4_PAPER_SUB"])

    _mini_chip(slide, x=1.12, y=6.72, w=3.28, h=0.54, head="heterogeneity", val="auth vs unauth, only where mode diverges", accent=GREEN, name="S4_NEXT_1")
    _mini_chip(slide, x=4.82, y=6.72, w=3.42, h=0.54, head="robustness", val="near-duplicate after exact duplicate core", accent=AMBER, name="S4_NEXT_2")
    _mini_chip(slide, x=8.68, y=6.72, w=2.98, h=0.54, head="optional later", val="viewer-local graph", accent=VIOLET, name="S4_NEXT_3")
    steps.append([*_mini_chip_names("S4_NEXT_1"), *_mini_chip_names("S4_NEXT_2"), *_mini_chip_names("S4_NEXT_3")])
    return steps


def _build_prs() -> tuple[Presentation, dict[int, list[list[str]]]]:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    blank = prs.slide_layouts[6]

    slide1 = prs.slides.add_slide(blank)
    steps1 = _slide1(slide1)

    slide2 = prs.slides.add_slide(blank)
    steps2 = _slide2(slide2)

    slide3 = prs.slides.add_slide(blank)
    steps3 = _slide3(slide3)

    slide4 = prs.slides.add_slide(blank)
    steps4 = _slide4(slide4)

    return prs, {1: steps1, 2: steps2, 3: steps3, 4: steps4}


def _iter_shape_infos_in_z_order(slide_root: etree._Element) -> list[tuple[int, str]]:
    infos: list[tuple[int, str]] = []
    sp_tree = slide_root.find(f"{_qn(PML_NS, 'cSld')}/{_qn(PML_NS, 'spTree')}")
    if sp_tree is None:
        return infos
    for child in sp_tree:
        local = etree.QName(child).localname
        if local in {"nvGrpSpPr", "grpSpPr"}:
            continue
        c_nv_pr = child.find(f".//{_qn(PML_NS, 'cNvPr')}")
        if c_nv_pr is None:
            continue
        raw = c_nv_pr.get("id")
        if raw is None:
            continue
        try:
            spid = int(raw)
        except ValueError:
            continue
        infos.append((spid, c_nv_pr.get("name") or ""))
    return infos


def _ensure_fade_transition(root: etree._Element) -> None:
    transition = root.find(_qn(PML_NS, "transition"))
    if transition is None:
        transition = etree.Element(_qn(PML_NS, "transition"), spd="slow")
        etree.SubElement(transition, _qn(PML_NS, "fade"))
        c_sld = root.find(_qn(PML_NS, "cSld"))
        if c_sld is not None:
            idx = list(root).index(c_sld)
            root.insert(idx + 1, transition)
        else:
            root.insert(0, transition)
        return
    transition.set("spd", transition.get("spd") or "slow")
    if transition.find(_qn(PML_NS, "fade")) is None:
        for child in list(transition):
            transition.remove(child)
        etree.SubElement(transition, _qn(PML_NS, "fade"))


def _build_click_timing_steps(*, steps: list[list[int]], effect_dur_ms: int) -> etree._Element:
    timing = etree.Element(_qn(PML_NS, "timing"))
    tn_lst = etree.SubElement(timing, _qn(PML_NS, "tnLst"))
    par = etree.SubElement(tn_lst, _qn(PML_NS, "par"))
    tm_root = etree.SubElement(par, _qn(PML_NS, "cTn"), id="1", dur="indefinite", restart="never", nodeType="tmRoot")
    tm_child = etree.SubElement(tm_root, _qn(PML_NS, "childTnLst"))
    seq = etree.SubElement(tm_child, _qn(PML_NS, "seq"), concurrent="1", nextAc="seek")
    main_ctn = etree.SubElement(seq, _qn(PML_NS, "cTn"), id="2", dur="indefinite", nodeType="mainSeq")
    main_child = etree.SubElement(main_ctn, _qn(PML_NS, "childTnLst"))

    animated_spids: list[int] = []
    next_id = 3
    for step in steps:
        if not step:
            continue
        animated_spids.extend(step)
        par_wrap = etree.SubElement(main_child, _qn(PML_NS, "par"))
        wrap = etree.SubElement(par_wrap, _qn(PML_NS, "cTn"), id=str(next_id), fill="hold")
        next_id += 1
        st_cond = etree.SubElement(wrap, _qn(PML_NS, "stCondLst"))
        etree.SubElement(st_cond, _qn(PML_NS, "cond"), delay="indefinite")
        cond_on_begin = etree.SubElement(st_cond, _qn(PML_NS, "cond"), evt="onBegin", delay="0")
        etree.SubElement(cond_on_begin, _qn(PML_NS, "tn"), val="2")
        wrap_child = etree.SubElement(wrap, _qn(PML_NS, "childTnLst"))
        par_hold = etree.SubElement(wrap_child, _qn(PML_NS, "par"))
        hold = etree.SubElement(par_hold, _qn(PML_NS, "cTn"), id=str(next_id), fill="hold")
        next_id += 1
        hold_cond = etree.SubElement(hold, _qn(PML_NS, "stCondLst"))
        etree.SubElement(hold_cond, _qn(PML_NS, "cond"), delay="0")
        hold_child = etree.SubElement(hold, _qn(PML_NS, "childTnLst"))
        par_effect = etree.SubElement(hold_child, _qn(PML_NS, "par"))
        click = etree.SubElement(
            par_effect,
            _qn(PML_NS, "cTn"),
            id=str(next_id),
            presetID="10",
            presetClass="entr",
            presetSubtype="0",
            fill="hold",
            grpId="0",
            nodeType="clickEffect",
        )
        next_id += 1
        click_cond = etree.SubElement(click, _qn(PML_NS, "stCondLst"))
        etree.SubElement(click_cond, _qn(PML_NS, "cond"), delay="0")
        click_child = etree.SubElement(click, _qn(PML_NS, "childTnLst"))
        for spid in step:
            set_el = etree.SubElement(click_child, _qn(PML_NS, "set"))
            c_bhvr = etree.SubElement(set_el, _qn(PML_NS, "cBhvr"))
            set_ctn = etree.SubElement(c_bhvr, _qn(PML_NS, "cTn"), id=str(next_id), dur="1", fill="hold")
            next_id += 1
            set_ctn_cond = etree.SubElement(set_ctn, _qn(PML_NS, "stCondLst"))
            etree.SubElement(set_ctn_cond, _qn(PML_NS, "cond"), delay="0")
            tgt_el = etree.SubElement(c_bhvr, _qn(PML_NS, "tgtEl"))
            etree.SubElement(tgt_el, _qn(PML_NS, "spTgt"), spid=str(spid))
            attr_name_lst = etree.SubElement(c_bhvr, _qn(PML_NS, "attrNameLst"))
            attr_name = etree.SubElement(attr_name_lst, _qn(PML_NS, "attrName"))
            attr_name.text = "style.visibility"
            to = etree.SubElement(set_el, _qn(PML_NS, "to"))
            etree.SubElement(to, _qn(PML_NS, "strVal"), val="visible")

            anim = etree.SubElement(click_child, _qn(PML_NS, "animEffect"), transition="in", filter="fade")
            anim_bhvr = etree.SubElement(anim, _qn(PML_NS, "cBhvr"))
            etree.SubElement(anim_bhvr, _qn(PML_NS, "cTn"), id=str(next_id), dur=str(int(effect_dur_ms)))
            next_id += 1
            anim_tgt = etree.SubElement(anim_bhvr, _qn(PML_NS, "tgtEl"))
            etree.SubElement(anim_tgt, _qn(PML_NS, "spTgt"), spid=str(spid))

    prev = etree.SubElement(seq, _qn(PML_NS, "prevCondLst"))
    cond_prev = etree.SubElement(prev, _qn(PML_NS, "cond"), evt="onPrev", delay="0")
    tgt_prev = etree.SubElement(cond_prev, _qn(PML_NS, "tgtEl"))
    etree.SubElement(tgt_prev, _qn(PML_NS, "sldTgt"))

    nxt = etree.SubElement(seq, _qn(PML_NS, "nextCondLst"))
    cond_next = etree.SubElement(nxt, _qn(PML_NS, "cond"), evt="onNext", delay="0")
    tgt_next = etree.SubElement(cond_next, _qn(PML_NS, "tgtEl"))
    etree.SubElement(tgt_next, _qn(PML_NS, "sldTgt"))

    bld_lst = etree.SubElement(timing, _qn(PML_NS, "bldLst"))
    seen: set[int] = set()
    for spid in animated_spids:
        if spid in seen:
            continue
        seen.add(spid)
        etree.SubElement(bld_lst, _qn(PML_NS, "bldP"), spid=str(spid), grpId="0", animBg="1")
    return timing


def _inject_steps(*, pptx_in: Path, pptx_out: Path, slide_steps: dict[int, list[list[str]]], effect_dur_ms: int) -> None:
    with zipfile.ZipFile(pptx_in, "r") as zin:
        with zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                if not info.filename.startswith("ppt/slides/slide") or not info.filename.endswith(".xml"):
                    zout.writestr(info, data)
                    continue
                slide_num = int(Path(info.filename).stem.replace("slide", ""))
                steps_by_name = slide_steps.get(slide_num)
                if not steps_by_name:
                    zout.writestr(info, data)
                    continue
                root = etree.fromstring(data)
                infos = _iter_shape_infos_in_z_order(root)
                name_to_spid = {name: spid for spid, name in infos if name}
                z_order = [spid for spid, _ in infos]
                steps: list[list[int]] = []
                for step in steps_by_name:
                    current = [name_to_spid[name] for name in step if name in name_to_spid]
                    if current:
                        current = [spid for spid in z_order if spid in current]
                        steps.append(current)
                for existing in list(root.findall(_qn(PML_NS, "timing"))):
                    root.remove(existing)
                if steps:
                    timing = _build_click_timing_steps(steps=steps, effect_dur_ms=effect_dur_ms)
                    _ensure_fade_transition(root)
                    transition = root.find(_qn(PML_NS, "transition"))
                    if transition is not None:
                        idx = list(root).index(transition)
                        root.insert(idx + 1, timing)
                    else:
                        root.append(timing)
                zout.writestr(info, etree.tostring(root, encoding="UTF-8", xml_declaration=False))


def build_deck(*, out_pre: Path, out_animated: Path, effect_dur_ms: int) -> None:
    out_pre.parent.mkdir(parents=True, exist_ok=True)
    out_animated.parent.mkdir(parents=True, exist_ok=True)
    prs, steps = _build_prs()
    prs.save(str(out_pre))
    _inject_steps(pptx_in=out_pre, pptx_out=out_animated, slide_steps=steps, effect_dur_ms=effect_dur_ms)


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a 4-slide observed top-K exposure deck with grouped click animations.")
    parser.add_argument("--out-pre", type=Path, default=Path("_build/observed_topk_exposure_4slides_preanim.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/observed_topk_exposure_4slides_animated.pptx"))
    parser.add_argument("--dur-ms", type=int, default=220)
    args = parser.parse_args()

    build_deck(out_pre=args.out_pre, out_animated=args.out, effect_dur_ms=int(args.dur_ms))
    print(f"OK: wrote {args.out_pre}")
    print(f"OK: wrote {args.out}")


if __name__ == "__main__":
    main()
