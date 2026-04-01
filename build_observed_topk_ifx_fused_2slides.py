#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import build_ifx_poster_svg as base
import build_observed_topk_svg_4slides as obs


def _card_bg_svg(*, w_px: int, h_px: int, accent_hex: str, title: str) -> str:
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=22,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="12" height="{h_px}" rx="22" ry="22" fill="#{accent_hex}" opacity="0.84"/>'
    body += base._svg_text(x=22, y=18, text=title, size_px=19, color=f"#{base.COLORS_HEX['ink']}", weight=900, baseline="hanging")
    body += f'<line x1="22" y1="46" x2="{w_px-18}" y2="46" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="1.5" stroke-opacity="0.75"/>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _small_matrix_svg(*, w_px: int, h_px: int, accent_hex: str) -> str:
    n = 7
    gap = 6
    cell = min((w_px - 28 - gap * (n - 1)) // n, (h_px - 28 - gap * (n - 1)) // n)
    x0 = 14
    y0 = 14
    weights = [0.24, 0.20, 0.16, 0.14, 0.11, 0.09, 0.06]
    max_w = max(weights)
    hi_r = 3
    body = f'<rect x="{x0-6}" y="{y0 + hi_r * (cell + gap) - 4}" width="{n * cell + (n-1) * gap + 12}" height="{cell + 8}" rx="10" ry="10" fill="#{accent_hex}" opacity="0.12"/>'
    for r in range(n):
        for c in range(n):
            x = x0 + c * (cell + gap)
            y = y0 + r * (cell + gap)
            p = weights[(c - r) % n]
            fill = obs._mix(base.COLORS_HEX["card2"], accent_hex, 0.14 + 0.70 * (p / max_w))
            stroke = obs._mix(base.COLORS_HEX["card_line"], accent_hex, 0.34)
            body += f'<rect x="{x}" y="{y}" width="{cell}" height="{cell}" rx="8" ry="8" fill="#{fill}" stroke="#{stroke}" stroke-width="1.5"/>'
            body += base._svg_text(x=x + cell // 2, y=y + cell // 2, text=f"{p:.2f}", size_px=11, color=f"#{base.COLORS_HEX['ink']}", weight=800, anchor="middle")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _small_vbars_svg(*, w_px: int, h_px: int, accent_hex: str) -> str:
    body = ""
    rows = 7
    bar_h = 16
    gap = 12
    x0 = 10
    y0 = 10
    track_w = w_px - 36
    for j in range(rows):
        frac = 1.0 / (1.0 + 0.18 * j)
        fill_w = int(track_w * frac)
        y = y0 + j * (bar_h + gap)
        body += f'<rect x="{x0}" y="{y}" width="{track_w}" height="{bar_h}" rx="8" ry="8" fill="#08111D" opacity="0.96"/>'
        body += f'<rect x="{x0}" y="{y}" width="{fill_w}" height="{bar_h}" rx="8" ry="8" fill="#{accent_hex}" opacity="0.90"/>'
        body += f'<circle cx="{x0 + fill_w}" cy="{y + bar_h//2}" r="6" fill="#{accent_hex}" opacity="0.90"/>'
        body += base._svg_text(x=w_px - 6, y=y + bar_h // 2, text=f"v{j+1}", size_px=11, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="end")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _mini_gap_heatmap_svg(*, w_px: int, h_px: int, accent_hex: str) -> str:
    rows = 5
    cols = 5
    gap = 8
    cell_w = (w_px - 20 - gap * (cols - 1)) // cols
    cell_h = (h_px - 26 - gap * (rows - 1)) // rows
    x0 = 10
    y0 = 12
    vals = [
        [0.10, 0.18, 0.22, 0.28, 0.36],
        [0.06, 0.12, 0.18, 0.26, 0.34],
        [0.04, 0.10, 0.16, 0.20, 0.28],
        [0.02, 0.06, 0.12, 0.18, 0.24],
        [0.00, 0.04, 0.08, 0.12, 0.16],
    ]
    body = ""
    for r in range(rows):
        for c in range(cols):
            heat = vals[r][c]
            x = x0 + c * (cell_w + gap)
            y = y0 + r * (cell_h + gap)
            fill = obs._mix(base.COLORS_HEX["card2"], accent_hex, 0.22 + 0.74 * heat)
            stroke = obs._mix(base.COLORS_HEX["card_line"], accent_hex, 0.28)
            body += f'<rect x="{x}" y="{y}" width="{cell_w}" height="{cell_h}" rx="8" ry="8" fill="#{fill}" stroke="#{stroke}" stroke-width="1.5"/>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_body(svg_text: str) -> str:
    match = re.search(r"^<svg[^>]*>(.*)</svg>$", svg_text, flags=re.DOTALL)
    if not match:
        raise ValueError("Expected an SVG root wrapper")
    return match.group(1)


def _risk_timeline_with_nodes_svg(*, w_px: int, h_px: int) -> str:
    body = _svg_body(obs._risk_timeline_base_svg(w_px=w_px, h_px=h_px))
    node_w = base._in_to_px(0.50)
    node_h = base._in_to_px(0.50)
    placements = [
        (base._in_to_px(0.40), base._in_to_px(0.36), "A", base.COLORS_HEX["cyan"], True),
        (base._in_to_px(1.20), base._in_to_px(0.36), "B", base.COLORS_HEX["green"], True),
        (base._in_to_px(2.42), base._in_to_px(0.36), "C", base.COLORS_HEX["amber"], False),
    ]
    for x_px, y_px, label, accent_hex, eligible in placements:
        node_svg = obs._risk_node_svg(
            label=label,
            accent_hex=accent_hex,
            eligible=eligible,
            w_px=node_w,
            h_px=node_h,
        )
        body += f'<g transform="translate({x_px},{y_px})">{_svg_body(node_svg)}</g>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _reinforcement_plot_svg(*, w_px: int, h_px: int) -> str:
    plot_y = base._in_to_px(0.26)
    plot_h = h_px - plot_y
    body = f'<g transform="translate(0,{plot_y})">{_svg_body(obs._plot_axes_svg(w_px=w_px, h_px=plot_h))}</g>'
    body += f'<g transform="translate(0,{plot_y})">{_svg_body(obs._plot_line_svg(values=[0.18, 0.34, 0.58, 0.78, 0.90], accent_hex=base.COLORS_HEX["cyan"], label="A", w_px=w_px, h_px=plot_h))}</g>'
    body += f'<g transform="translate(0,{plot_y})">{_svg_body(obs._plot_line_svg(values=[0.18, 0.22, 0.26, 0.29, 0.31], accent_hex=base.COLORS_HEX["green"], label="B", w_px=w_px, h_px=plot_h))}</g>'
    body += f'<g transform="translate(0,{plot_y})">{_svg_body(obs._plot_line_svg(values=[0.18, 0.15, 0.12, 0.10, 0.08], accent_hex=base.COLORS_HEX["amber"], label="C", w_px=w_px, h_px=plot_h))}</g>'
    body += f'<g transform="translate({base._in_to_px(0.16)},0)">{_svg_body(obs._metric_arrow_svg(w_px=base._in_to_px(1.12), h_px=base._in_to_px(0.64)))}</g>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _assert_no_nonstatic_overlap(tiles: list[base.SvgTile]) -> None:
    def intersects(a: base.SvgTile, b: base.SvgTile, eps: float = 1e-4) -> bool:
        if a.slide_num != b.slide_num:
            return False
        if a.static or b.static:
            return False
        return not (
            a.x_in + a.w_in <= b.x_in + eps
            or b.x_in + b.w_in <= a.x_in + eps
            or a.y_in + a.h_in <= b.y_in + eps
            or b.y_in + b.h_in <= a.y_in + eps
        )

    dynamic = [t for t in tiles if not t.static]
    for idx, a in enumerate(dynamic):
        for b in dynamic[idx + 1 :]:
            if intersects(a, b):
                raise ValueError(
                    f"Overlap detected on slide {a.slide_num}: {a.name} "
                    f"({a.x_in:.2f},{a.y_in:.2f},{a.w_in:.2f},{a.h_in:.2f}) "
                    f"vs {b.name} ({b.x_in:.2f},{b.y_in:.2f},{b.w_in:.2f},{b.h_in:.2f})"
                )


def _formula_tile(*, slide_num: int, name: str, expr: str, x_in: float, y_in: float, max_w_in: float, color_hex: str, fontsize_pt: float) -> base.SvgTile:
    return obs._formula_tile(
        slide_num=slide_num,
        name=name,
        expr=expr,
        x_in=x_in,
        y_in=y_in,
        max_w_in=max_w_in,
        color_hex=color_hex,
        fontsize_pt=fontsize_pt,
    )


def _build_tiles() -> list[base.SvgTile]:
    tiles: list[base.SvgTile] = []

    def add_bg(slide_num: int, title: str, subtitle: str) -> None:
        tiles.append(
            base.SvgTile(
                slide_num=slide_num,
                name=f"_BG_GRID_S{slide_num}",
                x_in=0.0,
                y_in=0.0,
                w_in=base.SLIDE_W_IN,
                h_in=base.SLIDE_H_IN,
                svg_text=base._svg_grid_bg(w_px=base._in_to_px(base.SLIDE_W_IN), h_px=base._in_to_px(base.SLIDE_H_IN)),
                static=True,
            )
        )
        tiles.append(
            base.SvgTile(
                slide_num=slide_num,
                name=f"STATIC_TITLE_S{slide_num}",
                x_in=base.SAFE_L_IN,
                y_in=0.28,
                w_in=base.SAFE_W_IN,
                h_in=0.92,
                svg_text=base._svg_title_block(
                    w_px=base._in_to_px(base.SAFE_W_IN),
                    h_px=base._in_to_px(0.92),
                    title=title,
                    subtitle=subtitle,
                ),
                static=True,
            )
        )

    # Slide 1
    add_bg(1, "Observed top-K exposure", "same content, fixed viewer, local risk set")
    tiles.extend(
        [
            base.SvgTile(slide_num=1, name="STATIC_S1_LEFT_BG", x_in=0.65, y_in=1.34, w_in=5.60, h_in=5.46, svg_text=_card_bg_svg(w_px=base._in_to_px(5.60), h_px=base._in_to_px(5.46), accent_hex=base.COLORS_HEX["cyan"], title="exposure"), static=True),
            base.SvgTile(slide_num=1, name="STATIC_S1_RT_BG", x_in=6.46, y_in=1.34, w_in=6.22, h_in=2.55, svg_text=_card_bg_svg(w_px=base._in_to_px(6.22), h_px=base._in_to_px(2.55), accent_hex=base.COLORS_HEX["amber"], title="observed top-K"), static=True),
            base.SvgTile(slide_num=1, name="STATIC_S1_RB_BG", x_in=6.46, y_in=4.09, w_in=6.22, h_in=2.71, svg_text=_card_bg_svg(w_px=base._in_to_px(6.22), h_px=base._in_to_px(2.71), accent_hex=base.COLORS_HEX["violet"], title="local risk set"), static=True),
        ]
    )
    tiles.extend(
        [
            base.SvgTile(slide_num=1, name="S1_PILL_THEORY", x_in=0.90, y_in=1.58, w_in=1.70, h_in=0.20, svg_text=obs._pill_svg("formula", base.COLORS_HEX["cyan"], w_in=1.70, h_in=0.20)),
            _formula_tile(slide_num=1, name="S1_EQ_THEORY", expr=r"$\mathrm{Exposure}(d_i\mid P)=\sum_j P_{i,j}v_j$", x_in=0.88, y_in=1.90, max_w_in=4.95, color_hex=base.COLORS_HEX["ink"], fontsize_pt=19.0),
            base.SvgTile(slide_num=1, name="S1_PILL_P", x_in=0.96, y_in=2.62, w_in=1.55, h_in=0.20, svg_text=obs._pill_svg("P", base.COLORS_HEX["cyan"], w_in=1.55, h_in=0.20)),
            base.SvgTile(slide_num=1, name="S1_MATRIX", x_in=0.96, y_in=2.96, w_in=2.55, h_in=2.55, svg_text=_small_matrix_svg(w_px=base._in_to_px(2.55), h_px=base._in_to_px(2.55), accent_hex=base.COLORS_HEX["cyan"])),
            base.SvgTile(slide_num=1, name="S1_ARROW_THEORY", x_in=3.60, y_in=4.02, w_in=0.30, h_in=0.26, svg_text=base._svg_arrow(w_px=base._in_to_px(0.30), h_px=base._in_to_px(0.26), accent_hex=base.COLORS_HEX["amber"], text="")),
            base.SvgTile(slide_num=1, name="S1_PILL_V", x_in=4.36, y_in=2.62, w_in=1.34, h_in=0.20, svg_text=obs._pill_svg("v", base.COLORS_HEX["green"], w_in=1.34, h_in=0.20)),
            base.SvgTile(slide_num=1, name="S1_VBARS", x_in=4.02, y_in=2.96, w_in=1.58, h_in=2.55, svg_text=_small_vbars_svg(w_px=base._in_to_px(1.58), h_px=base._in_to_px(2.55), accent_hex=base.COLORS_HEX["green"])),
            base.SvgTile(slide_num=1, name="S1_OUT", x_in=3.96, y_in=5.78, w_in=1.64, h_in=0.52, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(1.64), h_px=base._in_to_px(0.52), label="E_i", frac=0.66, accent_hex=base.COLORS_HEX["amber"], value="")),
            _formula_tile(slide_num=1, name="S1_EQ_OBS", expr=r"$O_{i q}=I_{i q}\,w(\mathrm{rank}_{i q})$", x_in=6.72, y_in=1.68, max_w_in=5.10, color_hex=base.COLORS_HEX["amber"], fontsize_pt=20.0),
            base.SvgTile(slide_num=1, name="S1_ROW_A", x_in=6.80, y_in=2.34, w_in=2.76, h_in=0.42, svg_text=base._svg_rank_row(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.42), rank=1, label="A", accent_hex=base.COLORS_HEX["cyan"])),
            base.SvgTile(slide_num=1, name="S1_ROW_B", x_in=6.80, y_in=2.90, w_in=2.76, h_in=0.42, svg_text=base._svg_rank_row(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.42), rank=7, label="B", accent_hex=base.COLORS_HEX["green"])),
            base.SvgTile(slide_num=1, name="S1_ROW_C", x_in=6.80, y_in=3.46, w_in=2.76, h_in=0.42, svg_text=base._svg_rank_row(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.42), rank=19, label="C", accent_hex=base.COLORS_HEX["amber"])),
            base.SvgTile(slide_num=1, name="S1_OBS_ARROW", x_in=9.70, y_in=2.88, w_in=0.74, h_in=0.34, svg_text=base._svg_arrow(w_px=base._in_to_px(0.74), h_px=base._in_to_px(0.34), accent_hex=base.COLORS_HEX["amber"], text="w")),
            base.SvgTile(slide_num=1, name="S1_BAR_A", x_in=10.62, y_in=2.26, w_in=1.42, h_in=0.48, svg_text=obs._rank_weight_bar_svg(label="A", frac=0.96, accent_hex=base.COLORS_HEX["cyan"], w_px=base._in_to_px(1.42), h_px=base._in_to_px(0.48))),
            base.SvgTile(slide_num=1, name="S1_BAR_B", x_in=10.62, y_in=2.84, w_in=1.42, h_in=0.48, svg_text=obs._rank_weight_bar_svg(label="B", frac=0.44, accent_hex=base.COLORS_HEX["green"], w_px=base._in_to_px(1.42), h_px=base._in_to_px(0.48))),
            base.SvgTile(slide_num=1, name="S1_BAR_C", x_in=10.62, y_in=3.42, w_in=1.42, h_in=0.48, svg_text=obs._rank_weight_bar_svg(label="C", frac=0.18, accent_hex=base.COLORS_HEX["amber"], w_px=base._in_to_px(1.42), h_px=base._in_to_px(0.48))),
            _formula_tile(slide_num=1, name="S1_EQ_RISK", expr=r"$R_{c q}=\{i\in c:\ \mathrm{indexed\_at}_i\leq \mathrm{captured\_at}_q\}$", x_in=6.72, y_in=4.34, max_w_in=5.20, color_hex=base.COLORS_HEX["violet"], fontsize_pt=18.0),
            base.SvgTile(slide_num=1, name="S1_CLUSTER", x_in=6.84, y_in=5.08, w_in=1.52, h_in=1.40, svg_text=base._svg_cluster_bubble(w_px=base._in_to_px(1.52), h_px=base._in_to_px(1.40), label="c", accent_hex=base.COLORS_HEX["cyan"])),
            base.SvgTile(slide_num=1, name="S1_TIMELINE", x_in=8.56, y_in=5.12, w_in=3.46, h_in=1.32, svg_text=_risk_timeline_with_nodes_svg(w_px=base._in_to_px(3.46), h_px=base._in_to_px(1.32))),
        ]
    )

    # Slide 2
    add_bg(2, "Duplicate-level exposure gaps", "aggregate, compare, reinforce")
    tiles.extend(
        [
            base.SvgTile(slide_num=2, name="STATIC_S2_LU_BG", x_in=0.65, y_in=1.34, w_in=7.08, h_in=3.48, svg_text=_card_bg_svg(w_px=base._in_to_px(7.08), h_px=base._in_to_px(3.48), accent_hex=base.COLORS_HEX["amber"], title="aggregate"), static=True),
            base.SvgTile(slide_num=2, name="STATIC_S2_LL_BG", x_in=0.65, y_in=5.02, w_in=7.08, h_in=1.78, svg_text=_card_bg_svg(w_px=base._in_to_px(7.08), h_px=base._in_to_px(1.78), accent_hex=base.COLORS_HEX["coral"], title="gap metrics"), static=True),
            base.SvgTile(slide_num=2, name="STATIC_S2_R_BG", x_in=7.96, y_in=1.34, w_in=4.72, h_in=5.46, svg_text=_card_bg_svg(w_px=base._in_to_px(4.72), h_px=base._in_to_px(5.46), accent_hex=base.COLORS_HEX["violet"], title="reinforcement"), static=True),
        ]
    )
    tiles.extend(
        [
            _formula_tile(slide_num=2, name="S2_EQ_AGG", expr=r"$O_{i t}=\sum_{q\in \mathcal{Q}_t} O_{i q}$", x_in=0.90, y_in=1.54, max_w_in=5.75, color_hex=base.COLORS_HEX["amber"], fontsize_pt=20.0),
            base.SvgTile(slide_num=2, name="S2_Q1", x_in=1.24, y_in=2.42, w_in=0.62, h_in=0.18, svg_text=obs._pill_svg("q1", base.COLORS_HEX["coral"], w_in=0.62, h_in=0.18)),
            base.SvgTile(slide_num=2, name="S2_Q2", x_in=2.06, y_in=2.42, w_in=0.62, h_in=0.18, svg_text=obs._pill_svg("q2", base.COLORS_HEX["coral"], w_in=0.62, h_in=0.18)),
            base.SvgTile(slide_num=2, name="S2_Q3", x_in=2.88, y_in=2.42, w_in=0.62, h_in=0.18, svg_text=obs._pill_svg("q3", base.COLORS_HEX["coral"], w_in=0.62, h_in=0.18)),
            base.SvgTile(slide_num=2, name="S2_Q4", x_in=3.70, y_in=2.42, w_in=0.62, h_in=0.18, svg_text=obs._pill_svg("q4", base.COLORS_HEX["coral"], w_in=0.62, h_in=0.18)),
            base.SvgTile(slide_num=2, name="S2_Q5", x_in=4.52, y_in=2.42, w_in=0.62, h_in=0.18, svg_text=obs._pill_svg("q5", base.COLORS_HEX["coral"], w_in=0.62, h_in=0.18)),
            base.SvgTile(slide_num=2, name="S2_ROW_A", x_in=0.96, y_in=2.78, w_in=4.92, h_in=0.52, svg_text=obs._heat_row_svg(label="A", values=[0.82, 0.64, 0.44, 0.20, 0.00], accent_hex=base.COLORS_HEX["cyan"], w_px=base._in_to_px(4.92), h_px=base._in_to_px(0.52))),
            base.SvgTile(slide_num=2, name="S2_ROW_B", x_in=0.96, y_in=3.42, w_in=4.92, h_in=0.52, svg_text=obs._heat_row_svg(label="B", values=[0.18, 0.46, 0.70, 0.58, 0.36], accent_hex=base.COLORS_HEX["green"], w_px=base._in_to_px(4.92), h_px=base._in_to_px(0.52))),
            base.SvgTile(slide_num=2, name="S2_ROW_C", x_in=0.96, y_in=4.06, w_in=4.92, h_in=0.52, svg_text=obs._heat_row_svg(label="C", values=[0.00, 0.00, 0.16, 0.40, 0.74], accent_hex=base.COLORS_HEX["amber"], w_px=base._in_to_px(4.92), h_px=base._in_to_px(0.52))),
            base.SvgTile(slide_num=2, name="S2_BAR_A", x_in=6.02, y_in=2.78, w_in=1.40, h_in=0.58, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(1.40), h_px=base._in_to_px(0.58), label="O_A,t", frac=0.80, accent_hex=base.COLORS_HEX["cyan"], value="2.10")),
            base.SvgTile(slide_num=2, name="S2_BAR_B", x_in=6.02, y_in=3.43, w_in=1.40, h_in=0.58, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(1.40), h_px=base._in_to_px(0.58), label="O_B,t", frac=0.87, accent_hex=base.COLORS_HEX["green"], value="2.28")),
            base.SvgTile(slide_num=2, name="S2_BAR_C", x_in=6.02, y_in=4.08, w_in=1.40, h_in=0.58, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(1.40), h_px=base._in_to_px(0.58), label="O_C,t", frac=0.50, accent_hex=base.COLORS_HEX["amber"], value="1.30")),
            base.SvgTile(slide_num=2, name="STATIC_S2_CARD_D", x_in=0.92, y_in=5.42, w_in=2.08, h_in=1.00, svg_text=_card_bg_svg(w_px=base._in_to_px(2.08), h_px=base._in_to_px(1.00), accent_hex=base.COLORS_HEX["amber"], title="dup gap"), static=True),
            base.SvgTile(slide_num=2, name="STATIC_S2_CARD_L", x_in=3.27, y_in=5.42, w_in=2.08, h_in=1.00, svg_text=_card_bg_svg(w_px=base._in_to_px(2.08), h_px=base._in_to_px(1.00), accent_hex=base.COLORS_HEX["cyan"], title="indiv fair"), static=True),
            base.SvgTile(slide_num=2, name="STATIC_S2_CARD_C", x_in=5.62, y_in=5.42, w_in=1.86, h_in=1.00, svg_text=_card_bg_svg(w_px=base._in_to_px(1.86), h_px=base._in_to_px(1.00), accent_hex=base.COLORS_HEX["violet"], title="ctx ΔE"), static=True),
            _formula_tile(slide_num=2, name="S2_EQ_D", expr=r"$D(p,p')=\left|E(p)-E(p')\right|$", x_in=1.06, y_in=5.86, max_w_in=1.76, color_hex=base.COLORS_HEX["amber"], fontsize_pt=14.0),
            _formula_tile(slide_num=2, name="S2_EQ_L", expr=r"$d\!\approx\!0\Rightarrow |E-E'|\leq \varepsilon$", x_in=3.42, y_in=5.86, max_w_in=1.74, color_hex=base.COLORS_HEX["ink"], fontsize_pt=13.4),
            _formula_tile(slide_num=2, name="S2_EQ_C", expr=r"$\Delta E=v(r_a)-v(r_b)$", x_in=5.78, y_in=5.86, max_w_in=1.52, color_hex=base.COLORS_HEX["violet"], fontsize_pt=13.0),
            _formula_tile(slide_num=2, name="S2_EQ_REIN", expr=r"$s_{i c,t+1}=\rho s_{i c t}+\gamma M^-_{i c t}+\cdots$", x_in=8.20, y_in=1.62, max_w_in=4.05, color_hex=base.COLORS_HEX["violet"], fontsize_pt=17.0),
            base.SvgTile(slide_num=2, name="S2_REIN_PLOT", x_in=8.18, y_in=2.08, w_in=4.14, h_in=2.60, svg_text=_reinforcement_plot_svg(w_px=base._in_to_px(4.14), h_px=base._in_to_px(2.60))),
            base.SvgTile(slide_num=2, name="S2_GAP_HEAT_BG", x_in=8.18, y_in=5.06, w_in=4.14, h_in=1.42, svg_text=_card_bg_svg(w_px=base._in_to_px(4.14), h_px=base._in_to_px(1.42), accent_hex=base.COLORS_HEX["coral"], title="gap field"), static=True),
            base.SvgTile(slide_num=2, name="S2_GAP_HEAT", x_in=9.02, y_in=5.56, w_in=2.46, h_in=0.70, svg_text=_mini_gap_heatmap_svg(w_px=base._in_to_px(2.46), h_px=base._in_to_px(0.70), accent_hex=base.COLORS_HEX["coral"])),
        ]
    )

    return tiles


def build_deck(*, out_pre: Path, out_svg: Path, out_animated: Path, effect_dur_ms: int, keep_intermediates: bool) -> None:
    out_pre = out_pre.resolve()
    out_svg = out_svg.resolve()
    out_animated = out_animated.resolve()
    out_pre.parent.mkdir(parents=True, exist_ok=True)
    out_svg.parent.mkdir(parents=True, exist_ok=True)
    out_animated.parent.mkdir(parents=True, exist_ok=True)

    tiles = _build_tiles()
    base._assert_tiles_within_slide_bounds(tiles)
    _assert_no_nonstatic_overlap(tiles)

    seen: set[str] = set()
    for tile in tiles:
        if tile.name in seen:
            raise ValueError(f"Duplicate tile name: {tile.name}")
        seen.add(tile.name)

    tmp_dir = out_pre.parent / "_observed_topk_ifx_fused_tmp"
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir, ignore_errors=True)
    svg_dir = tmp_dir / "svgs"
    png_dir = tmp_dir / "pngs"
    svg_dir.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)

    shape_to_svg: dict[str, Path] = {}
    shape_to_png: dict[str, Path] = {}
    for tile in tiles:
        svg_path = svg_dir / f"{tile.name}_s{tile.slide_num:02d}.svg"
        svg_path.write_text(tile.svg_text, encoding="utf-8")
        shape_to_svg[tile.name] = svg_path
        png_path = png_dir / f"{tile.name}_s{tile.slide_num:02d}.png"
        base._write_unique_png(out_path=png_path, seed=f"{tile.slide_num}:{tile.name}")
        shape_to_png[tile.name] = png_path

    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W_IN)
    prs.slide_height = Inches(base.SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    slides = [prs.slides.add_slide(blank) for _ in range(2)]
    for slide in slides:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = base._rgb(base.COLORS_HEX["bg"])

    for tile in tiles:
        slide = slides[tile.slide_num - 1]
        pic = slide.shapes.add_picture(
            str(shape_to_png[tile.name]),
            Inches(tile.x_in),
            Inches(tile.y_in),
            width=Inches(tile.w_in),
            height=Inches(tile.h_in),
        )
        pic.name = tile.name

    prs.save(str(out_pre))
    base._patch_pptx_replace_images_with_svg(pptx_in=out_pre, pptx_out=out_svg, shape_svg_map=shape_to_svg)

    repo_root = Path(__file__).resolve().parent
    slide2_dir = repo_root / "Slide2"
    sys.path.insert(0, str(slide2_dir))
    from pptx_click_animations import inject_click_reveals  # noqa: PLC0415

    exclude_by_slide = base._compute_exclude_spids_for_static_by_slide(out_svg)
    tmp_mid = out_animated.with_name(f"{out_animated.stem}__s1{out_animated.suffix}")
    res1 = inject_click_reveals(
        pptx_in=out_svg,
        pptx_out=tmp_mid,
        slide_nums={1},
        exclude_spids=exclude_by_slide.get(1, set()),
        effect_dur_ms=effect_dur_ms,
    )
    res2 = inject_click_reveals(
        pptx_in=tmp_mid,
        pptx_out=out_animated,
        slide_nums={2},
        exclude_spids=exclude_by_slide.get(2, set()),
        effect_dur_ms=effect_dur_ms,
    )
    try:
        tmp_mid.unlink()
    except OSError:
        pass

    clicks = {}
    clicks.update(res1.slide_click_effects)
    clicks.update(res2.slide_click_effects)
    print(f"OK: wrote {out_animated}")
    for slide_num, n in sorted(clicks.items()):
        print(f"OK: slide {slide_num} clickEffects={n}")
    print(f"OK: total clickEffects={sum(clicks.values())}")

    if not keep_intermediates:
        for path in [out_pre]:
            try:
                path.unlink()
            except OSError:
                pass


def main() -> None:
    parser = argparse.ArgumentParser(description="Fuse observed top-K deck with IFX poster into 2 SVG slides.")
    parser.add_argument("--out-pre", type=Path, default=Path("_build/observed_topk_ifx_fused_2slides_preanim.pptx"))
    parser.add_argument("--out-svg", type=Path, default=Path("_build/observed_topk_ifx_fused_2slides_raw.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/observed_topk_ifx_fused_2slides_animated.pptx"))
    parser.add_argument("--dur-ms", type=int, default=240)
    parser.add_argument("--keep-intermediates", action="store_true")
    args = parser.parse_args()

    build_deck(
        out_pre=args.out_pre,
        out_svg=args.out_svg,
        out_animated=args.out,
        effect_dur_ms=int(args.dur_ms),
        keep_intermediates=bool(args.keep_intermediates),
    )


if __name__ == "__main__":
    main()
