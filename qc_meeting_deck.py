#!/usr/bin/env python3
from __future__ import annotations

import argparse
import zipfile
from pathlib import Path

from lxml import etree
from pptx import Presentation

from qc_no_overlaps import check_no_overlaps


EMU_PER_INCH = 914400
PML_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


def _qn(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


def _check_slide_count(pptx_path: Path, *, expected: int) -> list[str]:
    prs = Presentation(str(pptx_path))
    if len(prs.slides) != expected:
        return [f"slide_count: expected={expected} actual={len(prs.slides)}"]
    return []


def _check_bounds(pptx_path: Path, *, margin_in: float) -> list[str]:
    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)
    margin = int(round(margin_in * EMU_PER_INCH))

    issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for sh in _iter_shapes(slide):
            w = int(getattr(sh, "width", 0))
            h = int(getattr(sh, "height", 0))
            if w <= 0 or h <= 0:
                continue
            left = int(getattr(sh, "left", 0))
            top = int(getattr(sh, "top", 0))
            right = left + w
            bottom = top + h
            if left < -margin or top < -margin or right > slide_w + margin or bottom > slide_h + margin:
                issues.append(
                    f"bounds: slide {slide_num:02d} id={int(getattr(sh,'shape_id',0))} "
                    f"name={(getattr(sh,'name','') or '').strip()} bbox=({left},{top})-({right},{bottom})"
                )
    return issues


def _check_safe_area(
    pptx_path: Path,
    *,
    safe_l_in: float,
    safe_r_in: float,
    safe_t_in: float,
    safe_b_in: float,
    tol_in: float = 0.00,
) -> list[str]:
    prs = Presentation(str(pptx_path))
    slide_w = int(prs.slide_width)
    slide_h = int(prs.slide_height)

    safe_l = int(round(safe_l_in * EMU_PER_INCH))
    safe_r = int(round(safe_r_in * EMU_PER_INCH))
    safe_t = int(round(safe_t_in * EMU_PER_INCH))
    safe_b = int(round(safe_b_in * EMU_PER_INCH))
    tol = int(round(tol_in * EMU_PER_INCH))

    x0 = safe_l - tol
    y0 = safe_t - tol
    x1 = slide_w - safe_r + tol
    y1 = slide_h - safe_b + tol

    issues: list[str] = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        for sh in _iter_shapes(slide):
            w = int(getattr(sh, "width", 0))
            h = int(getattr(sh, "height", 0))
            if w <= 0 or h <= 0:
                continue
            left = int(getattr(sh, "left", 0))
            top = int(getattr(sh, "top", 0))
            right = left + w
            bottom = top + h
            if left < x0 or top < y0 or right > x1 or bottom > y1:
                issues.append(
                    f"safe_area: slide {slide_num:02d} id={int(getattr(sh,'shape_id',0))} "
                    f"name={(getattr(sh,'name','') or '').strip()} bbox=({left},{top})-({right},{bottom})"
                )
    return issues


def _iter_slide_xmls(pptx_path: Path) -> dict[int, bytes]:
    out: dict[int, bytes] = {}
    with zipfile.ZipFile(pptx_path, "r") as zf:
        for info in zf.infolist():
            if not info.filename.startswith("ppt/slides/slide"):
                continue
            if not info.filename.endswith(".xml"):
                continue
            stem = info.filename.rsplit("/", 1)[-1]
            num_s = stem.removeprefix("slide").removesuffix(".xml")
            if not num_s.isdigit():
                continue
            out[int(num_s)] = zf.read(info.filename)
    return out


def _check_timing_present(pptx_path: Path, *, expected_slides: int) -> list[str]:
    xmls = _iter_slide_xmls(pptx_path)
    issues: list[str] = []
    for slide_num in range(1, expected_slides + 1):
        data = xmls.get(slide_num)
        if data is None:
            issues.append(f"timing: missing slide xml for slide {slide_num:02d}")
            continue
        root = etree.fromstring(data)
        if root.find(_qn(PML_NS, "timing")) is None:
            issues.append(f"timing: slide {slide_num:02d} missing <p:timing>")
    return issues


def _count_click_effects(pptx_path: Path) -> dict[int, int]:
    xmls = _iter_slide_xmls(pptx_path)
    out: dict[int, int] = {}
    ns = {"p": PML_NS}
    for slide_num, data in xmls.items():
        root = etree.fromstring(data)
        out[slide_num] = int(len(root.xpath(".//p:cTn[@nodeType='clickEffect']", namespaces=ns)))
    return out


def main() -> None:
    parser = argparse.ArgumentParser(description="QC for strict 6-slide meeting deck.")
    parser.add_argument("pptx", type=Path, nargs="?", default=Path("_build/bluesky_meeting_rqs_papers_data_6slides_strict_animated.pptx"))
    parser.add_argument("--expected-slides", type=int, default=6)
    parser.add_argument("--bounds-margin-in", type=float, default=0.02)
    parser.add_argument("--eq-mode", action="store_true", help="Whitelist EQ_STEP_* overlaps among themselves.")
    parser.add_argument("--eq-slide", type=int, default=0, help="If non-zero, require clickEffects >= --eq-target on this slide.")
    parser.add_argument("--eq-target", type=int, default=0)
    args = parser.parse_args()

    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        raise SystemExit(f"Missing pptx: {pptx_path}")

    errors: list[str] = []

    errors.extend(_check_slide_count(pptx_path, expected=int(args.expected_slides)))
    errors.extend(_check_bounds(pptx_path, margin_in=float(args.bounds_margin_in)))
    errors.extend(
        _check_safe_area(
            pptx_path,
            safe_l_in=0.70,
            safe_r_in=0.70,
            safe_t_in=0.45,
            safe_b_in=0.45,
            tol_in=0.00,
        )
    )
    errors.extend(_check_timing_present(pptx_path, expected_slides=int(args.expected_slides)))

    click_effects = _count_click_effects(pptx_path)
    s6 = click_effects.get(6, 0)
    if s6 < 80:
        errors.append(f"clickEffects: slide 06 expected>=80 actual={s6}")

    if int(args.eq_slide) and int(args.eq_target):
        eq = click_effects.get(int(args.eq_slide), 0)
        if eq < int(args.eq_target):
            errors.append(f"clickEffects: slide {int(args.eq_slide):02d} expected>={int(args.eq_target)} actual={eq}")

    errors.extend(check_no_overlaps(pptx_path=pptx_path, eq_mode=bool(args.eq_mode)))

    if errors:
        print("FAIL: QC errors")
        for line in errors:
            print(line)
        raise SystemExit(1)

    print("OK: QC passed")
    total = sum(click_effects.values())
    print(f"OK: total clickEffects={total} (slide06={s6})")


if __name__ == "__main__":
    main()

