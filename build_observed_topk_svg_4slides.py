#!/usr/bin/env python3
from __future__ import annotations

import argparse
import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import build_ifx_poster_enhanced_two_slide as theme
import build_ifx_poster_svg as base


theme._patch_theme()


def _mix(a: str, b: str, t: float) -> str:
    return theme._mix(a, b, t)


def _pill_svg(text: str, accent_hex: str, *, w_in: float = 2.55, h_in: float = 0.22) -> str:
    return base._svg_pill(w_px=base._in_to_px(w_in), h_px=base._in_to_px(h_in), text=text, accent_hex=accent_hex)


def _caption_svg(lines: list[str], *, accent_hex: str | None = None, w_in: float = 3.10, h_in: float = 0.72) -> str:
    return base._svg_caption(
        w_px=base._in_to_px(w_in),
        h_px=base._in_to_px(h_in),
        lines=lines,
        accent_hex=accent_hex,
    )


def _graph_frame_svg(*, w_px: int, h_px: int, accent_hex: str, title: str = "") -> str:
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=24,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="14" height="{h_px}" rx="24" ry="24" fill="#{accent_hex}" opacity="0.84"/>'
    body += f'<line x1="16" y1="1" x2="{w_px-18}" y2="1" stroke="#FFFFFF" stroke-opacity="0.10" stroke-width="2"/>'
    if title:
        body += base._svg_text(
            x=28,
            y=22,
            text=title,
            size_px=20,
            color=f"#{base.COLORS_HEX['ink']}",
            weight=900,
            baseline="hanging",
        )
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _rank_weight_bar_svg(*, label: str, frac: float, accent_hex: str, w_px: int, h_px: int) -> str:
    pad = 14
    frac = max(0.0, min(1.0, frac))
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=16,
        fill=f"#{base.COLORS_HEX['card2']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += base._svg_text(x=pad, y=int(h_px * 0.34), text=label, size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=800, baseline="middle")
    bx = pad
    by = int(h_px * 0.58)
    bw = w_px - 2 * pad
    bh = max(10, int(h_px * 0.20))
    fill_w = int(bw * frac)
    body += f'<rect x="{bx}" y="{by}" width="{bw}" height="{bh}" rx="10" ry="10" fill="#08111D" opacity="0.96"/>'
    body += f'<rect x="{bx}" y="{by}" width="{fill_w}" height="{bh}" rx="10" ry="10" fill="#{accent_hex}" opacity="0.92"/>'
    body += f'<circle cx="{bx + fill_w}" cy="{by + bh//2}" r="{max(4, bh//2 + 1)}" fill="#{accent_hex}" opacity="0.90"/>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _risk_timeline_base_svg(*, w_px: int, h_px: int) -> str:
    x0 = 52
    x1 = w_px - 48
    y = h_px // 2
    capture_x = int(x0 + 0.60 * (x1 - x0))
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=20,
        fill=f"#{base.COLORS_HEX['card2']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<line x1="{x0}" y1="{y}" x2="{x1}" y2="{y}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'
    body += f'<line x1="{capture_x}" y1="{y-72}" x2="{capture_x}" y2="{y+72}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="4"/>'
    body += f'<rect x="{capture_x-52}" y="{y-98}" width="104" height="28" rx="14" ry="14" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="2"/>'
    body += base._svg_text(x=capture_x, y=y - 84, text="capture q", size_px=14, color=f"#{base.COLORS_HEX['violet']}", weight=850, anchor="middle")
    body += base._svg_text(x=x0, y=y + 54, text="indexed_at", size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=800, baseline="hanging")
    body += base._svg_text(x=x1, y=y + 54, text="time", size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="end", baseline="hanging")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _risk_node_svg(*, label: str, accent_hex: str, eligible: bool, w_px: int = 66, h_px: int = 66) -> str:
    op = 0.95 if eligible else 0.35
    r = 22 if eligible else 18
    body = f'<circle cx="{w_px//2}" cy="{h_px//2-8}" r="{r}" fill="#{accent_hex}" opacity="{op:.2f}"/>'
    body += f'<circle cx="{w_px//2}" cy="{h_px//2-8}" r="{max(8, r-8)}" fill="#{base.COLORS_HEX["bg"]}" opacity="{op:.2f}"/>'
    body += base._svg_text(x=w_px // 2, y=h_px // 2 - 8, text=label, size_px=18, color=f"#{base.COLORS_HEX['ink']}", weight=900, anchor="middle")
    body += base._svg_text(x=w_px // 2, y=h_px - 8, text="in" if eligible else "out", size_px=13, color=f"#{accent_hex}", weight=850, anchor="middle")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _heat_row_svg(*, label: str, values: list[float], accent_hex: str, w_px: int, h_px: int) -> str:
    cell_w = 78
    gap = 10
    label_w = 34
    body = base._svg_text(x=16, y=h_px // 2, text=label, size_px=18, color=f"#{accent_hex}", weight=900, anchor="middle")
    x0 = label_w + 12
    for j, val in enumerate(values):
        x = x0 + j * (cell_w + gap)
        fill = _mix(base.COLORS_HEX["card2"], accent_hex, 0.14 + 0.70 * val)
        stroke = _mix(base.COLORS_HEX["card_line"], accent_hex, 0.34)
        body += f'<rect x="{x}" y="6" width="{cell_w}" height="{h_px-12}" rx="14" ry="14" fill="#{fill}" stroke="#{stroke}" stroke-width="2"/>'
        if val > 0:
            body += base._svg_text(x=x + cell_w // 2, y=h_px // 2, text=f"{val:.2f}", size_px=16, color=f"#{base.COLORS_HEX['ink']}", weight=850, anchor="middle")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _plot_axes_svg(*, w_px: int, h_px: int) -> str:
    x0 = 54
    y0 = h_px - 42
    x1 = w_px - 44
    y1 = 36
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=20,
        fill=f"#{base.COLORS_HEX['card2']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<line x1="{x0}" y1="{y1}" x2="{x0}" y2="{y0}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'
    body += f'<line x1="{x0}" y1="{y0}" x2="{x1}" y2="{y0}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'
    xs = [x0 + i * (x1 - x0) / 4 for i in range(5)]
    for i, x in enumerate(xs):
        body += f'<line x1="{int(x)}" y1="{y0}" x2="{int(x)}" y2="{y0+8}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="2"/>'
        body += base._svg_text(x=int(x), y=y0 + 12, text=f"t{i+1}", size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="middle", baseline="hanging")
    body += base._svg_text(x=x0, y=y1 - 12, text="share", size_px=14, color=f"#{base.COLORS_HEX['muted']}", weight=800, baseline="hanging")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _plot_line_svg(*, values: list[float], accent_hex: str, label: str, w_px: int, h_px: int) -> str:
    x0 = 54
    y0 = h_px - 42
    x1 = w_px - 44
    y1 = 36
    xs = [x0 + i * (x1 - x0) / 4 for i in range(5)]
    points = []
    for x, v in zip(xs, values, strict=True):
        y = y0 - v * (y0 - y1)
        points.append((x, y))
    poly = " ".join(f"{int(x)},{int(y)}" for x, y in points)
    body = f'<polyline fill="none" stroke="#{accent_hex}" stroke-width="4" points="{poly}"/>'
    for x, y in points:
        body += f'<circle cx="{int(x)}" cy="{int(y)}" r="6" fill="#{accent_hex}"/>'
    lx, ly = points[-1]
    body += base._svg_text(x=int(lx + 16), y=int(ly), text=label, size_px=18, color=f"#{accent_hex}", weight=900, baseline="middle")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _metric_arrow_svg(*, w_px: int, h_px: int) -> str:
    body = f'<rect x="0" y="0" width="74" height="28" rx="14" ry="14" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="2"/>'
    body += base._svg_text(x=37, y=14, text="M^-", size_px=16, color=f"#{base.COLORS_HEX['violet']}", weight=900, anchor="middle")
    body += f'<line x1="37" y1="28" x2="{w_px-18}" y2="{h_px-18}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="3" stroke-dasharray="8 6"/>'
    body += f'<polygon points="{w_px-26},{h_px-30} {w_px-10},{h_px-18} {w_px-30},{h_px-10}" fill="#{base.COLORS_HEX["violet"]}" opacity="0.72"/>'
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _formula_tile(*, slide_num: int, name: str, expr: str, x_in: float, y_in: float, max_w_in: float, color_hex: str, fontsize_pt: float) -> base.SvgTile:
    svg, w_in, h_in = base._mpl_math_svg(expr, color_hex=color_hex, fontsize_pt=fontsize_pt)
    w2, h2 = base._scale_to_fit(w_in=w_in, h_in=h_in, max_w_in=max_w_in)
    return base.SvgTile(slide_num=slide_num, name=name, x_in=x_in, y_in=y_in, w_in=w2, h_in=h2, svg_text=svg)


def _svg_topk_graph(*, w_px: int, h_px: int) -> str:
    pad = 18
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=24,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="14" height="{h_px}" rx="24" ry="24" fill="#{base.COLORS_HEX["cyan"]}" opacity="0.84"/>'
    body += f'<rect x="0" y="{int(h_px*0.60)}" width="14" height="{int(h_px*0.40)}" rx="24" ry="24" fill="#{base.COLORS_HEX["coral"]}" opacity="0.72"/>'
    body += base._svg_text(
        x=pad + 6,
        y=24,
        text="One query context q",
        size_px=16,
        color=f"#{base.COLORS_HEX['muted']}",
        weight=800,
        baseline="hanging",
    )
    body += base._svg_text(
        x=pad + 6,
        y=48,
        text="observed top-K panel",
        size_px=24,
        color=f"#{base.COLORS_HEX['ink']}",
        weight=900,
        baseline="hanging",
    )

    left_x = pad + 8
    row_w = int(w_px * 0.48)
    row_h = 58
    row_gap = 20
    row_y0 = 98
    ranks = [1, 7, 19]
    labels = ["A", "B", "C"]
    accents = [base.COLORS_HEX["cyan"], base.COLORS_HEX["green"], base.COLORS_HEX["amber"]]

    for idx, (rank, label, accent) in enumerate(zip(ranks, labels, accents, strict=True)):
        y = row_y0 + idx * (row_h + row_gap)
        badge_w = 62
        body += f'<rect x="{left_x}" y="{y}" width="{row_w}" height="{row_h}" rx="24" ry="24" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="2"/>'
        body += f'<rect x="{left_x}" y="{y}" width="{badge_w}" height="{row_h}" rx="24" ry="24" fill="#{accent}"/>'
        body += base._svg_text(x=left_x + badge_w // 2, y=y + row_h // 2, text=str(rank), size_px=24, color=f"#{base.COLORS_HEX['bg']}", weight=900, anchor="middle")
        body += base._svg_text(x=left_x + 88, y=y + row_h // 2, text=f"duplicate {label}", size_px=20, color=f"#{base.COLORS_HEX['ink']}", weight=800)

        bar_x = int(w_px * 0.64)
        bar_y = y + 16
        bar_w = int(w_px * 0.24)
        bar_h = 16
        frac = max(0.12, 1.0 / (1.0 + 0.18 * (rank - 1)))
        fill_w = int(bar_w * frac)
        body += f'<line x1="{left_x + row_w + 18}" y1="{y + row_h//2}" x2="{bar_x - 18}" y2="{y + row_h//2}" stroke="#{accent}" stroke-opacity="0.42" stroke-width="3"/>'
        body += f'<polygon points="{bar_x - 18},{y + row_h//2 - 8} {bar_x - 18},{y + row_h//2 + 8} {bar_x - 2},{y + row_h//2}" fill="#{accent}" opacity="0.60"/>'
        body += f'<rect x="{bar_x}" y="{bar_y}" width="{bar_w}" height="{bar_h}" rx="10" ry="10" fill="#08111D" opacity="0.96"/>'
        body += f'<rect x="{bar_x}" y="{bar_y}" width="{fill_w}" height="{bar_h}" rx="10" ry="10" fill="#{accent}" opacity="0.92"/>'
        body += f'<circle cx="{bar_x + fill_w}" cy="{bar_y + bar_h//2}" r="8" fill="#{accent}" opacity="0.90"/>'
        body += base._svg_text(x=bar_x, y=bar_y - 10, text="w(rank)", size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=750, baseline="hanging")

    body += f'<rect x="{int(w_px*0.61)}" y="{row_y0 - 22}" width="{int(w_px*0.29)}" height="32" rx="16" ry="16" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["amber"]}" stroke-width="2"/>'
    body += base._svg_text(
        x=int(w_px * 0.755),
        y=row_y0 - 6,
        text="opportunity proxy",
        size_px=15,
        color=f"#{base.COLORS_HEX['ink']}",
        weight=850,
        anchor="middle",
    )
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_risk_graph(*, w_px: int, h_px: int) -> str:
    pad = 18
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=24,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="14" height="{h_px}" rx="24" ry="24" fill="#{base.COLORS_HEX["cyan"]}" opacity="0.84"/>'
    body += base._svg_text(x=pad + 8, y=24, text="duplicate cluster c", size_px=24, color=f"#{base.COLORS_HEX['ink']}", weight=900, baseline="hanging")
    body += base._svg_text(x=pad + 8, y=52, text="only indexed versions enter the local risk set", size_px=15, color=f"#{base.COLORS_HEX['muted']}", weight=750, baseline="hanging")

    cx = 118
    cy = 182
    cr = 74
    body += f'<circle cx="{cx}" cy="{cy}" r="{cr}" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["cyan"]}" stroke-width="3"/>'
    body += f'<circle cx="{cx}" cy="{cy}" r="{cr-10}" fill="none" stroke="#{base.COLORS_HEX["cyan"]}" stroke-opacity="0.30" stroke-width="2" stroke-dasharray="8 8"/>'
    body += base._svg_text(x=cx, y=cy - 6, text="c", size_px=44, color=f"#{base.COLORS_HEX['ink']}", weight=900, anchor="middle")
    body += base._svg_text(x=cx, y=cy + 28, text="same text", size_px=16, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="middle")

    x0 = 246
    x1 = w_px - 48
    y_line = 182
    body += f'<line x1="{x0}" y1="{y_line}" x2="{x1}" y2="{y_line}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'

    capture_x = int(x0 + 0.58 * (x1 - x0))
    body += f'<line x1="{capture_x}" y1="{y_line - 104}" x2="{capture_x}" y2="{y_line + 96}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="4"/>'
    body += f'<rect x="{capture_x - 54}" y="{y_line - 128}" width="108" height="30" rx="15" ry="15" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="2"/>'
    body += base._svg_text(x=capture_x, y=y_line - 113, text="captured at q", size_px=14, color=f"#{base.COLORS_HEX['violet']}", weight=850, anchor="middle")

    nodes = [
        ("A", 0.12, base.COLORS_HEX["cyan"], True),
        ("B", 0.40, base.COLORS_HEX["green"], True),
        ("C", 0.82, base.COLORS_HEX["amber"], False),
    ]
    for label, frac, accent, eligible in nodes:
        x = int(x0 + frac * (x1 - x0))
        r = 22 if eligible else 18
        op = "0.95" if eligible else "0.35"
        body += f'<line x1="{cx + cr}" y1="{cy}" x2="{x - r - 8}" y2="{y_line}" stroke="#{accent}" stroke-opacity="0.22" stroke-width="2" stroke-dasharray="7 7"/>'
        body += f'<circle cx="{x}" cy="{y_line}" r="{r}" fill="#{accent}" opacity="{op}"/>'
        body += f'<circle cx="{x}" cy="{y_line}" r="{max(8, r-8)}" fill="#{base.COLORS_HEX["bg"]}" opacity="{op}"/>'
        body += base._svg_text(x=x, y=y_line, text=label, size_px=18, color=f"#{base.COLORS_HEX['ink']}", weight=900, anchor="middle")
        body += base._svg_text(x=x, y=y_line + 38, text="in" if eligible else "out", size_px=14, color=f"#{accent}", weight=850, anchor="middle")

    body += base._svg_text(x=x0, y=y_line + 86, text="indexed_at", size_px=14, color=f"#{base.COLORS_HEX['muted']}", weight=800, baseline="hanging")
    body += base._svg_text(x=x1, y=y_line + 86, text="time", size_px=14, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="end", baseline="hanging")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_aggregation_graph(*, w_px: int, h_px: int) -> str:
    pad = 18
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=24,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="14" height="{h_px}" rx="24" ry="24" fill="#{base.COLORS_HEX["amber"]}" opacity="0.82"/>'
    body += base._svg_text(x=pad + 8, y=24, text="aggregate across q in one hour", size_px=24, color=f"#{base.COLORS_HEX['ink']}", weight=900, baseline="hanging")

    grid_x = 48
    grid_y = 84
    cell_w = 88
    cell_h = 58
    rows = ["A", "B", "C"]
    cols = ["q1", "q2", "q3", "q4", "q5"]
    vals = [
        [0.82, 0.64, 0.44, 0.20, 0.00],
        [0.18, 0.46, 0.70, 0.58, 0.36],
        [0.00, 0.00, 0.16, 0.40, 0.74],
    ]
    accents = [base.COLORS_HEX["cyan"], base.COLORS_HEX["green"], base.COLORS_HEX["amber"]]

    for j, col in enumerate(cols):
        x = grid_x + j * (cell_w + 10) + cell_w // 2
        body += base._svg_text(x=x, y=grid_y - 18, text=col, size_px=14, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="middle")

    for i, row in enumerate(rows):
        y = grid_y + i * (cell_h + 12) + cell_h // 2
        body += base._svg_text(x=grid_x - 18, y=y, text=row, size_px=18, color=f"#{accents[i]}", weight=900, anchor="middle")
        for j, val in enumerate(vals[i]):
            x = grid_x + j * (cell_w + 10)
            fill = _mix(base.COLORS_HEX["card2"], accents[i], 0.12 + 0.66 * val)
            stroke = _mix(base.COLORS_HEX["card_line"], accents[i], 0.38)
            body += f'<rect x="{x}" y="{grid_y + i * (cell_h + 12)}" width="{cell_w}" height="{cell_h}" rx="14" ry="14" fill="#{fill}" stroke="#{stroke}" stroke-width="2"/>'
            if val > 0:
                body += base._svg_text(
                    x=x + cell_w // 2,
                    y=grid_y + i * (cell_h + 12) + cell_h // 2,
                    text=f"{val:.2f}",
                    size_px=16,
                    color=f"#{base.COLORS_HEX['ink']}",
                    weight=850,
                    anchor="middle",
                )

    sum_x = grid_x + len(cols) * (cell_w + 10) + 40
    body += base._svg_text(x=sum_x, y=grid_y - 18, text="O_it", size_px=16, color=f"#{base.COLORS_HEX['amber']}", weight=900, baseline="hanging")
    totals = [sum(r) for r in vals]
    max_total = max(totals)
    for i, total in enumerate(totals):
        y = grid_y + i * (cell_h + 12) + 15
        bar_w = 180
        fill_w = int(bar_w * (total / max_total))
        body += f'<rect x="{sum_x}" y="{y}" width="{bar_w}" height="26" rx="13" ry="13" fill="#08111D" opacity="0.96"/>'
        body += f'<rect x="{sum_x}" y="{y}" width="{fill_w}" height="26" rx="13" ry="13" fill="#{accents[i]}" opacity="0.90"/>'
        body += f'<circle cx="{sum_x + fill_w}" cy="{y + 13}" r="9" fill="#{accents[i]}" opacity="0.90"/>'
        body += base._svg_text(
            x=sum_x + bar_w + 48,
            y=y + 13,
            text=f"{total:.2f}",
            size_px=16,
            color=f"#{base.COLORS_HEX['ink']}",
            weight=850,
            anchor="middle",
        )

    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _svg_reinforcement_graph(*, w_px: int, h_px: int) -> str:
    pad = 18
    body = base._svg_rounded_rect(
        x=0,
        y=0,
        w=w_px,
        h=h_px,
        r=24,
        fill=f"#{base.COLORS_HEX['card']}",
        stroke=f"#{base.COLORS_HEX['card_line']}",
        stroke_w=2,
    )
    body += f'<rect x="0" y="0" width="14" height="{h_px}" rx="24" ry="24" fill="#{base.COLORS_HEX["violet"]}" opacity="0.84"/>'
    body += base._svg_text(x=pad + 8, y=24, text="within-content reinforcement", size_px=24, color=f"#{base.COLORS_HEX['ink']}", weight=900, baseline="hanging")

    x0 = 76
    y0 = h_px - 70
    x1 = w_px - 54
    y1 = 96
    body += f'<line x1="{x0}" y1="{y1}" x2="{x0}" y2="{y0}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'
    body += f'<line x1="{x0}" y1="{y0}" x2="{x1}" y2="{y0}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="3"/>'

    xs = [x0 + i * (x1 - x0) / 4 for i in range(5)]
    for i, x in enumerate(xs):
        body += f'<line x1="{int(x)}" y1="{y0}" x2="{int(x)}" y2="{y0 + 8}" stroke="#{base.COLORS_HEX["card_line"]}" stroke-width="2"/>'
        body += base._svg_text(x=int(x), y=y0 + 14, text=f"t{i+1}", size_px=13, color=f"#{base.COLORS_HEX['muted']}", weight=800, anchor="middle", baseline="hanging")

    lines = [
        ("A", base.COLORS_HEX["cyan"], [0.18, 0.34, 0.58, 0.78, 0.90]),
        ("B", base.COLORS_HEX["green"], [0.18, 0.22, 0.26, 0.29, 0.31]),
        ("C", base.COLORS_HEX["amber"], [0.18, 0.15, 0.12, 0.10, 0.08]),
    ]
    for label, accent, vals in lines:
        points = []
        for x, v in zip(xs, vals, strict=True):
            y = y0 - v * (y0 - y1)
            points.append((x, y))
        poly = " ".join(f"{int(x)},{int(y)}" for x, y in points)
        body += f'<polyline fill="none" stroke="#{accent}" stroke-width="4" points="{poly}"/>'
        for x, y in points:
            body += f'<circle cx="{int(x)}" cy="{int(y)}" r="6" fill="#{accent}"/>'
        lx, ly = points[-1]
        body += base._svg_text(x=int(lx + 16), y=int(ly), text=label, size_px=18, color=f"#{accent}", weight=900, baseline="middle")

    metric_x = int(xs[0])
    metric_y = int(y0 - 0.18 * (y0 - y1))
    body += f'<rect x="{metric_x - 38}" y="{metric_y - 58}" width="76" height="28" rx="14" ry="14" fill="#{base.COLORS_HEX["card2"]}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="2"/>'
    body += base._svg_text(x=metric_x, y=metric_y - 44, text="M^-", size_px=16, color=f"#{base.COLORS_HEX['violet']}", weight=900, anchor="middle")
    body += f'<line x1="{metric_x}" y1="{metric_y - 30}" x2="{int(xs[1])}" y2="{int(y0 - 0.34 * (y0 - y1)) - 8}" stroke="#{base.COLORS_HEX["violet"]}" stroke-width="3" stroke-dasharray="8 6"/>'
    body += f'<polygon points="{int(xs[1]) - 12},{int(y0 - 0.34 * (y0 - y1)) - 14} {int(xs[1]) - 2},{int(y0 - 0.34 * (y0 - y1)) - 2} {int(xs[1]) - 16},{int(y0 - 0.34 * (y0 - y1)) + 2}" fill="#{base.COLORS_HEX["violet"]}" opacity="0.70"/>'

    body += base._svg_text(x=x0, y=y1 - 16, text="opportunity share", size_px=14, color=f"#{base.COLORS_HEX['muted']}", weight=800, baseline="hanging")
    return base._svg_root(w_px=w_px, h_px=h_px, body=body)


def _build_tiles() -> list[base.SvgTile]:
    tiles: list[base.SvgTile] = []
    title_h = 0.92
    title_x = base.SAFE_L_IN
    title_w = base.SAFE_W_IN

    def add_bg(slide_num: int, title: str, subtitle: str) -> None:
        tiles.append(
            base.SvgTile(
                slide_num=slide_num,
                name=f"_BG_GRID_S{slide_num}",
                x_in=0.0,
                y_in=0.0,
                w_in=base.SLIDE_W_IN,
                h_in=base.SLIDE_H_IN,
                svg_text=base._svg_grid_bg(w_px=base._in_to_px(base.SLIDE_W_IN), h_px=base._in_to_px(base.SLIDE_H_IN)),
                static=True,
            )
        )
        tiles.append(
            base.SvgTile(
                slide_num=slide_num,
                name=f"STATIC_TITLE_S{slide_num}",
                x_in=title_x,
                y_in=0.28,
                w_in=title_w,
                h_in=title_h,
                svg_text=base._svg_title_block(
                    w_px=base._in_to_px(title_w),
                    h_px=base._in_to_px(title_h),
                    title=title,
                    subtitle=subtitle,
                ),
                static=True,
            )
        )

    add_bg(1, "Observed top-K opportunity", "formula -> panel -> weight")
    tiles.append(base.SvgTile(slide_num=1, name="STATIC_S1_FRAME", x_in=0.96, y_in=2.44, w_in=11.40, h_in=4.34, svg_text=_graph_frame_svg(w_px=base._in_to_px(11.40), h_px=base._in_to_px(4.34), accent_hex=base.COLORS_HEX["amber"], title="q -> observed top-K -> weight by rank"), static=True))
    tiles.append(base.SvgTile(slide_num=1, name="S1_PILL_FORMULA", x_in=1.04, y_in=1.46, w_in=2.50, h_in=0.22, svg_text=_pill_svg("formula", base.COLORS_HEX["amber"], w_in=2.50)))
    tiles.append(_formula_tile(slide_num=1, name="S1_EQ", expr=r"$O_{i q}=I_{i q}\cdot w(\mathrm{rank}_{i q})$", x_in=1.02, y_in=1.82, max_w_in=7.40, color_hex=base.COLORS_HEX["amber"], fontsize_pt=28.0))
    tiles.append(base.SvgTile(slide_num=1, name="S1_PILL_LIST", x_in=1.20, y_in=2.72, w_in=1.92, h_in=0.22, svg_text=_pill_svg("top-K", base.COLORS_HEX["coral"], w_in=1.92)))
    tiles.append(base.SvgTile(slide_num=1, name="S1_ROW_A", x_in=1.20, y_in=3.18, w_in=4.85, h_in=0.54, svg_text=base._svg_rank_row(w_px=base._in_to_px(4.85), h_px=base._in_to_px(0.54), rank=1, label="duplicate A", accent_hex=base.COLORS_HEX["cyan"])))
    tiles.append(base.SvgTile(slide_num=1, name="S1_ROW_B", x_in=1.20, y_in=3.82, w_in=4.85, h_in=0.54, svg_text=base._svg_rank_row(w_px=base._in_to_px(4.85), h_px=base._in_to_px(0.54), rank=7, label="duplicate B", accent_hex=base.COLORS_HEX["green"])))
    tiles.append(base.SvgTile(slide_num=1, name="S1_ROW_C", x_in=1.20, y_in=4.46, w_in=4.85, h_in=0.54, svg_text=base._svg_rank_row(w_px=base._in_to_px(4.85), h_px=base._in_to_px(0.54), rank=19, label="duplicate C", accent_hex=base.COLORS_HEX["amber"])))
    tiles.append(base.SvgTile(slide_num=1, name="S1_ARROW", x_in=6.30, y_in=3.70, w_in=1.15, h_in=0.52, svg_text=base._svg_arrow(w_px=base._in_to_px(1.15), h_px=base._in_to_px(0.52), accent_hex=base.COLORS_HEX["amber"], text="w(r)")))
    tiles.append(base.SvgTile(slide_num=1, name="S1_BAR_A", x_in=7.84, y_in=3.08, w_in=3.20, h_in=0.86, svg_text=_rank_weight_bar_svg(label="A", frac=0.96, accent_hex=base.COLORS_HEX["cyan"], w_px=base._in_to_px(3.20), h_px=base._in_to_px(0.86))))
    tiles.append(base.SvgTile(slide_num=1, name="S1_BAR_B", x_in=7.84, y_in=3.96, w_in=3.20, h_in=0.86, svg_text=_rank_weight_bar_svg(label="B", frac=0.44, accent_hex=base.COLORS_HEX["green"], w_px=base._in_to_px(3.20), h_px=base._in_to_px(0.86))))
    tiles.append(base.SvgTile(slide_num=1, name="S1_BAR_C", x_in=7.84, y_in=4.84, w_in=3.20, h_in=0.86, svg_text=_rank_weight_bar_svg(label="C", frac=0.18, accent_hex=base.COLORS_HEX["amber"], w_px=base._in_to_px(3.20), h_px=base._in_to_px(0.86))))
    tiles.append(base.SvgTile(slide_num=1, name="S1_TAKE", x_in=9.08, y_in=1.52, w_in=3.20, h_in=0.72, svg_text=_caption_svg(["opportunity != attention"], accent_hex=base.COLORS_HEX["amber"], w_in=3.20, h_in=0.72)))

    add_bg(2, "Local risk set", "cluster -> capture -> compare")
    tiles.append(base.SvgTile(slide_num=2, name="STATIC_S2_FRAME", x_in=0.96, y_in=2.44, w_in=11.40, h_in=4.34, svg_text=_graph_frame_svg(w_px=base._in_to_px(11.40), h_px=base._in_to_px(4.34), accent_hex=base.COLORS_HEX["cyan"], title="cluster first, capture second"), static=True))
    tiles.append(base.SvgTile(slide_num=2, name="S2_PILL_FORMULA", x_in=1.04, y_in=1.46, w_in=2.50, h_in=0.22, svg_text=_pill_svg("formula", base.COLORS_HEX["cyan"], w_in=2.50)))
    tiles.append(_formula_tile(slide_num=2, name="S2_EQ", expr=r"$R_{c q}=\{i\in c:\ \mathrm{indexed\_at}_i\leq \mathrm{captured\_at}_q\}$", x_in=1.02, y_in=1.82, max_w_in=8.15, color_hex=base.COLORS_HEX["cyan"], fontsize_pt=24.0))
    tiles.append(base.SvgTile(slide_num=2, name="S2_CLUSTER", x_in=1.20, y_in=3.14, w_in=2.40, h_in=2.18, svg_text=base._svg_cluster_bubble(w_px=base._in_to_px(2.40), h_px=base._in_to_px(2.18), label="cluster c", accent_hex=base.COLORS_HEX["cyan"])))
    tiles.append(base.SvgTile(slide_num=2, name="S2_TIMELINE", x_in=4.46, y_in=3.24, w_in=6.78, h_in=1.98, svg_text=_risk_timeline_base_svg(w_px=base._in_to_px(6.78), h_px=base._in_to_px(1.98))))
    tiles.append(base.SvgTile(slide_num=2, name="S2_NODE_A", x_in=5.08, y_in=3.70, w_in=0.72, h_in=0.72, svg_text=_risk_node_svg(label="A", accent_hex=base.COLORS_HEX["cyan"], eligible=True, w_px=base._in_to_px(0.72), h_px=base._in_to_px(0.72))))
    tiles.append(base.SvgTile(slide_num=2, name="S2_NODE_B", x_in=6.66, y_in=3.70, w_in=0.72, h_in=0.72, svg_text=_risk_node_svg(label="B", accent_hex=base.COLORS_HEX["green"], eligible=True, w_px=base._in_to_px(0.72), h_px=base._in_to_px(0.72))))
    tiles.append(base.SvgTile(slide_num=2, name="S2_NODE_C", x_in=9.10, y_in=3.70, w_in=0.72, h_in=0.72, svg_text=_risk_node_svg(label="C", accent_hex=base.COLORS_HEX["amber"], eligible=False, w_px=base._in_to_px(0.72), h_px=base._in_to_px(0.72))))
    tiles.append(base.SvgTile(slide_num=2, name="S2_TAKE", x_in=9.18, y_in=1.52, w_in=3.10, h_in=0.72, svg_text=_caption_svg(["local competition"], accent_hex=base.COLORS_HEX["cyan"], w_in=3.10, h_in=0.72)))

    add_bg(3, "Aggregate over q", "sum over q")
    tiles.append(base.SvgTile(slide_num=3, name="STATIC_S3_FRAME", x_in=0.96, y_in=2.44, w_in=11.40, h_in=4.34, svg_text=_graph_frame_svg(w_px=base._in_to_px(11.40), h_px=base._in_to_px(4.34), accent_hex=base.COLORS_HEX["amber"], title="heat by q, then sum by post"), static=True))
    tiles.append(base.SvgTile(slide_num=3, name="S3_PILL_FORMULA", x_in=1.04, y_in=1.46, w_in=2.50, h_in=0.22, svg_text=_pill_svg("formula", base.COLORS_HEX["amber"], w_in=2.50)))
    tiles.append(_formula_tile(slide_num=3, name="S3_EQ", expr=r"$O_{i t}=\sum_{q\in \mathcal{Q}_t} O_{i q}$", x_in=1.02, y_in=1.70, max_w_in=7.40, color_hex=base.COLORS_HEX["amber"], fontsize_pt=26.0))
    for idx, q in enumerate(["q1", "q2", "q3", "q4", "q5"]):
        tiles.append(base.SvgTile(slide_num=3, name=f"S3_Q_{idx+1}", x_in=2.02 + idx * 1.02, y_in=2.82, w_in=0.82, h_in=0.22, svg_text=_pill_svg(q, base.COLORS_HEX["coral"], w_in=0.82, h_in=0.22)))
    tiles.append(base.SvgTile(slide_num=3, name="S3_ROW_A", x_in=1.36, y_in=3.28, w_in=5.95, h_in=0.64, svg_text=_heat_row_svg(label="A", values=[0.82, 0.64, 0.44, 0.20, 0.00], accent_hex=base.COLORS_HEX["cyan"], w_px=base._in_to_px(5.95), h_px=base._in_to_px(0.64))))
    tiles.append(base.SvgTile(slide_num=3, name="S3_ROW_B", x_in=1.36, y_in=4.06, w_in=5.95, h_in=0.64, svg_text=_heat_row_svg(label="B", values=[0.18, 0.46, 0.70, 0.58, 0.36], accent_hex=base.COLORS_HEX["green"], w_px=base._in_to_px(5.95), h_px=base._in_to_px(0.64))))
    tiles.append(base.SvgTile(slide_num=3, name="S3_ROW_C", x_in=1.36, y_in=4.84, w_in=5.95, h_in=0.64, svg_text=_heat_row_svg(label="C", values=[0.00, 0.00, 0.16, 0.40, 0.74], accent_hex=base.COLORS_HEX["amber"], w_px=base._in_to_px(5.95), h_px=base._in_to_px(0.64))))
    tiles.append(base.SvgTile(slide_num=3, name="S3_BAR_A", x_in=8.16, y_in=3.16, w_in=2.76, h_in=0.90, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.90), label="O_A,t", frac=0.80, accent_hex=base.COLORS_HEX["cyan"], value="2.10")))
    tiles.append(base.SvgTile(slide_num=3, name="S3_BAR_B", x_in=8.16, y_in=4.02, w_in=2.76, h_in=0.90, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.90), label="O_B,t", frac=0.87, accent_hex=base.COLORS_HEX["green"], value="2.28")))
    tiles.append(base.SvgTile(slide_num=3, name="S3_BAR_C", x_in=8.16, y_in=4.88, w_in=2.76, h_in=0.90, svg_text=base._svg_scalar_bar(w_px=base._in_to_px(2.76), h_px=base._in_to_px(0.90), label="O_C,t", frac=0.50, accent_hex=base.COLORS_HEX["amber"], value="1.30")))
    tiles.append(base.SvgTile(slide_num=3, name="S3_TAKE", x_in=9.26, y_in=1.52, w_in=3.02, h_in=0.72, svg_text=_caption_svg(["sum over q"], accent_hex=base.COLORS_HEX["amber"], w_in=3.02, h_in=0.72)))

    add_bg(4, "Reinforcement", "feedback loop")
    tiles.append(base.SvgTile(slide_num=4, name="STATIC_S4_FRAME", x_in=0.96, y_in=2.44, w_in=11.40, h_in=4.34, svg_text=_graph_frame_svg(w_px=base._in_to_px(11.40), h_px=base._in_to_px(4.34), accent_hex=base.COLORS_HEX["violet"], title="if rho > 0, early opportunity compounds"), static=True))
    tiles.append(base.SvgTile(slide_num=4, name="S4_PILL_FORMULA", x_in=1.04, y_in=1.46, w_in=2.50, h_in=0.22, svg_text=_pill_svg("formula", base.COLORS_HEX["violet"], w_in=2.50)))
    tiles.append(_formula_tile(slide_num=4, name="S4_EQ", expr=r"$s_{i c,t+1}=\rho\, s_{i c t}+\gamma\, M^{-}_{i c t}+\cdots$", x_in=1.02, y_in=1.82, max_w_in=7.35, color_hex=base.COLORS_HEX["violet"], fontsize_pt=28.0))
    tiles.append(base.SvgTile(slide_num=4, name="S4_AXES", x_in=1.36, y_in=3.04, w_in=7.72, h_in=2.82, svg_text=_plot_axes_svg(w_px=base._in_to_px(7.72), h_px=base._in_to_px(2.82))))
    tiles.append(base.SvgTile(slide_num=4, name="S4_METRIC", x_in=1.74, y_in=2.86, w_in=2.18, h_in=1.08, svg_text=_metric_arrow_svg(w_px=base._in_to_px(2.18), h_px=base._in_to_px(1.08))))
    tiles.append(base.SvgTile(slide_num=4, name="S4_LINE_A", x_in=1.36, y_in=3.04, w_in=7.72, h_in=2.82, svg_text=_plot_line_svg(values=[0.18, 0.34, 0.58, 0.78, 0.90], accent_hex=base.COLORS_HEX["cyan"], label="A", w_px=base._in_to_px(7.72), h_px=base._in_to_px(2.82))))
    tiles.append(base.SvgTile(slide_num=4, name="S4_LINE_B", x_in=1.36, y_in=3.04, w_in=7.72, h_in=2.82, svg_text=_plot_line_svg(values=[0.18, 0.22, 0.26, 0.29, 0.31], accent_hex=base.COLORS_HEX["green"], label="B", w_px=base._in_to_px(7.72), h_px=base._in_to_px(2.82))))
    tiles.append(base.SvgTile(slide_num=4, name="S4_LINE_C", x_in=1.36, y_in=3.04, w_in=7.72, h_in=2.82, svg_text=_plot_line_svg(values=[0.18, 0.15, 0.12, 0.10, 0.08], accent_hex=base.COLORS_HEX["amber"], label="C", w_px=base._in_to_px(7.72), h_px=base._in_to_px(2.82))))
    tiles.append(base.SvgTile(slide_num=4, name="S4_TAKE", x_in=9.26, y_in=1.52, w_in=3.02, h_in=0.72, svg_text=_caption_svg(["rho > 0"], accent_hex=base.COLORS_HEX["violet"], w_in=3.02, h_in=0.72)))

    return tiles


def build_deck(*, out_pre: Path, out_svg: Path, out_animated: Path, effect_dur_ms: int, keep_intermediates: bool) -> None:
    out_pre = out_pre.resolve()
    out_svg = out_svg.resolve()
    out_animated = out_animated.resolve()
    out_pre.parent.mkdir(parents=True, exist_ok=True)
    out_svg.parent.mkdir(parents=True, exist_ok=True)
    out_animated.parent.mkdir(parents=True, exist_ok=True)

    tiles = _build_tiles()
    base._assert_tiles_within_slide_bounds(tiles)

    seen: set[str] = set()
    for tile in tiles:
        if tile.name in seen:
            raise ValueError(f"Duplicate tile name: {tile.name}")
        seen.add(tile.name)

    tmp_dir = out_pre.parent / "_observed_topk_svg_tmp"
    if tmp_dir.exists():
        shutil.rmtree(tmp_dir, ignore_errors=True)
    svg_dir = tmp_dir / "svgs"
    png_dir = tmp_dir / "pngs"
    svg_dir.mkdir(parents=True, exist_ok=True)
    png_dir.mkdir(parents=True, exist_ok=True)

    shape_to_svg: dict[str, Path] = {}
    shape_to_png: dict[str, Path] = {}
    for tile in tiles:
        svg_path = svg_dir / f"{tile.name}_s{tile.slide_num:02d}.svg"
        svg_path.write_text(tile.svg_text, encoding="utf-8")
        shape_to_svg[tile.name] = svg_path
        png_path = png_dir / f"{tile.name}_s{tile.slide_num:02d}.png"
        base._write_unique_png(out_path=png_path, seed=f"{tile.slide_num}:{tile.name}")
        shape_to_png[tile.name] = png_path

    prs = Presentation()
    prs.slide_width = Inches(base.SLIDE_W_IN)
    prs.slide_height = Inches(base.SLIDE_H_IN)
    blank = prs.slide_layouts[6]
    slides = [prs.slides.add_slide(blank) for _ in range(max(tile.slide_num for tile in tiles))]

    for slide in slides:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = base._rgb(base.COLORS_HEX["bg"])

    for tile in tiles:
        slide = slides[tile.slide_num - 1]
        pic = slide.shapes.add_picture(
            str(shape_to_png[tile.name]),
            Inches(tile.x_in),
            Inches(tile.y_in),
            width=Inches(tile.w_in),
            height=Inches(tile.h_in),
        )
        pic.name = tile.name

    prs.save(str(out_pre))
    base._patch_pptx_replace_images_with_svg(pptx_in=out_pre, pptx_out=out_svg, shape_svg_map=shape_to_svg)

    repo_root = Path(__file__).resolve().parent
    slide2_dir = repo_root / "Slide2"
    sys.path.insert(0, str(slide2_dir))
    from pptx_click_animations import inject_click_reveals  # noqa: PLC0415

    exclude_by_slide = base._compute_exclude_spids_for_static_by_slide(out_svg)
    current_in = out_svg
    temp_outputs: list[Path] = []
    total_clicks: dict[int, int] = {}

    slide_count = len(slides)
    for slide_num in range(1, slide_count + 1):
        current_out = out_animated if slide_num == slide_count else out_animated.with_name(f"{out_animated.stem}__s{slide_num:02d}{out_animated.suffix}")
        res = inject_click_reveals(
            pptx_in=current_in,
            pptx_out=current_out,
            slide_nums={slide_num},
            exclude_spids=exclude_by_slide.get(slide_num, set()),
            effect_dur_ms=effect_dur_ms,
        )
        total_clicks.update(res.slide_click_effects)
        if current_in != out_svg:
            temp_outputs.append(current_in)
        current_in = current_out

    for temp in temp_outputs:
        try:
            temp.unlink()
        except OSError:
            pass

    print(f"OK: wrote {out_animated}")
    for slide_num, clicks in sorted(total_clicks.items()):
        print(f"OK: slide {slide_num} clickEffects={clicks}")
    print(f"OK: total clickEffects={sum(total_clicks.values())}")

    if not keep_intermediates:
        for path in [out_pre]:
            try:
                path.unlink()
            except OSError:
                pass


def main() -> None:
    parser = argparse.ArgumentParser(description="Build a 4-slide SVG-style observed top-K deck.")
    parser.add_argument("--out-pre", type=Path, default=Path("_build/observed_topk_svg_4slides_preanim.pptx"))
    parser.add_argument("--out-svg", type=Path, default=Path("_build/observed_topk_svg_4slides_raw.pptx"))
    parser.add_argument("--out", type=Path, default=Path("_build/observed_topk_svg_4slides_animated.pptx"))
    parser.add_argument("--dur-ms", type=int, default=260)
    parser.add_argument("--keep-intermediates", action="store_true")
    args = parser.parse_args()

    build_deck(
        out_pre=args.out_pre,
        out_svg=args.out_svg,
        out_animated=args.out,
        effect_dur_ms=int(args.dur_ms),
        keep_intermediates=bool(args.keep_intermediates),
    )


if __name__ == "__main__":
    main()
