#!/usr/bin/env python3
from __future__ import annotations

import argparse
import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


# Canvas
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# Safe area (nothing may cross)
SAFE_L_IN = 0.70
SAFE_R_IN = 0.70
SAFE_T_IN = 0.45
SAFE_B_IN = 0.45

# Grid
GRID_COLS = 12
GUTTER_IN = 0.22
Y_STEP_IN = 0.12

# Cards
CARD_PAD_IN = 0.18

FONT_NAME = "Calibri"

COLORS_HEX = {
    "bg": "0B1320",
    "ink": "F7F9FC",
    "muted": "C9D2E3",
    "card": "162033",
    "card_line": "3A4D66",
    "cyan": "58B7E6",
    "amber": "F4C55D",
    "coral": "FF7D77",
    "green": "54C687",
}


def _rgb(hex6: str) -> RGBColor:
    return RGBColor.from_string(hex6)


SAFE_W_IN = SLIDE_W_IN - SAFE_L_IN - SAFE_R_IN
SAFE_H_IN = SLIDE_H_IN - SAFE_T_IN - SAFE_B_IN
COL_W_IN = (SAFE_W_IN - (GRID_COLS - 1) * GUTTER_IN) / GRID_COLS


def gx(col_idx: int) -> float:
    return SAFE_L_IN + col_idx * (COL_W_IN + GUTTER_IN)


def gw(col_span: int) -> float:
    if col_span <= 0:
        raise ValueError("col_span must be >= 1")
    return col_span * COL_W_IN + (col_span - 1) * GUTTER_IN


def gy(step: int) -> float:
    return SAFE_T_IN + step * Y_STEP_IN


def gh(steps: int) -> float:
    return steps * Y_STEP_IN


def _set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(COLORS_HEX["bg"])


def _style_text_frame(tf, *, margin_l: float = 0.18, margin_r: float = 0.18, margin_t: float = 0.10, margin_b: float = 0.10) -> None:
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(margin_l)
    tf.margin_right = Inches(margin_r)
    tf.margin_top = Inches(margin_t)
    tf.margin_bottom = Inches(margin_b)


def _set_text(
    shape,
    *,
    text: str,
    size_pt: int,
    bold: bool = False,
    color_hex: str = "F7F9FC",
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> None:
    tf = shape.text_frame
    tf.clear()
    _style_text_frame(tf)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    font = run.font
    font.name = FONT_NAME
    font.size = Pt(size_pt)
    font.bold = bold
    font.color.rgb = _rgb(color_hex)


def _add_title(slide, *, text: str, y: float, h: float, name: str = "TITLE") -> None:
    title = slide.shapes.add_textbox(Inches(gx(0)), Inches(y), Inches(gw(12)), Inches(h))
    title.name = name
    _set_text(title, text=text, size_pt=42, bold=True, color_hex=COLORS_HEX["ink"])


def _add_takeaway(slide, *, text: str, y: float, h: float, name: str = "TAKEAWAY") -> None:
    t = slide.shapes.add_textbox(Inches(gx(0)), Inches(y), Inches(gw(12)), Inches(h))
    t.name = name
    _set_text(t, text=text, size_pt=24, bold=False, color_hex=COLORS_HEX["muted"])


def _add_round_rect(slide, *, x: float, y: float, w: float, h: float, fill_hex: str, line_hex: str, name: str) -> object:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.name = name
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill_hex)
    shape.line.color.rgb = _rgb(line_hex)
    shape.line.width = Inches(0.02)
    return shape


def _add_rqs_rows(
    slide,
    *,
    y0: float,
    row_h: float,
    row_gap: float,
    accent_hex: str,
    items: list[str],
) -> None:
    if len(items) != 5:
        raise ValueError("RQS slide requires exactly 5 rows")

    for idx, slot_text in enumerate(items, start=1):
        y = y0 + (idx - 1) * (row_h + row_gap)

        num = _add_round_rect(
            slide,
            x=gx(0),
            y=y,
            w=gw(1),
            h=row_h,
            fill_hex=accent_hex,
            line_hex=accent_hex,
            name=f"RQ_NUM_{idx:02d}",
        )
        _style_text_frame(num.text_frame, margin_l=0.0, margin_r=0.0, margin_t=0.10, margin_b=0.10)
        _set_text(num, text=f"R{idx}", size_pt=20, bold=True, color_hex=COLORS_HEX["bg"], align=PP_ALIGN.CENTER)

        slot = _add_round_rect(
            slide,
            x=gx(1),
            y=y,
            w=gw(11),
            h=row_h,
            fill_hex=COLORS_HEX["card"],
            line_hex=COLORS_HEX["card_line"],
            name=f"RQ_SLOT_{idx:02d}",
        )
        _set_text(slot, text=slot_text, size_pt=20, bold=False, color_hex=COLORS_HEX["ink"])


def _add_cred_cards(slide, *, title: str, left_lines: list[str], right_lines: list[str]) -> None:
    _add_title(slide, text=title, y=gy(0), h=gh(6), name="TITLE")

    card_y = gy(8)
    card_h = gh(42)  # 5.04in

    left = _add_round_rect(
        slide,
        x=gx(0),
        y=card_y,
        w=gw(6),
        h=card_h,
        fill_hex=COLORS_HEX["card"],
        line_hex=COLORS_HEX["card_line"],
        name="CARD_LEFT",
    )
    _set_text(left, text="\n".join(left_lines), size_pt=19, bold=False, color_hex=COLORS_HEX["ink"])

    right = _add_round_rect(
        slide,
        x=gx(6),
        y=card_y,
        w=gw(6),
        h=card_h,
        fill_hex=COLORS_HEX["card"],
        line_hex=COLORS_HEX["card_line"],
        name="CARD_RIGHT",
    )
    _set_text(right, text="\n".join(right_lines), size_pt=19, bold=False, color_hex=COLORS_HEX["ink"])


def _add_artifacts_grid(
    slide,
    *,
    title: str,
    cards: list[tuple[str, list[tuple[str, str]]]],
) -> None:
    # Slide 6 layout: 2x2 cards, 10 label/value rows per card; label+value are separate shapes.
    if len(cards) != 4:
        raise ValueError("Need exactly 4 cards for a 2x2 grid")
    for _, rows in cards:
        if len(rows) != 10:
            raise ValueError("Each artifacts card must have exactly 10 rows")

    # Title (static; excluded from legacy animation injection later).
    title_shape = slide.shapes.add_textbox(Inches(gx(0)), Inches(gy(0)), Inches(gw(12)), Inches(gh(5)))
    title_shape.name = "TITLE"
    _set_text(title_shape, text=title, size_pt=40, bold=True, color_hex=COLORS_HEX["ink"])

    # Card grid geometry (snap to columns; y in 0.12 steps)
    grid_y = gy(6)  # directly below title
    card_h = gh(24)  # 2.88in
    row_gap = gh(1)  # 0.12in between top/bottom rows

    card_pos = [
        (gx(0), grid_y),  # TL
        (gx(6), grid_y),  # TR
        (gx(0), grid_y + card_h + row_gap),  # BL
        (gx(6), grid_y + card_h + row_gap),  # BR
    ]

    # Row layout inside a card
    header_h = gh(3)  # 0.36in
    header_gap = gh(1)  # 0.12in
    row_h = gh(2)  # 0.24in
    row_start_dy = header_h + header_gap

    pad_x = CARD_PAD_IN
    gap_x = 0.12
    label_w = 2.10

    for card_idx, ((card_title, rows), (cx, cy)) in enumerate(zip(cards, card_pos), start=1):
        bg = _add_round_rect(
            slide,
            x=cx,
            y=cy,
            w=gw(6),
            h=card_h,
            fill_hex=COLORS_HEX["card"],
            line_hex=COLORS_HEX["card_line"],
            name=f"_BG_CARD_{card_idx:02d}",
        )
        # Header text lives in the background shape (static).
        _style_text_frame(bg.text_frame, margin_l=0.18, margin_r=0.18, margin_t=0.12, margin_b=0.0)
        _set_text(bg, text=card_title, size_pt=18, bold=True, color_hex=COLORS_HEX["muted"])

        inner_x = round(cx + pad_x, 6)
        inner_w = round(gw(6) - 2 * pad_x, 6)
        value_w = inner_w - label_w - gap_x
        value_x = round(inner_x + label_w + gap_x, 6)

        for row_idx, (lab, val) in enumerate(rows, start=1):
            ry = round(cy + row_start_dy + (row_idx - 1) * row_h, 6)

            lab_sh = slide.shapes.add_textbox(Inches(inner_x), Inches(ry), Inches(label_w), Inches(row_h))
            lab_sh.name = f"S6_CARD{card_idx:02d}_R{row_idx:02d}_LABEL"
            _style_text_frame(lab_sh.text_frame, margin_l=0.0, margin_r=0.0, margin_t=0.0, margin_b=0.0)
            _set_text(lab_sh, text=lab, size_pt=16, bold=True, color_hex=COLORS_HEX["muted"])

            val_sh = slide.shapes.add_textbox(Inches(value_x), Inches(ry), Inches(value_w), Inches(row_h))
            val_sh.name = f"S6_CARD{card_idx:02d}_R{row_idx:02d}_VALUE"
            _style_text_frame(val_sh.text_frame, margin_l=0.0, margin_r=0.0, margin_t=0.0, margin_b=0.0)
            _set_text(val_sh, text=val, size_pt=16, bold=False, color_hex=COLORS_HEX["ink"])


@dataclass(frozen=True)
class ReportMetrics:
    generated_at_utc: str
    non_artifact_files: str
    non_artifact_size: str
    artifact_files: str
    artifact_size: str
    dir_size_gib: dict[str, str]
    hourly_runs: str
    hourly_rows_feed_items: str
    hourly_rows_post_labels: str
    hourly_rows_posts_first_seen: str
    wide_runs: str
    wide_rows_feed_items: str
    wide_rows_post_labels: str
    wide_rows_posts_first_seen: str
    labelerexp_runs: str
    labelerexp_rows_post_labels: str
    label_scan_rows: str
    label_scan_parts: str
    label_val_rows: dict[str, str]
    top_label_src_did: str
    top_label_src_rows: str
    control_state_db_size: str
    control_table_rows: dict[str, str]


def _re1(pattern: str, text: str, *, group: int = 1) -> str:
    m = re.search(pattern, text, flags=re.MULTILINE)
    if not m:
        raise ValueError(f"Missing pattern in REPORT.md: {pattern}")
    return m.group(group).strip()


def _parse_report(report_path: Path) -> ReportMetrics:
    raw = report_path.read_text(encoding="utf-8")

    generated_at_utc = _re1(r"Generated at \(UTC\): `([^`]+)`", raw)
    non_artifact_files = _re1(r"Non-artifact files: \*\*([0-9,]+)\*\*", raw)
    non_artifact_size = _re1(r"Non-artifact size: \*\*([0-9.]+ GiB)\*\*", raw)
    artifact_files = _re1(r"Artifact files\s*\(`\._\*`, `\.DS_Store`\): \*\*([0-9,]+)\*\*", raw)
    artifact_size = _re1(r"Artifact files\s*\(`\._\*`, `\.DS_Store`\): \*\*[0-9,]+\*\* \(([^)]+)\)", raw)

    dir_size_gib: dict[str, str] = {}
    for d in ["hourly", "effective_csv", "exports", "wide", "labelerexp"]:
        dir_size_gib[d] = _re1(rf"\| `{re.escape(d)}` \| [^|]+ \| ([0-9.]+) \|", raw)

    hourly_runs = _re1(r"### `hourly`\s+\n\s*- Runs: ([0-9,]+) \(", raw)
    hourly_rows_feed_items = _re1(r"Total `rows_feed_items`: \*\*([0-9,]+)\*\*", raw)
    hourly_rows_post_labels = _re1(r"Total `rows_post_labels`: \*\*([0-9,]+)\*\*", raw)
    hourly_rows_posts_first_seen = _re1(r"Total `rows_posts_first_seen`: \*\*([0-9,]+)\*\*", raw)

    wide_runs = _re1(r"### `wide`\s+\n\s*- Runs: ([0-9,]+) \(", raw)
    # Disambiguate the wide totals by anchoring in the wide section.
    wide_block = _re1(r"### `wide`\s+([\s\S]+?)### `labelerexp_hourly`", raw, group=1)
    wide_rows_feed_items = _re1(r"Total `rows_feed_items`: \*\*([0-9,]+)\*\*", wide_block)
    wide_rows_post_labels = _re1(r"Total `rows_post_labels`: \*\*([0-9,]+)\*\*", wide_block)
    wide_rows_posts_first_seen = _re1(r"Total `rows_posts_first_seen`: \*\*([0-9,]+)\*\*", wide_block)

    labelerexp_runs = _re1(r"### `labelerexp_hourly`\s+\n\s*- Runs: ([0-9,]+) \(", raw)
    labelerexp_block = _re1(r"### `labelerexp_hourly`\s+([\s\S]+?)## E\)", raw, group=1)
    labelerexp_rows_post_labels = _re1(r"Total `rows_post_labels`: \*\*([0-9,]+)\*\*", labelerexp_block)

    label_scan_rows = _re1(r"Scanned \*\*([0-9,]+)\*\* label rows across \*\*[0-9,]+\*\* CSV part files", raw)
    label_scan_parts = _re1(r"Scanned \*\*[0-9,]+\*\* label rows across \*\*([0-9,]+)\*\* CSV part files", raw)

    label_val_rows: dict[str, str] = {}
    for val in ["porn", "sexual", "nudity", "graphic-media", "!no-unauthenticated", "spam"]:
        label_val_rows[val] = _re1(rf"\| `{re.escape(val)}` \| ([0-9,]+) \|", raw)

    top_label_src_did = _re1(r"Top `label_src` \(labeler DID\):[\s\S]+?\| `(did:plc:[^`]+)` \|", raw)
    top_label_src_rows = _re1(r"Top `label_src` \(labeler DID\):[\s\S]+?\| `did:plc:[^`]+` \| ([0-9,]+) \|", raw)

    control_state_db_size = _re1(r"## K\) control_state\.db[\s\S]*?- Size: \*\*([0-9.]+ MiB)\*\*", raw)

    control_table_rows: dict[str, str] = {}
    for table in [
        "author_registry",
        "feed_catalog",
        "feed_generator_index_global",
        "feed_generator_index_parts",
        "feed_generator_index_repo_tasks",
        "post_registry",
        "queue_posts",
        "runs",
        "wide_sweep_tasks",
    ]:
        control_table_rows[table] = _re1(rf"\| `{re.escape(table)}` \| ([0-9,]+) \|", raw)

    return ReportMetrics(
        generated_at_utc=generated_at_utc,
        non_artifact_files=non_artifact_files,
        non_artifact_size=non_artifact_size,
        artifact_files=artifact_files,
        artifact_size=artifact_size,
        dir_size_gib=dir_size_gib,
        hourly_runs=hourly_runs,
        hourly_rows_feed_items=hourly_rows_feed_items,
        hourly_rows_post_labels=hourly_rows_post_labels,
        hourly_rows_posts_first_seen=hourly_rows_posts_first_seen,
        wide_runs=wide_runs,
        wide_rows_feed_items=wide_rows_feed_items,
        wide_rows_post_labels=wide_rows_post_labels,
        wide_rows_posts_first_seen=wide_rows_posts_first_seen,
        labelerexp_runs=labelerexp_runs,
        labelerexp_rows_post_labels=labelerexp_rows_post_labels,
        label_scan_rows=label_scan_rows,
        label_scan_parts=label_scan_parts,
        label_val_rows=label_val_rows,
        top_label_src_did=top_label_src_did,
        top_label_src_rows=top_label_src_rows,
        control_state_db_size=control_state_db_size,
        control_table_rows=control_table_rows,
    )


def build_deck(*, report_path: Path, out_pptx: Path) -> None:
    m = _parse_report(report_path)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    blank = prs.slide_layouts[6]

    # Slide 1
    s1 = prs.slides.add_slide(blank)
    _set_slide_bg(s1)
    _add_title(s1, text="Research questions: audit scope (R1-R5)", y=gy(0), h=gh(6), name="TITLE")
    _add_takeaway(s1, text="Takeaway: bias as measurable exposure + concentration outcomes.", y=gy(6), h=gh(3), name="TAKEAWAY")
    _add_rqs_rows(
        s1,
        y0=gy(10),
        row_h=gh(7),
        row_gap=gh(1),
        accent_hex=COLORS_HEX["cyan"],
        items=[
            "Discovery defaults: concentration in onboarding + /feeds surfaces.",
            "Infrastructure: hosting/provider concentration (domain + service DID).",
            "Exposure: winner-takes-all creators across parallel feed rankings.",
            "Safety: label coverage + disagreement by viewer_mode and labelers.",
            "Risk funnels: discovery-path feeds concentrate labeled risky content.",
        ],
    )

    # Slide 2
    s2 = prs.slides.add_slide(blank)
    _set_slide_bg(s2)
    _add_title(s2, text="Research questions: papers that define metrics", y=gy(0), h=gh(6), name="TITLE")
    _add_takeaway(s2, text="Takeaway: cite exposure/fair-ranking surveys to justify definitions.", y=gy(6), h=gh(3), name="TAKEAWAY")
    _add_rqs_rows(
        s2,
        y0=gy(10),
        row_h=gh(7),
        row_gap=gh(1),
        accent_hex=COLORS_HEX["amber"],
        items=[
            "Exposure weights w(r) + position bias (10.1145/3219819.3220088).",
            "Fair ranking survey: metrics + constraints (10.1145/3533379).",
            "Learning-to-rank fairness: evaluation framing (10.1145/3533380).",
            "Fairness in recommendation: stakeholders + pitfalls (10.1145/3610302).",
            "Fairness definitions are political: scope claims (Narayanan FAT* 2018).",
        ],
    )

    # Slide 3
    s3 = prs.slides.add_slide(blank)
    _set_slide_bg(s3)
    _add_title(s3, text="Research questions: what “bias” means here", y=gy(0), h=gh(6), name="TITLE")
    _add_takeaway(s3, text="Takeaway: separate self-selection from ranking/entrypoint allocation.", y=gy(6), h=gh(3), name="TAKEAWAY")
    _add_rqs_rows(
        s3,
        y0=gy(10),
        row_h=gh(7),
        row_gap=gh(1),
        accent_hex=COLORS_HEX["coral"],
        items=[
            "Provider/creator bias: different exposure at equal merit (audit outcomes).",
            "Content/opinion bias: topics/frames get different exposure probabilities.",
            "Chronological vs ranked feeds: mechanism differs; test rank-vs-age.",
            "Echo chambers: measure exposure homogeneity; don’t assume “bias”.",
            "Moderation labels: measure coverage/divergence, not ground-truth truth.",
        ],
    )

    # Slide 4
    s4 = prs.slides.add_slide(blank)
    _set_slide_bg(s4)
    _add_title(s4, text="Research questions: data answers (from REPORT.md)", y=gy(0), h=gh(6), name="TITLE")
    _add_takeaway(
        s4,
        text=f"Takeaway: {m.non_artifact_size} corpus; hourly feed_items={m.hourly_rows_feed_items} rows.",
        y=gy(6),
        h=gh(3),
        name="TAKEAWAY",
    )
    _add_rqs_rows(
        s4,
        y0=gy(10),
        row_h=gh(7),
        row_gap=gh(1),
        accent_hex=COLORS_HEX["green"],
        items=[
            f"R1 discovery: metadata JSONL + feed_catalog ({m.control_table_rows['feed_catalog']} rows).",
            f"R2 hosting: provider_domain/service_did in feed_catalog ({m.control_table_rows['feed_catalog']} rows).",
            f"R3 exposure: hourly feed_items ({m.hourly_rows_feed_items} rows; {m.hourly_runs} runs).",
            f"R4 safety: post_labels scanned ({m.label_scan_rows} rows; {m.label_scan_parts} parts).",
            f"R5 robustness: labelerexp post_labels ({m.labelerexp_rows_post_labels} rows; {m.labelerexp_runs} runs).",
        ],
    )

    # Slide 5
    s5 = prs.slides.add_slide(blank)
    _set_slide_bg(s5)
    _add_cred_cards(
        s5,
        title="Collection credibility: QC + reproducibility",
        left_lines=[
            "Run totals (REPORT.md)",
            f"hourly runs: {m.hourly_runs}",
            f"wide runs: {m.wide_runs}",
            f"labelerexp runs: {m.labelerexp_runs}",
            "Flagged hourly anomalies:",
            "2026-02-18/02; 2026-02-22/04; 2026-02-27/22",
        ],
        right_lines=[
            "Auditable artifacts",
            f"control_state.db size: {m.control_state_db_size}",
            f"feed_catalog rows: {m.control_table_rows['feed_catalog']}",
            f"post_registry rows: {m.control_table_rows['post_registry']}",
            "Join spine keys (REPORT.md):",
            "run_id; viewer_mode+vantage_id; feed_uri; post_uri(+post_cid); author_did",
        ],
    )

    # Slide 6
    s6 = prs.slides.add_slide(blank)
    _set_slide_bg(s6)
    _add_artifacts_grid(
        s6,
        title="Data artifacts collected: 2×2 audit table (Slide 6 micro-reveals)",
        cards=[
            (
                "Inventory (files + bytes)",
                [
                    ("report_utc", m.generated_at_utc),
                    ("files", m.non_artifact_files),
                    ("size", m.non_artifact_size),
                    ("artifact_files", m.artifact_files),
                    ("artifact_size", m.artifact_size),
                    ("hourly GiB", m.dir_size_gib["hourly"]),
                    ("effective_csv GiB", m.dir_size_gib["effective_csv"]),
                    ("exports GiB", m.dir_size_gib["exports"]),
                    ("wide GiB", m.dir_size_gib["wide"]),
                    ("labelerexp GiB", m.dir_size_gib["labelerexp"]),
                ],
            ),
            (
                "Run totals (progress.json)",
                [
                    ("hourly runs", m.hourly_runs),
                    ("hourly feed_items", m.hourly_rows_feed_items),
                    ("hourly post_labels", m.hourly_rows_post_labels),
                    ("hourly posts_first_seen", m.hourly_rows_posts_first_seen),
                    ("wide runs", m.wide_runs),
                    ("wide feed_items", m.wide_rows_feed_items),
                    ("wide post_labels", m.wide_rows_post_labels),
                    ("wide posts_first_seen", m.wide_rows_posts_first_seen),
                    ("labelerexp runs", m.labelerexp_runs),
                    ("labelerexp post_labels", m.labelerexp_rows_post_labels),
                ],
            ),
            (
                "Labels (scanned post_labels)",
                [
                    ("label rows", m.label_scan_rows),
                    ("CSV parts", m.label_scan_parts),
                    ("porn", m.label_val_rows["porn"]),
                    ("sexual", m.label_val_rows["sexual"]),
                    ("nudity", m.label_val_rows["nudity"]),
                    ("graphic-media", m.label_val_rows["graphic-media"]),
                    ("!no-unauth", m.label_val_rows["!no-unauthenticated"]),
                    ("spam", m.label_val_rows["spam"]),
                    ("top label_src", m.top_label_src_did),
                    ("top label_src rows", m.top_label_src_rows),
                ],
            ),
            (
                "control_state.db (SQLite)",
                [
                    ("db size", m.control_state_db_size),
                    ("author_registry", m.control_table_rows["author_registry"]),
                    ("feed_catalog", m.control_table_rows["feed_catalog"]),
                    ("post_registry", m.control_table_rows["post_registry"]),
                    ("runs", m.control_table_rows["runs"]),
                    ("wide_sweep_tasks", m.control_table_rows["wide_sweep_tasks"]),
                    ("index_repo_tasks", m.control_table_rows["feed_generator_index_repo_tasks"]),
                    ("index_parts", m.control_table_rows["feed_generator_index_parts"]),
                    ("index_global", m.control_table_rows["feed_generator_index_global"]),
                    ("queue_posts", m.control_table_rows["queue_posts"]),
                ],
            ),
        ],
    )

    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_pptx))


def main() -> None:
    parser = argparse.ArgumentParser(description="Build strict 6-slide meeting deck (pre-animation).")
    parser.add_argument("--report", type=Path, default=Path("data_v2_full/exports/inspection_2026-02-27/REPORT.md"))
    parser.add_argument("--out", type=Path, default=Path("_build/bluesky_meeting_rqs_papers_data_6slides_strict_preanim.pptx"))
    args = parser.parse_args()

    report = args.report.resolve()
    if not report.exists():
        raise SystemExit(f"Missing REPORT.md: {report}")

    build_deck(report_path=report, out_pptx=args.out.resolve())
    print(f"OK: wrote {args.out}")


if __name__ == "__main__":
    main()
