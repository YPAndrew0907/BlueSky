#!/usr/bin/env python3
from __future__ import annotations

import argparse
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation


EMU_PER_INCH = 914400


def _iter_shapes(slide):
    for shape in slide.shapes:
        yield shape
        if hasattr(shape, "shapes"):
            for child in shape.shapes:
                yield child


@dataclass(frozen=True)
class Box:
    slide_num: int
    shape_id: int
    name: str
    left: int
    top: int
    right: int
    bottom: int

    @property
    def width(self) -> int:
        return self.right - self.left

    @property
    def height(self) -> int:
        return self.bottom - self.top


def _shape_box(slide_num: int, shape) -> Box | None:
    w = int(getattr(shape, "width", 0))
    h = int(getattr(shape, "height", 0))
    if w <= 0 or h <= 0:
        return None
    left = int(getattr(shape, "left", 0))
    top = int(getattr(shape, "top", 0))
    return Box(
        slide_num=slide_num,
        shape_id=int(getattr(shape, "shape_id", 0)),
        name=(getattr(shape, "name", "") or "").strip(),
        left=left,
        top=top,
        right=left + w,
        bottom=top + h,
    )


def _is_thin_connector(box: Box, *, max_thickness_in: float = 0.03) -> bool:
    thr = int(round(max_thickness_in * EMU_PER_INCH))
    return box.width <= thr or box.height <= thr


def _overlap_area(a: Box, b: Box) -> tuple[int, int]:
    iw = min(a.right, b.right) - max(a.left, b.left)
    ih = min(a.bottom, b.bottom) - max(a.top, b.top)
    return iw, ih


def check_no_overlaps(
    *,
    pptx_path: Path,
    eq_mode: bool = False,
    ignore_bg_prefix: str = "_BG",
    ignore_eq_prefix: str = "EQ_STEP_",
) -> list[str]:
    prs = Presentation(str(pptx_path))
    issues: list[str] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        boxes: list[Box] = []
        for sh in _iter_shapes(slide):
            box = _shape_box(slide_num, sh)
            if box is None:
                continue
            if box.name.startswith(ignore_bg_prefix):
                continue
            if _is_thin_connector(box):
                continue
            boxes.append(box)

        for i in range(len(boxes)):
            a = boxes[i]
            for j in range(i + 1, len(boxes)):
                b = boxes[j]
                if eq_mode and a.name.startswith(ignore_eq_prefix) and b.name.startswith(ignore_eq_prefix):
                    # Whitelist stacked equation panels only among themselves.
                    continue
                iw, ih = _overlap_area(a, b)
                if iw > 0 and ih > 0:
                    issues.append(
                        f"slide {slide_num:02d}: overlap area={iw}x{ih} "
                        f"a(id={a.shape_id},name={a.name}) b(id={b.shape_id},name={b.name})"
                    )

    return issues


def main() -> None:
    parser = argparse.ArgumentParser(description="Strict overlap QC (positive-area overlaps).")
    parser.add_argument("pptx", type=Path, nargs="?", default=Path("_build/bluesky_meeting_rqs_papers_data_6slides_strict_animated.pptx"))
    parser.add_argument("--eq-mode", action="store_true", help="Whitelist overlaps among EQ_STEP_* shapes.")
    args = parser.parse_args()

    pptx_path = args.pptx.resolve()
    if not pptx_path.exists():
        raise SystemExit(f"Missing pptx: {pptx_path}")

    issues = check_no_overlaps(pptx_path=pptx_path, eq_mode=bool(args.eq_mode))
    if issues:
        print("FAIL: overlaps detected")
        for line in issues[:80]:
            print(line)
        if len(issues) > 80:
            print(f"... ({len(issues)} total)")
        raise SystemExit(1)

    print("OK: no overlaps")


if __name__ == "__main__":
    main()

